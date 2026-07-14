# CUI // SP-CTI
"""How much a document's word is worth, and which documents are copies.

Invented content — ICDEV is a public repo. The behaviours are the real ones.
"""
from __future__ import annotations

import openpyxl

from tests.bom import fixtures
from tools.bom import constants as C
from tools.bom.credibility import assess, load_config, propose_role
from tools.bom.derivative import find_derivatives
from tools.bom.extract_grid import extract_grid


def _wb(path, sheet_rows: list[list], title="Sheet1"):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = title
    for row in sheet_rows:
        ws.append(row)
    wb.save(path)
    return path


class TestCredibilityIsProposedNeverImposed:
    def test_an_assessment_is_never_binding_on_its_own(self, tmp_path):
        p = _wb(tmp_path / "a.xlsx", [["Item", "Cost"], ["Widget", 100]])
        a = assess(extract_grid(p))
        # A human's designation is what binds. Until then this is a suggestion
        # with a reason attached, and the reason is the point: somebody can
        # disagree with it specifically rather than just distrust the number.
        assert a.set_by == "ai_proposed"
        assert a.rationale

    def test_a_document_with_no_signal_stays_unknown(self, tmp_path):
        p = _wb(tmp_path / "bare.xlsx", [["a", "b"], ["c", "d"]])
        a = assess(extract_grid(p))
        # Silence is never confirmation. Nothing is trusted for having turned up.
        assert a.tier in ("unknown", "draft")
        assert C.CREDIBILITY_RANK[a.tier] >= C.CREDIBILITY_RANK["working"]


class TestTheVocabularyIsConfigurationNotCode:
    def test_the_customers_own_marker_is_honoured(self, tmp_path):
        """Users encode credibility by renaming files. Capture it; do not guess it.

        The lexicon lives in args/bom_credibility.yaml precisely so that a
        customer with a different convention adds a word rather than a patch.
        """
        cfg = load_config()
        assert "solid" in cfg["vocabulary"]["upgrade"]
        assert "wip" in cfg["vocabulary"]["downgrade"]

    def test_a_marked_final_file_outranks_a_marked_draft(self, tmp_path):
        rows = [["Item", "Part", "Qty", "Cost"], ["Switch", "SW-9500-16X", 1, 21000]]
        final = _wb(tmp_path / "network bom (solid) FINAL.xlsx", rows)
        draft = _wb(tmp_path / "network bom WIP draft v0.2.xlsx", rows)

        a, b = assess(extract_grid(final)), assess(extract_grid(draft))
        assert C.CREDIBILITY_RANK[a.tier] < C.CREDIBILITY_RANK[b.tier]

    def test_word_boundaries_are_respected(self, tmp_path):
        """"solid" must not fire on "solidarity"; "old" must not fire on "gold"."""
        p = _wb(tmp_path / "solidarity gold report.xlsx", [["Item", "Cost"], ["x", 1]])
        a = assess(extract_grid(p))
        names = {s.name for s in a.signals}
        assert "vocab:solid" not in names
        assert "vocab:old" not in names


class TestStructuralSignals:
    def test_live_formulas_are_worth_something(self, tmp_path):
        """Reproducible arithmetic can be argued with. A flat sheet is an assertion."""
        rich = fixtures.workbook_with_a_zeroed_line(tmp_path / "rich.xlsx")
        flat = _wb(tmp_path / "flat.xlsx", [["Item", "Qty", "Unit", "Extended"],
                                            ["Widget A", 2, 100, 200]])
        a, b = assess(extract_grid(rich)), assess(extract_grid(flat))
        assert a.score > b.score

    def test_a_script_written_document_is_marked_down(self, tmp_path):
        """Not a scandal — but a generated estimate is a different kind of evidence
        from one a person authored, and letting a machine's guess quietly acquire
        the authority of a human's judgement is how a bad number gets believed."""
        p = _wb(tmp_path / "gen.xlsx", [["Item", "Cost"], ["x", 1]])
        a = assess(extract_grid(p))
        assert any(s.name == "machine_generated" for s in a.signals)
        assert a.score < 0

    def test_a_heavily_downgraded_document_is_a_draft_not_an_unknown(self, tmp_path):
        """We know a great deal about it, none of it reassuring.

        'unknown' means no signal at all. Conflating the two would let a source we
        have every reason to distrust sit in the same bucket as one nobody has
        looked at yet.
        """
        p = _wb(tmp_path / "WIP draft rough v0.1 obsolete.xlsx", [["Item", "Cost"], ["x", 1]])
        a = assess(extract_grid(p))
        assert a.score < -5
        assert a.tier == "draft"


class TestRoleDetection:
    def test_serial_numbers_signal_an_inventory(self, tmp_path):
        """The distinction the whole engine turns on.

        A BOM CLAIMS things. An inventory IDENTIFIES INDIVIDUAL PHYSICAL UNITS.
        Serial numbers are the signature of the second — and that is how ground
        truth gets recognised without hardcoding anybody's filename. It generalises
        to any corpus in any industry.
        """
        p = _wb(tmp_path / "assets.xlsx", [
            ["Model", "Serial", "Warranty Ended"],
            ["Node A", "6HYXHB2", "2019-06-02"],
            ["Node B", "16H6FB2", "2019-03-22"],
            ["Node C", "33DBXG2", "2019-12-30"],
            ["Node D", "G62Q182", "2018-10-22"],
        ])
        role, conf, _ = propose_role(extract_grid(p), load_config())
        assert role == "inventory_truth"
        assert conf > 0.7

    def test_a_priced_list_is_a_claim_not_an_inventory(self, tmp_path):
        p = _wb(tmp_path / "bom.xlsx", [
            ["Item", "Part Number", "Qty", "Unit Price"],
            ["Switch", "SW-9500", 2, 21000],
            ["Router", "RT-8300", 1, 17500],
        ])
        role, _, _ = propose_role(extract_grid(p), load_config())
        assert role == "bom_claim"

    def test_an_inventory_is_worth_listening_to_whatever_it_is_called(self, tmp_path):
        """It is authoritative about COUNTS. That is what it is for.

        Role and tier are separate columns precisely because an inventory that can
        verify how many machines exist is still not authoritative about what they
        cost.
        """
        p = _wb(tmp_path / "scratch notes.xlsx", [
            ["Model", "Serial", "Warranty Ended"],
            ["Node A", "6HYXHB2", "2019-06-02"],
            ["Node B", "16H6FB2", "2019-03-22"],
            ["Node C", "33DBXG2", "2019-12-30"],
        ])
        a = assess(extract_grid(p))
        assert a.role == "inventory_truth"
        assert C.CREDIBILITY_RANK[a.tier] <= C.CREDIBILITY_RANK["working"]


class TestDerivatives:
    """The same document, twice, in two formats."""

    def test_a_copy_cannot_outrank_the_thing_it_copies(self, tmp_path):
        p = _wb(tmp_path / "print.pdf.xlsx", [["Item", "Cost"], ["x", 1]])
        a = assess(extract_grid(p), derivative_of="original.xlsx")
        assert a.tier == "derived"
        assert a.role == "derived"
        # No amount of confident vocabulary rescues a print of a spreadsheet: it
        # has lost the formulas, which is where the errors were hiding.
        assert C.CREDIBILITY_RANK["derived"] > C.CREDIBILITY_RANK["working"]

    def test_two_same_format_variants_are_never_treated_as_copies(self, tmp_path):
        """The most important guard in the module.

        Two same-format documents that are 99% identical are almost never one
        copied from the other. They are VARIANTS — the same design with one
        technology swapped for another, the same plan priced two ways. Excluding
        one from the rollups would silently DELETE AN OPTION the customer is still
        choosing between, and the deck would reach leadership missing an entire
        alternative with nothing to show it had ever existed.
        """
        common = [["Item", "Cost"], ["Core Switch", 21000], ["Firewall", 10500]]
        a = _wb(tmp_path / "design_option_a.xlsx", [*common, ["Platform A Licence", 48000]])
        b = _wb(tmp_path / "design_option_b.xlsx", [*common, ["Platform B Licence", 48000]])

        found = find_derivatives([extract_grid(a), extract_grid(b)])
        assert found == [], "a competing option was silently deleted as a duplicate"

    def test_a_copy_gets_exactly_one_attribution(self, tmp_path):
        """A document cannot be derived from two things at once.

        Emitting several attributions leaves whoever reads the register to pick one
        at random — and if the candidates are competing options, picking at random
        assigns the copy to the wrong funding decision.
        """
        rows = [["Item", "Part", "Cost"],
                ["Core Switch Assembly", "SW-9500-16X", 21000],
                ["Perimeter Firewall", "FW-2110-NGFW", 10500],
                ["Console Server Unit", "CS-16-PORT", 800]]
        wb = _wb(tmp_path / "source.xlsx", rows)

        from pptx import Presentation
        from pptx.util import Emu
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                box = s.shapes.add_textbox(
                    Emu(400_000 + c * 3_000_000), Emu(400_000 + r * 800_000),
                    Emu(2_800_000), Emu(600_000),
                )
                box.text_frame.text = str(val)
        deck = tmp_path / "printed.pptx"
        prs.save(deck)

        found = find_derivatives([extract_grid(wb), extract_grid(deck)])
        assert len(found) == 1
        d = found[0]
        assert d.derived == "printed.pptx"
        assert d.original == "source.xlsx"   # the copy with the structure wins
        assert d.overlap >= 0.8

    def test_numbers_are_compared_as_numbers_not_as_printed_text(self):
        """Regression, and the bug that made this module find nothing at all.

        A workbook stores 6400. Its PDF renders "$6,400.00". Stripping punctuation
        from the rendered string gives "640000" — which matches nothing — so two
        copies of the same document shared exactly ZERO rows and the duplicate
        sailed straight through.
        """
        from tools.bom.derivative import _token

        class _C:
            def __init__(self, text, num):
                self.value_text, self.value_num = text, num

        assert _token(_C("6400", 6400.0)) == _token(_C("$6,400.00", 6400.0))
