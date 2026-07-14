# CUI // SP-CTI
"""The workbook and the deck.

Invented content. ICDEV is a public repo.
"""
from __future__ import annotations

import io

import openpyxl
import pytest
from pptx import Presentation

from tools.bom.deck import NotReadyForLeadership, build_deck, freeze, render
from tools.bom.export_xlsx import export
from tools.bom.findings import Evidence, Finding
from tools.bom.lines import ExtractedLine
from tools.bom.pivot import Dataset, Row, pivot
from tools.bom.reconcile import Source


def _line(lid, desc, price, doc="bom.xlsx", cell="A1"):
    return ExtractedLine(
        line_id=lid, line_hash=f"h{lid}", source_document=doc, source_sheet="Sheet1",
        source_locator=cell, raw_text=desc, description=desc,
        part_number="SW-9500", manufacturer="Acme",
        qty=1, unit_price=price, extended_price=price,
    )


def _row(lid, price, *, committed=True, **dims):
    return Row(
        line_id=lid, cluster_id=f"c{lid}", description=f"item {lid}",
        dims=dims or {"manufacturer": "Acme", "price_basis": "msrp"},
        extended_price=price, qty=1, unit_price=price, committed=committed,
        excluded_reason="" if committed else "sources disagree",
    )


@pytest.fixture
def scenario():
    ds = Dataset(
        rows=[_row("a", 21000), _row("b", 9999, committed=False)],
        claim_sources={"bom.xlsx"},
    )
    lines = [_line("a", "Core switch", 21000), _line("b", "Firewall", 9999)]
    findings = [
        Finding(
            finding_type="unpriced_line_zeroed", kind="defect", severity="high",
            title="A line computes to zero", detail="An operand was never filled in.",
            impact_usd=None,
            evidence=[Evidence("bom.xlsx", "Sheet1", "E9", "chassis")],
        ),
        Finding(
            finding_type="hardcoded_rollup", kind="defect", severity="high",
            title="A typed number among formulas", impact_usd=192000.0,
            evidence=[Evidence("bom.xlsx", "Summary", "C15", "192000")],
        ),
    ]
    sources = {
        "bom.xlsx": Source("s1", credibility_tier="authoritative"),
        "print.pdf": Source("s2", credibility_tier="derived", role="derived"),
    }
    return ds, lines, findings, sources


class TestTheWorkbook:
    def test_it_opens(self, scenario):
        ds, lines, findings, sources = scenario
        wb = openpyxl.load_workbook(io.BytesIO(export(ds, findings, sources, lines)))
        assert wb.sheetnames == ["The Ask", "Bill of Materials", "Findings", "Sources"]

    def test_every_line_carries_its_provenance(self, scenario):
        """A reconciled BOM is only worth anything if a sceptic can take any number,
        follow it back to the cell it came from, and satisfy themselves it is real."""
        ds, lines, findings, sources = scenario
        wb = openpyxl.load_workbook(io.BytesIO(export(ds, findings, sources, lines)))
        ws = wb["Bill of Materials"]

        header = [c.value for c in ws[5]]
        assert "Source document" in header
        assert "Source cell" in header

        doc_col = header.index("Source document") + 1
        cell_col = header.index("Source cell") + 1
        assert ws.cell(6, doc_col).value == "bom.xlsx"
        assert ws.cell(6, cell_col).value == "Sheet1!A1"

    def test_findings_are_sorted_by_money(self, scenario):
        """The only ordering anybody with a budget has ever used."""
        ds, lines, findings, sources = scenario
        wb = openpyxl.load_workbook(io.BytesIO(export(ds, findings, sources, lines)))
        ws = wb["Findings"]
        impacts = [ws.cell(r, 3).value for r in range(6, 8)]
        assert impacts[0] == 192000.0     # priced first
        assert impacts[1] is None         # the one we refused to price

    def test_the_ask_states_what_is_excluded(self, scenario):
        ds, lines, findings, sources = scenario
        wb = openpyxl.load_workbook(io.BytesIO(export(ds, findings, sources, lines)))
        ws = wb["The Ask"]
        labels = [ws.cell(r, 1).value for r in range(1, 15)]
        blob = " ".join(str(x) for x in labels if x)
        assert "Committed" in blob
        assert "Open" in blob
        assert "NOTHING" in blob


class TestTheSnapshotIsFrozen:
    def test_the_same_evidence_gives_the_same_hash(self, scenario):
        """A figure quoted back at you six weeks later must be traceable to the
        exact state of the evidence that produced it."""
        ds, _lines, findings, _sources = scenario
        assert freeze(ds, findings).sha == freeze(ds, findings).sha

    def test_different_evidence_gives_a_different_hash(self, scenario):
        ds, _lines, findings, _sources = scenario
        moved = Dataset(rows=[_row("a", 22000)], claim_sources={"bom.xlsx"})
        assert freeze(ds, findings).sha != freeze(moved, findings).sha


class TestTheDeck:
    def test_it_renders_real_tables_and_a_talk_track(self, scenario, tmp_path):
        """The customer's own leadership deck had ZERO speaker notes across fifteen
        slides — we flagged it. Shipping a deck with the same hole would be
        embarrassing, and it nearly happened: the builder takes its table as
        bullets={headers, rows, footer}, and passing a bare list rendered the slide
        as the words "No table data.\""""
        ds, lines, findings, sources = scenario
        p = pivot(ds, rows="manufacturer", cols="price_basis")
        slides, _snap = build_deck(ds, findings, sources, pivot=p, project="Lab",
                                   audience="working")

        prs = Presentation(render(slides, theme="investment_deck", title="Lab"))
        tables = sum(
            1 for s in prs.slides for sh in s.shapes if getattr(sh, "has_table", False)
        )
        notes = sum(
            1 for s in prs.slides
            if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip()
        )
        assert tables >= 2
        assert notes >= len(prs.slides) - 1

    def test_the_tables_can_actually_be_READ(self, scenario):
        """The bug this exists for: every table rendered black-on-black.

        python-pptx creates tables with "Medium Style 2 - Accent 1", whose own
        text colours are applied at the RUN level and beat anything set on the
        paragraph. So the builder painted a dark fill, asked for white text, and
        PowerPoint drew the table style's dark text on top. Every table came out
        an empty box.

        Nothing about that is visible from Python. The text is in the XML,
        python-pptx reads it back happily, and a test asserting "the table has
        three rows and the right values" PASSES. It did pass. The deck went to the
        customer and the tables were blank.

        So this test does not ask whether the text is THERE. It asks whether a
        human could see it — contrast, at the run level, against the fill actually
        painted underneath.
        """
        ds, lines, findings, sources = scenario
        p = pivot(ds, rows="manufacturer", cols="price_basis")
        slides, _ = build_deck(ds, findings, sources, pivot=p, project="Lab",
                                   audience="working")
        prs = Presentation(render(slides, theme="investment_deck", title="Lab"))

        def luminance(rgb: str) -> float:
            def chan(c: float) -> float:
                c /= 255
                return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

            r, g, b = (int(str(rgb)[i:i + 2], 16) for i in (0, 2, 4))
            return 0.2126 * chan(r) + 0.7152 * chan(g) + 0.0722 * chan(b)

        def contrast(fg: str, bg: str) -> float:
            a, b = luminance(fg), luminance(bg)
            return (max(a, b) + 0.05) / (min(a, b) + 0.05)

        checked = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_table", False):
                    continue
                for row in shape.table.rows:
                    for cell in row.cells:
                        if not cell.text.strip():
                            continue
                        runs = cell.text_frame.paragraphs[0].runs
                        assert runs, f"no run on {cell.text!r} — colour cannot be set"

                        fg = runs[0].font.color.rgb
                        assert fg is not None, (
                            f"{cell.text!r} has no RUN-level colour; the table "
                            f"style will supply one and it will be the wrong one"
                        )
                        ratio = contrast(str(fg), str(cell.fill.fore_color.rgb))
                        assert ratio >= 4.5, (
                            f"{cell.text!r} renders at {ratio:.1f}:1 — invisible. "
                            f"WCAG AA needs 4.5:1."
                        )
                        checked += 1

        assert checked > 6, "no table cells were actually checked"

    def test_the_last_slide_is_an_explicit_outro(self, scenario):
        """The builder renders the LAST slide as an outro whatever type it claims to
        be. Without an explicit one it quietly eats a real slide — "Who said so" was
        being turned into a thank-you card, and nothing said so."""
        ds, _lines, findings, sources = scenario
        slides, _ = build_deck(ds, findings, sources, project="Lab", audience="working")
        assert slides[-1]["slide_type"] == "outro"

    def test_every_slide_cites_the_snapshot(self, scenario):
        ds, _lines, findings, sources = scenario
        slides, snap = build_deck(ds, findings, sources, project="Lab", audience="working")
        for s in slides:
            assert any(snap.sha in c["title"] for c in s["citations"])

    def test_an_unknown_theme_is_refused(self, scenario):
        ds, _lines, findings, sources = scenario
        slides, _ = build_deck(ds, findings, sources, project="Lab", audience="working")
        with pytest.raises(ValueError, match="unknown theme"):
            render(slides, theme="hot_pink")


class TestTheTwoAudiencesAnswerDifferentQuestions:
    """Leadership is being asked to fund an outcome. The workgroup is being asked to
    close the gaps. A deck that shows the working to an executive reads as hedging,
    and a deck that hides it from the workgroup is useless to them."""

    def _scenario(self):
        ds = Dataset(rows=[_row("switch", 36000)], claim_sources={"bom.xlsx"})
        f = Finding(finding_type="hardcoded_rollup", kind="defect", severity="high",
                    title="A typed-in total", impact_usd=1000,
                    evidence=[Evidence("bom.xlsx", "Summary", "C15", "")])
        return ds, [f]

    def test_the_working_deck_shows_what_we_found(self):
        ds, findings = self._scenario()
        slides, _ = build_deck(ds, findings, {}, project="Lab", audience="working")
        titles = [s["title"] for s in slides]
        assert any("found" in t.lower() for t in titles), titles

    def test_the_leadership_deck_does_NOT(self):
        """It is not that the findings are unimportant. It is that they are the
        reason this deck is allowed to state a number, not the deck's content."""
        ds, findings = self._scenario()
        slides, _ = build_deck(ds, findings, {}, project="Lab", audience="leadership")
        titles = [s["title"].lower() for s in slides]
        assert not any("found" in t for t in titles), titles
        assert not any("said so" in t for t in titles), titles
        assert not any("check this" in t for t in titles), titles

    def test_the_leadership_title_slide_leads_with_the_ask(self):
        ds, findings = self._scenario()
        slides, _ = build_deck(ds, findings, {}, project="Lab", audience="leadership")
        assert "$36,000" in slides[0]["bullets"][0]

    def test_an_unknown_audience_is_refused(self):
        ds, findings = self._scenario()
        with pytest.raises(ValueError, match="unknown audience"):
            build_deck(ds, findings, {}, project="Lab", audience="the board")


class TestALeadershipDeckCannotBePolishedPastNoNumber:
    """The one place the tool refuses rather than degrades.

    A polished deck states a number with confidence. There is no honest way to do
    that over four documents that price the same project differently — and a caveat
    in six-point type is not a fix, it is a disclaimer on a lie. The fix is a human
    nominating a source of record, which takes minutes because the reconciliation is
    already done.
    """

    def _competing(self):
        return Dataset(
            rows=[_row("switch", 36000)],
            claim_sources={"bom_a.xlsx", "bom_b.xlsx"},
        )

    def test_it_refuses(self):
        with pytest.raises(NotReadyForLeadership, match="2 sources"):
            build_deck(self._competing(), [], {}, project="Lab",
                       audience="leadership")

    def test_the_refusal_says_what_to_do_about_it(self):
        with pytest.raises(NotReadyForLeadership, match="source of record"):
            build_deck(self._competing(), [], {}, project="Lab",
                       audience="leadership")

    def test_the_working_deck_still_builds_and_states_it_honestly(self):
        """Refusing the leadership deck must not block the workgroup — they are the
        ones who can resolve it."""
        slides, _ = build_deck(self._competing(), [], {}, project="Lab",
                               audience="working")
        assert any("not yet a number" in s["title"] for s in slides)

    def test_one_source_of_record_unblocks_the_leadership_deck(self):
        ds = Dataset(rows=[_row("switch", 36000)], claim_sources={"bom_a.xlsx"})
        slides, snap = build_deck(ds, [], {}, project="Lab", audience="leadership")
        assert snap.is_a_total
        assert slides


class TestTheDeckRefusesToPrintATotalItDoesNotHave:
    def test_competing_claims_replace_the_ask_with_the_truth(self):
        """The room's instinct is to ask for 'the number'. The honest answer is that
        several documents disagree and somebody has to say which one governs."""
        ds = Dataset(
            rows=[_row("a", 21000)],
            claim_sources={"bom_one.xlsx", "bom_two.xlsx"},
        )
        slides, snap = build_deck(ds, [], {}, project="Lab", audience="working")

        assert snap.is_a_total is False
        ask = slides[1]
        assert "not yet a number" in ask["title"]
        assert any("competing estimates" in b.lower() for b in ask["bullets"])
        # And it names them, so nobody has to guess which documents are fighting.
        assert any("bom_one.xlsx" in b for b in ask["bullets"])

    def test_a_settled_dataset_states_the_figure(self):
        ds = Dataset(rows=[_row("a", 21000)], claim_sources={"bom.xlsx"})
        slides, snap = build_deck(ds, [], {}, project="Lab", audience="working")
        assert snap.is_a_total is True
        assert slides[1]["title"] == "The ask"
        assert any("21,000" in b for b in slides[1]["bullets"])


class TestWhatYouAlreadyOwn:
    def test_it_is_the_best_news_and_gets_its_own_slide(self):
        """Hardware already in the building costs nothing and is frequently the
        reason a team can start now instead of waiting a year. On a cost-sorted
        table it is the last row — which is a product failure, not a formatting one.
        """
        ds = Dataset(rows=[_row("a", 21000)], claim_sources={"bom.xlsx"})
        slides, _ = build_deck(
            ds, [], {}, project="Lab",
            owned_value=95000,
            owned_note="Nineteen machines, being repurposed rather than replaced.",
        )
        titles = [s["title"] for s in slides]
        assert "What we already own" in titles

    def test_it_is_omitted_when_there_is_nothing_to_say(self):
        ds = Dataset(rows=[_row("a", 21000)], claim_sources={"bom.xlsx"})
        slides, _ = build_deck(ds, [], {}, project="Lab", audience="working")
        assert "What we already own" not in [s["title"] for s in slides]
