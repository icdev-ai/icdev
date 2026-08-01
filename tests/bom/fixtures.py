# CUI // SP-CTI
"""Documents that misbehave in the specific ways real ones do.

These are BUILT, not checked in. Two reasons, and the second is the important
one.

First, hermetic: they run in CI, where no customer corpus exists and none should.

Second, and this is the whole point of the module: the failure modes this engine
hunts are *structural*. A formula multiplying a quantity by an empty cell yields
zero whether the item is a network switch or a marine gearbox. A worksheet that
reports A1:A1 can still be carrying an anchored image, in any workbook anyone has
ever made. A gap in the sheetId sequence means a sheet was deleted, always.

So the fixtures reproduce the *shapes* of the bugs, with invented content. An
engine that passes against these will find the same defects in documents nobody
here has ever seen — which is the actual requirement, and is a stronger claim
than "it works on the one corpus we had."
"""
from __future__ import annotations

import shutil
import zipfile
from pathlib import Path


def _inject_cached_values(path: Path, values: dict[str, dict[str, float]]) -> None:
    """Give formula cells the value Excel would have cached.

    openpyxl never RUNS a formula, so a workbook it writes has ``<f>`` and no
    ``<v>``. Load it with ``data_only=True`` and every formula cell comes back
    None. A workbook a human saved from Excel carries both.

    Both exist in the wild and the engine has to survive both — a script-written
    estimate is exactly the kind whose arithmetic most wants checking. So the
    fixtures make one of each rather than quietly only testing the easy one.
    """
    import re

    with zipfile.ZipFile(path) as src:
        items = [(i, src.read(i.filename)) for i in src.infolist()]

    # Keyed by SHEET, then by cell. Keying on the cell ref alone writes the same
    # cached value into every sheet that happens to use that address — B4 exists
    # on all of them — and a fixture that quietly lies produces tests that pass
    # against behaviour the engine does not have.
    order = [n for n in (i.filename for i, _ in items)
             if n.startswith("xl/worksheets/sheet") and n.endswith(".xml")]
    order.sort(key=lambda n: int(re.search(r"sheet(\d+)\.xml", n).group(1)))
    sheet_names = list(values)

    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as dst:
        for info, blob in items:
            if info.filename in order:
                idx = order.index(info.filename)
                name = sheet_names[idx] if idx < len(sheet_names) else None
                text = blob.decode("utf-8")
                for ref, val in (values.get(name) or {}).items():
                    # openpyxl emits an EMPTY <v></v> after the formula, so the
                    # cell already has the slot — it is just never filled, because
                    # nothing ever evaluated the formula.
                    #   <c r="D2"><f>B2*C2</f><v></v></c>
                    #     -> <c r="D2"><f>B2*C2</f><v>200</v></c>
                    text = re.sub(
                        rf'(<c r="{ref}"[^>]*>\s*<f>[^<]*</f>\s*)(?:<v>\s*</v>)?\s*</c>',
                        rf"\1<v>{val}</v></c>",
                        text,
                    )
                blob = text.encode("utf-8")
            dst.writestr(info, blob)


def workbook_with_a_zeroed_line(path: Path, *, cached: bool = True) -> Path:
    """A quantity, no unit price, and a formula that politely multiplies to zero.

    The line looks costed. It contributes nothing. The total is understated by
    whatever that item is worth, and the spreadsheet will never say so.

    ``cached=True`` mimics a workbook saved from Excel (formula AND the value it
    produced). ``cached=False`` mimics one a script wrote (formula, no value) —
    where the zero is not even visible, and the only evidence is a formula
    multiplying by a cell that does not exist.
    """
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "BOM"
    ws.append(["Item", "Qty", "Unit", "Extended"])

    ws.append(["Widget A", 2, 100, None])
    ws["D2"] = "=B2*C2"                      # 200 — a normal, healthy line

    ws.append(["Widget B", 1, None, None])   # unit price never entered
    ws["D3"] = "=B3*C3"                      # 0 — and nobody notices

    ws.append(["TOTAL", None, None, None])
    ws["D4"] = "=SUM(D2:D3)"

    wb.save(path)

    if cached:
        _inject_cached_values(path, {"BOM": {"D2": 200, "D3": 0, "D4": 200}})
    return path


def workbook_with_a_deleted_sheet(path: Path) -> Path:
    """sheetIds 1,2,4 — sheet 3 is gone.

    Excel never reuses a sheetId, so a hole in the run is a deletion. In a
    costing workbook that is a category somebody removed, and the surviving
    subtotals will not tell you whether its money went with it.

    openpyxl always writes a contiguous sequence, so the gap is punched into the
    saved package directly.
    """
    import openpyxl

    tmp = path.with_suffix(".seed.xlsx")
    wb = openpyxl.Workbook()
    wb.active.title = "Alpha"
    wb.create_sheet("Beta")
    wb.create_sheet("Gamma")
    wb.save(tmp)

    with zipfile.ZipFile(tmp) as src, zipfile.ZipFile(path, "w") as dst:
        for item in src.infolist():
            blob = src.read(item.filename)
            if item.filename == "xl/workbook.xml":
                text = blob.decode("utf-8")
                text = text.replace('sheetId="3"', 'sheetId="4"')
                blob = text.encode("utf-8")
            dst.writestr(item, blob)

    tmp.unlink()
    return path


def workbook_with_a_double_count(path: Path) -> Path:
    """The same licence, on two sheets, reaching the grand total twice.

    This is the distinction the whole formula graph exists for. The item appears
    on Networking and again on Simulation. A human even wrote "shared with the
    Networking sheet" next to the second occurrence — they KNEW. And both sheet
    subtotals still include it, so the money lands in the total twice.

    Note what makes it a bug: not that the item appears twice (a cross-reference
    is legitimate), but that two DIFFERENT subtotals consume it and both flow into
    the total. Had they fed the same subtotal, it would simply be a quantity of
    two, and nothing here would be wrong.
    """
    import openpyxl

    wb = openpyxl.Workbook()

    net = wb.active
    net.title = "Networking"
    net.append(["Item", "Cost", "Notes"])
    net.append(["Switch", 5000, ""])
    net.append(["Simulation Licence", 10000, ""])
    net["B4"] = "=SUM(B2:B3)"          # Networking subtotal = 15000

    sim = wb.create_sheet("Simulation")
    sim.append(["Item", "Cost", "Notes"])
    sim.append(["Test Harness", 2000, ""])
    sim.append(["Simulation Licence", 10000, "shared with Networking sheet"])
    sim["B4"] = "=SUM(B2:B3)"          # Simulation subtotal = 12000

    total = wb.create_sheet("Summary")
    total.append(["Sheet", "Subtotal"])
    total.append(["Networking", None])
    total["B2"] = "=Networking!B4"
    total.append(["Simulation", None])
    total["B3"] = "=Simulation!B4"
    total.append(["TOTAL", None])
    total["B4"] = "=SUM(B2:B3)"        # 27000 — of which 10000 is counted twice

    wb.save(path)
    _inject_cached_values(path, {
        "Networking": {"B4": 15000},
        "Simulation": {"B4": 12000},
        "Summary": {"B2": 15000, "B3": 12000, "B4": 27000},
    })
    return path


def workbook_with_a_hardcoded_rollup(path: Path) -> Path:
    """A summary block where one 'subtotal' is a number somebody typed.

    Its neighbours are live formulas. This one is a literal. Edit the sheets it
    claims to total and it will not move — and nothing on screen distinguishes it
    from the cells above and below.
    """
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Summary"
    ws.append(["Category", "Amount"])
    ws.append(["Compute", None])
    ws.append(["Networking", None])
    ws.append(["Storage", None])
    ws.append(["Security", None])
    ws.append(["TOTAL", None])

    ws["B2"] = "=SUM(Detail!B2:B3)"
    ws["B3"] = "=SUM(Detail!B4:B5)"
    ws["B4"] = 192000                    # <- typed in. Will never update.
    ws["B5"] = "=SUM(Detail!B6:B7)"
    ws["B6"] = "=SUM(B2:B5)"

    detail = wb.create_sheet("Detail")
    detail.append(["Item", "Cost"])
    for i in range(6):
        detail.append([f"Item {i}", 1000 * (i + 1)])

    wb.save(path)
    return path


def workbook_with_an_image_on_an_empty_sheet(path: Path, message: str) -> Path:
    """A worksheet with zero cells, carrying a picture.

    Excel reports the sheet as A1:A1. It opens blank. Nobody scrolling the
    workbook finds it. The picture is where somebody pasted a screenshot of the
    thing that actually constrains the project.
    """
    import openpyxl
    from openpyxl.drawing.image import Image as XLImage
    from PIL import Image, ImageDraw

    png = path.with_suffix(".png")
    img = Image.new("RGB", (760, 180), "white")
    draw = ImageDraw.Draw(img)
    # Large and plain: this has to survive OCR, and the test asserts on what
    # comes back out of the pixels.
    for i, line in enumerate(message.split("\n")):
        draw.text((14, 20 + i * 46), line, fill="black")
    img = img.resize((1520, 360), Image.LANCZOS)
    img.save(png)

    wb = openpyxl.Workbook()
    wb.active.title = "Data"
    wb["Data"]["A1"] = "something"
    blank = wb.create_sheet("Notes")     # not one cell written to it
    blank.add_image(XLImage(str(png)), "A1")
    wb.save(path)

    png.unlink(missing_ok=True)
    return path


def deck_with_a_table_that_is_not_a_table(path: Path) -> Path:
    """A grid drawn out of loose text boxes.

    ``shape.has_table`` is False for every one of them. This is what a
    script-generated deck looks like, and what you get when a human builds a
    table by duplicating a text box — so an extractor that only understands real
    tables reads the money slide as scattered words and loses the total in
    silence.
    """
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

    rows = [
        ("Category", "Amount"),
        ("Alpha", "$1,000"),
        ("Beta", "$2,500"),
        ("Total", "$3,500"),
    ]
    for r, (left_text, right_text) in enumerate(rows):
        # A couple of EMU of vertical jitter, because real decks are never
        # pixel-aligned and a clusterer that demands exactness is useless.
        jitter = 900 * (r % 2)
        for c, text in enumerate((left_text, right_text)):
            box = slide.shapes.add_textbox(
                Emu(500_000 + c * 3_000_000),
                Emu(500_000 + r * 800_000 + jitter),
                Emu(2_800_000),
                Emu(600_000),
            )
            box.text_frame.text = text

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "The long-lead item gates everything behind it. Order in week one or the "
        "schedule does not hold. Unit price for line 2 is still TBD."
    )
    prs.save(path)
    return path


def drawio_with_tabs(path: Path) -> Path:
    """A real .drawio file: <mxfile> wrapping one <diagram> per tab.

    The rack elevation is the only tab that says HOW MANY, so the tabs have to
    stay apart. It also draws more units than the inventory will turn out to
    verify — which is a question for a human, not a conclusion for a tool.
    """
    workers = "".join(
        f'<mxCell id="w{i}" value="NX-100 — Worker #{i}" style="rounded=0;" '
        f'vertex="1" parent="1"/>'
        for i in range(1, 13)
    )
    xml = f"""<mxfile host="test">
  <diagram name="Floor Plan" id="d1"><mxGraphModel><root>
    <mxCell id="0"/><mxCell id="1" parent="0"/>
    <mxCell id="f1" value="Room A" style="rounded=0;" vertex="1" parent="1"/>
  </root></mxGraphModel></diagram>
  <diagram name="Rack Elevation" id="d2"><mxGraphModel><root>
    <mxCell id="0"/><mxCell id="1" parent="0"/>
    {workers}
  </root></mxGraphModel></diagram>
</mxfile>"""
    path.write_text(xml, encoding="utf-8")
    return path


def copy_of(src: Path, dst: Path) -> Path:
    shutil.copy2(src, dst)
    return dst
