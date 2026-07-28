# CUI // SP-CTI
"""End-to-end / route-level V&V suite for the hardened /boundary ecosystem — task bdr-vv-1.

This suite asserts the full *user-visible* contract of the boundary canvas
surfaces that the bdr-* work hardened, exercising the real modules and the real
Flask blueprint routes (built via ``create_boundary_blueprint`` on a minimal
Flask app — the established pattern in tests/test_bdc_cato_readiness.py /
test_bdc_impact_panel.py / test_bdc_oscal_exporter.py / test_bdc_export_pptx.py).
Importing the full tools/dashboard/app.py probes localhost LLM servers and hangs
in the sandbox, so we never do that here.

Sections (map to the bdr-vv-1 acceptance items):
  (a) twin snapshot route      — status ok + persisted with real counts; DB-error -> 500
  (b) OSCAL export             — all 4 artifact types valid; unlinked project -> 400
  (c) cATO readiness           — green / amber / red / unknown bands; route wiring
  (d) boundary impact (ATO)    — context / assess GREEN+RED / alternatives / empty-state
  (e) supply-chain cross-link  — risk-score match / no-match / missing-table (graceful)
  (f) PPTX export              — valid pptx bytes (python-pptx round-trip)
  (g) template labels          — "Diagram Grade", "Boundary Impact (ATO)", "Not scored",
                                 no bare ambiguous "Assess" toolbar button

Placeholder normalization (bdr-sec-3): main-DB reads use PG-native ``%s``; the
storage layer translates for the SQLite init-fallback, so every seeded main DB
is driven through a translating StorageConnection.
"""
from __future__ import annotations

import io
import json
import sqlite3
from datetime import datetime, timedelta, timezone
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parents[1]


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def _sc(raw: sqlite3.Connection):
    """A translating StorageConnection over an in-memory sqlite conn that never
    closes it — production opens/closes get_connection() per call, and the shared
    ``:memory:`` DB must survive across those calls within one test."""
    from tools.db.storage import StorageConnection

    class _NoClose(StorageConnection):
        def close(self):  # keep shared in-memory conn alive
            pass

    return _NoClose(raw, "sqlite")


def _now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def _days_from_now(n: int) -> str:
    return (datetime.now(timezone.utc) + timedelta(days=n)).date().isoformat()


def _make_client(bp):
    """Wrap a boundary blueprint in a minimal authed Flask test client."""
    from flask import Flask

    app = Flask(__name__)
    app.secret_key = "test-secret"
    app.register_blueprint(bp, url_prefix="/boundary")
    client = app.test_client()
    with client.session_transaction() as sess:
        sess["user_id"] = "test-admin"
    return client


@pytest.fixture
def bdc_client(monkeypatch):
    monkeypatch.setenv("ICDEV_BOUNDARY_ENABLED", "true")
    from tools.boundary_canvas.blueprint import create_boundary_blueprint

    bp = create_boundary_blueprint()
    assert bp is not None
    return _make_client(bp)


# ===========================================================================
# (a) Digital-twin snapshot route — honest status + real counts / DB-error 500
# ===========================================================================

_TWIN_MAIN_DDL = """
CREATE TABLE IF NOT EXISTS project_controls (
    project_id TEXT, control_id TEXT, implementation_status TEXT
);
CREATE TABLE IF NOT EXISTS evidence (
    id INTEGER PRIMARY KEY AUTOINCREMENT, project_id TEXT, ref TEXT
);
CREATE TABLE IF NOT EXISTS compliance_snapshots (
    snapshot_id TEXT, project_id TEXT, framework_id TEXT, control_id TEXT,
    implementation_status TEXT, evidence_ref TEXT, taken_at TEXT
);
"""

_DESIGN_TWIN = "d-twin"


def _seed_twin_main(with_snapshots_table=True):
    raw = sqlite3.connect(":memory:")
    raw.row_factory = sqlite3.Row
    ddl = _TWIN_MAIN_DDL
    if not with_snapshots_table:
        # Drop the compliance_snapshots DDL statement -> INSERT will fail.
        ddl = _TWIN_MAIN_DDL.split("CREATE TABLE IF NOT EXISTS compliance_snapshots")[0]
    raw.executescript(ddl)
    for cid in ("AC-2", "AU-6"):
        raw.execute(
            "INSERT INTO project_controls (project_id, control_id, implementation_status) "
            "VALUES (?, ?, ?)",
            (_DESIGN_TWIN, cid, "satisfied"),
        )
    raw.execute("INSERT INTO evidence (project_id, ref) VALUES (?, ?)", (_DESIGN_TWIN, "ev-1"))
    raw.commit()
    return raw


def test_twin_snapshot_ok_persisted_real_counts(bdc_client, monkeypatch):
    raw = _seed_twin_main(with_snapshots_table=True)
    monkeypatch.setattr("tools.boundary_canvas.twin.get_connection", lambda *a, **k: _sc(raw))

    resp = bdc_client.post(f"/boundary/api/twin/{_DESIGN_TWIN}/snapshot", json={})
    assert resp.status_code == 201, resp.get_data(as_text=True)
    data = resp.get_json()
    assert data["status"] == "ok"
    assert data["persisted"] is True
    # Real, non-zeroed counts (not a fabricated success dict full of zeros).
    assert data["control_count"] == 2
    assert data["evidence_count"] == 1
    assert data["snapshot_id"]
    raw.close()


def test_twin_snapshot_db_error_returns_500_not_zeroed(bdc_client, monkeypatch):
    # No compliance_snapshots table -> INSERT fails -> honest error, route 500.
    raw = _seed_twin_main(with_snapshots_table=False)
    monkeypatch.setattr("tools.boundary_canvas.twin.get_connection", lambda *a, **k: _sc(raw))

    resp = bdc_client.post(f"/boundary/api/twin/{_DESIGN_TWIN}/snapshot", json={})
    assert resp.status_code == 500
    data = resp.get_json()
    assert data["status"] == "error"
    assert data["persisted"] is False
    assert data.get("error")
    raw.close()


# ===========================================================================
# (b) OSCAL cATO export — 4 valid artifact types + unlinked-project 400
# ===========================================================================

_OSCAL_DDL = """
CREATE TABLE IF NOT EXISTS projects (
    id TEXT PRIMARY KEY, name TEXT, description TEXT,
    impact_level TEXT DEFAULT 'IL5', classification TEXT DEFAULT 'CUI',
    type TEXT DEFAULT 'webapp', status TEXT DEFAULT 'under-development',
    ato_status TEXT DEFAULT 'none', cloud_environment TEXT DEFAULT 'aws-govcloud',
    created_by TEXT DEFAULT 'ICDEV', directory_path TEXT DEFAULT '',
    tech_stack_backend TEXT, tech_stack_frontend TEXT, tech_stack_database TEXT
);
CREATE TABLE IF NOT EXISTS project_controls (
    project_id TEXT, control_id TEXT, implementation_status TEXT,
    implementation_description TEXT, responsible_role TEXT, evidence_path TEXT, last_assessed TEXT
);
CREATE TABLE IF NOT EXISTS compliance_controls (id TEXT PRIMARY KEY, family TEXT, title TEXT, description TEXT);
CREATE TABLE IF NOT EXISTS poam_items (
    id INTEGER PRIMARY KEY AUTOINCREMENT, project_id TEXT, weakness_id TEXT,
    weakness_description TEXT, severity TEXT, source TEXT, control_id TEXT,
    status TEXT DEFAULT 'open', corrective_action TEXT, milestone_date TEXT,
    responsible_party TEXT, created_at TEXT
);
CREATE TABLE IF NOT EXISTS fedramp_assessments (
    project_id TEXT, control_id TEXT, baseline TEXT, status TEXT, implementation_status TEXT,
    evidence_description TEXT, evidence_path TEXT, notes TEXT, assessment_date TEXT, assessor TEXT
);
CREATE TABLE IF NOT EXISTS cmmc_assessments (
    project_id TEXT, practice_id TEXT, domain TEXT, level INTEGER, status TEXT,
    evidence_description TEXT, evidence_path TEXT, notes TEXT, nist_171_id TEXT,
    assessment_date TEXT, assessor TEXT
);
CREATE TABLE IF NOT EXISTS stig_findings (
    project_id TEXT, stig_id TEXT, finding_id TEXT, rule_id TEXT, severity TEXT, title TEXT,
    description TEXT, check_content TEXT, fix_text TEXT, status TEXT, comments TEXT,
    target_type TEXT, assessed_by TEXT, assessed_at TEXT
);
CREATE TABLE IF NOT EXISTS sbom_records (
    project_id TEXT, format TEXT, version TEXT, component_count INTEGER,
    vulnerability_count INTEGER, generated_at TEXT
);
"""

_OSCAL_PROJECT = "proj-vv-oscal"


@pytest.fixture
def oscal_db(tmp_path):
    from tools.db.storage import get_connection

    db_path = tmp_path / "icdev.db"
    conn = get_connection(db_path=str(db_path))
    conn.executescript(_OSCAL_DDL)
    conn.execute(
        "INSERT INTO projects (id, name, description, impact_level, status) VALUES (%s, %s, %s, %s, %s)",
        (_OSCAL_PROJECT, "VV OSCAL System", "Boundary VV.", "IL5", "operational"),
    )
    for cid, status in [("AC-2", "satisfied"), ("AU-6", "not_satisfied")]:
        conn.execute(
            "INSERT INTO project_controls (project_id, control_id, implementation_status, "
            "implementation_description, responsible_role) VALUES (%s, %s, %s, %s, %s)",
            (_OSCAL_PROJECT, cid, status, f"{cid} via baseline.", "ISSO"),
        )
    conn.commit()
    conn.close()
    return db_path


@pytest.mark.parametrize("artifact_type", ["ssp", "poam", "assessment_results", "component_definition"])
def test_oscal_all_four_types_valid(oscal_db, tmp_path, artifact_type):
    from tools.boundary_canvas import oscal_cato_exporter as ox

    result = ox.export_oscal_artifact(
        _OSCAL_PROJECT, artifact_type=artifact_type,
        output_dir=str(tmp_path / "out"), db_path=oscal_db,
    )
    assert result["status"] == "ok", result
    assert result["path"] and Path(result["path"]).exists()
    assert result["valid"] is True, result.get("errors")


def test_oscal_unlinked_project_route_400(bdc_client):
    # A design id with no matching projects row -> export_oscal returns an honest
    # error payload; the route surfaces it as 400 (never a phantom artifact/200).
    resp = bdc_client.get("/boundary/api/twin/no-such-design-vv/oscal-export?artifact_type=ssp")
    assert resp.status_code == 400
    data = resp.get_json()
    assert data["status"] == "error"
    assert data["path"] is None


# ===========================================================================
# (c) cATO readiness — green / amber / red / unknown bands + route wiring
# ===========================================================================

_READY_MAIN_DDL = """
CREATE TABLE IF NOT EXISTS compliance_twin_snapshots (
    id INTEGER PRIMARY KEY AUTOINCREMENT, snapshot_id TEXT, project_id TEXT, framework TEXT,
    control_id TEXT, implementation_status TEXT, evidence_ref TEXT, score REAL,
    assessor TEXT, notes TEXT, classification TEXT DEFAULT 'CUI // SP-CTI',
    created_at TEXT DEFAULT (datetime('now'))
);
CREATE TABLE IF NOT EXISTS compliance_twin_violations (
    id INTEGER PRIMARY KEY AUTOINCREMENT, snapshot_id TEXT, project_id TEXT, framework TEXT,
    control_id TEXT, violation_type TEXT, severity TEXT DEFAULT 'moderate', details TEXT,
    poam_id INTEGER, classification TEXT DEFAULT 'CUI // SP-CTI', created_at TEXT DEFAULT (datetime('now'))
);
CREATE TABLE IF NOT EXISTS poam_items (
    id INTEGER PRIMARY KEY AUTOINCREMENT, project_id TEXT, weakness_id TEXT, weakness_description TEXT,
    severity TEXT, source TEXT, control_id TEXT, status TEXT DEFAULT 'open', corrective_action TEXT,
    milestone_date TEXT, completion_date TEXT, responsible_party TEXT, resources_required TEXT,
    created_at TEXT DEFAULT (datetime('now'))
);
"""

_READY_CANVAS_DDL = """
CREATE TABLE IF NOT EXISTS bd_isa_tracker (
    id TEXT PRIMARY KEY, design_id TEXT, interconnection_id TEXT NOT NULL, isa_doc_id TEXT,
    status TEXT DEFAULT 'draft', expiry_date TEXT, isa_expiry_date TEXT, review_date TEXT,
    owner TEXT, notes TEXT, created_at TEXT DEFAULT CURRENT_TIMESTAMP, updated_at TEXT DEFAULT CURRENT_TIMESTAMP
);
"""


@pytest.fixture
def ready_main(tmp_path):
    from tools.db.storage import get_connection

    conn = get_connection(db_path=str(tmp_path / "ready_main.db"))
    conn.executescript(_READY_MAIN_DDL)
    conn.commit()
    yield conn
    conn.close()


@pytest.fixture
def ready_canvas():
    conn = sqlite3.connect(":memory:")
    conn.row_factory = sqlite3.Row
    conn.executescript(_READY_CANVAS_DDL)
    conn.commit()
    yield conn
    conn.close()


def _seed_ready_snapshot(conn, did, snap_id, controls, created_at=None):
    created_at = created_at or _now_iso()
    for cid, status, score in controls:
        conn.execute(
            "INSERT INTO compliance_twin_snapshots "
            "(snapshot_id, project_id, framework, control_id, implementation_status, score, created_at) "
            "VALUES (%s, %s, %s, %s, %s, %s, %s)",
            (snap_id, did, "FedRAMP Moderate", cid, status, score, created_at),
        )
    conn.commit()


def _seed_ready_violation(conn, did, control_id, poam_status):
    cur = conn.execute(
        "INSERT INTO poam_items (project_id, weakness_id, control_id, status, source) "
        "VALUES (%s, %s, %s, %s, %s)",
        (did, f"TWIN-{control_id}", control_id, poam_status, "test"),
    )
    poam_id = cur.lastrowid
    conn.execute(
        "INSERT INTO compliance_twin_violations "
        "(snapshot_id, project_id, framework, control_id, violation_type, poam_id) "
        "VALUES (%s, %s, %s, %s, %s, %s)",
        ("snap-x", did, "FedRAMP Moderate", control_id, "not_satisfied", poam_id),
    )
    conn.commit()


def _seed_ready_isa(conn, did, isa_id, expiry_days, status="active"):
    conn.execute(
        "INSERT INTO bd_isa_tracker (id, design_id, interconnection_id, status, expiry_date) "
        "VALUES (?, ?, ?, ?, ?)",
        (isa_id, did, f"ic-{isa_id}", status, _days_from_now(expiry_days)),
    )
    conn.commit()


def test_readiness_green_band(ready_main, ready_canvas):
    from tools.boundary_canvas import cato_readiness as cr

    did = "d-green"
    _seed_ready_snapshot(ready_main, did, "snap-g",
                         [("AC-2", "satisfied", 1.0), ("AU-6", "satisfied", 1.0)])
    _seed_ready_violation(ready_main, did, "AC-2", "closed")
    _seed_ready_isa(ready_canvas, did, "isa-1", 365)
    result = cr.compute_readiness(did, conn=ready_main, canvas_conn=ready_canvas)
    assert result["band"] == "green"
    assert result["score"] >= 80
    assert set(result["components"]) >= {"control_coverage", "poam", "isa_expiry", "freshness"}


def test_readiness_amber_band(ready_main, ready_canvas):
    from tools.boundary_canvas import cato_readiness as cr

    did = "d-amber"
    _seed_ready_snapshot(ready_main, did, "snap-a",
                         [("AC-2", "satisfied", 1.0), ("AU-6", "not_satisfied", 0.0)])
    for i in range(10):
        _seed_ready_violation(ready_main, did, f"CM-{i}", "open")
    _seed_ready_isa(ready_canvas, did, "isa-2", 30)
    result = cr.compute_readiness(did, conn=ready_main, canvas_conn=ready_canvas)
    assert result["band"] == "amber"
    assert 50 <= result["score"] < 80


def test_readiness_red_band(ready_main, ready_canvas):
    from tools.boundary_canvas import cato_readiness as cr

    did = "d-red"
    old = (datetime.now(timezone.utc) - timedelta(days=120)).isoformat()
    _seed_ready_snapshot(ready_main, did, "snap-r",
                         [("AC-2", "not_satisfied", 0.0), ("AU-6", "not_satisfied", 0.0)],
                         created_at=old)
    for i in range(20):
        _seed_ready_violation(ready_main, did, f"SC-{i}", "open")
    _seed_ready_isa(ready_canvas, did, "isa-3", -5)
    result = cr.compute_readiness(did, conn=ready_main, canvas_conn=ready_canvas)
    assert result["band"] == "red"
    assert result["score"] < 50


def test_readiness_unknown_when_empty(ready_main, ready_canvas):
    from tools.boundary_canvas import cato_readiness as cr

    result = cr.compute_readiness("d-empty", conn=ready_main, canvas_conn=ready_canvas)
    assert result["score"] is None
    assert result["band"] == "unknown"


def test_readiness_route_wiring(bdc_client, monkeypatch):
    canned = {
        "design_id": "d1", "project_id": "d1", "score": 72.5, "readiness_score": 72.5,
        "band": "amber", "components": {"control_coverage": {"score": 80.0}}, "weights": {},
    }
    monkeypatch.setattr("tools.boundary_canvas.cato_readiness.compute_readiness", lambda *a, **k: canned)
    resp = bdc_client.get("/boundary/api/cato/readiness/d1")
    assert resp.status_code == 200
    data = resp.get_json()
    assert data["score"] == 72.5
    assert data["band"] == "amber"
    assert "control_coverage" in data["components"]


# ===========================================================================
# (d) Boundary Impact (ATO) — context / assess GREEN+RED / alternatives / empty
# ===========================================================================

_IMPACT_DDL = """
-- Matches the `projects` shape used by the other fixture in this file and by
-- the real table. The two-column version worked only because these tests never
-- read past `name`; the first query that did would have failed here and not in
-- the suite above, for no reason a reader could see.
CREATE TABLE IF NOT EXISTS projects (
    id TEXT PRIMARY KEY, name TEXT, description TEXT,
    impact_level TEXT DEFAULT 'IL5', classification TEXT DEFAULT 'CUI',
    status TEXT DEFAULT 'under-development'
);
CREATE TABLE IF NOT EXISTS ato_system_registry (
    id TEXT PRIMARY KEY, project_id TEXT NOT NULL, system_name TEXT NOT NULL, system_acronym TEXT,
    ato_type TEXT, ato_date TEXT, ato_expiry TEXT, authorizing_official TEXT, accreditation_boundary TEXT,
    ssp_document_id INTEGER, impact_level TEXT, data_types TEXT, interconnections TEXT,
    baseline_controls TEXT, component_inventory TEXT, classification TEXT DEFAULT 'CUI',
    created_at TEXT DEFAULT (datetime('now')), updated_at TEXT DEFAULT (datetime('now'))
);
CREATE TABLE IF NOT EXISTS intake_requirements (
    id TEXT PRIMARY KEY, session_id TEXT, project_id TEXT, source_turn INTEGER, raw_text TEXT NOT NULL,
    refined_text TEXT, requirement_type TEXT DEFAULT 'functional', created_at TEXT DEFAULT (datetime('now'))
);
CREATE TABLE IF NOT EXISTS boundary_impact_assessments (
    id TEXT PRIMARY KEY, session_id TEXT, project_id TEXT NOT NULL, system_id TEXT NOT NULL,
    requirement_id TEXT, safe_item_id TEXT, impact_tier TEXT NOT NULL, impact_category TEXT NOT NULL,
    impact_description TEXT NOT NULL, affected_controls TEXT, affected_components TEXT,
    ssp_sections_impacted TEXT, remediation_required TEXT, alternative_approach TEXT,
    risk_score REAL DEFAULT 0.0, assessed_by TEXT DEFAULT 'icdev-requirements-analyst',
    assessed_at TEXT DEFAULT (datetime('now')), UNIQUE(requirement_id, system_id)
);
"""

_IMP_PROJECT = "proj-vv-imp"
_IMP_SYSTEM = "sys-vv-imp"
_IMP_REQ_GREEN = "req-vv-green"
_IMP_REQ_RED = "req-vv-red"


@pytest.fixture
def impact_db(tmp_path, monkeypatch):
    from tools.db.storage import get_connection
    from tools.requirements import boundary_analyzer

    db_file = tmp_path / "impact_main.db"
    conn = get_connection(db_path=str(db_file))
    conn.executescript(_IMPACT_DDL)
    conn.execute("INSERT INTO projects (id, name) VALUES (%s, %s)", (_IMP_PROJECT, "VV Impact"))
    conn.execute(
        "INSERT INTO ato_system_registry "
        "(id, project_id, system_name, ato_type, impact_level, classification, "
        " accreditation_boundary, data_types, interconnections, baseline_controls, component_inventory) "
        "VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)",
        (_IMP_SYSTEM, _IMP_PROJECT, "Mission System", "ato", "IL5", "CUI",
         "{}", "[]", "[]", '["AC-2", "AU-2"]', "[]"),
    )
    conn.execute(
        "INSERT INTO intake_requirements (id, session_id, project_id, raw_text, refined_text, requirement_type) "
        "VALUES (%s, %s, %s, %s, %s, %s)",
        (_IMP_REQ_GREEN, "sess-1", _IMP_PROJECT,
         "Apply a minor configuration change and patch to an existing internal component.", "", "functional"),
    )
    conn.execute(
        "INSERT INTO intake_requirements (id, session_id, project_id, raw_text, refined_text, requirement_type) "
        "VALUES (%s, %s, %s, %s, %s, %s)",
        (_IMP_REQ_RED, "sess-1", _IMP_PROJECT,
         "Process SECRET classification data and require boundary expansion to a new enclave.", "", "security"),
    )
    conn.commit()
    conn.close()
    monkeypatch.setattr(boundary_analyzer, "DB_PATH", db_file)
    return db_file


def test_impact_context_with_system(bdc_client, impact_db):
    resp = bdc_client.get("/boundary/api/impact/context?project_id=" + _IMP_PROJECT)
    assert resp.status_code == 200
    data = resp.get_json()
    assert data["empty"] is False
    assert len(data["systems"]) == 1
    req_ids = {r["id"] for r in data["requirements"]}
    assert {_IMP_REQ_GREEN, _IMP_REQ_RED} <= req_ids


def test_impact_context_empty_state(bdc_client, impact_db):
    resp = bdc_client.get("/boundary/api/impact/context?project_id=no-such")
    assert resp.status_code == 200
    data = resp.get_json()
    assert data["empty"] is True
    assert data["systems"] == []


def test_impact_assess_green_renders_tier_fields(bdc_client, impact_db):
    resp = bdc_client.post("/boundary/api/impact/assess",
                           json={"system_id": _IMP_SYSTEM, "requirement_id": _IMP_REQ_GREEN})
    assert resp.status_code == 200
    data = resp.get_json()
    assert data["impact_tier"] == "GREEN"
    assert isinstance(data["impact_score"], (int, float))
    assert "affected_controls" in data
    assert data["remediation_steps"]
    assert data["assessment_id"]


def test_impact_assess_red_yields_alternatives(bdc_client, impact_db):
    assess = bdc_client.post("/boundary/api/impact/assess",
                             json={"system_id": _IMP_SYSTEM, "requirement_id": _IMP_REQ_RED}).get_json()
    assert assess["impact_tier"] == "RED"
    resp = bdc_client.post("/boundary/api/impact/alternatives",
                           json={"assessment_id": assess["assessment_id"]})
    assert resp.status_code == 200
    data = resp.get_json()
    assert data["original_tier"] == "RED"
    assert len(data["alternatives"]) >= 1


def test_impact_alternatives_refused_for_green(bdc_client, impact_db):
    assess = bdc_client.post("/boundary/api/impact/assess",
                             json={"system_id": _IMP_SYSTEM, "requirement_id": _IMP_REQ_GREEN}).get_json()
    assert assess["impact_tier"] == "GREEN"
    resp = bdc_client.post("/boundary/api/impact/alternatives",
                           json={"assessment_id": assess["assessment_id"]})
    assert resp.status_code == 400


# ===========================================================================
# (e) Supply-chain cross-link matcher — risk-score match / no-match / missing-table
# ===========================================================================

class _FakeCanvasConn:
    """Minimal context-manager canvas connection returning canned rows by SQL.

    Mirrors tests/test_bdc_export_pptx.py::_FakeConn. init_db()'s DDL runs
    harmlessly through it (no-op executescript / execute), giving the test full
    control over whether ``bd_isa_tracker`` reads succeed or raise (missing-table).
    """

    def __init__(self, design_row, isas, isa_missing=False):
        self._design_row = design_row
        self._isas = isas
        self._isa_missing = isa_missing
        self._sql = ""

    def __enter__(self):
        return self

    def __exit__(self, *a):
        return False

    def execute(self, sql, params=()):
        self._sql = sql
        if (self._isa_missing and "bd_isa_tracker" in sql
                and sql.strip().upper().startswith("SELECT")):
            raise sqlite3.OperationalError("no such table: bd_isa_tracker")
        return self

    def executescript(self, sql):  # init_db schema load -> no-op
        return self

    def fetchone(self):
        if "boundary_designs" in self._sql:
            return self._design_row
        return None

    def fetchall(self):
        if "bd_isa_tracker" in self._sql:
            return self._isas
        return []

    def commit(self):
        pass

    def close(self):
        pass

    def cursor(self):
        return self


def _risk_client(monkeypatch, conn):
    monkeypatch.setenv("ICDEV_BOUNDARY_ENABLED", "true")
    # Patch BEFORE building the blueprint: the route captures get_connection as a
    # closure-local at build time (test_bdc_export_pptx.py pattern).
    monkeypatch.setattr("tools.boundary_canvas.db.init_db.get_connection", lambda *a, **k: conn)
    from tools.boundary_canvas.blueprint import create_boundary_blueprint

    bp = create_boundary_blueprint()
    assert bp is not None
    return _make_client(bp)


def _design_row(edges):
    return {"graph_json": json.dumps({"nodes": [], "edges": edges})}


def test_risk_score_supply_chain_match(monkeypatch):
    # A CRITICAL/HIGH interconnection -> supply_chain_risk cross-link populated.
    edges = [{"id": "e-cd", "data": {"interconnection_type": "cross-domain"}, "label": "Cross Domain Link"}]
    client = _risk_client(monkeypatch, _FakeCanvasConn(_design_row(edges), isas=[]))
    resp = client.get("/boundary/api/designs/d-match/risk-score")
    assert resp.status_code == 200
    data = resp.get_json()
    assert len(data["supply_chain_risk"]) >= 1
    assert data["scored_edges"][0]["risk_tier"] == "CRITICAL"


def test_risk_score_supply_chain_no_match(monkeypatch):
    # Only a LOW-risk (federation) interconnection -> no supply-chain cross-link.
    edges = [{"id": "e-fed", "data": {"interconnection_type": "federation"}, "label": "SSO"}]
    client = _risk_client(monkeypatch, _FakeCanvasConn(_design_row(edges), isas=[]))
    resp = client.get("/boundary/api/designs/d-nomatch/risk-score")
    assert resp.status_code == 200
    data = resp.get_json()
    assert data["supply_chain_risk"] == []
    assert data["scored_edges"][0]["risk_tier"] == "LOW"


def test_risk_score_missing_isa_table_degrades(monkeypatch):
    # bd_isa_tracker absent -> graceful 200 (no ISA overlay), not a 500.
    edges = [{"id": "e-cd", "data": {"interconnection_type": "cross-domain"}, "label": "Link"}]
    client = _risk_client(monkeypatch, _FakeCanvasConn(_design_row(edges), isas=[], isa_missing=True))
    resp = client.get("/boundary/api/designs/d-missing/risk-score")
    assert resp.status_code == 200
    data = resp.get_json()
    # Edge-derived supply-chain risk still computes without the ISA table.
    assert len(data["supply_chain_risk"]) >= 1


def test_risk_score_design_not_found(monkeypatch):
    client = _risk_client(monkeypatch, _FakeCanvasConn(None, isas=[]))
    resp = client.get("/boundary/api/designs/nope/risk-score")
    assert resp.status_code == 404


# ===========================================================================
# (f) PPTX export — valid pptx bytes (python-pptx round-trip)
# ===========================================================================

def test_pptx_export_round_trips():
    from pptx import Presentation

    from tools.boundary_canvas import export_pptx

    design = {
        "id": "d-pptx", "name": "VV Deck", "classification": "CUI",
        "graph_json": json.dumps({
            "nodes": [{"id": "n1", "type": "server", "label": "Web"}],
            "edges": [],
        }),
    }
    readiness = {"score": 74.0, "band": "amber", "components": {"control_coverage": {"score": 80.0}}}
    data = export_pptx.design_to_pptx(design, assessment={"grade": "B", "score": 82.5}, readiness=readiness)
    assert isinstance(data, (bytes, bytearray))
    assert data[:2] == b"PK"  # .pptx is a zip
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 3
    titles = [s.shapes.title.text for s in prs.slides if s.shapes.title is not None]
    assert "VV Deck" in titles


# ===========================================================================
# (g) Template-level label assertions
# ===========================================================================

_CANVAS_HTML = ROOT / "tools" / "dashboard" / "templates" / "boundary_canvas" / "canvas.html"
_TWIN_HTML = ROOT / "tools" / "dashboard" / "templates" / "boundary_canvas" / "twin.html"


def test_canvas_template_has_grade_and_impact_labels():
    html = _CANVAS_HTML.read_text(encoding="utf-8")
    assert "Diagram Grade" in html
    assert "Boundary Impact (ATO)" in html


def test_canvas_toolbar_has_no_bare_assess_button():
    html = _CANVAS_HTML.read_text(encoding="utf-8")
    # The old ambiguous top-level "Assess" toolbar button was renamed to
    # "Diagram Grade" — the button wired to canvasAssess() must carry the clear
    # label, not a bare "Assess". (The only remaining "Assess" is the scoped
    # button inside the titled Boundary Impact (ATO) modal.)
    grade_lines = [ln for ln in html.splitlines() if "canvasAssess()" in ln]
    assert grade_lines, "canvasAssess() toolbar button not found"
    assert any("Diagram Grade" in ln for ln in grade_lines)
    assert not any(">Assess<" in ln for ln in grade_lines)


def test_twin_template_handles_not_scored():
    html = _TWIN_HTML.read_text(encoding="utf-8")
    # Honest twin: a None score renders "Not scored", never a fabricated 0%.
    assert "Not scored" in html


# ===========================================================================
# Playwright screenshots (best-effort) — skipped when unavailable/sandboxed
# ===========================================================================

def test_playwright_screenshots_best_effort():
    """Attempt a live-server Playwright screenshot pass; skip gracefully if the
    playwright package/browsers are absent or the sandbox blocks localhost sockets.
    The route-level suite above is the required deliverable."""
    pytest.importorskip("playwright", reason="playwright python package not installed")
    # If the package ever lands, a live-server capture would go here; for now the
    # importorskip records an honest SKIP rather than a false pass.
    pytest.skip("playwright present but live-server capture not enabled in sandbox")
