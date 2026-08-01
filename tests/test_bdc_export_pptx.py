# CUI // SP-CTI
"""Tests for tools/boundary_canvas/export_pptx.py — task bdr-feat-5.

PPTX export for boundary designs. Coverage:
  * design_to_pptx returns valid .pptx bytes that python-pptx re-opens
  * expected slide count (title + diagram + compliance) and title text
  * compliance table populated from assessment + readiness
  * missing-readiness / missing-assessment case still yields a valid 3-slide deck
  * empty graph (no nodes) still produces a valid deck (placeholder diagram)
  * Flask route returns the bytes with pptx content-type + attachment filename
"""
from __future__ import annotations

import io
import json

from tools.boundary_canvas import export_pptx


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

def _design_fixture(name="Test Boundary", with_graph=True):
    graph = {
        "nodes": [
            {"id": "n1", "type": "server", "label": "Web Server"},
            {"id": "n2", "type": "database", "label": "Postgres DB"},
            {"id": "n3", "type": "firewall", "label": "Edge FW"},
        ],
        "edges": [
            {"source": "n3", "target": "n1", "type": "allows"},
            {"source": "n1", "target": "n2", "type": "queries"},
        ],
    } if with_graph else {"nodes": [], "edges": []}
    return {
        "id": "design-1",
        "name": name,
        "description": "A boundary for testing.",
        "classification": "CUI",
        "graph_json": json.dumps(graph),
    }


_ASSESSMENT = {
    "assessment_type": "full",
    "grade": "B",
    "score": 82.5,
    "cat1_findings": 0,
    "cat2_findings": 2,
    "cat3_findings": 5,
    "created_at": "2026-07-17 12:00:00",
}

_READINESS = {
    "design_id": "design-1",
    "project_id": "design-1",
    "score": 74.0,
    "readiness_score": 74.0,
    "band": "amber",
    "components": {
        "control_coverage": {"score": 80.0},
        "poam": {"score": 100.0},
        "isa_expiry": {"score": 50.0},
        "freshness": {"score": None},
    },
}


def _open(pptx_bytes):
    from pptx import Presentation

    assert isinstance(pptx_bytes, (bytes, bytearray))
    assert pptx_bytes[:2] == b"PK"  # .pptx is a zip
    return Presentation(io.BytesIO(pptx_bytes))


# ---------------------------------------------------------------------------
# Core export tests
# ---------------------------------------------------------------------------

def test_full_deck_round_trips():
    data = export_pptx.design_to_pptx(
        _design_fixture(), assessment=_ASSESSMENT, readiness=_READINESS
    )
    prs = _open(data)
    assert len(prs.slides) == 3
    # Title slide carries the design name.
    title_texts = [
        s.shapes.title.text for s in prs.slides if s.shapes.title is not None
    ]
    assert "Test Boundary" in title_texts
    assert "Boundary Diagram" in title_texts
    assert "Compliance & cATO Readiness" in title_texts


def test_compliance_table_has_assessment_and_readiness():
    data = export_pptx.design_to_pptx(
        _design_fixture(), assessment=_ASSESSMENT, readiness=_READINESS
    )
    prs = _open(data)
    compliance = prs.slides[2]
    tables = [sh for sh in compliance.shapes if sh.has_table]
    assert len(tables) == 1
    # Flatten all cell text.
    cells = [
        c.text
        for row in tables[0].table.rows
        for c in row.cells
    ]
    blob = " ".join(cells)
    assert "Grade" in blob and "B" in blob            # from assessment
    assert "cATO Readiness Score" in blob             # from readiness
    assert "AMBER" in blob


def test_missing_readiness_still_valid_deck():
    data = export_pptx.design_to_pptx(_design_fixture(), assessment=_ASSESSMENT, readiness=None)
    prs = _open(data)
    assert len(prs.slides) == 3
    # Assessment alone still produces a table.
    tables = [sh for sh in prs.slides[2].shapes if sh.has_table]
    assert len(tables) == 1


def test_missing_all_compliance_data_uses_placeholder():
    data = export_pptx.design_to_pptx(_design_fixture(), assessment=None, readiness=None)
    prs = _open(data)
    assert len(prs.slides) == 3
    compliance = prs.slides[2]
    assert not any(sh.has_table for sh in compliance.shapes)
    text = " ".join(
        sh.text_frame.text for sh in compliance.shapes if sh.has_text_frame
    )
    assert "No compliance assessment" in text


def test_readiness_with_none_score_treated_as_absent():
    unknown = {"score": None, "band": "unknown", "components": {}}
    data = export_pptx.design_to_pptx(_design_fixture(), assessment=None, readiness=unknown)
    prs = _open(data)
    # No usable rows -> placeholder, no table.
    assert not any(sh.has_table for sh in prs.slides[2].shapes)


def test_empty_graph_still_valid_deck():
    data = export_pptx.design_to_pptx(_design_fixture(with_graph=False))
    prs = _open(data)
    assert len(prs.slides) == 3
    diagram = prs.slides[1]
    text = " ".join(
        sh.text_frame.text for sh in diagram.shapes if sh.has_text_frame
    )
    assert "no nodes" in text.lower()


def test_graph_json_as_dict_accepted():
    design = _design_fixture()
    design["graph_json"] = json.loads(design["graph_json"])  # pass a dict, not a string
    data = export_pptx.design_to_pptx(design, assessment=_ASSESSMENT, readiness=_READINESS)
    prs = _open(data)
    assert len(prs.slides) == 3


# ---------------------------------------------------------------------------
# Route test (Flask test client) — mirrors tests/test_bdc_cato_readiness.py
#
# The route imports get_connection *inside* create_boundary_blueprint (it becomes
# a closure-local), so we patch init_db.get_connection BEFORE building the
# blueprint. This keeps the route test backend-agnostic (no real PG/SQLite writes)
# while still exercising the real query path, _row_to_dict, and Response wiring.
# ---------------------------------------------------------------------------

class _FakeConn:
    """A minimal context-manager connection returning canned rows by SQL."""

    def __init__(self, design_row, assessment_row=None):
        self._design_row = design_row
        self._assessment_row = assessment_row
        self._sql = ""

    def __enter__(self):
        return self

    def __exit__(self, *a):
        return False

    def execute(self, sql, params=()):
        self._sql = sql
        return self

    def fetchone(self):
        if "bd_assessments" in self._sql:
            return self._assessment_row
        if "boundary_designs" in self._sql:
            return self._design_row
        return None

    def close(self):
        pass


def _build_client(monkeypatch, conn):
    monkeypatch.setenv("ICDEV_BOUNDARY_ENABLED", "true")
    monkeypatch.setattr(
        "tools.boundary_canvas.db.init_db.get_connection", lambda *a, **k: conn
    )
    from flask import Flask
    from tools.boundary_canvas.blueprint import create_boundary_blueprint

    bp = create_boundary_blueprint()
    assert bp is not None
    app = Flask(__name__)
    app.secret_key = "test-secret"
    app.register_blueprint(bp, url_prefix="/boundary")
    client = app.test_client()
    with client.session_transaction() as sess:
        sess["user_id"] = "test-admin"
    return client


def test_route_returns_pptx_bytes(monkeypatch):
    design_row = {
        "id": "design-1",
        "name": "Route Design",
        "graph_json": json.dumps({"nodes": [{"id": "a", "label": "A"}], "edges": []}),
        "classification": "CUI",
    }
    client = _build_client(monkeypatch, _FakeConn(design_row, _ASSESSMENT))
    monkeypatch.setattr(
        "tools.boundary_canvas.cato_readiness.compute_readiness",
        lambda *a, **k: _READINESS,
    )

    resp = client.get("/boundary/api/export/design-1/pptx")
    assert resp.status_code == 200
    assert (
        resp.mimetype
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    cd = resp.headers.get("Content-Disposition", "")
    assert "attachment" in cd
    assert "Route_Design.pptx" in cd
    prs = _open(resp.data)
    assert len(prs.slides) == 3


def test_route_404_when_design_missing(monkeypatch):
    client = _build_client(monkeypatch, _FakeConn(None))
    resp = client.get("/boundary/api/export/nope/pptx")
    assert resp.status_code == 404


def test_route_requires_auth(monkeypatch):
    monkeypatch.setenv("ICDEV_BOUNDARY_ENABLED", "true")
    # No session user_id -> auth gate rejects before any DB access.
    from flask import Flask
    from tools.boundary_canvas.blueprint import create_boundary_blueprint

    app = Flask(__name__)
    app.secret_key = "s"
    app.register_blueprint(create_boundary_blueprint(), url_prefix="/boundary")
    resp = app.test_client().get("/boundary/api/export/d1/pptx")
    assert resp.status_code == 401
