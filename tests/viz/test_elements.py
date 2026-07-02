# CUI // SP-CTI
"""VIZ Epic G1 — positioned element model + WYSIWYG PPTX element rendering."""
from __future__ import annotations

from pptx import Presentation

from tools.viz.elements import Element, auto_layout, elements_to_dicts, elements_from_dicts
from tools.viz.spec import ChartSpec, Series
from tools.slides import pptx_builder

_PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
    "890000000a49444154789c6360000002000154a24f5f0000000049454e44ae426082"
)


def test_element_roundtrip():
    e = Element("text", 0.1, 0.2, 0.5, 0.3, z=2, payload={"text": "hi"},
                style={"fontSize": 24, "color": "#C8A951"}, id="t0")
    d = e.to_dict()
    assert d["type"] == "text" and d["x"] == 0.1 and d["z"] == 2
    back = Element.from_dict(d)
    assert back.to_dict() == d
    assert elements_from_dicts(elements_to_dicts([e]))[0].payload["text"] == "hi"


def test_from_dict_coerces_bad_values():
    e = Element.from_dict({"type": "image", "x": "nope", "w": None})
    assert e.x == 0.05 and e.w == 0.4   # fall back to defaults


def test_auto_layout_title_slide():
    els = auto_layout({"slide_type": "title", "title": "Hello"})
    assert len(els) == 2
    assert els[0].type == "text" and els[0].payload["text"] == "Hello"
    assert els[0].style["align"] == "center"


def test_auto_layout_chart_slide():
    chart = ChartSpec(title="C", chart_type="bar", categories=["a", "b"],
                      series=[Series("s", [1, 2])]).to_dict()
    els = auto_layout({"slide_type": "data", "title": "Metrics", "chart": chart})
    types = [e.type for e in els]
    assert "text" in types and "chart" in types
    chart_el = next(e for e in els if e.type == "chart")
    assert chart_el.payload["chart_type"] == "bar"
    # body element is below the title bar
    assert chart_el.y > 0.15


def test_auto_layout_bullets():
    els = auto_layout({"slide_type": "content", "title": "T", "bullets": ["one", "two"]})
    body = [e for e in els if e.type == "text" and "•" in e.payload.get("text", "")]
    assert body and "one" in body[0].payload["text"]


def test_pptx_element_slide_wysiwyg(tmp_path):
    img = tmp_path / "pic.png"
    img.write_bytes(_PNG)
    chart = ChartSpec(title="C", chart_type="column", categories=["x", "y"],
                      series=[Series("s", [3, 5])]).to_dict()
    slides = [{
        "slide_type": "content",
        "elements": [
            Element("text", 0.05, 0.04, 0.9, 0.12, z=2, payload={"text": "Freeform"},
                    style={"fontSize": 32, "bold": True, "color": "#C8A951"}).to_dict(),
            Element("chart", 0.05, 0.2, 0.5, 0.7, z=0, payload=chart).to_dict(),
            Element("image", 0.6, 0.2, 0.35, 0.5, z=1, payload={"src": str(img)}).to_dict(),
        ],
    }]
    path = pptx_builder.build(slides, title="Freeform Deck")
    prs = Presentation(path)
    assert len(prs.slides) == 1
    shapes = list(prs.slides[0].shapes)
    assert any(sh.has_chart for sh in shapes), "chart element rendered"
    # picture shape present (shape_type PICTURE == 13)
    assert any(getattr(sh, "shape_type", None) == 13 for sh in shapes), "image element rendered"
    # chart positioned at ~0.05*13.33in left (within tolerance)
    chart_shape = next(sh for sh in shapes if sh.has_chart)
    assert abs(chart_shape.left - int(0.05 * 13.33 * 914400)) < 914400  # within 1 inch


def test_pptx_shape_element(tmp_path):
    slides = [{"slide_type": "content", "elements": [
        Element("shape", 0.2, 0.2, 0.3, 0.3, payload={"shape": "ellipse"},
                style={"fill": "#C8A951", "stroke": "#FFFFFF", "strokeWidth": 2, "opacity": 1}).to_dict(),
        Element("shape", 0.2, 0.6, 0.5, 0.1, payload={"shape": "rectangle"},
                style={"fill": "#4A90D9", "cornerRadius": 8}).to_dict(),
    ]}]
    prs = Presentation(pptx_builder.build(slides, title="Shapes"))
    # auto-shapes present (shape_type AUTO_SHAPE == 1)
    autoshapes = [sh for sh in prs.slides[0].shapes if getattr(sh, "shape_type", None) == 1]
    assert len(autoshapes) >= 2


def test_pptx_text_element_honors_style(tmp_path):
    slides = [{"slide_type": "content", "elements": [
        Element("text", 0.1, 0.1, 0.8, 0.2, payload={"text": "Big Gold"},
                style={"fontSize": 40, "color": "#C8A951", "bold": True, "align": "center"}).to_dict()
    ]}]
    prs = Presentation(pptx_builder.build(slides, title="T"))
    tf_runs = [r for sh in prs.slides[0].shapes if sh.has_text_frame
               for p in sh.text_frame.paragraphs for r in p.runs]
    assert any(r.text == "Big Gold" and r.font.size and r.font.size.pt == 40 for r in tf_runs)
