# CUI // SP-CTI
"""Unit tests for the ICDEV™ Viz Kernel (tools/viz).

Each renderer must produce a non-empty, valid artifact for every spec kind,
fully offline (no network). PPTX must reopen; PNG must be a valid PNG; SVG/XML
must parse; Excalidraw must be valid JSON.
"""
from __future__ import annotations

import json
import xml.etree.ElementTree as ET
from pathlib import Path

import pytest

from tools.viz.spec import (
    ChartSpec, Series, TableSpec, DiagramSpec, KpiSpec, KpiTile,
    TimelineSpec, Milestone, spec_from_dict,
)

THEMES = ["midnight_executive", "govcon_proposal", "compliance_briefing"]


# ── fixtures ─────────────────────────────────────────────────────────────────

def _bar() -> ChartSpec:
    return ChartSpec(title="Compliance by Framework", chart_type="column",
                     categories=["FedRAMP", "CMMC", "NIST"],
                     series=[Series("Score", [88, 72, 95])], unit="%")


def _multi_line() -> ChartSpec:
    return ChartSpec(title="Burndown", chart_type="line",
                     categories=["W1", "W2", "W3", "W4"],
                     series=[Series("Planned", [40, 30, 20, 10]),
                             Series("Actual", [40, 34, 25, 14])])


def _pie() -> ChartSpec:
    return ChartSpec(title="Tasks", chart_type="donut",
                     categories=["Done", "Active", "Backlog"],
                     series=[Series("count", [12, 5, 8])])


def _gauge() -> ChartSpec:
    return ChartSpec(title="Readiness", chart_type="gauge",
                     series=[Series("v", [87])], unit="%", max_value=100)


def _table() -> TableSpec:
    return TableSpec(title="Canvas Status",
                     headers=["Canvas", "Status", "Coverage"],
                     rows=[["DIC", "Live", "94%"], ["ACE", "Partial", "61%"]])


def _diagram() -> DiagramSpec:
    return DiagramSpec(title="Pipeline",
                       nodes=[{"id": "a", "label": "Ingest"},
                              {"id": "b", "label": "Analyze"},
                              {"id": "c", "label": "Report"}],
                       edges=[{"source": "a", "target": "b", "label": "raw"},
                              {"source": "b", "target": "c", "label": "scored"}])


def _kpis() -> KpiSpec:
    return KpiSpec(title="Platform", tiles=[
        KpiTile("Canvases", "56", "+3"), KpiTile("Tests", "330", unit=" passing")])


# ── spec round-trip ──────────────────────────────────────────────────────────

@pytest.mark.parametrize("spec_fn", [_bar, _multi_line, _pie, _gauge, _table, _diagram, _kpis])
def test_spec_roundtrip(spec_fn):
    spec = spec_fn()
    d = spec.to_dict()
    rebuilt = spec_from_dict(d)
    assert rebuilt.to_dict() == d
    assert "kind" in d


def test_timeline_roundtrip():
    spec = TimelineSpec(title="Roadmap", milestones=[
        Milestone("Kernel", "Q1", "done"), Milestone("Web deck", "Q2", "in_progress")])
    assert spec_from_dict(spec.to_dict()).to_dict() == spec.to_dict()


def test_spec_from_dict_rejects_unknown():
    with pytest.raises(ValueError):
        spec_from_dict({"kind": "bogus"})


# ── PNG renderer ─────────────────────────────────────────────────────────────

@pytest.mark.parametrize("spec_fn", [_bar, _multi_line, _pie, _gauge])
def test_chart_to_png(spec_fn, tmp_path):
    from tools.viz.render_png import chart_to_png
    out = tmp_path / "c.png"
    path = chart_to_png(spec_fn(), out_path=str(out))
    data = Path(path).read_bytes()
    assert data[:8] == b"\x89PNG\r\n\x1a\n"  # PNG magic
    assert len(data) > 500


def test_diagram_to_png(tmp_path):
    from tools.viz.render_png import diagram_to_png
    out = tmp_path / "d.png"
    path = diagram_to_png(_diagram(), out_path=str(out))
    assert Path(path).read_bytes()[:8] == b"\x89PNG\r\n\x1a\n"


# ── SVG renderer ─────────────────────────────────────────────────────────────

@pytest.mark.parametrize("spec_fn", [_bar, _multi_line, _pie, _gauge])
def test_chart_to_svg_parses(spec_fn):
    from tools.viz.render_svg import chart_to_svg
    svg = chart_to_svg(spec_fn())
    assert svg.startswith("<svg") or svg.lstrip().startswith("<svg")
    ET.fromstring(svg)  # raises if malformed


def test_diagram_to_svg_parses():
    from tools.viz.render_svg import diagram_to_svg
    ET.fromstring(diagram_to_svg(_diagram()))


# ── HTML renderer ────────────────────────────────────────────────────────────

def test_html_fragments():
    from tools.viz.render_html import chart_to_html, table_to_html, kpis_to_html, diagram_to_html
    assert "<svg" in chart_to_html(_bar())
    assert "<table" in table_to_html(_table())
    assert "viz-kpis" in kpis_to_html(_kpis())
    # default: deterministic inline SVG (renders even in hidden slides)
    dh = diagram_to_html(_diagram())
    assert "<svg" in dh and "viz-diagram" in dh
    # opt-in: interactive Mermaid block with SVG fallback
    dh_mm = diagram_to_html(_diagram(), use_mermaid=True)
    assert "mermaid" in dh_mm and "flowchart" in dh_mm


# ── PPTX renderer ────────────────────────────────────────────────────────────

def _new_slide():
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


@pytest.mark.parametrize("spec_fn", [_bar, _multi_line, _pie])
def test_add_chart_reopens(spec_fn, tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    from tools.viz.render_pptx import add_chart
    prs, slide = _new_slide()
    add_chart(slide, spec_fn(), Inches(1), Inches(1), Inches(8), Inches(5))
    out = tmp_path / "chart.pptx"
    prs.save(str(out))
    reopened = Presentation(str(out))
    assert len(reopened.slides) == 1


def test_add_table_reopens(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    from tools.viz.render_pptx import add_table
    prs, slide = _new_slide()
    add_table(slide, _table(), Inches(1), Inches(1), Inches(8), Inches(3))
    out = tmp_path / "table.pptx"
    prs.save(str(out))
    assert len(Presentation(str(out)).slides) == 1


# ── draw.io + Excalidraw exporters ───────────────────────────────────────────

def test_to_drawio_parses():
    from tools.viz.render_diagram_export import to_drawio
    xml = to_drawio(_diagram())
    root = ET.fromstring(xml)
    assert root.tag == "mxGraphModel"


def test_to_excalidraw_valid_scene():
    from tools.viz.render_diagram_export import to_excalidraw
    scene = json.loads(to_excalidraw(_diagram()))
    assert scene["type"] == "excalidraw"
    assert scene["version"] == 2
    kinds = {el["type"] for el in scene["elements"]}
    assert "rectangle" in kinds and "text" in kinds and "arrow" in kinds
    # every arrow binds to real rectangles
    rect_ids = {el["id"] for el in scene["elements"] if el["type"] == "rectangle"}
    for el in scene["elements"]:
        if el["type"] == "arrow":
            assert el["startBinding"]["elementId"] in rect_ids
            assert el["endBinding"]["elementId"] in rect_ids


def test_excalidraw_deterministic():
    """No RNG — same spec must serialize identically (resume-safe)."""
    from tools.viz.render_diagram_export import to_excalidraw
    assert to_excalidraw(_diagram()) == to_excalidraw(_diagram())


# ── theming ──────────────────────────────────────────────────────────────────

@pytest.mark.parametrize("theme", THEMES)
def test_all_themes_render(theme, tmp_path):
    from tools.viz.render_png import chart_to_png
    from tools.viz.render_svg import chart_to_svg
    chart_to_png(_bar(), theme=theme, out_path=str(tmp_path / f"{theme}.png"))
    ET.fromstring(chart_to_svg(_bar(), theme=theme))
