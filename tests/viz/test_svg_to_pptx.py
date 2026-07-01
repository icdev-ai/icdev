# CUI // SP-CTI
"""Unit tests for the SVG -> native PPTX shape converter."""
import sys
from pathlib import Path
from unittest.mock import patch

import pytest

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from tools.viz import svg_to_pptx


def _new_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _svg(inner: str, viewbox: str = "0 0 200 100") -> str:
    return f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="{viewbox}">{inner}</svg>'


class TestPrimitives:
    def test_rect_becomes_freeform(self):
        _, slide = _new_slide()
        shapes = svg_to_pptx.render_svg_into_slide(
            slide, _svg('<rect x="0" y="0" width="50" height="30" fill="#ff0000"/>'),
            Inches(1), Inches(1), Inches(8), Inches(4),
        )
        assert len(shapes) == 1
        assert shapes[0].shape_type == MSO_SHAPE_TYPE.FREEFORM
        assert str(shapes[0].fill.fore_color.rgb) == "FF0000"

    def test_circle_and_ellipse_become_freeform(self):
        _, slide = _new_slide()
        shapes = svg_to_pptx.render_svg_into_slide(
            slide, _svg('<circle cx="10" cy="10" r="5" fill="blue"/>'
                        '<ellipse cx="50" cy="50" rx="8" ry="4" fill="green"/>'),
            Inches(1), Inches(1), Inches(8), Inches(4),
        )
        assert len(shapes) == 2
        assert all(s.shape_type == MSO_SHAPE_TYPE.FREEFORM for s in shapes)

    def test_line_is_open_freeform_with_stroke(self):
        _, slide = _new_slide()
        shapes = svg_to_pptx.render_svg_into_slide(
            slide, _svg('<line x1="0" y1="0" x2="10" y2="10" stroke="#000000" stroke-width="2"/>'),
            Inches(1), Inches(1), Inches(8), Inches(4),
        )
        assert len(shapes) == 1
        assert str(shapes[0].line.color.rgb) == "000000"

    def test_polygon_closed_polyline_open(self):
        _, slide = _new_slide()
        shapes = svg_to_pptx.render_svg_into_slide(
            slide, _svg('<polygon points="0,0 10,0 5,10" fill="orange"/>'
                        '<polyline points="0,0 10,0 5,10" stroke="black" fill="none"/>'),
            Inches(1), Inches(1), Inches(8), Inches(4),
        )
        assert len(shapes) == 2

    def test_path_with_cubic_curve_flattens_to_freeform(self):
        _, slide = _new_slide()
        shapes = svg_to_pptx.render_svg_into_slide(
            slide, _svg('<path d="M10,10 C20,0 40,0 50,10 Z" fill="purple"/>'),
            Inches(1), Inches(1), Inches(8), Inches(4),
        )
        assert len(shapes) == 1
        assert shapes[0].shape_type == MSO_SHAPE_TYPE.FREEFORM

    def test_path_with_multiple_subpaths_emits_multiple_shapes(self):
        _, slide = _new_slide()
        shapes = svg_to_pptx.render_svg_into_slide(
            slide, _svg('<path d="M0,0 L10,0 L5,10 Z M20,20 L30,20 L25,30 Z" fill="red"/>'),
            Inches(1), Inches(1), Inches(8), Inches(4),
        )
        assert len(shapes) == 2

    def test_text_becomes_textbox(self):
        _, slide = _new_slide()
        shapes = svg_to_pptx.render_svg_into_slide(
            slide, _svg('<text x="10" y="20" font-size="12" fill="#333333">Hello</text>'),
            Inches(1), Inches(1), Inches(8), Inches(4),
        )
        assert len(shapes) == 1
        assert shapes[0].shape_type == MSO_SHAPE_TYPE.TEXT_BOX
        assert shapes[0].text_frame.text == "Hello"

    def test_zero_size_rect_is_skipped(self):
        _, slide = _new_slide()
        shapes = svg_to_pptx.render_svg_into_slide(
            slide, _svg('<rect x="0" y="0" width="0" height="0" fill="red"/>'),
            Inches(1), Inches(1), Inches(8), Inches(4),
        )
        assert shapes == []


class TestTransformsAndGroups:
    def test_nested_group_translate_and_scale_compose(self):
        _, slide = _new_slide()
        svg = _svg(
            '<g transform="translate(10,10)">'
            '  <g transform="scale(2)">'
            '    <rect x="0" y="0" width="5" height="5" fill="red"/>'
            '  </g>'
            '</g>'
        )
        shapes = svg_to_pptx.render_svg_into_slide(slide, svg, Inches(0), Inches(0), Inches(2), Inches(1))
        assert len(shapes) == 1
        # translate(10,10) then scale(2) applied to (0,0)-(5,5) rect corners
        # should yield a bounding box starting at local (10,10), not (0,0).
        assert shapes[0].left > 0

    def test_defs_children_are_not_rendered(self):
        _, slide = _new_slide()
        svg = _svg(
            '<defs><pattern id="grid"><path d="M0,0 L10,10"/></pattern></defs>'
            '<rect x="0" y="0" width="10" height="10" fill="red"/>'
        )
        shapes = svg_to_pptx.render_svg_into_slide(slide, svg, Inches(0), Inches(0), Inches(2), Inches(1))
        assert len(shapes) == 1  # only the visible rect, not the defs' path

    def test_unsupported_gradient_fill_falls_back_to_flat_color(self):
        _, slide = _new_slide()
        svg = _svg('<rect x="0" y="0" width="10" height="10" fill="url(#bg)"/>')
        shapes = svg_to_pptx.render_svg_into_slide(slide, svg, Inches(0), Inches(0), Inches(2), Inches(1))
        assert len(shapes) == 1
        assert shapes[0].fill.fore_color.rgb is not None  # degraded to default gray, not raised


class TestFileAndErrorHandling:
    def test_embed_svg_file_reads_and_renders(self, tmp_path):
        _, slide = _new_slide()
        svg_path = tmp_path / "art.svg"
        svg_path.write_text(_svg('<rect x="0" y="0" width="10" height="10" fill="red"/>'), encoding="utf-8")
        shapes = svg_to_pptx.embed_svg_file(slide, str(svg_path), Inches(0), Inches(0), Inches(2), Inches(1))
        assert len(shapes) == 1

    def test_malformed_xml_raises(self):
        _, slide = _new_slide()
        with pytest.raises(Exception):
            svg_to_pptx.render_svg_into_slide(slide, "<svg><rect", Inches(0), Inches(0), Inches(2), Inches(1))

    def test_empty_viewbox_returns_no_shapes(self):
        _, slide = _new_slide()
        shapes = svg_to_pptx.render_svg_into_slide(
            slide, _svg('<rect x="0" y="0" width="10" height="10" fill="red"/>', viewbox="0 0 0 0"),
            Inches(0), Inches(0), Inches(2), Inches(1),
        )
        assert shapes == []


class TestPptxBuilderIntegration:
    def test_content_slide_routes_svg_through_native_embed(self, tmp_path):
        """.svg image_path must go through embed_svg_file, not add_picture (which can't rasterize SVG)."""
        from tools.slides import pptx_builder

        svg_path = tmp_path / "art.svg"
        svg_path.write_text(_svg('<rect x="0" y="0" width="10" height="10" fill="red"/>'), encoding="utf-8")

        with patch("tools.viz.svg_to_pptx.embed_svg_file") as mock_embed:
            mock_embed.return_value = []
            prs = Presentation()
            slide_data = {"slide_type": "content", "title": "T", "bullets": ["a"], "image_path": str(svg_path)}
            pptx_builder._build_content_slide(
                prs, slide_data, 1, pptx_builder.THEME_PALETTES["midnight_executive"], image_path=str(svg_path),
            )
            mock_embed.assert_called_once()

    def test_content_slide_png_still_uses_add_picture(self, tmp_path):
        from tools.slides import pptx_builder

        png_path = tmp_path / "art.png"
        # Minimal 1x1 PNG.
        png_path.write_bytes(bytes.fromhex(
            "89504e470d0a1a0a0000000d4948445200000001000000010802000000907753"
            "de0000000c4944415478da6360606060000000050001a5f645400000000049454e44ae426082"
        ))
        with patch("tools.viz.svg_to_pptx.embed_svg_file") as mock_embed:
            prs = Presentation()
            slide_data = {"slide_type": "content", "title": "T", "bullets": ["a"], "image_path": str(png_path)}
            pptx_builder._build_content_slide(
                prs, slide_data, 1, pptx_builder.THEME_PALETTES["midnight_executive"], image_path=str(png_path),
            )
            mock_embed.assert_not_called()

    def test_svg_art_slide_type_builds_native_shapes(self):
        from tools.slides import pptx_builder

        svg = _svg('<rect x="0" y="0" width="10" height="10" fill="red"/>')
        # A single-slide deck coerces slide 0 to "title" (build()'s `i == 0`
        # override) — use a 3-slide deck so svg_art lands in the middle.
        slides = [
            {"slide_type": "title", "title": "Intro"},
            {"slide_type": "svg_art", "title": "Art", "svg_code": svg},
            {"slide_type": "outro", "title": "Thanks"},
        ]
        path = pptx_builder.build(slides, title="SVG Art Smoke")
        try:
            prs = Presentation(path)
            shape_types = [s.shape_type for s in prs.slides[1].shapes]
            assert MSO_SHAPE_TYPE.FREEFORM in shape_types
        finally:
            Path(path).unlink(missing_ok=True)
