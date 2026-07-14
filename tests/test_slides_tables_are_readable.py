# CUI // SP-CTI
"""Every table ICDEV renders must be legible on a screen.

This is not a style preference. python-pptx creates tables with the "Medium Style
2 - Accent 1" table style, and that style supplies its own text colours AT THE RUN
LEVEL — which beats anything set on the paragraph. So a builder that paints a dark
fill and asks for white text gets the table style's DARK text drawn on that dark
fill, and the table renders as an empty box.

None of that is visible from Python. The text sits in the XML, python-pptx reads it
back happily, and a test asserting "the table has the right rows and the right
values" passes cleanly. It did pass. The deck reached the customer with every table
blank.

So these tests do not ask whether the text is present. They ask whether a person
could READ it — run-level colour, against the fill actually painted underneath, at
a contrast ratio a human eye can resolve.
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from tools.slides.constants import THEMES
from tools.slides.pptx_builder import build

_SLIDES = [
    {"slide_type": "title", "title": "Deck"},
    {
        "slide_type": "table",
        "title": "A table",
        "bullets": {
            "headers": ["Severity", "Impact", "What we found"],
            "rows": [
                ["HIGH", "$192,000", "A typed number in a block of formulas"],
                ["MEDIUM", "—", "Still unfinished"],
            ],
            "footer": ["1 high, 1 medium", "", ""],
        },
        "speaker_notes": "notes",
    },
    {"slide_type": "outro", "title": "End"},
]

# WCAG AA for body text. Below this, somebody at the back of the room cannot read
# the slide — which for a table of figures means the slide is not there.
_MIN_CONTRAST = 4.5


def _luminance(rgb: str) -> float:
    def chan(c: float) -> float:
        c /= 255
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

    r, g, b = (int(str(rgb)[i:i + 2], 16) for i in (0, 2, 4))
    return 0.2126 * chan(r) + 0.7152 * chan(g) + 0.0722 * chan(b)


def _contrast(fg: str, bg: str) -> float:
    a, b = _luminance(fg), _luminance(bg)
    return (max(a, b) + 0.05) / (min(a, b) + 0.05)


def _tables(path: str):
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                yield shape.table


@pytest.mark.parametrize("theme", THEMES)
def test_table_text_is_legible_in_every_theme(theme):
    path = build(_SLIDES, theme=theme, title="legibility")

    checked = 0
    for tbl in _tables(path):
        for row in tbl.rows:
            for cell in row.cells:
                if not cell.text.strip():
                    continue

                runs = cell.text_frame.paragraphs[0].runs
                assert runs, f"[{theme}] {cell.text!r}: no run — no colour can be set on it"

                fg = runs[0].font.color.rgb
                assert fg is not None, (
                    f"[{theme}] {cell.text!r}: no RUN-level colour. The table style "
                    f"will supply one, and it will not be the one you painted the "
                    f"cell for."
                )

                ratio = _contrast(str(fg), str(cell.fill.fore_color.rgb))
                assert ratio >= _MIN_CONTRAST, (
                    f"[{theme}] {cell.text!r} renders at {ratio:.1f}:1 "
                    f"(#{fg} on #{cell.fill.fore_color.rgb}) — effectively invisible."
                )
                checked += 1

    assert checked >= 9, f"[{theme}] no table cells were actually inspected"


def test_the_imposed_table_style_is_stripped():
    """"Medium Style 2 - Accent 1" is what does the damage.

    Removing it — and the banding flags that came with it — is what makes the
    builder's own fills and fonts the only formatting in play.
    """
    path = build(_SLIDES, theme="investment_deck", title="style")

    for tbl in _tables(path):
        pr = tbl._tbl.find(qn("a:tblPr"))
        style = pr.find(qn("a:tableStyleId"))
        assert style is not None
        # "No Style, No Grid" — the one built-in style that imposes nothing.
        assert style.text == "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"
        for flag in ("firstRow", "bandRow", "firstCol", "bandCol"):
            assert pr.get(flag) == "0", f"{flag} still on — the style would come back"


def test_a_table_still_carries_its_data():
    """The obvious half. Kept, but it is the half that already passed while the
    tables were blank — which is the entire lesson."""
    path = build(_SLIDES, theme="investment_deck", title="data")
    tbl = next(_tables(path))
    assert tbl.cell(0, 0).text == "Severity"
    assert tbl.cell(1, 1).text == "$192,000"


class TestATableMustFitOnTheSlide:
    """add_table() takes a height as a HINT, not a constraint.

    If the text in a cell wraps, PowerPoint grows the row to fit it and keeps
    growing — so a long description in a narrow column silently pushes the bottom
    of the table off the bottom of the slide. Asking for a smaller row height does
    not help: it is a request, and PowerPoint declines it.

    The only reliable fix is to put less in the table.
    """

    _BIG = [
        {"slide_type": "title", "title": "T"},
        {
            "slide_type": "table",
            "title": "A very long table",
            "bullets": {
                "headers": ["Manufacturer", "MSRP", "ROM", "Unknown", "Total"],
                "rows": [
                    [f"Manufacturer with a rather long name number {i}",
                     f"${i * 1000:,}", "—", f"${i * 2000:,}", f"${i * 3000:,}"]
                    for i in range(1, 30)
                ],
                "footer": ["Total", "$1", "$2", "$3", "$4"],
            },
        },
        {"slide_type": "outro", "title": "End"},
    ]

    def test_thirty_rows_do_not_run_off_the_bottom(self):
        from pptx.util import Inches

        path = build(self._BIG, theme="investment_deck", title="overflow")
        prs = Presentation(path)
        slide_h = prs.slide_height

        for shape in (sh for s in prs.slides for sh in s.shapes
                      if getattr(sh, "has_table", False)):
            bottom = shape.top + shape.height
            assert bottom <= slide_h - Inches(0.25), (
                f"the table ends {(bottom - slide_h) / 914400:.2f}in past the bottom "
                f"of the slide"
            )

    def test_what_was_dropped_is_DECLARED(self):
        """A table that silently shows the first twelve of thirty rows is worse
        than one that says it is doing so. Silent truncation reads as completeness.
        """
        path = build(self._BIG, theme="investment_deck", title="overflow")
        tbl = next(_tables(path))
        cells = [c.text for r in tbl.rows for c in r.cells]
        assert any("more row" in c for c in cells), (
            "rows were dropped and nothing said so"
        )

    def test_the_header_and_the_footer_are_never_dropped(self):
        """The header says what the columns mean; the footer usually carries the
        total. A table missing either is not a shorter table, it is a broken one."""
        path = build(self._BIG, theme="investment_deck", title="overflow")
        tbl = next(_tables(path))
        assert tbl.cell(0, 0).text.startswith("Manufacturer")
        last = [tbl.cell(len(tbl.rows) - 1, c).text for c in range(len(tbl.columns))]
        assert any("more row" in x for x in last) or last[0] == "Total"
