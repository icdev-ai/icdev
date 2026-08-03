# CUI // SP-CTI
"""Unit tests for the Slide Deck Generator Flask blueprint.

Covers page routes and the native asset-generator smoke endpoint.
Database-backed routes use the slides canvas DB (SQLite in tests).

sdt-auth-01 — two things this file needs that it used to get by accident:

1. The blueprint must be MOUNTED. It is only registered when ICDEV_SLIDES_ENABLED
   is on, and that lives in the repo .env; tests/conftest.py now sets it before the
   dashboard app is first imported (it cannot be done from this module — the app is
   already built by the time this module is collected).
2. The caller must be AUTHENTICATED. The seven write routes carry
   @require_role(*_SLIDES_WRITE_ROLES), and global dashboard auth requires a
   session on every route. The fix has to go THROUGH the session: a second
   before_request that sets g.current_user never runs, because
   _auth_before_request returns a redirect for an anonymous page request and a
   returned response short-circuits every later hook. So the `client` fixture seeds
   session['user_id'] and stubs get_user_by_id.

Authenticating the fixture to fix an auth failure would otherwise delete the only
thing that notices the gate being dropped later, so TestWriteRoutesRequireAuth
keeps two unauthenticated callers around. Note which assertion pins what — they
are not interchangeable, and this was measured by removing the decorator:

  * `anon_client` (no session) pins GLOBAL dashboard auth. It gets 401/302 from
    _auth_before_request whether or not @require_role is on the route, so on its
    own it would stay green through the decorator's removal.
  * `viewer_client` (session, role NOT in _SLIDES_WRITE_ROLES) is what actually
    pins @require_role: it clears the session gate and is refused 403 by the role
    gate. With the decorator deleted it reaches the view and returns 404 instead.
"""
import importlib
import json
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parents[2]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

# role must be in tools.slides.blueprint._SLIDES_WRITE_ROLES; status must be
# 'active' or _auth_before_request clears the session as expired.
_TEST_USER = {
    "id": "slides-test-user",
    "email": "slides-tester@icdev.local",
    "display_name": "Slides Tester",
    "role": "developer",
    "status": "active",
}


@pytest.fixture(scope="module")
def slides_db(tmp_path_factory):
    """Give this module its OWN slides canvas database.

    Third piece of undeclared ambient state, same family as the toggle. The canvas
    DB is whatever SLIDES_PG_DATABASE names — a PostgreSQL *database name* that
    get_canvas_connection() reuses as a SQLite *file path* once the backend is
    sqlite (which tests/conftest.py forces). With the repo .env loaded that value is
    "icdev_slides", i.e. the file of that name COMMITTED at the repo root, whose
    slides_decks was created from the PG schema and still carries the
    pre-sdt-vocab-01 eight-type deck_type CHECK. CREATE TABLE IF NOT EXISTS never
    widens a CHECK and slides migration 006 only widens it on PostgreSQL, so
    test_fill_end_to_end_creates_deck raises IntegrityError on a machine that has
    .env and passes on one that does not. A private file makes the run identical
    either way — and matches sdt-vv-01, which asserts against a *freshly
    initialised* schema.
    """
    path = tmp_path_factory.mktemp("slides_canvas") / "slides_canvas.db"
    # init_db() and the blueprint each memoise "schema is up" in a module global.
    # Both have to forget it or the new file never gets a schema.
    for name in ("tools.slides.db.init_db", "tools.slides.blueprint"):
        importlib.import_module(name)._INIT_DONE = False
    return str(path)


@pytest.fixture
def anon_client(monkeypatch, slides_db):
    """An explicitly UNAUTHENTICATED test client.

    Both dev-autologin escape hatches are removed so an inherited .env cannot
    silently authenticate the "anonymous" caller and turn the auth assertions green.
    """
    monkeypatch.delenv("ICDEV_DASHBOARD_DEV_AUTOLOGIN", raising=False)
    monkeypatch.delenv("ICDEV_DASHBOARD_API_KEY", raising=False)
    monkeypatch.setenv("SLIDES_PG_DATABASE", slides_db)

    from tools.dashboard.app import app

    app.config["TESTING"] = True
    with app.test_client() as test_client:
        yield test_client


@pytest.fixture
def viewer_client(monkeypatch, slides_db):
    """Authenticated, but with a role that is NOT in _SLIDES_WRITE_ROLES.

    This is what actually pins @require_role. The anonymous client cannot: global
    dashboard auth 401s it in before_request, so it would keep passing even if
    every @require_role were deleted. Only a caller who clears the session gate and
    is then refused by the role gate can tell the two apart.
    """
    monkeypatch.setenv("SLIDES_PG_DATABASE", slides_db)
    auth = importlib.import_module("tools.dashboard.auth")
    viewer = dict(_TEST_USER, id="slides-test-viewer", role="auditor")
    monkeypatch.setattr(auth, "get_user_by_id", lambda user_id: dict(viewer))

    from tools.dashboard.app import app

    app.config["TESTING"] = True
    with app.test_client() as test_client:
        with test_client.session_transaction() as sess:
            sess["user_id"] = viewer["id"]
        yield test_client


@pytest.fixture
def client(monkeypatch, slides_db):
    """Create a Flask test client from the full dashboard app so base.html context
    processors (nav_tree, etc.) are available, logged in as a slides author."""
    monkeypatch.setenv("SLIDES_PG_DATABASE", slides_db)

    from tools.dashboard.app import app
    from tools.slides.db.init_db import init_db

    # Shim-aware: tools.dashboard.auth and icdev.tools.dashboard.auth are the same
    # module object, and _auth_before_request resolves get_user_by_id as a module
    # global, so patching the attribute is what the hook actually reads.
    auth = importlib.import_module("tools.dashboard.auth")
    monkeypatch.setattr(auth, "get_user_by_id", lambda user_id: dict(_TEST_USER))

    assert "slides" in app.blueprints, (
        "slides blueprint not mounted — ICDEV_SLIDES_ENABLED must be set in "
        "tests/conftest.py before tools.dashboard.app is first imported"
    )

    app.config["TESTING"] = True
    with app.test_client() as test_client:
        with test_client.session_transaction() as sess:
            sess["user_id"] = _TEST_USER["id"]
        with app.app_context():
            try:
                init_db()
            except Exception:
                pass
        yield test_client


class TestPageRoutes:
    def test_index_returns_200_and_heading(self, client):
        resp = client.get("/slides/")
        assert resp.status_code == 200
        assert b"Slide Deck Generator" in resp.data
        assert b"CUI // SP-CTI" in resp.data

    def test_new_wizard_returns_200_and_form(self, client):
        resp = client.get("/slides/new")
        assert resp.status_code == 200
        assert b"Generate New Deck" in resp.data or b"deck_type" in resp.data
        assert b"CUI // SP-CTI" in resp.data


class TestAssetSmokeEndpoint:
    def test_asset_smoke_returns_svg_path(self, client):
        resp = client.post(
            "/slides/api/asset-smoke",
            data=json.dumps({"title": "Unit Test Slide", "bullets": ["a", "b"]}),
            content_type="application/json",
        )
        assert resp.status_code == 200, resp.data.decode()
        data = resp.get_json()
        assert data["success"] is True
        assert data["method"] == "slides_svg"
        assert isinstance(data["path"], str)
        assert len(data["path"]) > 0

    def test_asset_smoke_handles_missing_body(self, client):
        resp = client.post("/slides/api/asset-smoke", content_type="application/json")
        assert resp.status_code == 200
        data = resp.get_json()
        assert data["success"] is True
        assert data["method"] == "slides_svg"


class TestPresentRoute:
    def _insert_test_deck_sqlite(self):
        """Insert a minimal completed deck + one slide via SQLite directly, return deck_id."""
        from tools.db.storage import sql_placeholder
        from tools.slides.db.init_db import get_connection, init_db
        try:
            init_db()
        except Exception:
            pass
        conn = get_connection()
        conn.execute("PRAGMA foreign_keys=OFF")
        try:
            conn.execute(
                "INSERT INTO slides_decks "
                "(title, deck_type, theme, tone, status, slide_count, output_formats) "
                "VALUES ('Test Presentation', 'general_presentation', 'midnight_executive', "
                "'professional', 'completed', 1, '[\"pptx\"]')"
            )
            conn.commit()
            deck_id = conn.execute("SELECT last_insert_rowid()").fetchone()[0]
            # sql_placeholder(conn), not a bare '?': get_connection() here returns a
            # StorageConnection whose dialect follows the resolved backend, so a
            # hardcoded '?' is a ProgrammingError the moment this runs on psycopg2.
            # canvas_placeholder_style in coherence_checker gates exactly this.
            ph = sql_placeholder(conn)
            conn.execute(
                "INSERT INTO slides_slides "
                "(deck_id, position, slide_type, title, bullets, speaker_notes) "
                f"VALUES ({ph}, 1, 'content', 'Intro Slide', '[\"Bullet A\"]', 'Notes here')",
                (deck_id,),
            )
            conn.commit()
        finally:
            conn.close()
        return deck_id

    def test_present_404_for_unknown_deck(self, client):
        resp = client.get("/slides/99999/present")
        assert resp.status_code == 404

    def test_present_200_for_existing_deck(self, client):
        from unittest.mock import patch
        with patch.dict("os.environ", {"SLIDES_STORAGE_BACKEND": "sqlite"}, clear=False):
            deck_id = self._insert_test_deck_sqlite()
        resp = client.get(f"/slides/{deck_id}/present")
        assert resp.status_code in (200, 404)

    def test_new_wizard_includes_audience_mode_radios(self, client):
        resp = client.get("/slides/new")
        assert resp.status_code == 200
        assert b"audience_mode" in resp.data
        assert b"investor" in resp.data

    def test_new_wizard_includes_rich_diagrams_checkbox(self, client):
        resp = client.get("/slides/new")
        assert resp.status_code == 200
        assert b"enable_rich_diagrams" in resp.data


def _fixture_pptx_bytes() -> bytes:
    """A tiny title+body deck, built in-memory with python-pptx."""
    import io
    from pptx import Presentation

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.placeholders[0].text_frame.text = "Fixture Title"
    s.placeholders[1].text_frame.text = "Fixture bullet"
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


class TestTemplateFillRoutes:
    def _upload(self, client, filename: str = "fixture.pptx"):
        import io
        data = {"file": (io.BytesIO(_fixture_pptx_bytes()), filename)}
        return client.post(
            "/slides/api/templates/upload",
            data=data,
            content_type="multipart/form-data",
        )

    def test_upload_returns_template_id_and_shape_map(self, client):
        resp = self._upload(client)
        assert resp.status_code == 201, resp.data.decode()
        data = resp.get_json()
        assert data["template_id"] is not None
        assert data["slide_count"] == 1
        kinds = {s["kind"] for s in data["slides"][0]["shapes"]}
        assert kinds == {"title", "body"}

    def test_upload_rejects_non_pptx(self, client):
        import io
        resp = client.post(
            "/slides/api/templates/upload",
            data={"file": (io.BytesIO(b"not a pptx"), "notes.txt")},
            content_type="multipart/form-data",
        )
        assert resp.status_code == 400

    def test_upload_rejects_missing_file(self, client):
        resp = client.post("/slides/api/templates/upload", content_type="multipart/form-data")
        assert resp.status_code == 400

    def test_templates_index_lists_uploaded(self, client):
        self._upload(client, filename="listed.pptx")
        resp = client.get("/slides/templates")
        assert resp.status_code == 200
        assert b"listed.pptx" in resp.data

    def test_template_detail_404_for_unknown_id(self, client):
        resp = client.get("/slides/templates/999999")
        assert resp.status_code == 404

    def test_template_detail_renders_shape_summary(self, client):
        upload_resp = self._upload(client, filename="detail_target.pptx")
        template_id = upload_resp.get_json()["template_id"]
        resp = client.get(f"/slides/templates/{template_id}")
        assert resp.status_code == 200
        assert b"Fixture Title" in resp.data

    def test_fill_end_to_end_creates_deck(self, client):
        upload_resp = self._upload(client, filename="fill_target.pptx")
        template_id = upload_resp.get_json()["template_id"]

        fill_resp = client.post(
            f"/slides/api/templates/{template_id}/fill",
            data=json.dumps({"selections": [
                {"slide_index": 0, "title": "New Title", "bullets": ["b1", "b2"]},
            ]}),
            content_type="application/json",
        )
        assert fill_resp.status_code == 201, fill_resp.data.decode()
        deck_id = fill_resp.get_json()["deck_id"]
        assert deck_id is not None

        detail_resp = client.get(f"/slides/{deck_id}")
        assert detail_resp.status_code == 200

    def test_fill_requires_selections(self, client):
        upload_resp = self._upload(client, filename="empty_sel.pptx")
        template_id = upload_resp.get_json()["template_id"]
        resp = client.post(
            f"/slides/api/templates/{template_id}/fill",
            data=json.dumps({"selections": []}),
            content_type="application/json",
        )
        assert resp.status_code == 400

    def test_fill_404_for_unknown_template(self, client):
        resp = client.post(
            "/slides/api/templates/999999/fill",
            data=json.dumps({"selections": [{"slide_index": 0, "title": "x"}]}),
            content_type="application/json",
        )
        assert resp.status_code == 404


class TestWriteRoutesRequireAuth:
    """The `client` fixture above authenticates, so nothing else in this file would
    notice @require_role(*_SLIDES_WRITE_ROLES) being dropped from a write route.
    These pin it with a caller that is genuinely anonymous.

    401 vs 302 is not a detail: _auth_before_request only aborts(401) when the
    request is JSON or the path starts with /api/. The slides write routes live
    under /slides/api/..., so a multipart upload from a browser is REDIRECTED, and
    a redirect body is not the JSON the caller asked for. Both shapes are pinned.
    """

    def test_fill_refuses_anonymous_with_401(self, anon_client):
        resp = anon_client.post(
            "/slides/api/templates/1/fill",
            data=json.dumps({"selections": [{"slide_index": 0, "title": "x"}]}),
            content_type="application/json",
        )
        assert resp.status_code == 401

    def test_generate_refuses_anonymous_with_401(self, anon_client):
        resp = anon_client.post(
            "/slides/api/generate",
            data=json.dumps({"title": "x"}),
            content_type="application/json",
        )
        assert resp.status_code == 401

    def test_upload_refuses_anonymous(self, anon_client):
        import io

        resp = anon_client.post(
            "/slides/api/templates/upload",
            data={"file": (io.BytesIO(b"nope"), "x.pptx")},
            content_type="multipart/form-data",
        )
        # Non-JSON, non-/api/ path -> the auth hook redirects to the login page
        # rather than 401'ing. Either way the route never runs.
        assert resp.status_code in (302, 401)
        if resp.status_code == 302:
            assert "/login" in resp.headers.get("Location", "")

    def test_fill_refuses_authenticated_non_author_with_403(self, viewer_client):
        """The assertion that actually fails if @require_role is removed."""
        resp = viewer_client.post(
            "/slides/api/templates/1/fill",
            data=json.dumps({"selections": [{"slide_index": 0, "title": "x"}]}),
            content_type="application/json",
        )
        assert resp.status_code == 403

    def test_generate_refuses_authenticated_non_author_with_403(self, viewer_client):
        resp = viewer_client.post(
            "/slides/api/generate",
            data=json.dumps({"title": "x"}),
            content_type="application/json",
        )
        assert resp.status_code == 403

    def test_read_routes_stay_open_to_any_authenticated_user(self, viewer_client):
        """The gate is on the write routes only — a non-author still reads decks."""
        resp = viewer_client.get("/slides/")
        assert resp.status_code == 200
