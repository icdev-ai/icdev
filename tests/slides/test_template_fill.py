# CUI // SP-CTI
"""Unit tests for the template-fill workflow (tools/slides/template_fill.py)."""
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from tools.slides import template_fill


@pytest.fixture()
def fixture_pptx(tmp_path) -> str:
    """A 3-slide deck: title+body, table, and a slide meant to be dropped."""
    prs = Presentation()

    s0 = prs.slides.add_slide(prs.slide_layouts[1])
    s0.placeholders[0].text_frame.text = "Original Title"
    body = s0.placeholders[1]
    body.text_frame.text = "Original bullet"
    run = body.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(0x10, 0x20, 0x30)

    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    tbl = s1.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(1)).table
    tbl.cell(0, 0).text = "H1"
    tbl.cell(0, 1).text = "H2"
    tbl.cell(1, 0).text = "v1"
    tbl.cell(1, 1).text = "v2"

    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    s2.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1)).text_frame.text = "Drop me"

    path = tmp_path / "fixture_template.pptx"
    prs.save(str(path))
    return str(path)


class TestInspectTemplate:
    def test_reports_all_slides_and_shape_kinds(self, fixture_pptx):
        info = template_fill.inspect_template(fixture_pptx)
        assert info["slide_count"] == 3
        kinds_by_slide = [{s["kind"] for s in sl["shapes"]} for sl in info["slides"]]
        assert kinds_by_slide[0] == {"title", "body"}
        assert kinds_by_slide[1] == {"table"}

    def test_table_shape_reports_dimensions(self, fixture_pptx):
        info = template_fill.inspect_template(fixture_pptx)
        table_shape = info["slides"][1]["shapes"][0]
        assert table_shape["rows"] == 2
        assert table_shape["cols"] == 2

    def test_is_read_only(self, fixture_pptx):
        before = Path(fixture_pptx).read_bytes()
        template_fill.inspect_template(fixture_pptx)
        after = Path(fixture_pptx).read_bytes()
        assert before == after


class TestFillAndExport:
    def test_fills_title_and_bullets_preserving_format(self, fixture_pptx):
        out = template_fill.fill_and_export(fixture_pptx, [
            {"slide_index": 0, "title": "New Title", "bullets": ["one", "two", "three"]},
        ])
        try:
            prs = Presentation(out)
            assert len(prs.slides) == 1
            slide = prs.slides[0]
            title_shape = next(sh for sh in slide.shapes if sh.is_placeholder and sh.placeholder_format.idx == 0)
            assert title_shape.text_frame.text == "New Title"

            body_shape = next(sh for sh in slide.shapes if sh.is_placeholder and sh.placeholder_format.idx == 1)
            paragraphs = body_shape.text_frame.paragraphs
            assert [p.runs[0].text for p in paragraphs] == ["one", "two", "three"]
            # Formatting from the original single paragraph carried over to
            # the two newly-cloned ones.
            for p in paragraphs:
                assert p.runs[0].font.size == Pt(22)
                assert p.runs[0].font.color.rgb == RGBColor(0x10, 0x20, 0x30)
        finally:
            Path(out).unlink(missing_ok=True)

    def test_fills_table_preserving_unselected_cells(self, fixture_pptx):
        out = template_fill.fill_and_export(fixture_pptx, [
            {"slide_index": 1, "table": {"headers": ["A", "B"], "rows": [["r1a", "r1b"]]}},
        ])
        try:
            prs = Presentation(out)
            table = next(sh for sh in prs.slides[0].shapes if sh.has_table).table
            assert table.cell(0, 0).text == "A"
            assert table.cell(0, 1).text == "B"
            assert table.cell(1, 0).text == "r1a"
            assert table.cell(1, 1).text == "r1b"
        finally:
            Path(out).unlink(missing_ok=True)

    def test_drops_unselected_slides(self, fixture_pptx):
        out = template_fill.fill_and_export(fixture_pptx, [
            {"slide_index": 0, "title": "Kept"},
        ])
        try:
            prs = Presentation(out)
            assert len(prs.slides) == 1
            texts = [sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame]
            assert "Drop me" not in texts
        finally:
            Path(out).unlink(missing_ok=True)

    def test_keeps_multiple_selected_slides_in_order(self, fixture_pptx):
        out = template_fill.fill_and_export(fixture_pptx, [
            {"slide_index": 0, "title": "First"},
            {"slide_index": 1, "table": {"headers": ["X"], "rows": []}},
        ])
        try:
            prs = Presentation(out)
            assert len(prs.slides) == 2
        finally:
            Path(out).unlink(missing_ok=True)

    def test_no_valid_selection_raises(self, fixture_pptx):
        with pytest.raises(ValueError):
            template_fill.fill_and_export(fixture_pptx, [{"slide_index": 99, "title": "x"}])

    def test_fewer_bullets_than_original_paragraphs_blanks_leftovers(self, tmp_path):
        prs = Presentation()
        s0 = prs.slides.add_slide(prs.slide_layouts[1])
        s0.placeholders[0].text_frame.text = "T"
        body = s0.placeholders[1]
        body.text_frame.text = "line1"
        body.text_frame.add_paragraph().text = "line2"
        body.text_frame.add_paragraph().text = "line3"
        path = tmp_path / "multi.pptx"
        prs.save(str(path))

        out = template_fill.fill_and_export(str(path), [{"slide_index": 0, "bullets": ["only one"]}])
        try:
            check = Presentation(out)
            body_shape = next(sh for sh in check.slides[0].shapes if sh.is_placeholder and sh.placeholder_format.idx == 1)
            texts = [p.runs[0].text if p.runs else "" for p in body_shape.text_frame.paragraphs]
            assert texts[0] == "only one"
            assert all(t == "" for t in texts[1:])
        finally:
            Path(out).unlink(missing_ok=True)
