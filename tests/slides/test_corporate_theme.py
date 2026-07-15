# CUI // SP-CTI
"""The light corporate status theme, and the fallbacks that keep it from breaking
the dark themes it lives beside.

The whole risk of adding the first LIGHT theme to a builder written for dark ones
is that a colour the dark themes never set — a card fill, the text on the header
band — silently defaults to something that renders black-on-black or blue-on-navy.
So these tests assert two things: the light theme is actually legible, and every
dark theme still resolves the new keys to its old behaviour.
"""
from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[2]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from tools.slides import pptx_builder as B
from tools.slides.constants import THEME_PALETTES


def _hex(s: str):
    """'RRGGBB' → (r, g, b)."""
    return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _lum(rgb) -> float:
    return 0.2126 * rgb[0] + 0.7152 * rgb[1] + 0.0722 * rgb[2]


def _contrast(fg, bg) -> float:
    def chan(c):
        c /= 255
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

    def L(rgb):
        return 0.2126 * chan(rgb[0]) + 0.7152 * chan(rgb[1]) + 0.0722 * chan(rgb[2])

    a, b = L(fg), L(bg)
    return (max(a, b) + 0.05) / (min(a, b) + 0.05)


class TestTheThemeExists:
    def test_it_is_registered(self):
        assert "corporate_status" in THEME_PALETTES

    def test_it_is_actually_light(self):
        """A 'light corporate' theme with a dark page would be neither."""
        assert _lum(THEME_PALETTES["corporate_status"]["bg"]) >= 200


class TestLegibility:
    """Nothing the theme puts on a surface may be unreadable on it. WCAG AA is
    4.5:1; the loud failure this guards against is black-on-black."""

    def _p(self):
        return THEME_PALETTES["corporate_status"]

    def test_body_text_reads_on_the_card(self):
        p = self._p()
        card = p.get("card", p["dark"])
        assert _contrast(B._on_card_text(p), card) >= 4.5

    def test_the_band_title_reads_on_the_navy_band(self):
        """This is the trap: the builder used to hardcode the blue accent as the
        title colour, and blue on navy is unreadable. The light theme must override
        to white."""
        p = self._p()
        assert _contrast(B._band_text(p), p["dark"]) >= 4.5

    def test_body_and_subtext_read_on_the_page(self):
        p = self._p()
        assert _contrast(p["text"], p["bg"]) >= 4.5
        assert _contrast(p["subtext"], p["bg"]) >= 4.5


class TestDarkThemesAreUnchanged:
    """The point of the _opt() fallback: a dark theme sets none of the new keys and
    must behave exactly as before."""

    DARK = [t for t in THEME_PALETTES if t != "corporate_status"]

    def test_card_fill_falls_back_to_dark(self):
        for name in self.DARK:
            p = THEME_PALETTES[name]
            assert tuple(B._card_fill(p)) == tuple(B._rgb(p, "dark")), name

    def test_band_text_falls_back_to_accent(self):
        for name in self.DARK:
            p = THEME_PALETTES[name]
            assert tuple(B._band_text(p)) == tuple(B._rgb(p, "accent")), name

    def test_a_theme_without_a_rotation_reuses_its_accent(self):
        for name in self.DARK:
            p = THEME_PALETTES[name]
            if "rotation" not in p:
                rot = B._rotation(p)
                assert len(rot) == 1
                assert tuple(rot[0]) == tuple(B._rgb(p, "accent")), name


class TestRotation:
    def test_the_corporate_theme_cycles_four_accents(self):
        rot = B._rotation(THEME_PALETTES["corporate_status"])
        assert len(rot) == 4
        assert len({tuple(c) for c in rot}) == 4   # all distinct


class TestItRendersEndToEnd:
    """A palette that passes every unit check can still throw at draw time. Build a
    real deck and confirm each slide type produces shapes."""

    def _deck(self):
        return [
            {"slide_type": "title", "title": "Status Update",
             "subtitle": "July 2026", "speaker_notes": "sub"},
            {"slide_type": "card_grid", "title": "Focus Areas",
             "bullets": [
                 {"label": str(i), "title": f"Area {i}",
                  "body": "body text here", "meta": "SME: name"}
                 for i in range(1, 7)
             ]},
            {"slide_type": "content", "title": "Overview",
             "bullets": ["one", "two", "three"]},
            {"slide_type": "outro", "title": "Close", "bullets": ["done"]},
        ]

    def test_it_builds_a_file_with_every_slide(self):
        from pptx import Presentation

        path = B.build(self._deck(), theme="corporate_status", title="T")
        prs = Presentation(path)
        assert len(prs.slides) == 4
        for s in prs.slides:
            assert len(s.shapes) > 0

    def test_the_card_grid_draws_six_distinct_accent_stripes(self):
        """The signature. Six cards should not all wear the same colour."""
        from pptx import Presentation

        path = B.build(self._deck(), theme="corporate_status", title="T")
        prs = Presentation(path)
        grid = prs.slides[1]
        fills = set()
        for sh in grid.shapes:
            try:
                if sh.fill.type is not None and sh.height < B.Inches(0.12):
                    fills.add(str(sh.fill.fore_color.rgb))
            except Exception:
                pass
        # At least the four rotation colours appear among the thin stripes.
        assert len(fills) >= 4


class TestRoadmap:
    """The phased timeline — numbered circles on a spine, boxes alternating above
    and below. It must not be slide 0 or the last slide, or the builder's title/
    outro override at those positions runs instead (true for every slide type)."""

    def _deck(self, n=4):
        return [
            {"slide_type": "title", "title": "Deck"},
            {"slide_type": "roadmap", "title": "Roadmap",
             "phases": [
                 {"label": f"Phase {i}", "title": f"Phase {i}",
                  "body": "what happens here", "date": f"Est. M{i}"}
                 for i in range(1, n + 1)
             ]},
            {"slide_type": "outro", "title": "End"},
        ]

    def _ovals(self, slide):
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        out = []
        for sh in slide.shapes:
            try:
                if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and sh.auto_shape_type == 9:
                    out.append(sh)
            except Exception:
                pass
        return out

    def test_one_numbered_circle_per_phase(self):
        from pptx import Presentation

        prs = Presentation(B.build(self._deck(4), theme="corporate_status", title="T"))
        road = prs.slides[1]
        ovals = self._ovals(road)
        assert len(ovals) == 4
        assert {o.text_frame.text for o in ovals} == {"1", "2", "3", "4"}

    def test_phases_alternate_above_and_below_the_spine(self):
        """Boxes on one side only would overlap and be unreadable. Odd-indexed
        phases sit on the opposite side from even ones."""
        from pptx import Presentation

        prs = Presentation(B.build(self._deck(4), theme="corporate_status", title="T"))
        road = prs.slides[1]
        ovals = sorted(self._ovals(road), key=lambda o: o.left)
        line_y = ovals[0].top
        # Find each phase box (tall rounded rects) and bucket by side of the line.
        tops = sorted({sh.top for sh in road.shapes
                       if sh.height > B.Inches(1.0) and sh.width > B.Inches(1.5)})
        assert any(t < line_y for t in tops), "no box above the line"
        assert any(t > line_y for t in tops), "no box below the line"

    def test_circles_cycle_the_rotation_on_the_light_theme(self):
        from pptx import Presentation

        prs = Presentation(B.build(self._deck(4), theme="corporate_status", title="T"))
        colors = {str(o.fill.fore_color.rgb) for o in self._ovals(prs.slides[1])}
        assert len(colors) == 4   # blue / purple / green / amber, all distinct

    def test_the_date_line_is_legible_on_its_box(self):
        """The green accent as small text on its own pale tint is only 2.8:1 — so
        the builder DARKENS the accent for the label. Assert the colour it actually
        uses clears the 3:1 bar for large/bold text, for every rotation colour."""
        p = THEME_PALETTES["corporate_status"]
        for color in B._rotation(p):
            tint = B._tint(color, 0.86)
            date_c = B._shade(color)          # what the builder draws on a light box
            assert _contrast(date_c, tint) >= 3.0, str(color)

    def test_no_phases_is_a_message_not_a_crash(self):
        from pptx import Presentation

        deck = [{"slide_type": "title", "title": "D"},
                {"slide_type": "roadmap", "title": "Empty", "phases": []},
                {"slide_type": "outro", "title": "E"}]
        prs = Presentation(B.build(deck, theme="corporate_status", title="T"))
        assert len(prs.slides) == 3

    def test_it_caps_at_five_phases(self):
        from pptx import Presentation

        prs = Presentation(B.build(self._deck(8), theme="corporate_status", title="T"))
        assert len(self._ovals(prs.slides[1])) == 5

    def test_it_renders_on_a_dark_theme_too(self):
        from pptx import Presentation

        prs = Presentation(B.build(self._deck(3), theme="midnight_executive", title="T"))
        assert len(self._ovals(prs.slides[1])) == 3


class TestTablesDoNotClip:
    """A native PowerPoint table grows rows to fit their text, so enough rows push
    the table off the bottom of the slide — the last rows and the footer simply
    gone. These lock the fit: readable headers, wrapping on, the table inside the
    slide, and any overflow DECLARED rather than dropped off the edge."""

    def _deck(self, nrows, cols=5):
        rows = [
            [f"Line item {i} with a fairly long descriptive name that wraps"]
            + [f"c{c}v{i}" for c in range(cols - 1)]
            for i in range(1, nrows + 1)
        ]
        return [
            {"slide_type": "title", "title": "D"},
            {"slide_type": "table", "title": f"{nrows} rows",
             "bullets": {"headers": [f"Col {c}" for c in range(cols)],
                         "rows": rows,
                         "footer": ["Total"] + ["—"] * (cols - 1)}},
            {"slide_type": "outro", "title": "E"},
        ]

    def _table(self, prs):
        return next(sh for sh in prs.slides[1].shapes if sh.has_table)

    def test_the_header_is_readable_not_blue_on_navy(self):
        """The bug: header text was the blue accent on the navy header fill — a
        few percent of contrast. It must be the band colour (white on a light
        theme)."""
        from pptx import Presentation

        prs = Presentation(B.build(self._deck(5), theme="corporate_status", title="T"))
        tbl = self._table(prs).table
        p = THEME_PALETTES["corporate_status"]
        hdr = str(tbl.cell(0, 0).text_frame.paragraphs[0].font.color.rgb)
        fill = str(tbl.cell(0, 0).fill.fore_color.rgb)
        assert hdr == str(B._band_text(p))
        assert _contrast(_hex(hdr), _hex(fill)) >= 4.5

    def test_cells_wrap_instead_of_clipping(self):
        from pptx import Presentation

        prs = Presentation(B.build(self._deck(5), theme="corporate_status", title="T"))
        tbl = self._table(prs).table
        assert tbl.cell(1, 0).text_frame.word_wrap is True

    def test_a_big_table_stays_inside_the_slide(self):
        from pptx import Presentation

        for nrows in (5, 12, 30):
            prs = Presentation(B.build(self._deck(nrows), theme="corporate_status",
                                       title="T"))
            shp = self._table(prs)
            assert shp.top + shp.height <= B.H, f"{nrows} rows overflow the slide"

    def test_the_font_shrinks_as_rows_grow(self):
        from pptx import Presentation

        def body_pt(nrows):
            prs = Presentation(B.build(self._deck(nrows), theme="corporate_status",
                                       title="T"))
            return self._table(prs).table.cell(1, 0).text_frame.paragraphs[0].font.size.pt

        assert body_pt(5) > body_pt(20)

    def test_overflow_is_declared_not_silently_dropped(self):
        from pptx import Presentation

        prs = Presentation(B.build(self._deck(40), theme="corporate_status", title="T"))
        slide = prs.slides[1]
        note = [sh for sh in slide.shapes
                if sh.has_text_frame and "more row" in sh.text_frame.text]
        assert note, "capped rows must be announced, not silently cut off-slide"

    def test_the_footer_row_always_survives(self):
        """A total that falls off the bottom is worse than no total. Even when body
        rows are capped, the footer must still be in the table."""
        from pptx import Presentation

        prs = Presentation(B.build(self._deck(40), theme="corporate_status", title="T"))
        tbl = self._table(prs).table
        assert tbl.cell(len(tbl.rows) - 1, 0).text_frame.text == "Total"
