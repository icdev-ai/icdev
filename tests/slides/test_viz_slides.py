# CUI // SP-CTI
"""Tests for VIZ Epic B — rich slide builders + deterministic data mapper."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from tools.slides import pptx_builder, viz_mapper
from tools.viz.spec import ChartSpec, Series, TableSpec, KpiSpec, KpiTile, DiagramSpec


def _open(path: str) -> Presentation:
    assert Path(path).exists()
    return Presentation(path)


def test_build_dispatches_all_viz_slide_types(tmp_path):
    slides = [
        {"slide_type": "title", "title": "VIZ Deck", "bullets": []},
        {"slide_type": "agenda", "title": "Agenda",
         "bullets": ["Overview", "Metrics", "Roadmap"]},
        {"slide_type": "data", "title": "KPIs",
         "kpis": KpiSpec(title="KPIs", tiles=[KpiTile("Canvases", "56"),
                                              KpiTile("Tests", "330")]).to_dict()},
        {"slide_type": "data", "title": "Completion",
         "chart": ChartSpec(title="Completion", chart_type="bar",
                            categories=["A", "B"], series=[Series("pct", [80, 40])]).to_dict()},
        {"slide_type": "data", "title": "Status",
         "table": TableSpec(headers=["Canvas", "Status"],
                            rows=[["DIC", "Live"]]).to_dict()},
        {"slide_type": "content", "title": "Architecture",
         "diagram": DiagramSpec(title="Flow", nodes=[{"id": "a", "label": "X"},
                                {"id": "b", "label": "Y"}],
                                edges=[{"source": "a", "target": "b"}]).to_dict()},
        {"slide_type": "quote", "title": "Q", "bullets": ["A system that builds systems."]},
        {"slide_type": "outro", "title": "Thank You", "bullets": ["Contact us"]},
    ]
    path = pptx_builder.build(slides, title="VIZ Test Deck")
    prs = _open(path)
    assert len(prs.slides) == len(slides)
    # the chart slide carries a native chart graphic-frame
    assert any(sh.has_chart for slide in prs.slides for sh in slide.shapes)
    # the table slide carries a native table
    assert any(sh.has_table for slide in prs.slides for sh in slide.shapes)


def test_viz_mapper_kanban_real_numbers():
    raw = {
        "kanban": {
            "total_projects": 2,
            "in_progress_tasks": 5,
            "backlog_tasks": 9,
            "summary": "2 projects active.",
            "projects": [
                {"key": "viz", "name": "Viz", "progress_pct": 30, "in_progress": 3, "backlog": 4},
                {"key": "dic", "name": "DIC", "progress_pct": 90, "in_progress": 2, "backlog": 5},
            ],
        }
    }
    out = viz_mapper.build_data_slides(raw, max_slides=4)
    assert out, "mapper should produce data slides from kanban"
    # KPIs now live in a Tableau-style dashboard slide's tiles
    dash = next(s for s in out if s.get("dashboard"))
    kpi_tile = next(t["spec"] for t in dash["dashboard"]["tiles"]
                    if t["spec"].get("kind") == "kpis")
    vals = {t["label"]: t["value"] for t in kpi_tile["tiles"]}
    assert vals["In Progress"] == "5"
    assert dash.get("insight"), "dashboard slide carries a Story-Point insight"
    # progress chart is sorted descending by real pct (DIC=90 before viz=30)
    chart = next(s for s in out if s.get("chart") and "Completion" in s["title"])
    assert chart["chart"]["series"][0]["values"][0] == 90.0
    assert chart.get("insight"), "focused chart slide carries an insight"


def test_viz_mapper_empty_safe():
    assert viz_mapper.build_data_slides({}) == []
    assert viz_mapper.build_data_slides(None) == []
