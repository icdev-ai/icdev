# CUI // SP-CTI
"""Table slides must emit integer EMU coordinates.

`_build_table_slide` sized rows as ``min(Inches(0.5), (H - Inches(2.0)) / num_rows)``.
``Inches()`` returns an ``Emu`` (an ``int`` subclass), but ``/`` is true division, so
once ``num_rows`` exceeded 11 the ``min()`` picked the float branch and the frame
height reached the XML as ``cy="5029200.0"``.

OOXML ``ST_PositiveCoordinate`` derives from ``xsd:long`` — a decimal literal is not
schema-valid, and python-pptx raises ``ValueError`` the moment anything reads that
shape's ``.height``. The failure is silent at write time: the file is produced, the
build reports success, and the damage only surfaces when the deck is opened.

Regression guard for the Slides canvas (`tools/slides/pptx_builder.py`), which is
mirrored to `icdev/tools/slides/pptx_builder.py`.
"""
from __future__ import annotations

import re
import zipfile
from pathlib import Path

import pytest

from tools.slides.pptx_builder import build

#: cx/cy/x/y carrying a decimal point anywhere in the package.
FLOAT_COORD = re.compile(rb'\s(?:cx|cy|x|y)="-?\d+\.\d+"')


def _table_deck(num_data_rows: int) -> list[dict]:
    return [
        {"slide_type": "title", "title": "Fixture"},
        {
            "slide_type": "table",
            "title": f"{num_data_rows}-row table",
            "bullets": {
                "headers": ["Col A", "Col B"],
                "rows": [[f"r{i}", "value"] for i in range(num_data_rows)],
                "footer": ["Total", str(num_data_rows)],
            },
        },
        {"slide_type": "outro", "title": "End"},
    ]


def _float_coords(path: str) -> list[bytes]:
    with zipfile.ZipFile(path) as zf:
        return [
            match
            for info in zf.infolist()
            if info.filename.endswith(".xml")
            for match in FLOAT_COORD.findall(zf.read(info.filename))
        ]


@pytest.fixture
def built_deck(tmp_path):
    """build() writes to its own output dir; move each artifact into tmp_path."""
    created: list[Path] = []

    def _build(slides: list[dict], title: str) -> str:
        src = Path(build(slides, title=title))
        created.append(src)
        dest = tmp_path / src.name
        src.replace(dest)
        created[-1] = dest
        return str(dest)

    yield _build

    for p in created:
        p.unlink(missing_ok=True)


# 12 rows total (10 data + header + footer) is the first size where
# (H - 2in) / num_rows drops below Inches(0.5) and the float branch wins.
@pytest.mark.parametrize("num_data_rows", [1, 4, 9, 10, 25])
def test_table_slide_emits_integer_emu(built_deck, num_data_rows):
    path = built_deck(_table_deck(num_data_rows), f"emu-{num_data_rows}")
    offenders = _float_coords(path)
    assert not offenders, (
        f"{num_data_rows}-row table wrote non-integer EMU coordinates "
        f"(schema-invalid, PowerPoint reports the file as needing repair): {offenders}"
    )


def test_tall_table_height_is_readable(built_deck):
    """python-pptx must be able to read back the frame it just wrote."""
    from pptx import Presentation

    path = built_deck(_table_deck(25), "emu-readback")
    prs = Presentation(path)
    frames = [sh for slide in prs.slides for sh in slide.shapes if sh.has_table]
    assert frames, "fixture produced no table"
    for frame in frames:
        # Raised ValueError: invalid literal for int() ... '5029200.0' before the fix.
        assert isinstance(frame.height, int)
        assert frame.height > 0


def test_tall_table_stays_on_slide(built_deck):
    """The height cap must survive the int coercion — no running off the canvas."""
    from pptx import Presentation
    from pptx.util import Inches

    path = built_deck(_table_deck(25), "emu-bounds")
    prs = Presentation(path)
    frame = next(sh for slide in prs.slides for sh in slide.shapes if sh.has_table)
    assert frame.top + frame.height <= Inches(7.5), "table overruns the slide"
