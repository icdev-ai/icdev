# CUI // SP-CTI
"""Building inside somebody else's corporate master.

The master is the customer's: their fonts, their palette, their footer, and — in a
regulated environment — their classification banners, baked into parallel layout
families. None of that can be checked in here, so the fixture BUILDS a master with
the same shape: three marking families, four roles each.

That is a stronger test than the real file would be. It proves the module reads a
master it has never seen, which is the only property that makes it reusable.
"""
from __future__ import annotations

import re
import zipfile

import pytest
from pptx import Presentation
from pptx.util import Inches

from tools.slides.brand_deck import build_branded, families, layout_roles

_SLIDE_PART = re.compile(r"^ppt/slides/slide\d+\.xml$")


@pytest.fixture
def master(tmp_path):
    """A corporate master, invented: 3 marking families x 4 layout roles.

    python-pptx cannot author slide layouts, so this is built by RENAMING the
    layouts of the default template — which is exactly the structure a real master
    has, and gives us names to match against.
    """
    prs = Presentation()
    names = [
        "NoMarking_CoverSlide_White", "NoMarking_Title and Content_White",
        "NoMarking_TitleOnly_White", "NoMarking_BreakSlide_White",
        "CUI_CoverSlide_White", "CUI_Title and Content_White",
        "CUI_TitleOnly_White", "CUI_BreakSlide_White",
        "Proprietary_CoverSlide_White", "Proprietary_Title and Content_White",
        "Proprietary_TitleOnly_White",
    ]
    for layout, name in zip(prs.slide_layouts, names):
        layout.name = name

    # A master normally ships with example slides. They are somebody else's deck
    # and must not survive into ours.
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "SOMEBODY ELSE'S DECK"

    p = tmp_path / "master.pptx"
    prs.save(p)
    return str(p)


class TestReadingAMasterItHasNeverSeen:
    def test_it_finds_the_marking_families(self, master):
        assert families(master) == ["CUI", "NoMarking", "Proprietary"]

    def test_each_family_resolves_every_role(self, master):
        roles = layout_roles(master, "CUI")
        assert roles["cover"].startswith("CUI_Cover")
        assert roles["content"].startswith("CUI_Title and Content")
        assert roles["table"].startswith("CUI_TitleOnly")
        assert roles["section"].startswith("CUI_BreakSlide")

    def test_an_unknown_family_is_refused_rather_than_guessed(self, master):
        """Silently falling back to an unmarked layout would be a spill."""
        with pytest.raises(ValueError, match="no layouts"):
            layout_roles(master, "SECRET")


class TestTheMarkingIsChosenNotTyped:
    """The banner comes from the LAYOUT. That makes classifying the deck a single
    parameter, set once from the classification of the data, rather than twelve
    text boxes somebody has to remember."""

    @pytest.mark.parametrize("marking", ["CUI", "NoMarking", "Proprietary"])
    def test_every_slide_lands_in_the_requested_family(self, master, tmp_path, marking):
        out = tmp_path / f"{marking}.pptx"
        build_branded(
            master,
            [
                {"slide_type": "title", "title": "T", "bullets": ["a"]},
                {"slide_type": "content", "title": "C", "bullets": ["b"]},
            ],
            marking=marking,
            out=str(out),
        )
        prs = Presentation(str(out))
        assert prs.slides
        for s in prs.slides:
            assert s.slide_layout.name.startswith(marking + "_")

    def test_the_master_is_never_modified(self, master, tmp_path):
        before = open(master, "rb").read()
        build_branded(master, [{"slide_type": "title", "title": "T"}],
                      marking="CUI", out=str(tmp_path / "o.pptx"))
        assert open(master, "rb").read() == before


class TestWeBuildOurDeckNotTheirs:
    def test_the_templates_own_slides_are_dropped(self, master, tmp_path):
        out = tmp_path / "o.pptx"
        build_branded(master, [{"slide_type": "title", "title": "Ours"}],
                      marking="CUI", out=str(out))
        prs = Presentation(str(out))
        assert len(prs.slides) == 1
        text = " ".join(
            sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame
        )
        assert "SOMEBODY ELSE" not in text

    def test_the_templates_slides_are_not_still_INSIDE_THE_FILE(
        self, master, tmp_path,
    ):
        """The bug this exists for, and PowerPoint will never show it to you.

        Removing a slide id from the presentation's list stops the slide being
        DISPLAYED; the slide part stays in the .pptx package. The deck opens
        looking exactly right, and the file you hand to leadership still contains
        the previous author's slides — their pricing, their decision matrix —
        recoverable by anyone who unzips it.

        So this test reads the ZIP, not the render. Every check that goes through
        python-pptx's object model passes while the content is still there.
        """
        out = tmp_path / "o.pptx"
        build_branded(
            master,
            [{"slide_type": "title", "title": "Ours"},
             {"slide_type": "content", "title": "Also ours", "bullets": ["x"]}],
            marking="CUI", out=str(out),
        )

        with zipfile.ZipFile(out) as z:
            parts = [n for n in z.namelist() if _SLIDE_PART.match(n)]
            assert len(parts) == len(set(parts)), (
                f"duplicate slide parts in the package: {sorted(parts)}"
            )
            assert len(parts) == 2, (
                f"expected 2 slide parts, found {len(parts)}: {sorted(parts)} — "
                f"the template's own slides are still in the file"
            )
            blob = b"".join(z.read(n) for n in parts)

        assert b"SOMEBODY ELSE" not in blob

    def test_unfilled_placeholders_are_removed(self, master, tmp_path):
        """An empty placeholder renders as 'Click to add text' in presentation
        mode — the single most common way a generated deck announces itself."""
        out = tmp_path / "o.pptx"
        build_branded(master, [{"slide_type": "content", "title": "Only a title"}],
                      marking="CUI", out=str(out))
        prs = Presentation(str(out))
        empty = [
            ph for ph in prs.slides[0].placeholders
            if ph.has_text_frame and not ph.text_frame.text.strip()
        ]
        assert not empty


class TestContentIsNeverSilentlyLost:
    """The worst bug available in a generated artifact: it still looks finished."""

    def test_bullets_survive_a_layout_that_has_no_body_placeholder(
        self, master, tmp_path,
    ):
        """A section-break layout is a title and nothing else. Hand it bullets and
        they vanish — the deck renders, the slide looks deliberate, and the
        sentence you most wanted the room to hear is simply not there."""
        out = tmp_path / "o.pptx"
        build_branded(
            master,
            [{"slide_type": "outro", "title": "The close",
              "bullets": ["The 12 servers we already own do not wait on a PO."]}],
            marking="CUI",
            out=str(out),
        )
        text = " ".join(
            sh.text_frame.text
            for sh in Presentation(str(out)).slides[0].shapes
            if sh.has_text_frame
        )
        assert "do not wait on a PO" in text

    def test_a_table_reports_the_rows_it_could_not_fit(self, master, tmp_path):
        out = tmp_path / "o.pptx"
        build_branded(
            master,
            [{
                "slide_type": "table", "title": "Big",
                "bullets": {
                    "headers": ["A", "B"],
                    "rows": [[str(i), "x"] for i in range(40)],
                    "footer": [],
                },
            }],
            marking="CUI", out=str(out),
        )
        slide = Presentation(str(out)).slides[0]
        text = " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
        assert "further rows" in text, "a truncated table must say so"

    def test_a_table_renders_its_data(self, master, tmp_path):
        out = tmp_path / "o.pptx"
        build_branded(
            master,
            [{
                "slide_type": "table", "title": "Dates",
                "bullets": {
                    "headers": ["Milestone", "Target"],
                    "rows": [["Lab", "2026-12-31"]],
                    "footer": [],
                },
            }],
            marking="CUI", out=str(out),
        )
        slide = Presentation(str(out)).slides[0]
        tbl = next(sh for sh in slide.shapes if getattr(sh, "has_table", False)).table
        assert tbl.cell(0, 0).text == "Milestone"
        assert tbl.cell(1, 1).text == "2026-12-31"

    def test_a_table_fits_on_the_slide(self, master, tmp_path):
        out = tmp_path / "o.pptx"
        build_branded(
            master,
            [{
                "slide_type": "table", "title": "Big",
                "bullets": {"headers": ["A", "B", "C"],
                            "rows": [[str(i)] * 3 for i in range(30)], "footer": []},
            }],
            marking="CUI", out=str(out),
        )
        prs = Presentation(str(out))
        tbl = next(sh for sh in prs.slides[0].shapes if getattr(sh, "has_table", False))
        assert tbl.top + tbl.height <= prs.slide_height - Inches(0.2)

    def test_speaker_notes_survive(self, master, tmp_path):
        out = tmp_path / "o.pptx"
        build_branded(
            master,
            [{"slide_type": "content", "title": "T", "bullets": ["x"],
              "speaker_notes": "Do not soften this."}],
            marking="CUI", out=str(out),
        )
        notes = Presentation(str(out)).slides[0].notes_slide.notes_text_frame.text
        assert "Do not soften this" in notes


class TestItReportsWhatItGuessed:
    def test_a_role_with_no_matching_layout_is_named(self, master, tmp_path):
        """Proprietary has no BreakSlide in the fixture. The deck still builds —
        a deck rendered in the wrong layout is fixable, one that did not render is
        not — but nobody should discover the substitution in the room."""
        r = build_branded(
            master, [{"slide_type": "outro", "title": "Close"}],
            marking="Proprietary", out=str(tmp_path / "o.pptx"),
        )
        assert "section" in r["unmatched"]

    def test_a_clean_resolution_reports_nothing_unmatched(self, master, tmp_path):
        r = build_branded(
            master, [{"slide_type": "outro", "title": "Close"}],
            marking="CUI", out=str(tmp_path / "o.pptx"),
        )
        assert r["unmatched"] == []
        assert r["marking"] == "CUI"
