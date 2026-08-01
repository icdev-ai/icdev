# CUI // SP-CTI
"""Build a deck inside somebody else's corporate master.

`pptx_builder` draws slides from scratch under one of nine ICDEV themes. That is
right for an internal artifact and wrong for anything that leaves the building: a
deck going to leadership has to look like it came from the organisation, because a
deck that does not is read as a draft no matter what it says.

The organisation already has the answer. A corporate template carries its fonts,
its palette, its logo, its footer, and — in a regulated environment — its
CLASSIFICATION MARKINGS, baked into the slide layouts by the people whose job that
is. So we do not approximate it. We open it and build inside it.

`template_fill.fill_and_export` fills the slides a template already HAS. This
module does the other half: it reads the master's LAYOUTS and adds new slides from
them, so a deck of any length inherits the brand without anyone hand-editing a
theme file.

**Markings are chosen, not typed.** A master in a defence environment ships parallel
layout families — ``NoMarking_*``, ``ProprietaryInfo_*``, ``CUI_*`` — that differ
only in the banner. Picking the family is therefore the same act as classifying the
deck, and it is done ONCE, from the classification of the data, rather than by a
human remembering to change a text box on twelve slides. Getting that wrong is a
spill; a tool that makes it a parameter is a tool that can be audited.

Public API::

    families(master)              -> ["CUI", "NoMarking", "ProprietaryInfo"]
    layout_roles(master, marking) -> {"cover": "...", "content": "...", ...}
    build_branded(master, slides, marking=..., out=...) -> str   # a .pptx path
"""
from __future__ import annotations

import os
import re
import tempfile
from typing import Any, Iterable

from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches, Pt

from tools.slides.pptx_builder import (
    _clip,
    _fit_font,
    _fit_rows,
    _neutralize_table_style,
)

# The roles a deck needs, and how a corporate master tends to name them. Matched
# case-insensitively, first pattern wins, so the list is an ordered preference.
#
# These are PATTERNS, not a hardcoded vocabulary: a master that calls its content
# layout something else adds a pattern here, or passes an explicit map. Nothing in
# this module knows what a particular customer's layouts are called.
LAYOUT_PATTERNS: dict[str, tuple[str, ...]] = {
    "cover":   (r"cover.*white", r"cover.*general", r"coverslide", r"^title slide$"),
    "content": (r"title and content.*white", r"title and content", r"^content$"),
    "table":   (r"titleonly.*white", r"titleonly", r"title only",
                r"title and content.*white", r"title and content"),
    "section": (r"breakslide.*white", r"breakslide", r"section"),
}

# A layout family is a marking. "CUI_Title and Content_Black" is the CUI family;
# the same layout without the prefix is unmarked. Anything that does not look like
# a family prefix falls into "" — the master's own default.
_FAMILY = re.compile(r"^(?:\d+_)?([A-Za-z]+)_")

# Layout families that mean "no marking at all". Everything else is treated as a
# real marking, so an unknown family is never silently downgraded to unmarked.
_UNMARKED = ("nomarking",)


def _family_of(name: str) -> str:
    m = _FAMILY.match(name or "")
    return m.group(1) if m else ""


def families(master: str) -> list[str]:
    """The marking families this master offers, e.g. CUI / ProprietaryInfo."""
    prs = Presentation(master)
    return sorted({
        f for f in (_family_of(lo.name) for lo in prs.slide_layouts) if f
    })


def _layouts_in_family(prs, marking: str) -> list[Any]:
    want = (marking or "").lower()
    return [
        lo for lo in prs.slide_layouts
        if _family_of(lo.name).lower() == want
    ]


def layout_roles(master: str, marking: str = "") -> dict[str, str]:
    """Which layout plays which role, for this marking. Names, for inspection."""
    prs = Presentation(master)
    return {
        role: lo.name
        for role, lo in _resolve_roles(prs, marking).items()
    }


def _resolve_roles(prs, marking: str) -> dict[str, Any]:
    pool = _layouts_in_family(prs, marking) if marking else list(prs.slide_layouts)
    if not pool:
        raise ValueError(
            f"master has no layouts in the {marking!r} family; "
            f"available: {sorted({_family_of(lo.name) for lo in prs.slide_layouts})}"
        )

    out: dict[str, Any] = {}
    for role, patterns in LAYOUT_PATTERNS.items():
        for pat in patterns:
            hit = next((lo for lo in pool if re.search(pat, lo.name, re.I)), None)
            if hit is not None:
                out[role] = hit
                break
        if role not in out:
            # Fall back to the first layout in the family rather than failing: a
            # deck rendered in the wrong layout is fixable; a deck that did not
            # render is not. But say so — see build_branded's `unmatched`.
            out[role] = pool[0]
    return out


# ── Filling ──────────────────────────────────────────────────────────────────

def _placeholders(slide) -> tuple[Any | None, list[Any]]:
    """The title placeholder, and every body placeholder, in reading order."""
    title = None
    bodies = []
    for ph in slide.placeholders:
        t = ph.placeholder_format.type
        idx = ph.placeholder_format.idx
        if title is None and (idx == 0 or "TITLE" in str(t)):
            title = ph
        elif ph.has_text_frame:
            bodies.append(ph)
    return title, bodies


def _set_title(slide, text: str) -> None:
    title, _ = _placeholders(slide)
    if title is None or not text:
        return
    tf = title.text_frame
    tf.text = text
    # Do not touch size or colour: the master already decided, and overriding it is
    # exactly the thing this module exists not to do.


def _set_bullets(slide, prs, bullets: Iterable[str]) -> bool:
    """Fill the body placeholder — or make one, rather than lose the content.

    Not every layout has a body. A section-break layout in particular is usually a
    title and nothing else, so handing it bullets means they vanish: the deck
    renders, the slide looks deliberate, and the sentence you most wanted the room
    to hear is simply not there. Silent content loss in a generated artifact is the
    worst bug available, because the artifact still looks finished.
    """
    items = [str(b) for b in bullets if str(b).strip()]
    if not items:
        return False

    _, bodies = _placeholders(slide)
    if bodies:
        tf = bodies[0].text_frame
        tf.clear()
        for i, text in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.level = 0
        return True

    box = slide.shapes.add_textbox(
        Inches(0.9), Inches(2.6),
        prs.slide_width - Inches(1.8), Inches(2.4),
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(18)
        run.font.color.theme_color = MSO_THEME_COLOR.TEXT_1
    return True


def _drop_empty_placeholders(slide) -> None:
    """An unfilled placeholder renders as 'Click to add text' in presentation mode.

    It is the single most common way a generated deck announces that it was
    generated.
    """
    for ph in list(slide.placeholders):
        if ph.has_text_frame and not ph.text_frame.text.strip():
            ph._element.getparent().remove(ph._element)


def _add_table(slide, prs, spec: dict) -> None:
    """A table, sized to the slide and coloured from the MASTER's theme.

    Header fill and body text use theme colours (ACCENT_1 / TEXT_1) rather than
    literal RGB, so the table inherits whatever palette the corporate master
    defines — including a palette nobody here has ever seen.
    """
    headers = [str(h) for h in (spec.get("headers") or [])]
    all_rows = [[str(c) for c in r] for r in (spec.get("rows") or [])]
    footer = spec.get("footer") or []
    if not headers:
        return

    rows, footer, dropped = _fit_rows(all_rows, headers, footer)
    n_cols = len(headers)
    n_rows = len(rows) + 1 + (1 if footer else 0)
    font_pt = _fit_font(n_rows, n_cols)

    left, top = Inches(0.6), Inches(1.7)
    width = prs.slide_width - Inches(1.2)
    height = prs.slide_height - top - Inches(0.9)

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = shape.table
    _neutralize_table_style(tbl)

    row_h = int(height / max(1, n_rows))
    for r in tbl.rows:
        r.height = row_h

    def _write(cell, text: str, *, header: bool) -> None:
        cell.text = _clip(text, n_cols, font_pt)
        para = cell.text_frame.paragraphs[0]
        for run in para.runs or [para.add_run()]:
            run.font.size = Pt(font_pt)
            run.font.bold = header
            # Theme colours, so the master's palette wins.
            run.font.color.theme_color = (
                MSO_THEME_COLOR.BACKGROUND_1 if header else MSO_THEME_COLOR.TEXT_1
            )
        if header:
            cell.fill.solid()
            cell.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        else:
            cell.fill.background()

    for c, h in enumerate(headers):
        _write(tbl.cell(0, c), h, header=True)
    for r, row in enumerate(rows, start=1):
        for c in range(n_cols):
            _write(tbl.cell(r, c), row[c] if c < len(row) else "", header=False)

    if footer:
        last = len(rows) + 1
        for c in range(n_cols):
            _write(tbl.cell(last, c),
                   str(footer[c]) if c < len(footer) else "", header=False)
            tbl.cell(last, c).text_frame.paragraphs[0].runs[0].font.italic = True

    if dropped:
        # Say what was left out. A table that silently keeps the first twelve rows
        # is a table that lies about the other thirty-five.
        box = slide.shapes.add_textbox(
            left, prs.slide_height - Inches(0.85), width, Inches(0.3),
        )
        p = box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"{dropped} further rows — see the workbook"
        run.font.size = Pt(9)
        run.font.italic = True
        run.font.color.theme_color = MSO_THEME_COLOR.TEXT_1


def build_branded(
    master: str,
    slides: list[dict],
    *,
    marking: str = "",
    out: str | None = None,
) -> dict[str, Any]:
    """Render `slides` into a copy of `master`, using its layouts.

    Returns {path, marking, layouts, unmatched}. ``unmatched`` names any role that
    fell back to a default layout — a deck can ship with one, but nobody should
    discover it in the room.

    The master itself is never modified.
    """
    if not os.path.exists(master):
        raise FileNotFoundError(master)

    prs = Presentation(master)

    # Start from the master's LAYOUTS, not its slides. Whatever content the
    # template shipped with is somebody else's deck.
    #
    # Removing the id from _sldIdLst is not enough, and the failure is invisible:
    # the slide stops being SHOWN while its part stays in the package, so the file
    # you hand to leadership still contains the previous author's slides — their
    # pricing, their decision matrix — recoverable by anyone who unzips it. The
    # deck looks right in PowerPoint the entire time.
    #
    # drop_rel() severs the relationship, and an unrelated part is not serialised.
    for sld in list(prs.slides._sldIdLst):
        rId = sld.rId
        prs.slides._sldIdLst.remove(sld)
        prs.part.drop_rel(rId)

    roles = _resolve_roles(prs, marking)
    unmatched = [
        role for role, lo in roles.items()
        if not any(re.search(p, lo.name, re.I) for p in LAYOUT_PATTERNS[role])
    ]

    for spec in slides:
        kind = str(spec.get("slide_type") or "content")
        bullets = spec.get("bullets")
        is_table = isinstance(bullets, dict)

        if kind == "title":
            role = "cover"
        elif kind == "outro":
            role = "section"
        elif is_table:
            role = "table"
        else:
            role = "content"

        slide = prs.slides.add_slide(roles[role])
        _set_title(slide, str(spec.get("title") or ""))

        if is_table:
            _add_table(slide, prs, bullets)
        elif bullets:
            _set_bullets(slide, prs, bullets)

        notes = spec.get("speaker_notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)

        _drop_empty_placeholders(slide)

    path = out or os.path.join(
        tempfile.gettempdir(), f"branded_{abs(hash(master)) % 10**8}.pptx",
    )
    prs.save(path)

    return {
        "path": os.path.abspath(path),
        "marking": marking or "(master default)",
        "layouts": {r: lo.name for r, lo in roles.items()},
        "unmatched": unmatched,
        "slides": len(slides),
    }


__all__ = [
    "LAYOUT_PATTERNS",
    "build_branded",
    "families",
    "layout_roles",
]
