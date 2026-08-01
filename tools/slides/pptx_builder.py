# CUI // SP-CTI
"""PPTX Builder — assembles slide decks via python-pptx.

Reuses all primitives from tools/presentations/generate_exec_deck.py
(_blank, _bg, _rect, _box, _title, _footer, _gold_bar, _notes, _card).
Extends with multi-theme support and image embedding.

Themes: midnight_executive | govcon_proposal | compliance_briefing
"""
from __future__ import annotations

import hashlib
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from tools.slides.constants import THEME_PALETTES, DEFAULT_THEME

# ── Canvas dimensions (16:9 widescreen) ──────────────────────────────────────
W  = Inches(13.33)
H  = Inches(7.50)
LM = Inches(0.55)
CW = W - LM * 2

_ICDEV_ROOT = Path(__file__).resolve().parents[3]
_OUTPUT_DIR = _ICDEV_ROOT / "tools" / "presentations" / "slides"


def _rgb(palette: dict, key: str) -> RGBColor:
    r, g, b = palette[key]
    return RGBColor(r, g, b)


def _opt(palette: dict, key: str, fallback: str) -> RGBColor:
    """A colour that a theme MAY define, falling back to one it must.

    This is what lets a light theme coexist with the dark ones without a flag on
    every call site: a dark theme never sets "card" or "band_text", so it inherits
    the old behaviour; a light theme sets them and gets white cards and a white
    band title. No theme has to know the others exist.
    """
    return _rgb(palette, key if key in palette else fallback)


def _card_fill(palette: dict) -> RGBColor:
    # Dark themes fill cards with "dark"; a light theme fills them white via "card".
    return _opt(palette, "card", "dark")


def _band_text(palette: dict) -> RGBColor:
    # Title colour ON the navy header band. Blue-on-navy is unreadable, so a light
    # theme overrides to white; dark themes keep the accent they always used.
    return _opt(palette, "band_text", "accent")


def _on_card_text(palette: dict) -> RGBColor:
    """Body colour that is legible on the card fill, whichever way it goes.

    A white card needs dark text; a navy card needs light text. Decide from the
    card's own luminance rather than from a per-theme flag, so a new theme cannot
    forget to set it and ship black-on-black.
    """
    r, g, b = palette.get("card", palette["dark"])
    luminance = 0.2126 * r + 0.7152 * g + 0.0722 * b
    return _rgb(palette, "text") if luminance < 140 else _rgb(palette, "dark")


def _rotation(palette: dict) -> list[RGBColor]:
    """The per-card accent cycle. A theme without one just reuses its accent."""
    rot = palette.get("rotation")
    if not rot:
        return [_rgb(palette, "accent")]
    return [RGBColor(*c) for c in rot]


def _is_light(palette: dict) -> bool:
    r, g, b = palette["bg"]
    return (0.2126 * r + 0.7152 * g + 0.0722 * b) >= 140


def _tint(color: RGBColor, amount: float = 0.85) -> RGBColor:
    """Mix a colour toward white — the pale phase-box fills in the reference deck.

    Computed, not per-theme: a light tint of whatever accent a phase carries, so a
    theme that adds a fifth rotation colour gets a matching box for free.
    """
    return RGBColor(
        round(color[0] + (255 - color[0]) * amount),
        round(color[1] + (255 - color[1]) * amount),
        round(color[2] + (255 - color[2]) * amount),
    )


def _shade(color: RGBColor, amount: float = 0.35) -> RGBColor:
    """Mix a colour toward black. A rotation accent as a pale box fill is bright —
    green in particular clears barely 2.8:1 as small text on its own tint — so the
    date LABEL on the box is a darkened version of the accent, legible by
    construction rather than by luck."""
    return RGBColor(
        round(color[0] * (1 - amount)),
        round(color[1] * (1 - amount)),
        round(color[2] * (1 - amount)),
    )


# ── Primitives (identical to generate_exec_deck.py) ──────────────────────────

def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color: RGBColor) -> None:
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def _rect(slide, l, t, w, h, fill: RGBColor, line=None, lw=Pt(1)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s


def _box(
    slide, l, t, w, h, text: str = "", size: int = 14,
    bold: bool = False, italic: bool = False,
    color: RGBColor = None, align=PP_ALIGN.LEFT, wrap: bool = True
):
    if color is None:
        color = RGBColor(0xFF, 0xFF, 0xFF)
    shape = slide.shapes.add_textbox(l, t, w, h)
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf


def _accent_bar(slide, palette: dict, top=Inches(1.55), h=Inches(0.04)) -> None:
    _rect(slide, LM, top, CW, h, _rgb(palette, "accent"))


def _footer(slide, n: int, palette: dict, text: str = "ICDEV™  ·  A System That Builds Systems") -> None:
    dark = _rgb(palette, "dark")
    _box(slide, LM, H - Inches(0.32), CW - Inches(0.6), Inches(0.28),
         text, size=8, color=dark)
    _box(slide, W - Inches(0.8), H - Inches(0.32), Inches(0.7), Inches(0.28),
         str(n), size=9, color=dark, align=PP_ALIGN.RIGHT)


def _citation_footer(slide, palette: dict, citations: list[dict]) -> None:
    """Render short source citations at the bottom of a content slide."""
    if not citations:
        return
    subtext = _rgb(palette, "subtext")
    lines = []
    for i, src in enumerate(citations[:3], start=1):
        title = (src.get("title") or "Source")[:60]
        url = src.get("url", "")
        line = f"[{i}] {title}"
        if url:
            line += f" — {url[:80]}"
        lines.append(line)
    text = " | ".join(lines)
    _box(slide, LM, H - Inches(0.58), CW, Inches(0.26),
         text, size=7, color=subtext, wrap=True)


def _notes(slide, text: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    tf.text = text


def _card(slide, l, t, w, h, heading: str, body: str, palette: dict,
          accent: RGBColor | None = None) -> None:
    accent = accent or _rgb(palette, "accent")
    subtext = _rgb(palette, "subtext")
    border = _opt(palette, "border", "accent")
    body_c = _on_card_text(palette)
    # White card with a hairline border and a coloured LEFT accent bar — the
    # reference deck's signature. On a dark theme the fill is "dark" and the effect
    # is a subtly bordered panel; same code, both worlds.
    _rect(slide, l, t, w, h, _card_fill(palette), border)
    _rect(slide, l, t, Inches(0.08), h, accent)          # left accent bar
    pad = Inches(0.18)
    _box(slide, l + pad, t + Inches(0.10), w - pad * 2, Inches(0.38),
         heading, size=13, bold=True, color=body_c)
    _box(slide, l + pad, t + Inches(0.50), w - pad * 2, h - Inches(0.62),
         body, size=10, color=subtext, wrap=True)


# ── Slide Builders ────────────────────────────────────────────────────────────

def _build_title_slide(prs: Presentation, slide_data: dict, n: int, palette: dict) -> None:
    s = _blank(prs)
    accent = _rgb(palette, "accent")

    # A cover wants a deep field behind a bright title. Dark themes already have
    # that in "bg"; a light theme's "bg" is white, which makes a flat cover — so a
    # light theme covers in its navy "dark" instead and sets the title white. Both
    # are decided from luminance, not a flag, so a new theme cannot get it wrong.
    br, bgc, bb = palette["bg"]
    is_light = (0.2126 * br + 0.7152 * bgc + 0.0722 * bb) >= 140
    cover_bg = _rgb(palette, "dark") if is_light else _rgb(palette, "bg")
    title_c = _band_text(palette) if is_light else accent
    text_c = _band_text(palette) if is_light else _rgb(palette, "text")
    subtext = _opt(palette, "subtext", "subtext")

    _bg(s, cover_bg)

    # Top + bottom accent bars. The bottom uses the last rotation colour (amber in
    # the corporate theme) for the reference's two-tone rule; falls back to accent.
    rot = _rotation(palette)
    _rect(s, 0, 0, W, Inches(0.1), accent)
    _rect(s, 0, H - Inches(0.1), W, Inches(0.1), rot[-1])

    title = slide_data.get("title", "ICDEV™")
    _box(s, LM, Inches(1.0), CW, Inches(2.0),
         title, size=44, bold=True, color=title_c, align=PP_ALIGN.CENTER)
    _rect(s, Inches(3.8), Inches(3.2), Inches(5.73), Inches(0.03), accent)

    subtitle = slide_data.get("speaker_notes", "")[:120]
    tagline = slide_data.get("subtitle") or "ICDEV™  ·  A System That Builds Systems"
    _box(s, LM, Inches(3.35), CW, Inches(0.5),
         tagline, size=14, color=text_c, align=PP_ALIGN.CENTER)
    if subtitle:
        _box(s, LM, Inches(3.9), CW, Inches(0.4),
             subtitle[:100], size=12, color=subtext, align=PP_ALIGN.CENTER, italic=True)


def _build_content_slide(
    prs: Presentation, slide_data: dict, n: int, palette: dict, image_path: str | None = None
) -> None:
    s = _blank(prs)
    _bg(s, _rgb(palette, "bg"))
    accent = _rgb(palette, "accent")
    dark = _rgb(palette, "dark")

    # Left accent stripe
    _rect(s, 0, 0, Inches(0.12), H, accent)

    # Title bar — full-width navy band, title in the theme's band colour (white on
    # a light theme, accent on the dark themes).
    _rect(s, Inches(0.12), 0, W - Inches(0.12), Inches(0.72), dark)
    title = slide_data.get("title", "")[:80]
    _box(s, Inches(0.24), Inches(0.12), CW, Inches(0.55),
         title, size=22, bold=True, color=_band_text(palette))

    # Accent underline
    _accent_bar(s, palette, top=Inches(0.72), h=Inches(0.04))

    bullets = slide_data.get("bullets", [])
    if image_path and Path(image_path).exists():
        # Two-column layout: bullets left, image right
        col_w = CW * 0.55
        _add_bullets(s, bullets, Inches(0.24), Inches(0.85), col_w, palette)
        img_l = LM + col_w + Inches(0.15)
        img_w = CW - col_w - Inches(0.1)
        img_h = Inches(5.6)
        try:
            if str(image_path).lower().endswith(".svg"):
                # add_picture cannot rasterize SVG (no PIL SVG support) — embed
                # as native, editable shapes instead of a flattened picture.
                from tools.viz import svg_to_pptx
                svg_to_pptx.embed_svg_file(s, image_path, img_l, Inches(0.85), img_w, img_h)
            else:
                s.shapes.add_picture(image_path, img_l, Inches(0.85), img_w, img_h)
        except Exception:
            pass
    else:
        _add_bullets(s, bullets, Inches(0.24), Inches(0.85), CW, palette)

    _citation_footer(s, palette, slide_data.get("citations", []))
    _footer(s, n, palette)
    notes_text = slide_data.get("speaker_notes", "")
    if notes_text:
        _notes(s, notes_text)


def _add_bullets(slide, bullets: list[str], l, t, w, palette: dict) -> None:
    text_c = _rgb(palette, "text")
    subtext = _rgb(palette, "subtext")
    if not bullets:
        return
    shape = slide.shapes.add_textbox(l, t, w, Inches(5.6))
    tf = shape.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets[:5]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(6) if i > 0 else Pt(0)
        run = p.add_run()
        run.text = f"▸  {bullet}"
        run.font.size = Pt(16)
        run.font.color.rgb = text_c if i == 0 else subtext


def _try_mmdc_render(mermaid_code: str) -> str | None:
    """Try rendering Mermaid diagram to PNG via mmdc CLI. Returns temp PNG path or None."""
    import subprocess
    import tempfile
    import os
    tmp_dir = tempfile.gettempdir()
    in_path = os.path.join(tmp_dir, f"icdev_mmdc_in_{id(mermaid_code)}.mmd")
    out_path = os.path.join(tmp_dir, f"icdev_mmdc_out_{id(mermaid_code)}.png")
    try:
        with open(in_path, "w", encoding="utf-8") as f:
            f.write(mermaid_code)
        result = subprocess.run(
            ["mmdc", "-i", in_path, "-o", out_path, "-b", "transparent", "-w", "1200"],
            timeout=15, capture_output=True,
        )
        if result.returncode == 0 and Path(out_path).exists():
            return out_path
    except Exception:
        pass
    finally:
        try:
            os.unlink(in_path)
        except Exception:
            pass
    return None


def _build_mermaid_slide(prs: Presentation, slide_data: dict, n: int, palette: dict) -> None:
    s = _blank(prs)
    _bg(s, _rgb(palette, "bg"))
    accent = _rgb(palette, "accent")
    dark = _rgb(palette, "dark")
    subtext = _rgb(palette, "subtext")

    _rect(s, 0, 0, Inches(0.12), H, accent)
    _rect(s, Inches(0.12), 0, W - Inches(0.12), Inches(0.72), dark)
    title = slide_data.get("title", "")[:80]
    _box(s, Inches(0.24), Inches(0.12), CW, Inches(0.55),
         title, size=22, bold=True, color=_band_text(palette))
    _accent_bar(s, palette, top=Inches(0.72), h=Inches(0.04))

    mermaid_code = slide_data.get("mermaid_code") or ""
    img_path = _try_mmdc_render(mermaid_code) if mermaid_code else None

    if img_path:
        try:
            s.shapes.add_picture(img_path, LM, Inches(0.9), CW, Inches(5.2))
        except Exception:
            img_path = None

    if not img_path:
        # Fallback: monospace text block with raw Mermaid code
        _box(s, LM, Inches(0.9), CW, Inches(4.0),
             mermaid_code[:800] if mermaid_code else "(No diagram generated)",
             size=9, color=subtext, wrap=True)
        _box(s, LM, Inches(5.0), CW, Inches(0.4),
             "ℹ️  Open the web presentation viewer or HTML export to see the interactive diagram.",
             size=11, color=accent, italic=True, wrap=True)

    _footer(s, n, palette)
    notes_text = slide_data.get("speaker_notes", "")
    if notes_text:
        _notes(s, notes_text)


def _build_svg_slide(prs: Presentation, slide_data: dict, n: int, palette: dict) -> None:
    """Full-slide vector art, rendered as native/editable PPTX shapes (not a picture)."""
    s = _blank(prs)
    _bg(s, _rgb(palette, "bg"))
    accent = _rgb(palette, "accent")
    dark = _rgb(palette, "dark")
    subtext = _rgb(palette, "subtext")

    _rect(s, 0, 0, Inches(0.12), H, accent)
    _rect(s, Inches(0.12), 0, W - Inches(0.12), Inches(0.72), dark)
    title = slide_data.get("title", "")[:80]
    _box(s, Inches(0.24), Inches(0.12), CW, Inches(0.55),
         title, size=22, bold=True, color=_band_text(palette))
    _accent_bar(s, palette, top=Inches(0.72), h=Inches(0.04))

    svg_code = slide_data.get("svg_code") or ""
    rendered = False
    if svg_code:
        try:
            from tools.viz import svg_to_pptx
            shapes = svg_to_pptx.render_svg_into_slide(s, svg_code, LM, Inches(0.9), CW, Inches(5.6))
            rendered = bool(shapes)
        except Exception:
            rendered = False

    if not rendered:
        _box(s, LM, Inches(2.8), CW, Inches(0.6),
             "(No vector art generated)", size=14, color=subtext, align=PP_ALIGN.CENTER)

    _footer(s, n, palette)
    notes_text = slide_data.get("speaker_notes", "")
    if notes_text:
        _notes(s, notes_text)


def _build_three_placeholder_slide(prs: Presentation, slide_data: dict, n: int, palette: dict) -> None:
    s = _blank(prs)
    _bg(s, _rgb(palette, "bg"))
    accent = _rgb(palette, "accent")
    dark = _rgb(palette, "dark")
    subtext = _rgb(palette, "subtext")

    _rect(s, 0, 0, Inches(0.12), H, accent)
    _rect(s, Inches(0.12), 0, W - Inches(0.12), Inches(0.72), dark)
    title = slide_data.get("title", "")[:80]
    _box(s, Inches(0.24), Inches(0.12), CW, Inches(0.55),
         title, size=22, bold=True, color=_band_text(palette))
    _accent_bar(s, palette, top=Inches(0.72), h=Inches(0.04))

    cfg = slide_data.get("three_scene_config") or {}
    preset = cfg.get("preset", "neural_network") if isinstance(cfg, dict) else "neural_network"

    _rect(s, LM, Inches(1.1), CW, Inches(4.2), dark)
    _box(s, LM + Inches(0.3), Inches(2.0), CW - Inches(0.6), Inches(0.6),
         "🌐  3D Animation", size=28, bold=True, color=accent, align=PP_ALIGN.CENTER)
    _box(s, LM + Inches(0.3), Inches(2.7), CW - Inches(0.6), Inches(0.5),
         f"Scene: {preset}", size=14, color=subtext, align=PP_ALIGN.CENTER, italic=True)

    _box(s, LM, Inches(5.4), CW, Inches(0.5),
         "Interactive 3D animation — open the web presentation viewer to experience this slide.",
         size=11, color=accent, italic=True, wrap=True)

    _footer(s, n, palette)
    notes_text = slide_data.get("speaker_notes", "")
    if notes_text:
        _notes(s, notes_text)


def _build_excalidraw_placeholder_slide(prs: Presentation, slide_data: dict, n: int, palette: dict) -> None:
    s = _blank(prs)
    _bg(s, _rgb(palette, "bg"))
    accent = _rgb(palette, "accent")
    dark = _rgb(palette, "dark")
    subtext = _rgb(palette, "subtext")

    _rect(s, 0, 0, Inches(0.12), H, accent)
    _rect(s, Inches(0.12), 0, W - Inches(0.12), Inches(0.72), dark)
    title = slide_data.get("title", "")[:80]
    _box(s, Inches(0.24), Inches(0.12), CW, Inches(0.55),
         title, size=22, bold=True, color=_band_text(palette))
    _accent_bar(s, palette, top=Inches(0.72), h=Inches(0.04))

    _rect(s, LM, Inches(1.1), CW, Inches(4.2), dark)
    _box(s, LM + Inches(0.3), Inches(2.0), CW - Inches(0.6), Inches(0.6),
         "✏️  Hand-Drawn Diagram", size=24, bold=True, color=accent, align=PP_ALIGN.CENTER)
    _box(s, LM + Inches(0.3), Inches(2.7), CW - Inches(0.6), Inches(0.8),
         slide_data.get("speaker_notes", "")[:200],
         size=12, color=subtext, align=PP_ALIGN.CENTER, wrap=True)

    _box(s, LM, Inches(5.4), CW, Inches(0.5),
         "Interactive sketch — open the web presentation viewer to experience this slide.",
         size=11, color=accent, italic=True, wrap=True)

    _footer(s, n, palette)
    notes_text = slide_data.get("speaker_notes", "")
    if notes_text:
        _notes(s, notes_text)


def _build_roadmap_slide(prs: Presentation, slide_data: dict, n: int, palette: dict) -> None:
    """A phased timeline: a horizontal spine, numbered circles, and phase boxes
    that alternate above and below the line.

    This is the reference deck's milestone slide. Each phase takes the next colour
    in the rotation, so the circles and boxes march blue → purple → green → amber;
    the boxes are a pale tint of that colour on a light theme, or the card fill on
    a dark one, so the same layout reads on either.

    slide_data["phases"] = [{"label","title","body","date"}], up to 5.
    """
    s = _blank(prs)
    _bg(s, _rgb(palette, "bg"))
    accent = _rgb(palette, "accent")
    dark = _rgb(palette, "dark")
    subtext = _rgb(palette, "subtext")
    rotation = _rotation(palette)
    light = _is_light(palette)

    _rect(s, 0, 0, Inches(0.12), H, accent)
    _rect(s, Inches(0.12), 0, W - Inches(0.12), Inches(0.72), dark)
    _box(s, Inches(0.24), Inches(0.12), CW, Inches(0.55),
         slide_data.get("title", "")[:80], size=22, bold=True, color=_band_text(palette))
    _accent_bar(s, palette, top=Inches(0.72), h=Inches(0.04))

    phases = [p for p in (slide_data.get("phases") or []) if isinstance(p, dict)][:5]
    if not phases:
        _box(s, LM, Inches(1.5), CW, Inches(2.0),
             "No phases — pass slide_data['phases'] = [{label,title,body,date}].",
             size=14, color=subtext)
        _footer(s, n, palette)
        return

    # The spine, centred vertically in the body area.
    line_y = Inches(4.0)
    line_h = Inches(0.06)
    spine = _rgb(palette, "dark") if light else _rgb(palette, "subtext")
    _rect(s, LM, line_y, CW, line_h, spine)
    # Arrowhead at the right end.
    _rect(s, W - LM, line_y - Inches(0.09), Inches(0.18), Inches(0.24), spine)

    seg = CW / len(phases)
    circle_d = Inches(0.5)
    box_w = min(Inches(2.9), seg - Inches(0.25))
    box_h = Inches(1.55)
    gap = Inches(0.55)   # circle-to-box vertical gap

    for i, ph in enumerate(phases):
        color = rotation[i % len(rotation)]
        cx = LM + seg * i + seg / 2
        circle_l = cx - circle_d / 2
        circle_t = line_y + line_h / 2 - circle_d / 2

        above = (i % 2 == 0)   # alternate: 1 above, 2 below, 3 above ...

        # Connector stub between the circle and its box.
        stub_x = cx - Inches(0.01)
        if above:
            box_t = circle_t - gap - box_h
            _rect(s, stub_x, box_t + box_h, Inches(0.02), gap, color)
        else:
            box_t = circle_t + circle_d + gap
            _rect(s, stub_x, circle_t + circle_d, Inches(0.02), gap, color)

        # The phase box: tinted fill on a light theme, card fill on a dark one.
        box_l = cx - box_w / 2
        fill = _tint(color, 0.86) if light else _card_fill(palette)
        _rect(s, box_l, box_t, box_w, box_h, fill, color, lw=Pt(1.5))
        pad = Inches(0.14)
        title_c = _rgb(palette, "dark") if light else _on_card_text(palette)
        _box(s, box_l + pad, box_t + Inches(0.10), box_w - pad * 2, Inches(0.35),
             str(ph.get("title", ""))[:40], size=13, bold=True, color=title_c)
        _box(s, box_l + pad, box_t + Inches(0.48), box_w - pad * 2, Inches(0.62),
             str(ph.get("body", ""))[:90], size=10, color=subtext, wrap=True)
        date = str(ph.get("date", ""))[:28]
        if date:
            # Darken the accent for the label on a light (tinted) box; on a dark
            # theme the bright accent already reads against the dark card fill.
            date_c = _shade(color) if light else color
            _box(s, box_l + pad, box_t + box_h - Inches(0.30), box_w - pad * 2,
                 Inches(0.26), date, size=10, bold=True, color=date_c)

        # The numbered circle sits ON TOP of the spine, drawn last so it wins.
        circ = s.shapes.add_shape(9, circle_l, circle_t, circle_d, circle_d)  # 9 = oval
        circ.fill.solid()
        circ.fill.fore_color.rgb = color
        circ.line.color.rgb = _rgb(palette, "bg")
        circ.line.width = Pt(2.5)
        ctf = circ.text_frame
        ctf.word_wrap = False
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        crun = cp.add_run()
        crun.text = str(ph.get("label", i + 1)).replace("Phase", "").strip() or str(i + 1)
        crun.font.size = Pt(16)
        crun.font.bold = True
        crun.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    _footer(s, n, palette)
    notes_text = slide_data.get("speaker_notes", "")
    if notes_text:
        _notes(s, notes_text)


def _build_card_grid_slide(prs: Presentation, slide_data: dict, n: int, palette: dict) -> None:
    """3-column card grid (investment overview / capability comparison)."""
    s = _blank(prs)
    _bg(s, _rgb(palette, "bg"))
    accent = _rgb(palette, "accent")
    dark = _rgb(palette, "dark")
    subtext = _rgb(palette, "subtext")
    rotation = _rotation(palette)

    _rect(s, 0, 0, Inches(0.12), H, accent)
    _rect(s, Inches(0.12), 0, W - Inches(0.12), Inches(0.72), dark)
    title = slide_data.get("title", "")[:80]
    _box(s, Inches(0.24), Inches(0.12), CW, Inches(0.55),
         title, size=22, bold=True, color=_band_text(palette))
    _accent_bar(s, palette, top=Inches(0.72), h=Inches(0.04))

    bullets = slide_data.get("bullets", [])
    # Expect bullets to be card dicts or plain strings
    cards = []
    for b in bullets:
        if isinstance(b, dict):
            cards.append(b)
        elif isinstance(b, str):
            cards.append({"label": "", "title": b, "body": "", "meta": ""})

    if not cards:
        _box(s, LM, Inches(1.5), CW, Inches(3.0),
             "No card data — bullets should be card objects.", size=14, color=subtext)
        _footer(s, n, palette)
        return

    cols = 3
    rows = (len(cards) + cols - 1) // cols
    card_w = CW / cols - Inches(0.12)
    card_h = min(Inches(2.0), (H - Inches(1.4)) / rows - Inches(0.1))
    start_y = Inches(0.85)
    pad = Inches(0.1)

    for ci, card in enumerate(cards[:9]):  # max 9 cards (3x3)
        row, col = divmod(ci, cols)
        cl = LM + col * (card_w + Inches(0.12))
        ct = start_y + row * (card_h + Inches(0.08))

        # An explicit per-card colour wins; otherwise cycle the theme's rotation so
        # adjacent cards differ (blue → purple → green → amber) instead of every
        # card wearing the same accent.
        card_accent_hex = card.get("accent_color", "")
        if card_accent_hex and card_accent_hex.startswith("#") and len(card_accent_hex) >= 7:
            try:
                ca = RGBColor(int(card_accent_hex[1:3], 16),
                              int(card_accent_hex[3:5], 16),
                              int(card_accent_hex[5:7], 16))
            except Exception:
                ca = rotation[ci % len(rotation)]
        else:
            ca = rotation[ci % len(rotation)]

        border = _opt(palette, "border", "accent")
        title_c = _on_card_text(palette)
        _rect(s, cl, ct, card_w, card_h, _card_fill(palette), border)
        _rect(s, cl, ct, card_w, Inches(0.06), ca)  # top accent bar

        label = str(card.get("label", ""))[:10]
        title_text = str(card.get("title", ""))[:50]
        body_text = str(card.get("body", ""))[:120]
        meta_text = str(card.get("meta", ""))[:60]

        if label:
            _box(s, cl + pad, ct + Inches(0.10), card_w - pad * 2, Inches(0.28),
                 label, size=11, bold=True, color=ca)
        body_top = ct + Inches(0.10 + (0.28 if label else 0))
        if title_text:
            _box(s, cl + pad, body_top, card_w - pad * 2, Inches(0.28),
                 title_text, size=11, bold=True, color=title_c)
            body_top += Inches(0.30)
        if body_text:
            _box(s, cl + pad, body_top, card_w - pad * 2, card_h - (body_top - ct) - Inches(0.28),
                 body_text, size=10, color=subtext, wrap=True)
        if meta_text:
            _box(s, cl + pad, ct + card_h - Inches(0.24), card_w - pad * 2, Inches(0.2),
                 meta_text, size=9, bold=True, color=ca, wrap=False)

    _footer(s, n, palette)
    notes_text = slide_data.get("speaker_notes", "")
    if notes_text:
        _notes(s, notes_text)


def _table_fit(n_body_rows: int, has_hf: int, avail_h) -> tuple[int, int]:
    """Choose a font size that lets the rows FIT, and how many body rows fit at it.

    A native PowerPoint table grows each row to fit its wrapped text — so a fixed
    row height is only a minimum, and enough rows push the table straight off the
    bottom of the slide, where the last rows and the footer are simply gone. That
    is the "cutoff" nobody put there on purpose.

    So instead of a fixed height, pick the largest font at which every row fits,
    and if even the smallest font cannot hold them all, say how many were dropped
    rather than letting them fall off the edge.
    """
    avail = avail_h / 914400.0   # EMU → inches
    for pt in (12, 11, 10, 9, 8):
        # A body cell may wrap to ~2 lines; header/footer to 1. Row height ≈ two
        # lines of this font plus the tight cell margins we set below.
        body_row_in = (pt * 1.28 * 2) / 72.0 + 0.10
        hf_row_in = (pt * 1.28) / 72.0 + 0.10
        max_body = int((avail - has_hf * hf_row_in) / body_row_in)
        if max_body >= n_body_rows:
            return pt, n_body_rows
    # Smallest font, capped — caller appends a "+N more" note.
    pt = 8
    body_row_in = (pt * 1.28 * 2) / 72.0 + 0.10
    hf_row_in = (pt * 1.28) / 72.0 + 0.10
    max_body = max(1, int((avail - has_hf * hf_row_in) / body_row_in))
    return pt, max_body
# "No Style, No Grid" — the one built-in table style that imposes nothing.
_NO_TABLE_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"

# How many rows fit between the title band and the footer without PowerPoint
# growing them past the bottom of the slide. Measured, not derived: a 10pt row with
# one line of text is about 0.30", and there is roughly 5.4" to play with.
_MAX_TABLE_ROWS = 14


def _fit_rows(all_rows, headers, footer):
    """Trim a table to what will physically fit. Returns (rows, footer, dropped).

    The header and the footer are never dropped — the header says what the columns
    mean and the footer usually carries the total, and a table missing either is
    not a shorter table, it is a broken one.
    """
    has_header = bool(headers)
    has_footer = bool(footer)

    body_start = 1 if has_header else 0
    body_end = len(all_rows) - (1 if has_footer else 0)
    body = all_rows[body_start:body_end]

    room = _MAX_TABLE_ROWS - (1 if has_header else 0) - (1 if has_footer else 0)
    dropped = max(0, len(body) - room)
    if dropped:
        body = body[:room]

    out = ([all_rows[0]] if has_header else []) + body
    if has_footer:
        out.append(all_rows[-1])
    return out, (all_rows[-1] if has_footer else None), dropped


def _fit_font(num_rows: int, num_cols: int) -> int:
    """Smaller type for busier tables. A wall of 12pt text is a wall."""
    if num_rows > 10 or num_cols > 5:
        return 9
    if num_rows > 7 or num_cols > 4:
        return 10
    return 12


def _clip(text: str, num_cols: int, font_pt: int) -> str:
    """Keep a cell to roughly one line.

    A cell that wraps to three lines makes PowerPoint grow the row, and a few of
    those push the table off the slide. The full text is in the workbook; the slide
    is for the argument, not the archive.
    """
    usable_in = 12.0 / max(num_cols, 1)          # the content width, in inches
    chars = int(usable_in * (96 / (font_pt * 0.62)))
    return text if len(text) <= chars else text[: max(chars - 1, 8)].rstrip() + "…"


def _readable(fg, bg, fallback):
    """The more legible of two colours against this background.

    Two of the shipped themes put their accent on a dark header fill at under
    4.5:1, which is below what a person at the back of a room can resolve. Rather
    than hand-tuning palettes, pick whichever of the two candidate colours a human
    can actually read — the theme keeps its character and the table stays legible.
    """
    def lum(c):
        def chan(x):
            x /= 255
            return x / 12.92 if x <= 0.03928 else ((x + 0.055) / 1.055) ** 2.4
        return 0.2126 * chan(c[0]) + 0.7152 * chan(c[1]) + 0.0722 * chan(c[2])

    def ratio(a, b):
        la, lb = lum(a), lum(b)
        return (max(la, lb) + 0.05) / (min(la, lb) + 0.05)

    bg_t = (bg[0], bg[1], bg[2])
    if ratio((fg[0], fg[1], fg[2]), bg_t) >= 4.5:
        return fg
    return fallback


def _neutralize_table_style(tbl) -> None:
    """Stop PowerPoint's default table style from overriding our colours.

    python-pptx creates every table with "Medium Style 2 - Accent 1", which brings
    its own banding and — fatally — its own text colours. Those are applied at the
    run level and beat anything set on the paragraph, so a builder that paints a
    dark fill and asks for white text gets the table style's DARK text on that dark
    fill, and the table renders as an empty box.

    Nothing about that failure is visible from Python: the text is present in the
    XML, python-pptx reads it back happily, and every automated check passes. It is
    only wrong on a screen.

    So: drop the style, and turn off the banding and first-row emphasis that came
    with it. Explicit formatting is then the only formatting there is.
    """
    from pptx.oxml.ns import qn

    tbl_pr = tbl._tbl.find(qn("a:tblPr"))
    if tbl_pr is None:
        return

    for style_id in tbl_pr.findall(qn("a:tableStyleId")):
        tbl_pr.remove(style_id)

    style = tbl_pr.makeelement(qn("a:tableStyleId"), {})
    style.text = _NO_TABLE_STYLE
    tbl_pr.append(style)

    # These flags tell PowerPoint to apply the style's special formatting to the
    # header row and to alternate rows. With no style they are meaningless; with
    # one they are how the colours come back.
    tbl_pr.set("firstRow", "0")
    tbl_pr.set("bandRow", "0")
    tbl_pr.set("firstCol", "0")
    tbl_pr.set("bandCol", "0")


def _build_table_slide(prs: Presentation, slide_data: dict, n: int, palette: dict) -> None:
    """Render a data table using python-pptx's native table shape."""
    s = _blank(prs)
    _bg(s, _rgb(palette, "bg"))
    accent = _rgb(palette, "accent")
    dark = _rgb(palette, "dark")
    subtext = _rgb(palette, "subtext")
    text_c = _rgb(palette, "text")
    band_c = _band_text(palette)          # readable on the navy band AND the cells

    _rect(s, 0, 0, Inches(0.12), H, accent)
    _rect(s, Inches(0.12), 0, W - Inches(0.12), Inches(0.72), dark)
    title = slide_data.get("title", "")[:80]
    _box(s, Inches(0.24), Inches(0.12), CW, Inches(0.55),
         title, size=22, bold=True, color=band_c)
    _accent_bar(s, palette, top=Inches(0.72), h=Inches(0.04))

    tbl_data = slide_data.get("bullets") or {}
    if isinstance(tbl_data, str):
        import json as _json
        try:
            tbl_data = _json.loads(tbl_data)
        except Exception:
            tbl_data = {}

    headers = tbl_data.get("headers", []) if isinstance(tbl_data, dict) else []
    rows = list(tbl_data.get("rows", []) if isinstance(tbl_data, dict) else [])
    footer = tbl_data.get("footer", []) if isinstance(tbl_data, dict) else []

    if not (headers or rows or footer):
        _box(s, LM, Inches(1.5), CW, Inches(1.0), "No table data.", size=14, color=subtext)
        _footer(s, n, palette)
        return

    top = Inches(0.95)
    # Emu subclasses int, but `/` is true division: a height that reaches the XML
    # as cy="5029200.0" is not a valid ST_PositiveCoordinate (xsd:long), so
    # PowerPoint reports the deck as needing repair and python-pptx raises on
    # .height — with no error at write time. This expression is int-only, but
    # the guard is kept explicit so a future divisor here cannot silently
    # reintroduce the corruption main fixed.
    avail_h = Emu(int(H - top - Inches(0.5)))   # leave room for the page footer
    has_hf = (1 if headers else 0) + (1 if footer else 0)
    font_pt, max_body = _table_fit(len(rows), has_hf, avail_h)

    # Cap rows to what fits, and SAY what was cut instead of clipping it off-slide.
    dropped = 0
    if len(rows) > max_body:
        dropped = len(rows) - max_body
        rows = rows[:max_body]

    all_rows = ([headers] if headers else []) + rows + ([footer] if footer else [])
    num_rows = len(all_rows)
    num_cols = max(len(r) for r in all_rows)

    tbl_shape = s.shapes.add_table(num_rows, num_cols, LM, top, CW, avail_h)
    tbl = tbl_shape.table
    tbl.first_row = bool(headers)         # let the table style bar the header row

    _neutralize_table_style(tbl)

    for ri, row_data in enumerate(all_rows):
        is_header = ri == 0 and bool(headers)
        is_footer = ri == len(all_rows) - 1 and bool(footer)
        band = is_header or is_footer
        for ci in range(num_cols):
            cell = tbl.cell(ri, ci)
            cell.text = str(row_data[ci]) if ci < len(row_data) else ""
            # Tight margins + explicit wrap: less forced wrapping, shorter rows,
            # and no horizontal clipping of a long cell.
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            tf = cell.text_frame
            tf.word_wrap = True

            fill = dark if band else _rgb(palette, "bg")
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill

            # Banded rows take the theme accent — unless a human cannot read it
            # against this fill, in which case they get the body colour. Two of
            # the shipped themes put their accent on the dark header at under
            # 4.5:1, and a header nobody can read is not a header.
            colour = _readable(band_c, fill, text_c) if band else text_c
            size = Pt(font_pt - 1 if band else font_pt)

            para = tf.paragraphs[0]
            para.font.size = size
            para.font.bold = band
            para.font.color.rgb = colour

            # Set the colour on the RUN as well, not only the paragraph.
            #
            # A paragraph-level colour is a DEFAULT, and a PowerPoint table's own
            # style supplies run-level formatting that beats it. Without this the
            # builder painted a dark fill, set light paragraph text, and
            # PowerPoint rendered the table style's dark text on top: every table
            # came out black-on-black. The data was all in the XML, which is why
            # it survived — reading a .pptx back with python-pptx shows you the
            # text and tells you nothing about whether a human can SEE it.
            for run in para.runs:
                run.font.size = size
                run.font.bold = band
                run.font.color.rgb = colour

    if dropped:
        _box(s, LM, H - Inches(0.5), CW, Inches(0.22),
             f"+ {dropped} more row(s) — full table in the workbook",
             size=9, italic=True, color=subtext)

    _footer(s, n, palette)
    notes_text = slide_data.get("speaker_notes", "")
    if notes_text:
        _notes(s, notes_text)


def _build_outro_slide(prs: Presentation, slide_data: dict, n: int, palette: dict) -> None:
    s = _blank(prs)
    _bg(s, _rgb(palette, "bg"))
    accent = _rgb(palette, "accent")
    subtext = _rgb(palette, "subtext")

    _rect(s, 0, 0, W, Inches(0.1), accent)
    _rect(s, 0, H - Inches(0.1), W, Inches(0.1), accent)

    title = slide_data.get("title", "Thank You")
    _box(s, LM, Inches(1.2), CW, Inches(1.5),
         title, size=40, bold=True, color=accent, align=PP_ALIGN.CENTER)
    _rect(s, Inches(3.8), Inches(2.8), Inches(5.73), Inches(0.03), accent)

    bullets = slide_data.get("bullets", [])
    for i, bullet in enumerate(bullets[:3]):
        _box(s, LM, Inches(3.0 + i * 0.7), CW, Inches(0.55),
             bullet, size=16, color=subtext, align=PP_ALIGN.CENTER)

    _footer(s, n, palette)
    notes_text = slide_data.get("speaker_notes", "")
    if notes_text:
        _notes(s, notes_text)


# ── Main Builder ──────────────────────────────────────────────────────────────

def build(
    slides: list[dict],
    theme: str = DEFAULT_THEME,
    title: str = "ICDEV™ Presentation",
) -> str:
    """Assemble slides into a .pptx file. Returns absolute file path."""
    palette = THEME_PALETTES.get(theme, THEME_PALETTES[DEFAULT_THEME])
    prs = _new_prs()

    _OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    for i, slide_data in enumerate(slides):
        n = i + 1
        slide_type = slide_data.get("slide_type", "content")
        image_path = slide_data.get("image_path")

        if slide_type == "title" or i == 0:
            _build_title_slide(prs, slide_data, n, palette)
        elif slide_type == "outro" or i == len(slides) - 1:
            _build_outro_slide(prs, slide_data, n, palette)
        elif slide_type == "mermaid_diagram":
            _build_mermaid_slide(prs, slide_data, n, palette)
        elif slide_type == "svg_art":
            _build_svg_slide(prs, slide_data, n, palette)
        elif slide_type == "three_animation":
            _build_three_placeholder_slide(prs, slide_data, n, palette)
        elif slide_type == "excalidraw_sketch":
            _build_excalidraw_placeholder_slide(prs, slide_data, n, palette)
        elif slide_type == "card_grid":
            _build_card_grid_slide(prs, slide_data, n, palette)
        elif slide_type == "roadmap":
            _build_roadmap_slide(prs, slide_data, n, palette)
        elif slide_type == "table":
            _build_table_slide(prs, slide_data, n, palette)
        else:
            _build_content_slide(prs, slide_data, n, palette, image_path)

    # Generate filename with timestamp
    ts = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
    slug = hashlib.sha256(title.encode()).hexdigest()[:8]
    filename = f"{ts}_{slug}.pptx"
    out_path = _OUTPUT_DIR / filename
    prs.save(str(out_path))
    return str(out_path)
