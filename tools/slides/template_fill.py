# CUI // SP-CTI
"""Template-fill workflow — reuse a customer-supplied .pptx design as-is.

Ingests an uploaded .pptx (e.g. a federal agency's own proposal template),
lets the caller pick which of its slides to reuse, fills new content into
those slides in place (preserving every shape/theme/master exactly), and
drops the rest — no LLM step, deterministic python-pptx only.

Public API:
  inspect_template(pptx_path) -> dict        # read-only shape map
  fill_and_export(pptx_path, selections) -> str  # -> new .pptx path
"""
from __future__ import annotations

import copy
import hashlib
from datetime import datetime, timezone

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from tools.slides.pptx_builder import _OUTPUT_DIR

_TITLE_PLACEHOLDER_TYPES = (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
_BODY_PLACEHOLDER_TYPES = (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE)
_PREVIEW_LEN = 80


def _classify_shape(shape) -> str | None:
    """Classify a shape into a fillable "kind", or None if not fillable."""
    if shape.has_table:
        return "table"
    if shape.has_chart:
        return "chart"
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if shape.is_placeholder:
        ph_type = shape.placeholder_format.type
        if ph_type in _TITLE_PLACEHOLDER_TYPES:
            return "title"
        if ph_type in _BODY_PLACEHOLDER_TYPES:
            return "body"
        return None  # date/footer/slide-number placeholders — not content fields
    if shape.has_text_frame and shape.text_frame.text.strip():
        return "body"
    return None


def _shape_preview(shape, kind: str) -> dict:
    info: dict = {"shape_id": shape.shape_id, "kind": kind, "name": shape.name}
    if kind in ("title", "body"):
        text = shape.text_frame.text.strip()
        info["preview_text"] = text[:_PREVIEW_LEN]
    elif kind == "table":
        table = shape.table
        info["rows"] = len(table.rows)
        info["cols"] = len(table.columns)
        info["preview_text"] = table.cell(0, 0).text.strip()[:_PREVIEW_LEN] if info["rows"] and info["cols"] else ""
    elif kind == "chart":
        info["preview_text"] = str(shape.chart.chart_type)
    elif kind == "picture":
        info["preview_text"] = shape.name
    return info


def inspect_template(pptx_path: str) -> dict:
    """Read-only: enumerate every slide's fillable shapes.

    Returns {slide_count, slides: [{index, shapes: [{shape_id, kind, name,
    preview_text, rows?, cols?}]}]}.
    """
    prs = Presentation(pptx_path)
    slides = []
    for index, slide in enumerate(prs.slides):
        shapes = []
        for shape in slide.shapes:
            kind = _classify_shape(shape)
            if kind is None:
                continue
            shapes.append(_shape_preview(shape, kind))
        slides.append({"index": index, "shapes": shapes})
    return {"slide_count": len(prs.slides), "slides": slides}


# ── Content substitution (format-preserving) ─────────────────────────────────

def _replace_text_preserve_format(text_frame, new_text: str) -> None:
    """Overwrite only the first run's text — keeps its font/size/color."""
    paragraphs = text_frame.paragraphs
    if not paragraphs:
        text_frame.text = new_text
        return
    first_para = paragraphs[0]
    if first_para.runs:
        first_para.runs[0].text = new_text
        for run in first_para.runs[1:]:
            run.text = ""
    else:
        first_para.text = new_text


def _replace_bullets_preserve_format(text_frame, bullets: list[str]) -> None:
    """Rewrite as one paragraph per bullet, reusing the first paragraph's XML
    (via deep-copy) as the style template for any additional paragraphs so
    indentation/bullet-glyph formatting carries over, not just run font."""
    if not bullets:
        return
    paragraphs = list(text_frame.paragraphs)
    template_p = paragraphs[0]._p if paragraphs else None

    if template_p is not None:
        while len(text_frame.paragraphs) < len(bullets):
            text_frame._txBody.append(copy.deepcopy(template_p))

    paragraphs = text_frame.paragraphs
    for i, bullet in enumerate(bullets):
        para = paragraphs[i]
        if para.runs:
            para.runs[0].text = bullet
            for run in para.runs[1:]:
                run.text = ""
        else:
            para.text = bullet

    # Fewer bullets than existing paragraphs — blank the leftovers rather
    # than risk XML surgery to delete them.
    for para in paragraphs[len(bullets):]:
        if para.runs:
            para.runs[0].text = ""
            for run in para.runs[1:]:
                run.text = ""
        else:
            para.text = ""


def _fill_table(table, content: dict) -> None:
    headers = content.get("headers") or []
    rows = content.get("rows") or []
    all_rows = ([headers] if headers else []) + rows
    for ri in range(min(len(all_rows), len(table.rows))):
        row_data = all_rows[ri]
        for ci in range(min(len(row_data), len(table.columns))):
            cell = table.cell(ri, ci)
            _replace_text_preserve_format(cell.text_frame, str(row_data[ci]))


def _fill_chart(chart, content: dict) -> None:
    categories = content.get("categories") or []
    series = content.get("series") or []
    if not categories or not series:
        return
    data = CategoryChartData()
    data.categories = categories
    for s in series:
        name = s.get("name", "")
        values = s.get("values", [])
        data.add_series(name, values)
    chart.replace_data(data)


def _fill_slide(slide, sel: dict) -> None:
    title_shape = body_shape = table_shape = chart_shape = None
    for shape in slide.shapes:
        kind = _classify_shape(shape)
        if kind == "title" and title_shape is None:
            title_shape = shape
        elif kind == "body" and body_shape is None:
            body_shape = shape
        elif kind == "table" and table_shape is None:
            table_shape = shape
        elif kind == "chart" and chart_shape is None:
            chart_shape = shape

    title = sel.get("title")
    if title is not None and title_shape is not None:
        _replace_text_preserve_format(title_shape.text_frame, title)

    bullets = sel.get("bullets")
    if bullets and body_shape is not None:
        _replace_bullets_preserve_format(body_shape.text_frame, bullets)

    table = sel.get("table")
    if table and table_shape is not None:
        _fill_table(table_shape.table, table)

    chart_data = sel.get("chart")
    if chart_data and chart_shape is not None:
        _fill_chart(chart_shape.chart, chart_data)

    notes = sel.get("notes")
    if notes is not None:
        slide.notes_slide.notes_text_frame.text = notes


def fill_and_export(pptx_path: str, selections: list[dict]) -> str:
    """Fill selected slides in place, drop the rest, save to a new .pptx.

    Each selection is {slide_index, title?, bullets?, table?, chart?, notes?}
    where slide_index refers to the original 0-based index from
    inspect_template(). Returns the absolute path of the new file.
    """
    prs = Presentation(pptx_path)
    total = len(prs.slides)
    selected_by_index = {
        int(sel["slide_index"]): sel
        for sel in selections
        if 0 <= int(sel.get("slide_index", -1)) < total
    }
    if not selected_by_index:
        raise ValueError("selections must include at least one valid slide_index")

    for idx, sel in selected_by_index.items():
        _fill_slide(prs.slides[idx], sel)

    xml_slides = prs.slides._sldIdLst
    slide_id_list = list(xml_slides)
    for idx, sld in enumerate(slide_id_list):
        if idx not in selected_by_index:
            rId = sld.rId
            xml_slides.remove(sld)
            prs.part.drop_rel(rId)

    _OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    ts = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
    slug = hashlib.sha256(str(pptx_path).encode()).hexdigest()[:8]
    out_path = _OUTPUT_DIR / f"template_fill_{ts}_{slug}.pptx"
    prs.save(str(out_path))
    return str(out_path)
