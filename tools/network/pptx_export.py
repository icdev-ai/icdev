# CUI // SP-CTI
"""ICDEV™ Network Design Canvas -- PPTX topology export (ndc-brg-04).

Renders an NDC topology as an editable PowerPoint deck through the
``tools/viz`` presentation layer (pure-Python ``python-pptx``; fully air-gap,
no network, no npm). It reuses existing viz/slides primitives rather than
hand-writing raw pptx XML:

  * ``tools.slides.pptx_builder`` — deck/slide primitives (``_new_prs``,
    ``_blank``, ``_rect``, ``_box`` and the 16:9 canvas dimensions).
  * ``tools.viz.render_svg.diagram_to_svg`` + ``tools.viz.svg_to_pptx`` —
    the topology graph is rendered as a ``DiagramSpec`` → SVG → native,
    editable PPTX vector shapes.
  * ``tools.viz.render_pptx.add_table`` — a native, editable device-inventory
    table (not a flattened image).

Deck structure (3 slides), each carrying a classification banner derived from
the topology's ``classification`` column (markings are mandatory on every
ICDEV artifact — see ``tools/network/pdf_export.py`` for the sibling PDF
banner behavior this mirrors):

  1. Title     — topology name, classification marking, device/link counts.
  2. Diagram   — nodes/edges as native PPTX shapes (spring/grid layout).
  3. Inventory — device table (Name / Type / Vendor / Model / Mgmt IP).

Public API::

    export_topology_pptx(topo_id, out_path=None) -> bytes | str | None

Returns ``bytes`` (when ``out_path`` is None) or the output path string.
Returns ``None`` when the topology does not exist so callers can answer 404.
Raises :class:`PptxDependencyError` when ``python-pptx`` is unavailable so the
route can answer a clean 501 instead of leaking a stack trace.
"""
from __future__ import annotations

import io
from datetime import datetime, timezone

from tools.logging.icdev_logger import get_logger

logger = get_logger(__name__)

# Node object types that are drawing/annotation artifacts, not real devices.
_NON_DEVICE_TYPES = {"draw-rect", "zone", "boundary", "text-annotation", "label", "group", "link"}

# Classification banner background colors (R, G, B). Mirrors pdf_export intent:
# UNCLASSIFIED/PUBLIC green, CUI red, SECRET+ darker reds.
_CLS_BANNER_COLORS: dict[str, tuple[int, int, int]] = {
    "UNCLASSIFIED": (0x0D, 0x3B, 0x1E),
    "CUI // SP-CTI": (0xB4, 0x1E, 0x1E),
    "CUI": (0xB4, 0x1E, 0x1E),
    "SECRET": (0x8B, 0x00, 0x00),
    "TOP SECRET": (0x5B, 0x0A, 0x0A),
    "TOP SECRET//SCI": (0x3A, 0x0A, 0x4B),
}

_THEME = "midnight_executive"


class PptxDependencyError(RuntimeError):
    """python-pptx (viz PPTX renderer dependency) is unavailable."""


def _require_pptx() -> None:
    """Assert the python-pptx rendering stack is importable.

    Raises :class:`PptxDependencyError` (never a bare ``ImportError``) so callers
    degrade to a clean 501 instead of leaking a trace. Isolated into its own
    function so a missing optional dependency has a single detection point — and
    so tests can simulate it without fighting the import cache.
    """
    try:
        import pptx  # noqa: F401
        import tools.viz.render_pptx  # noqa: F401
        import tools.viz.svg_to_pptx  # noqa: F401
    except ImportError as exc:
        raise PptxDependencyError(f"python-pptx unavailable: {exc}") from exc


def _now_str() -> str:
    return datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")


def _marking(classification: str | None) -> str:
    """Normalize a topology ``classification`` value to a banner marking."""
    c = (classification or "CUI").strip().upper()
    if c in ("", "PUBLIC", "UNCLASSIFIED", "U", "UNCLASS"):
        return "UNCLASSIFIED"
    if c in ("CUI", "CONTROLLED", "FOUO", "CUI//SP-CTI", "CUI // SP-CTI"):
        return "CUI // SP-CTI"
    return c


def _device_nodes(graph: dict) -> list[dict]:
    """Real device nodes only (drawing/annotation shapes filtered out)."""
    out = []
    for n in graph.get("nodes", []):
        if not isinstance(n, dict):
            continue
        if str(n.get("type", "")) in _NON_DEVICE_TYPES:
            continue
        out.append(n)
    return out


def _node_field(node: dict, *keys: str) -> str:
    """First non-empty value across the node, its ``config`` and ``meta``."""
    cfg = node.get("config") or {}
    meta = node.get("meta") or {}
    for src in (node, cfg, meta):
        if not isinstance(src, dict):
            continue
        for k in keys:
            v = src.get(k)
            if v:
                return str(v)
    return ""


# ── Public API ─────────────────────────────────────────────────────────────────

def export_topology_pptx(topo_id: str, out_path: str | None = None):
    """Export an NDC topology as a PPTX deck via the tools/viz layer.

    Args:
        topo_id: Topology primary key.
        out_path: Optional file path. When given, the deck is written there and
            the path is returned; otherwise the deck bytes are returned.

    Returns:
        ``bytes`` or the ``out_path`` string, or ``None`` if the topology is
        missing.

    Raises:
        PptxDependencyError: python-pptx is not importable.
    """
    # Detect a missing optional dependency up front (monkeypatchable seam), then
    # import the pptx-dependent stack.
    _require_pptx()
    try:
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt  # noqa: F401  (Pt kept for parity/future use)
        from tools.slides.pptx_builder import CW, H, LM, W, _blank, _box, _new_prs, _rect
        from tools.viz.render_pptx import add_table
        from tools.viz.render_svg import diagram_to_svg
        from tools.viz.spec import DiagramSpec, TableSpec
        from tools.viz import svg_to_pptx
    except ImportError as exc:  # python-pptx (or a viz renderer dep) missing
        raise PptxDependencyError(f"python-pptx unavailable: {exc}") from exc

    from pptx.dml.color import RGBColor

    # ── Load topology read-only (reuses the parsed-graph cache) ─────────────────
    from tools.network.blueprint_helpers import _row_to_dict, get_parsed_graph
    from tools.network.db.init_db import get_connection
    from tools.db.storage import sql_placeholder

    conn = get_connection()
    try:
        ph = sql_placeholder(conn)
        row = conn.execute(
            f"SELECT * FROM topologies WHERE id={ph}", (topo_id,)
        ).fetchone()
        if row is None:
            return None
        topo = _row_to_dict(row)
        graph = get_parsed_graph(conn, topo_id) or {"nodes": [], "edges": []}
    finally:
        try:
            conn.close()
        except Exception:
            pass

    name = str(topo.get("name") or topo_id)
    marking = _marking(topo.get("classification"))
    devices = _device_nodes(graph)
    edges = [e for e in graph.get("edges", []) if isinstance(e, dict)]

    # ── Build the deck ──────────────────────────────────────────────────────────
    prs = _new_prs()
    banner_rgb = RGBColor(*_CLS_BANNER_COLORS.get(marking, _CLS_BANNER_COLORS["CUI"]))
    white = RGBColor(0xFF, 0xFF, 0xFF)
    band_h = Inches(0.24)

    def _cls_banner(slide) -> None:
        """Top + bottom classification bands on every slide (mandatory markings)."""
        _rect(slide, 0, 0, W, band_h, banner_rgb)
        _box(slide, 0, Inches(0.02), W, Inches(0.20), marking,
             size=9, bold=True, color=white, align=PP_ALIGN.CENTER)
        _rect(slide, 0, H - band_h, W, band_h, banner_rgb)
        _box(slide, 0, H - Inches(0.22), W, Inches(0.20), marking,
             size=9, bold=True, color=white, align=PP_ALIGN.CENTER)

    from tools.slides.constants import THEME_PALETTES, DEFAULT_THEME
    palette = THEME_PALETTES.get(_THEME, THEME_PALETTES[DEFAULT_THEME])

    def _rgb(key: str) -> RGBColor:
        r, g, b = palette[key]
        return RGBColor(r, g, b)

    # Slide 1 — Title
    s1 = _blank(prs)
    s1.background.fill.solid()
    s1.background.fill.fore_color.rgb = _rgb("bg")
    _cls_banner(s1)
    _box(s1, LM, Inches(2.3), CW, Inches(1.4), name,
         size=40, bold=True, color=_rgb("accent"), align=PP_ALIGN.CENTER)
    _rect(s1, Inches(3.8), Inches(3.75), Inches(5.73), Inches(0.03), _rgb("accent"))
    _box(s1, LM, Inches(3.95), CW, Inches(0.5),
         "ICDEV™ Network Design Canvas — Topology Export",
         size=15, color=_rgb("text"), align=PP_ALIGN.CENTER)
    _box(s1, LM, Inches(4.55), CW, Inches(0.9),
         f"Classification: {marking}\n"
         f"Devices: {len(devices)}    Links: {len(edges)}\n"
         f"Generated: {_now_str()}",
         size=13, color=_rgb("subtext"), align=PP_ALIGN.CENTER)

    # Slide 2 — Topology diagram (native, editable vector shapes)
    s2 = _blank(prs)
    s2.background.fill.solid()
    s2.background.fill.fore_color.rgb = _rgb("bg")
    _cls_banner(s2)
    _rect(s2, 0, band_h, W, Inches(0.6), _rgb("dark"))
    _box(s2, Inches(0.24), band_h + Inches(0.08), CW, Inches(0.45),
         "Network Topology", size=22, bold=True, color=_rgb("accent"))
    _MAX_DIAGRAM_NODES = 40
    diagram_devices = devices[:_MAX_DIAGRAM_NODES]
    id_set = {str(n.get("id", n.get("label", ""))) for n in diagram_devices}
    spec = DiagramSpec(
        title="",
        nodes=[
            {
                "id": str(n.get("id", n.get("label", f"n{i}"))),
                "label": str(n.get("label") or n.get("id") or f"node{i}"),
                "type": str(n.get("type", "")),
            }
            for i, n in enumerate(diagram_devices)
        ],
        edges=[
            {"source": str(e.get("source", "")), "target": str(e.get("target", ""))}
            for e in edges
            if str(e.get("source", "")) in id_set and str(e.get("target", "")) in id_set
        ],
        layout="spring",
    )
    rendered = False
    if spec.nodes:
        try:
            svg = diagram_to_svg(spec, _THEME)
            shapes = svg_to_pptx.render_svg_into_slide(
                s2, svg, LM, Inches(1.05), CW, Inches(5.55)
            )
            rendered = bool(shapes)
        except Exception as exc:  # never leak a raster/vector failure to the caller
            logger.warning("topology diagram render failed for %s: %s", topo_id, exc)
            rendered = False
    if not rendered:
        _box(s2, LM, Inches(3.0), CW, Inches(0.6),
             "(No devices to diagram)", size=16,
             color=_rgb("subtext"), align=PP_ALIGN.CENTER)
    if len(devices) > _MAX_DIAGRAM_NODES:
        _box(s2, LM, H - Inches(0.5), CW, Inches(0.24),
             f"Diagram shows first {_MAX_DIAGRAM_NODES} of {len(devices)} devices — "
             f"see inventory slide for the full list.",
             size=9, italic=True, color=_rgb("subtext"), align=PP_ALIGN.CENTER)

    # Slide 3 — Device inventory (native, editable table)
    s3 = _blank(prs)
    s3.background.fill.solid()
    s3.background.fill.fore_color.rgb = _rgb("bg")
    _cls_banner(s3)
    _rect(s3, 0, band_h, W, Inches(0.6), _rgb("dark"))
    _box(s3, Inches(0.24), band_h + Inches(0.08), CW, Inches(0.45),
         "Device Inventory", size=22, bold=True, color=_rgb("accent"))
    _MAX_TABLE_ROWS = 24
    table_rows = [
        [
            (n.get("label") or n.get("id") or "")[:36],
            str(n.get("type", ""))[:20],
            _node_field(n, "vendor", "manufacturer")[:18],
            _node_field(n, "model")[:22],
            _node_field(n, "ip", "mgmt_ip", "management_ip")[:20],
        ]
        for n in devices[:_MAX_TABLE_ROWS]
    ]
    tspec = TableSpec(
        title="Device Inventory",
        headers=["Device Name", "Type", "Vendor", "Model", "Mgmt IP"],
        rows=table_rows,
    )
    try:
        add_table(s3, tspec, LM, Inches(1.05), CW, Inches(5.4), _THEME)
    except Exception as exc:
        logger.warning("inventory table render failed for %s: %s", topo_id, exc)
        _box(s3, LM, Inches(3.0), CW, Inches(0.6),
             "(No device inventory available)", size=16,
             color=_rgb("subtext"), align=PP_ALIGN.CENTER)
    if len(devices) > _MAX_TABLE_ROWS:
        _box(s3, LM, H - Inches(0.5), CW, Inches(0.24),
             f"... and {len(devices) - _MAX_TABLE_ROWS} more devices",
             size=9, italic=True, color=_rgb("subtext"), align=PP_ALIGN.CENTER)

    logger.info(
        "NDC PPTX export: topo=%s classification=%s devices=%d links=%d",
        topo_id, marking, len(devices), len(edges),
    )

    # ── Emit ────────────────────────────────────────────────────────────────────
    if out_path:
        import pathlib

        pathlib.Path(out_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))
        return str(out_path)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
