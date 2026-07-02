# CUI // SP-CTI
"""SVG -> native PPTX shape converter.

Deterministic, dependency-free (stdlib xml.etree.ElementTree) conversion of a
bounded SVG subset into native python-pptx DrawingML shapes (freeform
polygons/polylines + textboxes) so vector art lands in a .pptx as real,
clickable/editable shapes instead of a flattened raster picture.

Supported SVG: <svg viewBox>, <g transform="translate()|scale()|matrix()|rotate()">,
<rect>, <circle>, <ellipse>, <line>, <polyline>, <polygon>,
<path> (M/L/H/V/C/S/Q/T/Z, absolute + relative; A/arc degrades to a straight
line to the arc endpoint), <text>/<tspan>. Curves are flattened to line
segments because python-pptx's FreeformBuilder only supports straight edges.
Unsupported constructs (gradients, masks, CSS classes, <defs> contents) fall
back to a flat color or are skipped rather than raising.
"""
from __future__ import annotations

import math
import re
import xml.etree.ElementTree as ET
from dataclasses import dataclass
from pathlib import Path

from defusedxml.ElementTree import fromstring as _safe_xml_fromstring

from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor

_BEZIER_STEPS = 16

_DEFAULT_FILL = RGBColor(0x80, 0x80, 0x80)

_NAMED_COLORS: dict[str, tuple[int, int, int] | None] = {
    "black": (0, 0, 0), "white": (255, 255, 255), "red": (255, 0, 0),
    "green": (0, 128, 0), "blue": (0, 0, 255), "gray": (128, 128, 128),
    "grey": (128, 128, 128), "navy": (0, 0, 128), "gold": (255, 215, 0),
    "silver": (192, 192, 192), "orange": (255, 165, 0), "yellow": (255, 255, 0),
    "purple": (128, 0, 128), "teal": (0, 128, 128), "none": None,
    "transparent": None,
}

Matrix = tuple[float, float, float, float, float, float]
IDENTITY: Matrix = (1.0, 0.0, 0.0, 1.0, 0.0, 0.0)

_TRANSFORM_RE = re.compile(r"(translate|scale|matrix|rotate)\s*\(([^)]*)\)")
_PATH_TOKEN_RE = re.compile(r"([MmLlHhVvCcSsQqTtAaZz])|(-?\d*\.?\d+(?:[eE][-+]?\d+)?)")


# ── Geometry helpers ──────────────────────────────────────────────────────────

def _mat_mul(m1: Matrix, m2: Matrix) -> Matrix:
    """Compose m1 . m2 — applying the result to a point equals m1(m2(point))."""
    a1, b1, c1, d1, e1, f1 = m1
    a2, b2, c2, d2, e2, f2 = m2
    return (
        a1 * a2 + c1 * b2,
        b1 * a2 + d1 * b2,
        a1 * c2 + c1 * d2,
        b1 * c2 + d1 * d2,
        a1 * e2 + c1 * f2 + e1,
        b1 * e2 + d1 * f2 + f1,
    )


def _apply(m: Matrix, x: float, y: float) -> tuple[float, float]:
    a, b, c, d, e, f = m
    return (a * x + c * y + e, b * x + d * y + f)


def _parse_transform(attr: str | None) -> Matrix:
    if not attr:
        return IDENTITY
    m = IDENTITY
    for kind, args in _TRANSFORM_RE.findall(attr):
        nums = [float(v) for v in re.split(r"[,\s]+", args.strip()) if v]
        if kind == "translate" and nums:
            tx = nums[0]
            ty = nums[1] if len(nums) > 1 else 0.0
            m = _mat_mul(m, (1.0, 0.0, 0.0, 1.0, tx, ty))
        elif kind == "scale" and nums:
            sx = nums[0]
            sy = nums[1] if len(nums) > 1 else sx
            m = _mat_mul(m, (sx, 0.0, 0.0, sy, 0.0, 0.0))
        elif kind == "matrix" and len(nums) == 6:
            m = _mat_mul(m, tuple(nums))  # type: ignore[arg-type]
        elif kind == "rotate" and nums:
            rad = math.radians(nums[0])
            cos_, sin_ = math.cos(rad), math.sin(rad)
            rot: Matrix = (cos_, sin_, -sin_, cos_, 0.0, 0.0)
            if len(nums) >= 3:
                cx, cy = nums[1], nums[2]
                m = _mat_mul(m, (1.0, 0.0, 0.0, 1.0, cx, cy))
                m = _mat_mul(m, rot)
                m = _mat_mul(m, (1.0, 0.0, 0.0, 1.0, -cx, -cy))
            else:
                m = _mat_mul(m, rot)
    return m


def _ellipse_points(cx: float, cy: float, rx: float, ry: float, segments: int = 48) -> list[tuple[float, float]]:
    return [
        (cx + rx * math.cos(2 * math.pi * i / segments), cy + ry * math.sin(2 * math.pi * i / segments))
        for i in range(segments)
    ]


def _cubic_points(p0, p1, p2, p3, steps: int = _BEZIER_STEPS) -> list[tuple[float, float]]:
    pts = []
    for i in range(1, steps + 1):
        t = i / steps
        mt = 1 - t
        x = mt ** 3 * p0[0] + 3 * mt ** 2 * t * p1[0] + 3 * mt * t ** 2 * p2[0] + t ** 3 * p3[0]
        y = mt ** 3 * p0[1] + 3 * mt ** 2 * t * p1[1] + 3 * mt * t ** 2 * p2[1] + t ** 3 * p3[1]
        pts.append((x, y))
    return pts


def _quad_points(p0, p1, p2, steps: int = _BEZIER_STEPS) -> list[tuple[float, float]]:
    pts = []
    for i in range(1, steps + 1):
        t = i / steps
        mt = 1 - t
        x = mt ** 2 * p0[0] + 2 * mt * t * p1[0] + t ** 2 * p2[0]
        y = mt ** 2 * p0[1] + 2 * mt * t * p1[1] + t ** 2 * p2[1]
        pts.append((x, y))
    return pts


# ── Path ("d" attribute) parsing ─────────────────────────────────────────────

def _tokenize_path(d: str) -> list[str | float]:
    tokens: list[str | float] = []
    for cmd_tok, num_tok in _PATH_TOKEN_RE.findall(d):
        if cmd_tok:
            tokens.append(cmd_tok)
        elif num_tok:
            tokens.append(float(num_tok))
    return tokens


def _parse_path(d: str) -> list[tuple[list[tuple[float, float]], bool]]:
    """Return a list of (points, closed) sub-paths, curves flattened to lines."""
    tokens = _tokenize_path(d)
    subpaths: list[tuple[list[tuple[float, float]], bool]] = []
    cur: list[tuple[float, float]] = []
    x = y = start_x = start_y = 0.0
    prev_ctrl: tuple[float, float] | None = None
    prev_cmd = ""
    cmd: str | None = None
    i, n = 0, len(tokens)

    def take(k: int) -> list[float]:
        nonlocal i
        vals = tokens[i:i + k]
        i += k
        return vals if all(isinstance(v, float) for v in vals) and len(vals) == k else []

    def flush(closed: bool) -> None:
        nonlocal cur
        if len(cur) >= 2:
            subpaths.append((cur, closed))
        cur = []

    while i < n:
        if isinstance(tokens[i], str):
            cmd = tokens[i]  # type: ignore[assignment]
            i += 1
        elif cmd is None:
            break
        else:
            cmd = "L" if cmd == "M" else ("l" if cmd == "m" else cmd)

        upper = cmd.upper()
        rel = cmd.islower()

        if upper == "M":
            vals = take(2)
            if not vals:
                break
            nx, ny = vals
            if rel:
                nx += x; ny += y
            if cur:
                flush(False)
            x = start_x = nx
            y = start_y = ny
            cur = [(x, y)]
            prev_ctrl = None
        elif upper == "L":
            vals = take(2)
            if not vals:
                break
            nx, ny = vals
            if rel:
                nx += x; ny += y
            x, y = nx, ny
            cur.append((x, y))
            prev_ctrl = None
        elif upper == "H":
            vals = take(1)
            if not vals:
                break
            nx = vals[0] + x if rel else vals[0]
            x = nx
            cur.append((x, y))
            prev_ctrl = None
        elif upper == "V":
            vals = take(1)
            if not vals:
                break
            ny = vals[0] + y if rel else vals[0]
            y = ny
            cur.append((x, y))
            prev_ctrl = None
        elif upper == "C":
            vals = take(6)
            if not vals:
                break
            x1, y1, x2, y2, nx, ny = vals
            if rel:
                x1 += x; y1 += y; x2 += x; y2 += y; nx += x; ny += y
            cur.extend(_cubic_points((x, y), (x1, y1), (x2, y2), (nx, ny)))
            prev_ctrl = (x2, y2)
            x, y = nx, ny
        elif upper == "S":
            vals = take(4)
            if not vals:
                break
            x2, y2, nx, ny = vals
            if rel:
                x2 += x; y2 += y; nx += x; ny += y
            if prev_ctrl and prev_cmd.upper() in ("C", "S"):
                x1, y1 = 2 * x - prev_ctrl[0], 2 * y - prev_ctrl[1]
            else:
                x1, y1 = x, y
            cur.extend(_cubic_points((x, y), (x1, y1), (x2, y2), (nx, ny)))
            prev_ctrl = (x2, y2)
            x, y = nx, ny
        elif upper == "Q":
            vals = take(4)
            if not vals:
                break
            x1, y1, nx, ny = vals
            if rel:
                x1 += x; y1 += y; nx += x; ny += y
            cur.extend(_quad_points((x, y), (x1, y1), (nx, ny)))
            prev_ctrl = (x1, y1)
            x, y = nx, ny
        elif upper == "T":
            vals = take(2)
            if not vals:
                break
            nx, ny = vals
            if rel:
                nx += x; ny += y
            if prev_ctrl and prev_cmd.upper() in ("Q", "T"):
                x1, y1 = 2 * x - prev_ctrl[0], 2 * y - prev_ctrl[1]
            else:
                x1, y1 = x, y
            cur.extend(_quad_points((x, y), (x1, y1), (nx, ny)))
            prev_ctrl = (x1, y1)
            x, y = nx, ny
        elif upper == "A":
            # Arcs are unsupported — degrade to a straight line to the endpoint.
            vals = take(7)
            if not vals:
                break
            nx = vals[5] + x if rel else vals[5]
            ny = vals[6] + y if rel else vals[6]
            x, y = nx, ny
            cur.append((x, y))
            prev_ctrl = None
        elif upper == "Z":
            x, y = start_x, start_y
            flush(True)
            cur = [(x, y)]
            prev_ctrl = None
        else:
            break

        prev_cmd = cmd

    flush(False)
    return subpaths


# ── Style parsing ─────────────────────────────────────────────────────────────

def _num(value, default: float = 0.0) -> float:
    if value is None:
        return default
    if isinstance(value, (int, float)):
        return float(value)
    match = re.match(r"-?\d*\.?\d+(?:[eE][-+]?\d+)?", str(value).strip())
    return float(match.group(0)) if match else default


def _parse_color(value: str | None) -> RGBColor | None:
    if not value:
        return None
    value = value.strip().lower()
    if value in ("none", "transparent"):
        return None
    if value.startswith("#"):
        hexv = value[1:]
        if len(hexv) == 3:
            hexv = "".join(c * 2 for c in hexv)
        if len(hexv) >= 6:
            try:
                return RGBColor(int(hexv[0:2], 16), int(hexv[2:4], 16), int(hexv[4:6], 16))
            except ValueError:
                return _DEFAULT_FILL
        return _DEFAULT_FILL
    m = re.match(r"rgb\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)", value)
    if m:
        return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    if value in _NAMED_COLORS:
        rgb = _NAMED_COLORS[value]
        return RGBColor(*rgb) if rgb else None
    return _DEFAULT_FILL  # unsupported (url(#gradient), css var, etc.)


def _style(el: ET.Element) -> dict[str, str]:
    style: dict[str, str] = {}
    for key in ("fill", "stroke", "stroke-width", "opacity", "fill-opacity", "font-size", "font-weight", "text-anchor"):
        v = el.get(key)
        if v is not None:
            style[key] = v
    inline = el.get("style")
    if inline:
        for part in inline.split(";"):
            if ":" in part:
                k, v = part.split(":", 1)
                style[k.strip()] = v.strip()
    return style


def _tag(el: ET.Element) -> str:
    t = el.tag
    return t.split("}", 1)[1] if t.startswith("{") else t


# ── Draw-op intermediate representation ───────────────────────────────────────

@dataclass
class _DrawOp:
    kind: str  # "poly" | "text"
    points: list[tuple[float, float]] | None = None
    closed: bool = False
    fill: RGBColor | None = None
    stroke: RGBColor | None = None
    stroke_width_pt: float = 1.0
    text: str | None = None
    x: float = 0.0
    y: float = 0.0
    font_size_pt: float = 18.0
    text_anchor: str = "start"
    bold: bool = False


def _walk(el: ET.Element, matrix: Matrix, ops: list[_DrawOp]) -> None:
    tag = _tag(el)

    if tag in ("svg", "g"):
        local_m = _mat_mul(matrix, _parse_transform(el.get("transform")))
        for child in el:
            _walk(child, local_m, ops)
        return

    local_m = _mat_mul(matrix, _parse_transform(el.get("transform")))
    style = _style(el)
    fill = _parse_color(style.get("fill", "#808080"))
    stroke = _parse_color(style.get("stroke"))
    stroke_w = _num(style.get("stroke-width"), 1.0)

    if tag == "rect":
        x, y = _num(el.get("x")), _num(el.get("y"))
        w, h = _num(el.get("width")), _num(el.get("height"))
        if w <= 0 or h <= 0:
            return
        corners = [(x, y), (x + w, y), (x + w, y + h), (x, y + h)]
        ops.append(_DrawOp("poly", [_apply(local_m, px, py) for px, py in corners], True, fill, stroke, stroke_w))

    elif tag == "circle":
        cx, cy, r = _num(el.get("cx")), _num(el.get("cy")), _num(el.get("r"))
        if r <= 0:
            return
        pts = [_apply(local_m, px, py) for px, py in _ellipse_points(cx, cy, r, r)]
        ops.append(_DrawOp("poly", pts, True, fill, stroke, stroke_w))

    elif tag == "ellipse":
        cx, cy = _num(el.get("cx")), _num(el.get("cy"))
        rx, ry = _num(el.get("rx")), _num(el.get("ry"))
        if rx <= 0 or ry <= 0:
            return
        pts = [_apply(local_m, px, py) for px, py in _ellipse_points(cx, cy, rx, ry)]
        ops.append(_DrawOp("poly", pts, True, fill, stroke, stroke_w))

    elif tag == "line":
        x1, y1 = _num(el.get("x1")), _num(el.get("y1"))
        x2, y2 = _num(el.get("x2")), _num(el.get("y2"))
        pts = [_apply(local_m, x1, y1), _apply(local_m, x2, y2)]
        ops.append(_DrawOp("poly", pts, False, None, stroke or fill, stroke_w))

    elif tag in ("polyline", "polygon"):
        raw = el.get("points", "")
        nums = [float(v) for v in re.split(r"[,\s]+", raw.strip()) if v]
        coords = list(zip(nums[0::2], nums[1::2]))
        if len(coords) < 2:
            return
        pts = [_apply(local_m, px, py) for px, py in coords]
        ops.append(_DrawOp("poly", pts, tag == "polygon", fill, stroke, stroke_w))

    elif tag == "path":
        for sub_pts, closed in _parse_path(el.get("d", "")):
            pts = [_apply(local_m, px, py) for px, py in sub_pts]
            if len(pts) >= 2:
                ops.append(_DrawOp("poly", pts, closed, fill, stroke, stroke_w))

    elif tag == "text":
        text = "".join(el.itertext()).strip()
        if not text:
            return
        px, py = _apply(local_m, _num(el.get("x")), _num(el.get("y")))
        fsize = _num(style.get("font-size"), 18.0)
        weight = (style.get("font-weight") or "").lower()
        ops.append(_DrawOp(
            "text", fill=fill or _DEFAULT_FILL, text=text, x=px, y=py,
            font_size_pt=fsize, text_anchor=style.get("text-anchor", "start"),
            bold=weight in ("bold", "700", "800", "900"),
        ))

    # Unknown tags (defs, clipPath, mask, linearGradient, pattern, ...) are
    # intentionally not recursed into — their children are definitions, not
    # visible geometry.


# ── Shape emission ─────────────────────────────────────────────────────────────

def _emit_poly(slide, op: _DrawOp):
    pts = op.points
    if not pts or len(pts) < 2:
        return None
    start_x, start_y = pts[0]
    fb = slide.shapes.build_freeform(start_x, start_y, scale=1.0)
    fb.add_line_segments(pts[1:], close=op.closed)
    shape = fb.convert_to_shape()
    if op.fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = op.fill
    else:
        shape.fill.background()
    if op.stroke is not None:
        shape.line.color.rgb = op.stroke
        shape.line.width = Pt(max(op.stroke_width_pt, 0.25))
    else:
        shape.line.fill.background()
    return shape


def _emit_text(slide, op: _DrawOp):
    if not op.text:
        return None
    font_pt = max(op.font_size_pt, 6.0)
    box_w = max(int(Pt(font_pt * 0.62 * len(op.text))), int(Pt(20)))
    box_h = int(Pt(font_pt * 1.6))
    left = int(op.x)
    top = int(op.y) - int(Pt(font_pt))
    if op.text_anchor == "middle":
        left -= box_w // 2
    elif op.text_anchor == "end":
        left -= box_w
    box = slide.shapes.add_textbox(Emu(max(left, 0)), Emu(max(top, 0)), Emu(box_w), Emu(box_h))
    tf = box.text_frame
    tf.word_wrap = False
    run = tf.paragraphs[0].add_run()
    run.text = op.text
    run.font.size = Pt(font_pt)
    run.font.bold = op.bold
    if op.fill is not None:
        run.font.color.rgb = op.fill
    return box


# ── Public API ──────────────────────────────────────────────────────────────

def render_svg_into_slide(slide, svg_source: str, left, top, width, height) -> list:
    """Parse `svg_source` and emit native shapes into `slide` at the given EMU box.

    Returns the list of created python-pptx shapes (may be empty on an
    unparseable/empty SVG — this never raises on malformed geometry, only on
    malformed XML, matching how callers already wrap graphics generation in
    try/except).
    """
    root = _safe_xml_fromstring(svg_source)
    vb = root.get("viewBox")
    if vb:
        parts = [float(v) for v in re.split(r"[,\s]+", vb.strip()) if v]
        vx, vy, vw, vh = (parts + [0.0, 0.0, 100.0, 100.0])[:4]
    else:
        vx, vy = 0.0, 0.0
        vw, vh = _num(root.get("width"), 100.0), _num(root.get("height"), 100.0)
    if vw <= 0 or vh <= 0:
        return []

    left, top, width, height = int(left), int(top), int(width), int(height)
    scale_x = width / vw
    scale_y = height / vh
    base: Matrix = (scale_x, 0.0, 0.0, scale_y, left - vx * scale_x, top - vy * scale_y)

    ops: list[_DrawOp] = []
    _walk(root, base, ops)

    shapes = []
    for op in ops:
        shape = _emit_poly(slide, op) if op.kind == "poly" else _emit_text(slide, op)
        if shape is not None:
            shapes.append(shape)
    return shapes


def embed_svg_file(slide, svg_path: str, left, top, width, height) -> list:
    """Read `svg_path` and render it natively into `slide`. See `render_svg_into_slide`."""
    svg_source = Path(svg_path).read_text(encoding="utf-8")
    return render_svg_into_slide(slide, svg_source, left, top, width, height)
