"""
ICDEV™ Executive Presentation Generator
Produces: tools/presentations/output/ICDEV_Executive_Deck.pptx
Theme: Dark-navy Gov/DoD executive style
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0A, 0x16, 0x28)   # slide background
GOLD   = RGBColor(0xC8, 0xA9, 0x51)   # accents, titles
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)   # body text
LGRAY  = RGBColor(0xE0, 0xE6, 0xF0)   # sub-text
MGRAY  = RGBColor(0x1E, 0x3A, 0x5F)   # divider / card backgrounds
DGOLD  = RGBColor(0xA0, 0x82, 0x35)   # dark gold for hover-contrast
GREEN  = RGBColor(0x2E, 0xCC, 0x71)   # success / positive
RED    = RGBColor(0xE7, 0x4C, 0x3C)   # alert
BLUE   = RGBColor(0x2E, 0x86, 0xC1)   # info

W  = Inches(13.33)
H  = Inches(7.50)
LM = Inches(0.55)   # left margin
RM = Inches(0.55)   # right margin
TM = Inches(0.40)   # top margin
CW = W - LM - RM    # content width


# ── Helpers ────────────────────────────────────────────────────────────────

def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])   # blank layout


def _bg(slide, color: RGBColor = NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, left, top, width, height,
         text="", size=18, bold=False, italic=False,
         color=WHITE, align=PP_ALIGN.LEFT,
         bg_color=None, word_wrap=True):
    shape = slide.shapes.add_textbox(left, top, width, height)
    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
    tf = shape.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf


def _para(tf, text, size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def _rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _gold_bar(slide, top=Inches(1.55), height=Inches(0.04)):
    _rect(slide, LM, top, CW, height, GOLD)


def _slide_num(slide, n: int):
    _box(slide, W - Inches(0.8), H - Inches(0.35), Inches(0.7), Inches(0.3),
         str(n), size=9, color=MGRAY, align=PP_ALIGN.RIGHT)


def _footer(slide, text="ICDEV™  |  FORGE Framework  |  IL4+ Ready  |  Apache 2.0"):
    _box(slide, LM, H - Inches(0.32), CW, Inches(0.28),
         text, size=8, color=MGRAY)


def _title_area(slide, title, subtitle=""):
    _box(slide, LM, TM, CW, Inches(1.05),
         title, size=32, bold=True, color=GOLD)
    _gold_bar(slide)
    if subtitle:
        _box(slide, LM, Inches(1.65), CW, Inches(0.45),
             subtitle, size=14, color=LGRAY, italic=True)


def _card(slide, left, top, width, height, heading, body, head_size=13, body_size=11):
    _rect(slide, left, top, width, height, MGRAY, GOLD)
    inner_pad = Inches(0.12)
    tf = _box(slide, left + inner_pad, top + inner_pad,
              width - inner_pad * 2, Inches(0.34),
              heading, size=head_size, bold=True, color=GOLD)
    _box(slide, left + inner_pad, top + Inches(0.38),
         width - inner_pad * 2, height - Inches(0.50),
         body, size=body_size, color=LGRAY, word_wrap=True)


def _bullet_block(slide, items, left, top, width, height,
                  bullet="▸", size=14, gap=Inches(0.38)):
    y = top
    for item in items:
        _box(slide, left, y, width, gap,
             f"{bullet}  {item}", size=size, color=WHITE, word_wrap=True)
        y += gap


# ── Individual Slide Builders ───────────────────────────────────────────────

def slide_01_cover(prs):
    s = _blank(prs); _bg(s)
    # Thick gold bar top
    _rect(s, 0, 0, W, Inches(0.12), GOLD)
    # Main title
    _box(s, LM, Inches(1.0), CW, Inches(1.8),
         "ICDEV™", size=80, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    _box(s, LM, Inches(2.7), CW, Inches(0.7),
         "The AI Lab That Builds AI Labs", size=28, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
    _box(s, LM, Inches(3.4), CW, Inches(0.5),
         "Adopt vs. Build  ·  Executive Briefing", size=16,
         color=LGRAY, align=PP_ALIGN.CENTER, italic=True)
    # Divider line
    _rect(s, Inches(3.5), Inches(4.05), Inches(6.33), Inches(0.03), GOLD)
    # Stats strip
    stats = ["42 Compliance Frameworks", "15 AI Agents", "12 Design Canvases",
             "530+ Deterministic Tools", "IL4+ Ready"]
    col_w = CW / len(stats)
    for i, stat in enumerate(stats):
        _box(s, LM + col_w * i, Inches(4.15), col_w, Inches(0.5),
             stat, size=11, color=LGRAY, align=PP_ALIGN.CENTER)
    _box(s, LM, Inches(4.7), CW, Inches(0.35),
         "May 2026  ·  CONFIDENTIAL", size=10, color=MGRAY, align=PP_ALIGN.CENTER)
    # Thick gold bar bottom
    _rect(s, 0, H - Inches(0.12), W, Inches(0.12), GOLD)


def slide_02_problem(prs):
    s = _blank(prs); _bg(s)
    _title_area(s, "The Problem Every Government Program Faces")
    _footer(s); _slide_num(s, 2)

    items = [
        ("⏱  12–18 Months",
         "Programs spend over a year and millions of dollars before writing a single line of production code — "
         "just to get an Authority to Operate (ATO)."),
        ("🔒  Compliance Discovered Late",
         "Security and compliance are identified after development — forcing costly redesigns, missed deadlines, "
         "and emergency POAM entries that follow the system for its entire lifecycle."),
        ("🧠  Expert Bottlenecks",
         "One ISSO. One security architect. One CI/CD engineer. Any one of them on leave or departed "
         "can stall an entire program. Knowledge lives in people, not in systems."),
    ]
    y = Inches(2.05)
    for icon_label, body in items:
        _rect(s, LM, y, CW, Inches(1.35), MGRAY, GOLD)
        _box(s, LM + Inches(0.18), y + Inches(0.1), Inches(2.8), Inches(0.45),
             icon_label, size=15, bold=True, color=GOLD)
        _box(s, LM + Inches(0.18), y + Inches(0.52), CW - Inches(0.36), Inches(0.75),
             body, size=12, color=LGRAY, word_wrap=True)
        y += Inches(1.50)

    _box(s, LM, Inches(6.65), CW, Inches(0.40),
         "These are not edge cases. They are the default experience for government software programs today.",
         size=11, italic=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_03_question(prs):
    s = _blank(prs); _bg(s)
    _footer(s); _slide_num(s, 3)

    _box(s, LM, TM, CW, Inches(0.8),
         "The Decision", size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    _gold_bar(s, Inches(0.85), Inches(0.04))

    _box(s, LM, Inches(1.1), CW, Inches(0.9),
         "Build from scratch — or adopt a platform that's already solved it?",
         size=22, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

    col_w = Inches(5.8)
    gap   = Inches(0.3)
    left1 = LM
    left2 = LM + col_w + gap

    # Build column
    _rect(s, left1, Inches(2.1), col_w, Inches(4.5), MGRAY, RED)
    _box(s, left1 + Inches(0.15), Inches(2.2), col_w - Inches(0.3), Inches(0.5),
         "Build From Scratch", size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)
    build_items = [
        "12–18 months before first ATO",
        "Compliance mapped after the code is written",
        "Expert knowledge locked in individuals",
        "Each new system = new ATO from zero",
        "No crosswalk between frameworks",
        "Legacy modernization effectively impossible",
        "No AI governance built in",
    ]
    y = Inches(2.75)
    for item in build_items:
        _box(s, left1 + Inches(0.15), y, col_w - Inches(0.3), Inches(0.38),
             f"✗  {item}", size=11, color=LGRAY)
        y += Inches(0.42)

    # Adopt column
    _rect(s, left2, Inches(2.1), col_w, Inches(4.5), MGRAY, GREEN)
    _box(s, left2 + Inches(0.15), Inches(2.2), col_w - Inches(0.3), Inches(0.5),
         "Adopt ICDEV™", size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    adopt_items = [
        "ATO in 1–3 months",
        "Compliance automated on first build",
        "15 AI agents replace expert dependency",
        "42-framework crosswalk — implement once",
        "Living applications that self-improve",
        "7Rs modernization pipeline built in",
        "OMB M-25-21 · NIST AI RMF baked in",
    ]
    y = Inches(2.75)
    for item in adopt_items:
        _box(s, left2 + Inches(0.15), y, col_w - Inches(0.3), Inches(0.38),
             f"✓  {item}", size=11, color=LGRAY)
        y += Inches(0.42)


def slide_04_what_is_icdev(prs):
    s = _blank(prs); _bg(s)
    _title_area(s, "What is ICDEV™?", "Intelligent Certified Development — a meta-builder for government software")
    _footer(s); _slide_num(s, 4)

    _box(s, LM, Inches(2.15), CW, Inches(0.55),
         "Plain English Requirements  →  ATO-Ready System",
         size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    _box(s, LM, Inches(2.7), CW, Inches(0.45),
         '"One developer built ICDEV™.  Imagine what your team could do."',
         size=14, italic=True, color=LGRAY, align=PP_ALIGN.CENTER)

    stats = [
        ("42", "Compliance\nFrameworks"),
        ("15", "AI Agents"),
        ("12", "Design\nCanvases"),
        ("530+", "Deterministic\nTools"),
        ("588", "Database\nTables"),
        ("6", "Supported\nLanguages"),
    ]
    box_w = CW / len(stats) - Inches(0.1)
    y = Inches(3.25)
    for i, (num, label) in enumerate(stats):
        x = LM + i * (box_w + Inches(0.1))
        _rect(s, x, y, box_w, Inches(2.5), MGRAY, GOLD)
        _box(s, x, y + Inches(0.25), box_w, Inches(0.9),
             num, size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        _box(s, x, y + Inches(1.15), box_w, Inches(0.8),
             label, size=11, color=LGRAY, align=PP_ALIGN.CENTER)

    _box(s, LM, Inches(6.1), CW, Inches(0.45),
         "FORGE Framework  ·  ANVIL Workflow  ·  Apache 2.0  ·  AWS GovCloud  ·  IL4+",
         size=11, color=MGRAY, align=PP_ALIGN.CENTER)


def slide_05_forge(prs):
    s = _blank(prs); _bg(s)
    _title_area(s, "The FORGE Framework", "Six layers that separate probabilistic AI reasoning from deterministic execution")
    _footer(s); _slide_num(s, 5)

    layers = [
        ("Goals",         "Process definitions\nWhat to achieve,\nwhich tools, outputs"),
        ("Orchestration", "AI reasoning layer\nClaude decides order,\nhandles errors"),
        ("Tools",         "530+ Python scripts\nOne job each\nDeterministic"),
        ("Args",          "YAML/JSON config\nChange behavior\nwithout editing code"),
        ("Context",       "Static reference\nCompliance catalogs\nTone rules"),
        ("Hard Prompts",  "Reusable LLM\ninstruction\ntemplates"),
    ]
    box_w = CW / len(layers) - Inches(0.08)
    y = Inches(2.05)
    for i, (name, desc) in enumerate(layers):
        x = LM + i * (box_w + Inches(0.08))
        _rect(s, x, y, box_w, Inches(2.8), MGRAY, GOLD)
        _box(s, x, y + Inches(0.15), box_w, Inches(0.55),
             name, size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        _rect(s, x + Inches(0.1), y + Inches(0.72), box_w - Inches(0.2), Inches(0.03), GOLD)
        _box(s, x + Inches(0.08), y + Inches(0.82), box_w - Inches(0.16), Inches(1.85),
             desc, size=10, color=LGRAY, align=PP_ALIGN.CENTER)

    _box(s, LM, Inches(5.10), CW, Inches(0.55),
         "Why It Matters", size=14, bold=True, color=GOLD)
    _box(s, LM, Inches(5.65), CW, Inches(0.90),
         "LLMs are probabilistic. At 90% accuracy per step, a 5-step workflow degrades to ~59% end-to-end (0.9⁵). "
         "FORGE confines AI reasoning to orchestration only — all execution is deterministic Python. "
         "Result: >95% reliable, auditable, air-gap-safe delivery.",
         size=12, color=LGRAY, word_wrap=True)


def slide_06_anvil(prs):
    s = _blank(prs); _bg(s)
    _title_area(s, "The ANVIL Workflow", "Five-phase TDD cycle — every feature, every time, with compliance baked in")
    _footer(s); _slide_num(s, 6)

    phases = [
        ("A\nrchitect",    "Design & acceptance\ncriteria. System\narchitecture defined."),
        ("N\navigate",     "Map to existing tools\nand patterns. No\nnew code without cause."),
        ("V\nerify",       "Write failing tests\nfirst. RED phase.\nSpec is the test."),
        ("I\nntegrate",    "Generate implementation.\nGREEN phase. Tests\nmust pass."),
        ("L\naunch",       "Refactor, security scan,\ncompliance map,\nmerge to main."),
    ]
    arrow_w = Inches(0.35)
    box_w   = (CW - arrow_w * 4) / 5
    y_box   = Inches(2.05)
    colors  = [BLUE, BLUE, RED, GREEN, GOLD]

    x = LM
    for i, (name, desc) in enumerate(phases):
        _rect(s, x, y_box, box_w, Inches(3.2), MGRAY, colors[i])
        _box(s, x, y_box + Inches(0.12), box_w, Inches(0.7),
             name.replace("\n", ""), size=16, bold=True, color=colors[i], align=PP_ALIGN.CENTER)
        _rect(s, x + Inches(0.1), y_box + Inches(0.85), box_w - Inches(0.2), Inches(0.03), colors[i])
        _box(s, x + Inches(0.08), y_box + Inches(0.98), box_w - Inches(0.16), Inches(2.0),
             desc, size=10, color=LGRAY, align=PP_ALIGN.CENTER)
        x += box_w
        if i < 4:
            _box(s, x, y_box + Inches(1.4), arrow_w, Inches(0.5),
                 "→", size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
            x += arrow_w

    _box(s, LM, Inches(5.55), CW, Inches(0.45),
         "Compliance baked in at every phase — not bolted on after.",
         size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    _box(s, LM, Inches(6.05), CW, Inches(0.55),
         "Every ANVIL run produces: working code · passing tests · SBOM · compliance evidence · security scan report",
         size=11, color=LGRAY, align=PP_ALIGN.CENTER)


def slide_07_canvases(prs):
    s = _blank(prs); _bg(s)
    _title_area(s, "12 Design Canvases", "Design visually. Query naturally. Comply automatically.")
    _footer(s); _slide_num(s, 7)

    canvases = [
        ("NDC",  "Network Design",        "Cloud topology, FedRAMP/IL/CMMC mapping, COTS cost modeling"),
        ("SDC",  "Security Design",        "STRIDE threat modeling, MITRE ATT&CK, SSP/POAM generator"),
        ("PDC",  "Pipeline Design",        "Visual CI/CD builder, SLSA assessment, OWASP coverage"),
        ("BDC",  "Boundary Design",        "ATO boundary, ISA lifecycle, PPS matrix auto-generation"),
        ("DDC",  "Data Design",            "Data classification zones, PII/PHI/CUI tracking, lineage"),
        ("ODC",  "Observability Design",   "Detection coverage, Sigma rules, MITRE ATT&CK detection"),
        ("IDC",  "Infrastructure Design",  "IaC resource design, 6 CSPs, Terraform generation"),
        ("AADC", "Agentic AI Design",      "7 solution packs, 40+ node types, AI risk register"),
        ("QDC",  "Quality Design",         "Code quality gates, UQS scoring, NIST SA-11 mapping"),
        ("MDC",  "Migration Design",       "7R assessment, strangler fig mapping, ATO bridge"),
        ("AIMC", "AI/ML Model Design",     "Foundation model catalog, DoD RAI compliance, IL assessment"),
        ("MSN",  "Mission Canvas",         "Strategic mission planning, operational design, COA support"),
    ]

    cols, rows = 4, 3
    box_w = (CW - Inches(0.15) * (cols - 1)) / cols
    box_h = Inches(1.32)
    y_start = Inches(2.05)

    for i, (code, name, desc) in enumerate(canvases):
        col = i % cols
        row = i // cols
        x = LM + col * (box_w + Inches(0.15))
        y = y_start + row * (box_h + Inches(0.1))
        _rect(s, x, y, box_w, box_h, MGRAY, GOLD)
        _box(s, x + Inches(0.1), y + Inches(0.08), box_w - Inches(0.2), Inches(0.38),
             f"{code}  ·  {name}", size=11, bold=True, color=GOLD)
        _box(s, x + Inches(0.1), y + Inches(0.5), box_w - Inches(0.2), Inches(0.72),
             desc, size=9, color=LGRAY, word_wrap=True)

    _box(s, LM, Inches(6.67), CW, Inches(0.38),
         "All 12 canvases: Natural-language IQE queries · Knowledge graph indexing · Air-gap safe · No CDN dependencies",
         size=10, color=MGRAY, align=PP_ALIGN.CENTER)


def slide_08_chat(prs):
    s = _blank(prs); _bg(s)
    _title_area(s, "Unified Chat — ICDEV in Action",
                 "Not a chatbot. A structured intake engine with AI governance baked in.")
    _footer(s); _slide_num(s, 8)

    # 3-pane layout diagram using rectangles
    left_w   = Inches(2.8)
    center_w = Inches(5.4)
    right_w  = Inches(3.6)
    pane_h   = Inches(4.0)
    pane_y   = Inches(2.05)
    gap      = Inches(0.12)

    x_l = LM
    x_c = x_l + left_w + gap
    x_r = x_c + center_w + gap

    _rect(s, x_l, pane_y, left_w, pane_h, MGRAY, GOLD)
    _rect(s, x_c, pane_y, center_w, pane_h, MGRAY, BLUE)
    _rect(s, x_r, pane_y, right_w, pane_h, MGRAY, GREEN)

    # Left pane content
    _box(s, x_l + Inches(0.1), pane_y + Inches(0.1), left_w - Inches(0.2), Inches(0.38),
         "Contexts + Use Cases", size=11, bold=True, color=GOLD)
    left_items = ["+ New Chat", "▸ General Modernization", "▸ Budget Sprint", "▸ Doc Refresh",
                  "▸ Custom Context 1", "▸ Custom Context 2"]
    y = pane_y + Inches(0.55)
    for item in left_items:
        _box(s, x_l + Inches(0.1), y, left_w - Inches(0.2), Inches(0.38),
             item, size=10, color=LGRAY)
        y += Inches(0.40)

    # Center pane
    _box(s, x_c + Inches(0.1), pane_y + Inches(0.1), center_w - Inches(0.2), Inches(0.38),
         "Message Stream", size=11, bold=True, color=BLUE)
    msgs = [("User", "What's the migration path for our legacy COBOL system?"),
            ("ICDEV™", "Based on your inventory, I recommend a Replatform strategy (R4)..."),
            ("User", "What does that mean for our current ATO?"),
            ("ICDEV™", "I've pre-loaded 12 ATO bridge requirements. Readiness: 73%.")]
    y = pane_y + Inches(0.55)
    for who, msg in msgs:
        c = GOLD if who == "User" else GREEN
        _box(s, x_c + Inches(0.1), y, center_w - Inches(0.2), Inches(0.22),
             who, size=9, bold=True, color=c)
        _box(s, x_c + Inches(0.1), y + Inches(0.22), center_w - Inches(0.2), Inches(0.42),
             msg, size=9, color=LGRAY, word_wrap=True)
        y += Inches(0.75)
    _rect(s, x_c + Inches(0.1), pane_y + pane_h - Inches(0.55),
          center_w - Inches(0.2), Inches(0.38), NAVY, BLUE)
    _box(s, x_c + Inches(0.15), pane_y + pane_h - Inches(0.50),
         center_w - Inches(0.3), Inches(0.30),
         "Send a message...  [↑ Send]", size=9, color=MGRAY)

    # Right pane
    _box(s, x_r + Inches(0.1), pane_y + Inches(0.1), right_w - Inches(0.2), Inches(0.38),
         "RICOAS Readiness", size=11, bold=True, color=GREEN)
    right_items = ["Readiness Score: 73 / 100", "Completeness: ✓ 8/12 reqs",
                   "Clarity: ✓ High", "Compliance: ⚠ ATO Bridge",
                   "AI Governance: ✓ Active",
                   "──────────────────",
                   "AI Governance Pillars",
                   "Transparency: ✓",
                   "Accountability: ✓",
                   "Fairness: ⚠ Pending"]
    y = pane_y + Inches(0.55)
    for item in right_items:
        _box(s, x_r + Inches(0.1), y, right_w - Inches(0.2), Inches(0.34),
             item, size=9, color=LGRAY)
        y += Inches(0.36)

    # Callouts below
    callouts = [
        ("Typed Contexts", "Every chat carries project metadata, use case, and compliance scope"),
        ("AI Governance", "6-pillar posture sidebar: Transparency, Accountability, Fairness, Safety, XAI, Privacy"),
        ("Intervention", "Operator injects mid-stream directions to guide the agent during processing"),
    ]
    cw = CW / 3 - Inches(0.1)
    y = Inches(6.15)
    for i, (title, body) in enumerate(callouts):
        x = LM + i * (cw + Inches(0.1))
        _box(s, x, y, cw, Inches(0.28), title, size=10, bold=True, color=GOLD)
        _box(s, x, y + Inches(0.28), cw, Inches(0.40), body, size=9, color=LGRAY, word_wrap=True)


def slide_09_use_cases(prs):
    s = _blank(prs); _bg(s)
    _title_area(s, "Pre-Built Use Cases — Start in Minutes",
                 "Expert AI personas, pre-loaded requirements, and canvas wiring — out of the box")
    _footer(s); _slide_num(s, 9)

    cases = [
        (
            "⚙  General Modernization",
            "BADGE: AI Boost  ·  12 Pre-Loaded Requirements",
            [
                "7Rs framework (Retire → Re-architect)",
                "Asset inventory + migration phasing",
                "ATO continuity bridge maintained",
                "Canvas wiring: Migration, Digital Twin,",
                "  Supply Chain, Compliance",
                "Output: phased migration plan with",
                "  rollback runbooks",
            ]
        ),
        (
            "$  Year-End Budget Sprint",
            "BADGE: RICOAS  ·  14 Pre-Loaded Requirements",
            [
                '"Spend it or lose it" procurement',
                "BOM + IGCE + SOW generated in days",
                "Tier 1 (execution-ready) vs Tier 2 (backup)",
                "FAR/DFARS + NDAA §889 compliance",
                "Canvas wiring: CPMP, Digital Twin,",
                "  Supply Chain, Kanban",
                "Output: CSV/XLSX BOM with lead times",
            ]
        ),
        (
            "📄  Crowd-Sourced Doc Refresh",
            "BADGE: RAG + KG  ·  10 Pre-Loaded Requirements",
            [
                "Staleness detection + crowd voting",
                "AI rebuild with role-weighted edits",
                "Immutable audit trail per NIST AU",
                "Auto-indexed to Knowledge Graph",
                "  within 24 hrs of approval",
                "Canvas wiring: Knowledge, RAG,",
                "  Compliance, Audit",
            ]
        ),
    ]

    cw = (CW - Inches(0.2)) / 3
    for i, (title, badge, items) in enumerate(cases):
        x = LM + i * (cw + Inches(0.1))
        y = Inches(2.05)
        _rect(s, x, y, cw, Inches(4.6), MGRAY, GOLD)
        _box(s, x + Inches(0.12), y + Inches(0.1), cw - Inches(0.24), Inches(0.5),
             title, size=13, bold=True, color=GOLD)
        _box(s, x + Inches(0.12), y + Inches(0.62), cw - Inches(0.24), Inches(0.32),
             badge, size=9, color=GOLD, italic=True)
        _rect(s, x + Inches(0.12), y + Inches(0.95), cw - Inches(0.24), Inches(0.02), GOLD)
        yy = y + Inches(1.05)
        for item in items:
            _box(s, x + Inches(0.12), yy, cw - Inches(0.24), Inches(0.40),
                 f"• {item}", size=10, color=LGRAY, word_wrap=True)
            yy += Inches(0.44)

    _box(s, LM, Inches(6.78), CW, Inches(0.32),
         "Add new use cases by editing args/use_cases.yaml — zero Python required. Auto-seeds RICOAS requirements and wires canvases.",
         size=10, color=MGRAY, align=PP_ALIGN.CENTER)


def slide_10_digital_twin(prs):
    s = _blank(prs); _bg(s)
    _title_area(s, "Digital Program Twin",
                 "Simulate before you commit. Three Courses of Action. Confidence-banded outcomes.")
    _footer(s); _slide_num(s, 10)

    dims = [
        ("Architecture", "Components, API surface,\ncoupling, data flow"),
        ("Compliance",   "Control coverage,\nPOAM projection, boundary tier"),
        ("Supply Chain", "New dependencies,\nvendor risk, SBOM delta"),
        ("Schedule",     "PERT estimates,\nMonte Carlo confidence"),
        ("Cost",         "T-shirt roll-up,\ninfrastructure delta"),
        ("Risk",         "Compound score,\nrisk interactions"),
    ]
    dw = (CW - Inches(0.25)) / 6
    y = Inches(2.05)
    for i, (dim, desc) in enumerate(dims):
        x = LM + i * (dw + Inches(0.05))
        _rect(s, x, y, dw, Inches(1.55), MGRAY, BLUE)
        _box(s, x + Inches(0.06), y + Inches(0.1), dw - Inches(0.12), Inches(0.42),
             dim, size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        _box(s, x + Inches(0.06), y + Inches(0.55), dw - Inches(0.12), Inches(0.85),
             desc, size=9, color=LGRAY, align=PP_ALIGN.CENTER)

    _box(s, LM, Inches(3.73), CW, Inches(0.35),
         "Monte Carlo: 10,000 iterations  →  P10 (optimistic) / P50 (likely) / P80 (management reserve) / P90 (conservative)",
         size=11, color=GOLD, align=PP_ALIGN.CENTER, bold=True)

    coas = [
        ("Speed COA",           "1–2 PIs",    "Lowest", "HIGH",   RED,    "P1 only. Fast. Technical debt."),
        ("Balanced COA ★",      "2–3 PIs",    "Medium", "MOD",    GREEN,  "P1+P2. Recommended. Best tradeoff."),
        ("Comprehensive COA",   "3–5 PIs",    "Highest","LOW",    LGRAY,  "Full scope. Complete coverage."),
    ]
    tw = CW / 3 - Inches(0.1)
    y = Inches(4.2)
    for i, (name, timeline, cost, risk, clr, desc) in enumerate(coas):
        x = LM + i * (tw + Inches(0.1))
        _rect(s, x, y, tw, Inches(2.05), MGRAY, clr)
        _box(s, x + Inches(0.1), y + Inches(0.08), tw - Inches(0.2), Inches(0.40),
             name, size=12, bold=True, color=clr, align=PP_ALIGN.CENTER)
        for j, (label, val) in enumerate([("Timeline", timeline), ("Cost", cost), ("Risk", risk)]):
            _box(s, x + Inches(0.1), y + Inches(0.55) + j * Inches(0.35),
                 tw - Inches(0.2), Inches(0.32),
                 f"{label}: {val}", size=10, color=LGRAY)
        _box(s, x + Inches(0.1), y + Inches(1.6), tw - Inches(0.2), Inches(0.35),
             desc, size=9, color=LGRAY, italic=True)

    _box(s, LM, Inches(6.45), CW, Inches(0.35),
         "No commercial vendor offers program-level Digital Twin semantics with NIST/FedRAMP verdicts. ICDEV does.",
         size=11, color=GOLD, align=PP_ALIGN.CENTER, italic=True)


def slide_11_compliance(prs):
    s = _blank(prs); _bg(s)
    _title_area(s, "42-Framework Compliance Superpower",
                 "Implement once. Satisfy many. Dual-hub crosswalk eliminates duplicate assessments.")
    _footer(s); _slide_num(s, 11)

    # Big stat
    _box(s, LM, Inches(2.05), CW, Inches(0.65),
         "Implement AC-2 (Account Management) once → automatically satisfies 30+ frameworks",
         size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Two hub boxes
    hub_w = Inches(4.0)
    hub_h = Inches(3.2)
    hub_y = Inches(2.85)

    # US Hub
    _rect(s, LM, hub_y, hub_w, hub_h, MGRAY, BLUE)
    _box(s, LM + Inches(0.1), hub_y + Inches(0.1), hub_w - Inches(0.2), Inches(0.42),
         "US HUB — NIST 800-53 Rev 5", size=12, bold=True, color=BLUE)
    us = ["FedRAMP (Mod/High/20x)", "CMMC Level 2 & 3", "NIST 800-171 Rev 2",
          "FIPS 199/200", "CNSSI 1253 (IL6)", "DoD CSSP (DI 8530.01)",
          "DoD MOSA (10 U.S.C. §4401)", "NIST SP 800-207 Zero Trust"]
    for j, fw in enumerate(us):
        _box(s, LM + Inches(0.12), hub_y + Inches(0.62) + j * Inches(0.31),
             hub_w - Inches(0.24), Inches(0.28), f"✓  {fw}", size=10, color=LGRAY)

    # Arrow
    arr_x = LM + hub_w + Inches(0.15)
    _box(s, arr_x, hub_y + Inches(1.3), Inches(1.3), Inches(0.6),
         "↔", size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    _box(s, arr_x, hub_y + Inches(1.9), Inches(1.3), Inches(0.45),
         "Crosswalk\nEngine", size=9, color=GOLD, align=PP_ALIGN.CENTER)

    # Intl Hub
    ih_x = arr_x + Inches(1.35)
    _rect(s, ih_x, hub_y, hub_w, hub_h, MGRAY, GREEN)
    _box(s, ih_x + Inches(0.1), hub_y + Inches(0.1), hub_w - Inches(0.2), Inches(0.42),
         "INTL HUB — ISO/IEC 27001:2022", size=12, bold=True, color=GREEN)
    intl = ["ISO/IEC 42001:2023 (AI Mgmt)", "EU AI Act (Annex III)",
            "HIPAA Security Rule", "HITRUST CSF v11",
            "PCI DSS v4.0", "SOC 2 Type II",
            "CJIS Security Policy", "NIST AI RMF 1.0 + AI 600-1"]
    for j, fw in enumerate(intl):
        _box(s, ih_x + Inches(0.12), hub_y + Inches(0.62) + j * Inches(0.31),
             hub_w - Inches(0.24), Inches(0.28), f"✓  {fw}", size=10, color=LGRAY)

    _box(s, LM, Inches(6.22), CW, Inches(0.45),
         "ATO in 1–3 months.  Not 12–18.  SSP, POAM, STIG, SBOM, OSCAL generated automatically on first build.",
         size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_12_agents(prs):
    s = _blank(prs); _bg(s)
    _title_area(s, "15-Agent Autonomous Ecosystem",
                 "Agents coordinate so your team doesn't have to.")
    _footer(s); _slide_num(s, 12)

    tiers = [
        ("CORE TIER", [("Orchestrator :8443", "Task routing, DAG execution, dispatcher mode"),
                       ("Architect :8444", "ANVIL A/T phases, system design, M-ANVIL for models")]),
        ("DOMAIN TIER", [
            ("Builder :8445",          "TDD code gen — 6 languages"),
            ("Compliance :8446",       "ATO: SSP, POAM, STIG, SBOM, OSCAL, eMASS"),
            ("Security :8447",         "SAST, deps, secrets, containers, ATLAS"),
            ("Infrastructure :8448",   "Terraform, Ansible, K8s, 6 CSPs"),
            ("MBSE :8451",             "SysML, DOORS NG, digital thread, DES"),
            ("Modernization :8452",    "7R assessment, migration, code gen"),
            ("Requirements :8453",     "Intake, gap detection, SAFe, readiness"),
            ("Supply Chain :8454",     "SBOM, CVE triage, SCRM, §889"),
            ("Simulation :8455",       "Digital Twin, Monte Carlo, COA gen"),
            ("DevSecOps+ZTA :8457",    "Pipeline security, NIST 800-207, policy-as-code"),
            ("Gateway :8458",          "Remote commands, 8-gate security chain"),
        ]),
        ("SUPPORT TIER", [("Knowledge :8449", "Self-healing patterns, ML recommendations"),
                          ("Monitor :8450",   "Logs, metrics, alerts, SLA tracking")]),
    ]

    tier_colors = [GOLD, BLUE, GREEN]
    tier_h      = [Inches(0.85), Inches(3.2), Inches(0.85)]
    y = Inches(2.0)

    for t_idx, (tier_name, agents) in enumerate(tiers):
        clr = tier_colors[t_idx]
        th  = tier_h[t_idx]
        _rect(s, LM, y, CW, th, MGRAY, clr)
        _box(s, LM + Inches(0.1), y + Inches(0.08), Inches(1.8), Inches(0.38),
             tier_name, size=11, bold=True, color=clr)

        aw = (CW - Inches(0.2) - Inches(1.9)) / len(agents)
        ax = LM + Inches(1.9)
        for name, desc in agents:
            _box(s, ax, y + Inches(0.08), aw - Inches(0.05), Inches(0.32),
                 name, size=9, bold=True, color=clr)
            _box(s, ax, y + Inches(0.40), aw - Inches(0.05), th - Inches(0.50),
                 desc, size=8, color=LGRAY, word_wrap=True)
            ax += aw

        y += th + Inches(0.1)

    _box(s, LM, Inches(6.32), CW, Inches(0.5),
         "A2A protocol: JSON-RPC 2.0 over mutual TLS  ·  All agents publish /.well-known/agent.json  "
         "·  Self-healing daemon: 30-min failure triage  ·  Confidence ≥ 0.7 required for auto-patch",
         size=10, color=MGRAY, align=PP_ALIGN.CENTER)


def slide_13_academy(prs):
    s = _blank(prs); _bg(s)
    _title_area(s, "FORGE Academy — Building AI-Ready Teams",
                 "Gamified, role-based AI training designed for government and defense organizations")
    _footer(s); _slide_num(s, 13)

    # Stats row
    stats = [("12", "Role Tracks"), ("75", "Missions"), ("165", "Mission Steps"),
             ("5", "Ranks"), ("3", "Certifications")]
    sw = CW / len(stats) - Inches(0.1)
    for i, (num, lbl) in enumerate(stats):
        x = LM + i * (sw + Inches(0.1))
        _rect(s, x, Inches(2.05), sw, Inches(1.1), MGRAY, GOLD)
        _box(s, x, Inches(2.1), sw, Inches(0.55),
             num, size=32, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        _box(s, x, Inches(2.65), sw, Inches(0.4),
             lbl, size=10, color=LGRAY, align=PP_ALIGN.CENTER)

    # Rank progression
    ranks = ["Recruit\n0 XP", "Operative\n500 XP", "Specialist\n2,000 XP",
             "Architect\n5,000 XP", "Sensei\n10,000+ XP"]
    rw = (CW - Inches(0.6)) / 5
    for i, rank in enumerate(ranks):
        x = LM + i * (rw + Inches(0.15))
        clr = [LGRAY, LGRAY, BLUE, GREEN, GOLD][i]
        _rect(s, x, Inches(3.35), rw, Inches(0.75), MGRAY, clr)
        _box(s, x, Inches(3.38), rw, Inches(0.7),
             rank, size=10, color=clr, align=PP_ALIGN.CENTER)
        if i < 4:
            _box(s, x + rw + Inches(0.02), Inches(3.55), Inches(0.12), Inches(0.40),
                 "→", size=14, color=GOLD, align=PP_ALIGN.CENTER)

    # Two columns
    col_w = (CW - Inches(0.2)) / 2
    _box(s, LM, Inches(4.25), col_w, Inches(0.38),
         "Role Tracks (12)", size=13, bold=True, color=GOLD)
    roles = ["DevOps · DataOps · SecOps", "SWE/Architect · NetOps · SRE",
             "ISSO · ISSM · CISO", "PM · Analyst · Leadership"]
    for j, r in enumerate(roles):
        _box(s, LM, Inches(4.65) + j * Inches(0.38),
             col_w, Inches(0.35), f"• {r}", size=11, color=LGRAY)

    x2 = LM + col_w + Inches(0.2)
    _box(s, x2, Inches(4.25), col_w, Inches(0.38),
         "Certifications + Alignment", size=13, bold=True, color=GOLD)
    certs = ["Foundation (2yr) — Tier 1 + role Tier 2 + 75% score",
             "Practitioner (2yr) — AADC ≥80 + 1 GameDay",
             "Expert (3yr) — Capstone + top-50% GameDay finish",
             "DoD AI Workforce Framework + NICE KSA crosswalk"]
    for j, c in enumerate(certs):
        _box(s, x2, Inches(4.65) + j * Inches(0.38),
             col_w, Inches(0.35), f"• {c}", size=11, color=LGRAY)

    _box(s, LM, Inches(6.45), CW, Inches(0.40),
         "Auto-updating curriculum: Genesis daemon detects proven patterns → auto-generates new missions → human review → activated",
         size=11, color=GOLD, align=PP_ALIGN.CENTER)


def slide_14_gameday(prs):
    s = _blank(prs); _bg(s)
    _title_area(s, "AI GameDay — Where Training Meets Reality",
                 "Live, AI-scored tabletop exercises. Competitive. NIST-compliant. After-action ready.")
    _footer(s); _slide_num(s, 14)

    teams = [
        ("RED TEAM", "Threat",       RED,   "Scout · Threat Analyst\nExploit Engineer\nRed Orchestrator"),
        ("BLUE TEAM","Defense",      BLUE,  "SOC Analyst · Security Arch\nIR Responder\nBlue Orchestrator"),
        ("GOLD TEAM","Innovation",   GOLD,  "Researcher · Builder\nEvaluator\nGold Orchestrator"),
        ("GREEN TEAM","Compliance",  GREEN, "NIST Auditor · Risk Assessor\nPolicy Advisor\nGreen Orchestrator"),
    ]
    tw = (CW - Inches(0.3)) / 4
    for i, (name, role, clr, members) in enumerate(teams):
        x = LM + i * (tw + Inches(0.1))
        _rect(s, x, Inches(2.05), tw, Inches(1.9), MGRAY, clr)
        _box(s, x + Inches(0.08), Inches(2.12), tw - Inches(0.16), Inches(0.38),
             name, size=11, bold=True, color=clr, align=PP_ALIGN.CENTER)
        _box(s, x + Inches(0.08), Inches(2.52), tw - Inches(0.16), Inches(0.32),
             role, size=10, color=GOLD, align=PP_ALIGN.CENTER, italic=True)
        _box(s, x + Inches(0.08), Inches(2.85), tw - Inches(0.16), Inches(0.95),
             members, size=9, color=LGRAY, align=PP_ALIGN.CENTER)

    # Injects
    _box(s, LM, Inches(4.1), CW, Inches(0.40),
         "Operation CIPHER FORGE — 5 Injects (4 Hours Total)", size=13, bold=True, color=GOLD)
    injects = [
        ("Signal Cluster",      "15 min", "SIGINT threat assessment"),
        ("COA Posture",         "40 min", "Force posture recommendation"),
        ("Ransomware Cascade",  "Seq",    "IR: contain, eradicate, recover"),
        ("Fine-Tune Sprint",    "Seq",    "ML: generate + evaluate training pairs"),
        ("War Council Brief",   "Final",  "Executive synthesis of all learnings"),
    ]
    iw = CW / len(injects) - Inches(0.06)
    for i, (name, dur, desc) in enumerate(injects):
        x = LM + i * (iw + Inches(0.06))
        _rect(s, x, Inches(4.6), iw, Inches(1.35), MGRAY, BLUE)
        _box(s, x + Inches(0.08), Inches(4.68), iw - Inches(0.16), Inches(0.38),
             name, size=10, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        _box(s, x + Inches(0.08), Inches(5.06), iw - Inches(0.16), Inches(0.28),
             dur, size=10, color=GOLD, align=PP_ALIGN.CENTER)
        _box(s, x + Inches(0.08), Inches(5.35), iw - Inches(0.16), Inches(0.52),
             desc, size=9, color=LGRAY, align=PP_ALIGN.CENTER, word_wrap=True)

    _box(s, LM, Inches(6.1), CW, Inches(0.32),
         "Scoring: Adversarial 40% · Innovation 25% · Compliance 20% · Training Pairs 15%  "
         "·  AI-scored · Server-verified receipts · Leaderboard live",
         size=10, color=LGRAY, align=PP_ALIGN.CENTER)
    _box(s, LM, Inches(6.45), CW, Inches(0.35),
         "Output: Auto-generated AAR · Kanban debrief card · Fine-tune pairs → model improvement  "
         "·  NIST AC-2, AU-2, SI-10 compliant",
         size=10, color=MGRAY, align=PP_ALIGN.CENTER)


def slide_15_build_vs_adopt(prs):
    s = _blank(prs); _bg(s)
    _title_area(s, "Build vs. Adopt — The Numbers Don't Lie",
                 "Nine dimensions. ICDEV wins on every axis.")
    _footer(s); _slide_num(s, 15)

    rows = [
        ("Dimension",            "Build From Scratch",           "Adopt ICDEV™"),
        ("ATO Timeline",         "12–18 months",                 "1–3 months"),
        ("Compliance",           "Manual, post-hoc, incomplete", "42 frameworks, auto-crosswalk"),
        ("AI Agents",            "None (hire specialists)",      "15 coordinating agents"),
        ("Digital Twin",         "Not available",                "6-dim simulation, 3 COAs, Monte Carlo"),
        ("AI Governance",        "Build from scratch (months)",  "OMB M-25-21 + NIST AI RMF built in"),
        ("Legacy Modernization", "No path. Greenfield only.",    "7Rs + cross-language + ATO bridge"),
        ("Audit Trail",          "Manual / ad hoc",              "Append-only, NIST AU, day one"),
        ("Scalability",          "Re-architect per program",     "Generated apps inherit full FORGE/ANVIL"),
        ("Cost",                 "$5M–10M+ to replicate",        "$570K–880K Year 1 (open-source base)"),
    ]

    col_widths = [Inches(2.5), Inches(4.8), Inches(4.8)]
    col_colors = [MGRAY, MGRAY, MGRAY]
    header_colors = [MGRAY, MGRAY, GREEN]

    y = Inches(2.0)
    for r_idx, row in enumerate(rows):
        x = LM
        for c_idx, (cell, cw) in enumerate(zip(row, col_widths)):
            is_header = r_idx == 0
            bg = [MGRAY, NAVY, MGRAY][c_idx] if not is_header else [MGRAY, MGRAY, GREEN][c_idx]
            tc = GOLD if is_header else (WHITE if c_idx == 0 else (LGRAY if c_idx == 1 else WHITE))
            _rect(s, x, y, cw - Inches(0.02), Inches(0.42), bg, None)
            _box(s, x + Inches(0.1), y + Inches(0.04), cw - Inches(0.14), Inches(0.35),
                 cell, size=10 if not is_header else 11, bold=is_header or c_idx == 0,
                 color=tc)
            x += cw
        y += Inches(0.43)


def slide_16_timeline(prs):
    s = _blank(prs); _bg(s)
    _title_area(s, "From Intake to ATO in 1–3 Months",
                 "A single auditable pipeline — intake through Authority to Operate.")
    _footer(s); _slide_num(s, 16)

    steps = [
        ("RICOAS\nIntake",     "Requirements\ngap detection\nreadiness score"),
        ("Digital\nTwin",      "6-dim simulation\n3 COAs\nMonte Carlo"),
        ("Design\nCanvases",   "12 visual builders\nNL queries\ncompliance map"),
        ("ANVIL\nBuild",       "TDD: RED→GREEN\nSecurity scan\nSBOM generated"),
        ("Compliance\nAuto",   "42 frameworks\nSSP+POAM+STIG\nOSCAL+eMASS"),
        ("Deploy\n+ ATO",      "K8s + IaC\nATO package\nSigned"),
    ]
    sw = (CW - Inches(0.5)) / len(steps)
    y = Inches(2.2)
    for i, (step, desc) in enumerate(steps):
        x = LM + i * (sw + Inches(0.1))
        clr = [BLUE, BLUE, GOLD, GREEN, GREEN, GOLD][i]
        _rect(s, x, y, sw, Inches(1.55), MGRAY, clr)
        _box(s, x + Inches(0.06), y + Inches(0.1), sw - Inches(0.12), Inches(0.55),
             step, size=11, bold=True, color=clr, align=PP_ALIGN.CENTER)
        _box(s, x + Inches(0.06), y + Inches(0.68), sw - Inches(0.12), Inches(0.78),
             desc, size=9, color=LGRAY, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            _box(s, x + sw + Inches(0.01), y + Inches(0.65), Inches(0.08), Inches(0.4),
                 "→", size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Old vs new
    _rect(s, LM, Inches(4.05), CW / 2 - Inches(0.1), Inches(0.7), MGRAY, RED)
    _box(s, LM + Inches(0.1), Inches(4.1), CW / 2 - Inches(0.3), Inches(0.55),
         "Without ICDEV™:  12–18 months  ·  $2M–5M+  ·  Expert-dependent", size=13, color=RED)
    _rect(s, LM + CW / 2 + Inches(0.1), Inches(4.05), CW / 2 - Inches(0.1), Inches(0.7), MGRAY, GREEN)
    _box(s, LM + CW / 2 + Inches(0.2), Inches(4.1), CW / 2 - Inches(0.3), Inches(0.55),
         "With ICDEV™:  1–3 months  ·  $570K–880K  ·  15 agents", size=13, color=GREEN)

    # PI timeline
    _box(s, LM, Inches(4.95), CW, Inches(0.40), "Program Increment Timeline:", size=13, bold=True, color=GOLD)
    pis = [
        ("Speed COA",         "1–2 PIs", "P1 only, MVP delivery"),
        ("Balanced COA  ★",   "2–3 PIs", "P1+P2, recommended"),
        ("Comprehensive COA", "3–5 PIs", "Full scope"),
    ]
    pw = CW / 3 - Inches(0.1)
    for i, (name, pi, desc) in enumerate(pis):
        x = LM + i * (pw + Inches(0.1))
        _rect(s, x, Inches(5.42), pw, Inches(0.75), MGRAY, [LGRAY, GREEN, LGRAY][i])
        _box(s, x + Inches(0.1), Inches(5.46), pw - Inches(0.2), Inches(0.30),
             name, size=11, bold=True, color=[LGRAY, GREEN, LGRAY][i])
        _box(s, x + Inches(0.1), Inches(5.76), pw - Inches(0.2), Inches(0.28),
             f"{pi}  ·  {desc}", size=10, color=LGRAY)


def slide_17_lab_vision(prs):
    s = _blank(prs); _bg(s)
    _title_area(s, "The AI Lab Vision",
                 "A living lab that trains your people, hardens your systems, and accelerates every program.")
    _footer(s); _slide_num(s, 17)

    # Central ICDEV platform box
    cx = LM + CW / 2 - Inches(1.9)
    _rect(s, cx, Inches(2.3), Inches(3.8), Inches(1.6), MGRAY, GOLD)
    _box(s, cx + Inches(0.15), Inches(2.4), Inches(3.5), Inches(0.55),
         "ICDEV™ Platform", size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    _box(s, cx + Inches(0.15), Inches(2.97), Inches(3.5), Inches(0.75),
         "FORGE · ANVIL · 12 Canvases · 15 Agents\nChat + Use Cases · Digital Twin · 42 Frameworks",
         size=10, color=LGRAY, align=PP_ALIGN.CENTER)

    satellites = [
        (LM,                         Inches(2.0),  "FORGE Academy",
         "12 roles · 75 missions\n4 cohorts × 25 devs\n8 weeks each", BLUE),
        (LM,                         Inches(4.2),  "AI GameDay",
         "4 teams · 5 injects\nQuarterly tournaments\nAAR + certification", GREEN),
        (W - LM - Inches(3.5),       Inches(2.0),  "Cloud Infrastructure",
         "AWS GovCloud / Azure Gov\nK8s · 15 agents\nIL4+ ready", GOLD),
        (W - LM - Inches(3.5),       Inches(4.2),  "100+ Developers",
         "AI-trained workforce\nCertified Sensei bench\nProgram acceleration", LGRAY),
    ]

    for sx, sy, title, desc, clr in satellites:
        _rect(s, sx, sy, Inches(3.3), Inches(1.55), MGRAY, clr)
        _box(s, sx + Inches(0.1), sy + Inches(0.1), Inches(3.1), Inches(0.42),
             title, size=13, bold=True, color=clr)
        _box(s, sx + Inches(0.1), sy + Inches(0.55), Inches(3.1), Inches(0.85),
             desc, size=10, color=LGRAY)
        # Arrow toward center (simplified — just a dash)
        if sx == LM:
            _box(s, sx + Inches(3.32), sy + Inches(0.55), Inches(0.5), Inches(0.45),
                 "→", size=18, bold=True, color=GOLD)
        else:
            _box(s, sx - Inches(0.55), sy + Inches(0.55), Inches(0.5), Inches(0.45),
                 "←", size=18, bold=True, color=GOLD)

    _box(s, LM, Inches(6.1), CW, Inches(0.60),
         "Every program powered by ICDEV™ generates its own ATO package, trains its own team, and "
         "improves the shared knowledge base — compounding returns across the enterprise.",
         size=12, color=LGRAY, align=PP_ALIGN.CENTER, word_wrap=True)


def slide_18_the_ask(prs):
    s = _blank(prs); _bg(s)
    _title_area(s, "What We're Asking For",
                 "Three specific approvals to stand up the ICDEV™ AI Lab.")
    _footer(s); _slide_num(s, 18)

    asks = [
        ("1.  Infrastructure Approval",
         "2-node Kubernetes cluster  ·  64 vCPU, 256 GB RAM, 10 TB NVMe\n"
         "AWS Bedrock (Claude) or Azure OpenAI Gov access (IL4+)\n"
         "Isolated lab VLAN + CAC auth  ·  GitLab EE + JFrog Artifactory\n"
         "Timeline: 30-day procurement",
         BLUE),
        ("2.  Personnel (4 FTEs)",
         "AI Lab Director — strategy, stakeholder liaison, program integration\n"
         "Platform / MLOps Engineer — ICDEV deployment, agent health, Bedrock\n"
         "Training Lead — Academy curriculum, cohort scheduling, GameDay facilitation\n"
         "ISSO — lab ATO, STIG compliance, audit trail review",
         GREEN),
        ("3.  Budget Authority — Year 1",
         "Infrastructure (cloud + on-prem K8s): $150K–300K\n"
         "Personnel (4 FTEs, blended rate): $300K–400K\n"
         "Academy (4 cohorts) + GameDay (4 events): $120K–160K\n"
         "ICDEV™ platform: $0 (Apache 2.0 open source)\n"
         "Total: $570K–860K  against a Year-1 return of $4M–6M+",
         GOLD),
    ]

    y = Inches(2.05)
    for title, body, clr in asks:
        _rect(s, LM, y, CW, Inches(1.5), MGRAY, clr)
        _box(s, LM + Inches(0.15), y + Inches(0.1), CW - Inches(0.3), Inches(0.45),
             title, size=15, bold=True, color=clr)
        _box(s, LM + Inches(0.15), y + Inches(0.58), CW - Inches(0.3), Inches(0.82),
             body, size=11, color=LGRAY, word_wrap=True)
        y += Inches(1.62)


def slide_19_why_now(prs):
    s = _blank(prs); _bg(s)
    _title_area(s, "Why Now — The Mandate Is Here",
                 "Federal AI policy and DoD strategy are not optional. ICDEV is already mapped to all of them.")
    _footer(s); _slide_num(s, 19)

    mandates = [
        ("OMB M-25-21",
         "High-Impact AI Transparency",
         "Requires model cards, system cards, AI inventory, confabulation detection, fairness assessment.\n"
         "ICDEV status: BUILT — AI Transparency canvas + cross-framework audit.",
         GOLD),
        ("OMB M-26-04",
         "Unbiased AI",
         "Requires fairness testing, bias detection, and accountability plans for all federal AI systems.\n"
         "ICDEV status: BUILT — Fairness assessor + appeals + CAIO designation workflow.",
         BLUE),
        ("DoD AI Strategy",
         "AI-Ready Workforce",
         "Requires AI literacy across all personnel grades; NICE KSA alignment.\n"
         "ICDEV status: BUILT — FORGE Academy is DoD AI Workforce Framework crosswalked.",
         GREEN),
        ("EO on AI / NIST AI RMF",
         "Govern, Map, Measure, Manage",
         "All high-impact AI systems must implement GOVERN, MAP, MEASURE, MANAGE functions.\n"
         "ICDEV status: BUILT — NIST AI RMF 1.0 + AI 600-1 + OWASP Agentic AI integrated.",
         LGRAY),
    ]

    cw2 = (CW - Inches(0.15)) / 2
    positions = [(LM, Inches(2.05)), (LM + cw2 + Inches(0.15), Inches(2.05)),
                 (LM, Inches(4.2)),  (LM + cw2 + Inches(0.15), Inches(4.2))]

    for (x, y), (ref, title, body, clr) in zip(positions, mandates):
        _rect(s, x, y, cw2, Inches(1.95), MGRAY, clr)
        _box(s, x + Inches(0.12), y + Inches(0.1), cw2 - Inches(0.24), Inches(0.38),
             ref, size=13, bold=True, color=clr)
        _box(s, x + Inches(0.12), y + Inches(0.5), cw2 - Inches(0.24), Inches(0.3),
             title, size=11, color=GOLD, italic=True)
        _box(s, x + Inches(0.12), y + Inches(0.82), cw2 - Inches(0.24), Inches(1.0),
             body, size=9, color=LGRAY, word_wrap=True)

    _box(s, LM, Inches(6.35), CW, Inches(0.45),
         "Your organization doesn't have to figure this out. We already did.",
         size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_20_next_steps(prs):
    s = _blank(prs); _bg(s)
    _rect(s, 0, 0, W, Inches(0.12), GOLD)
    _rect(s, 0, H - Inches(0.12), W, Inches(0.12), GOLD)
    _footer(s); _slide_num(s, 20)

    _box(s, LM, TM, CW, Inches(0.75),
         "Next Steps", size=40, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    _gold_bar(s, Inches(0.8), Inches(0.04))

    steps = [
        ("Day 1–30",  "Approve AI Lab stand-up",
         "Authorize infrastructure procurement (K8s cluster + cloud LLM access).\n"
         "Designate AI Lab Director and ISSO. Begin FY procurement action."),
        ("Day 31–60", "Deploy ICDEV™ + Onboard First Program",
         "ICDEV deployed and health-checked. First program team onboarded to canvases.\n"
         "RICOAS intake session completed. Digital Twin run. ATO package auto-generated."),
        ("Day 61–90", "Launch FORGE Academy Cohort 1 + AI GameDay #1",
         "25-developer Cohort 1 begins 8-week track.\n"
         "AI GameDay #1 executed — 4 teams, Operation CIPHER FORGE, live AAR."),
    ]

    y = Inches(1.3)
    for date, title, body in steps:
        _rect(s, LM, y, CW, Inches(1.6), MGRAY, GOLD)
        _box(s, LM + Inches(0.15), y + Inches(0.1), Inches(1.5), Inches(0.42),
             date, size=14, bold=True, color=GOLD)
        _box(s, LM + Inches(1.7), y + Inches(0.1), CW - Inches(1.85), Inches(0.42),
             title, size=14, bold=True, color=WHITE)
        _box(s, LM + Inches(0.15), y + Inches(0.56), CW - Inches(0.3), Inches(0.88),
             body, size=12, color=LGRAY, word_wrap=True)
        y += Inches(1.72)

    _rect(s, LM, Inches(6.28), CW, Inches(0.55), MGRAY, GOLD)
    _box(s, LM + Inches(0.15), Inches(6.33), CW - Inches(0.3), Inches(0.45),
         "One platform.  One decision.  Infinite programs.",
         size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ── Main ────────────────────────────────────────────────────────────────────

def main():
    out_dir = Path(__file__).parent / "output"
    out_dir.mkdir(exist_ok=True)

    prs = _new_prs()

    slide_01_cover(prs)
    slide_02_problem(prs)
    slide_03_question(prs)
    slide_04_what_is_icdev(prs)
    slide_05_forge(prs)
    slide_06_anvil(prs)
    slide_07_canvases(prs)
    slide_08_chat(prs)
    slide_09_use_cases(prs)
    slide_10_digital_twin(prs)
    slide_11_compliance(prs)
    slide_12_agents(prs)
    slide_13_academy(prs)
    slide_14_gameday(prs)
    slide_15_build_vs_adopt(prs)
    slide_16_timeline(prs)
    slide_17_lab_vision(prs)
    slide_18_the_ask(prs)
    slide_19_why_now(prs)
    slide_20_next_steps(prs)

    out_path = out_dir / "ICDEV_Executive_Deck.pptx"
    prs.save(str(out_path))
    print(f"[OK] Saved: {out_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
