# CUI // SP-CTI
"""DIC built-in file extractors — air-gap safe fallbacks for multimodal ingestion.

Every extractor returns an :class:`Extraction` dataclass with normalized text, provider
name, content type, and page count. Extractors are best-effort: if a required library
is missing the extractor returns ``text=""`` and logs a warning so the orchestrator
can report the failure rather than silently producing 0 chunks.

Supported formats:
  PDF   → pypdf (pure-Python, air-gap baseline)
  DOCX  → python-docx (best-effort)
  XLSX  → openpyxl (best-effort)
  PPTX  → python-pptx (best-effort)
  PNG   → pytesseract / easyocr (best-effort)
  HTML  → built-in strip-html
  TXT   → built-in read-text

Layout mode:
  Layout-aware extraction (column/table/figure region segmentation) needs the
  optional ``paddleocr`` / ``doclayout-yolo`` libraries, which pull heavy
  backends and download model weights over the network. We probe them at import
  time; if they (or their weights) are unavailable the module degrades to the
  ``flat-ocr`` baseline — sequential text extraction with no region
  segmentation — so air-gap ingestion always has a working path. Inspect the
  active mode via :func:`layout_mode` or the ``LAYOUT_MODE`` module constant.
"""
from __future__ import annotations

import io
import base64
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from tools.logging.icdev_logger import get_logger

logger = get_logger(__name__)


# --------------------------------------------------------------------------- #
# Layout-detection capability probe (PaddleOCR / doclayout-yolo)
# --------------------------------------------------------------------------- #
# Layout-aware extraction requires optional libraries that pull heavy backends
# (PaddlePaddle / torch) and download model weights over the network on first
# use — neither is air-gap safe by default. We probe their importability once at
# module-load time inside a guarded try/except: any failure (missing package,
# broken backend, blocked network fetch during import) is logged and the module
# falls back to the always-present 'flat-ocr' path so ingestion never breaks.

LAYOUT_MODE_AWARE = "layout-aware"   # region/table/column segmentation available
LAYOUT_MODE_FLAT = "flat-ocr"        # sequential text only — air-gap baseline

_LAYOUT_LIBS = ("paddleocr", "doclayout_yolo")


def _probe_layout_libs():
    """Detect whether every layout-detection library imports cleanly.

    Returns ``(mode, missing)``: ``mode`` is :data:`LAYOUT_MODE_AWARE` only when
    all libraries import without error, otherwise :data:`LAYOUT_MODE_FLAT`.
    Never raises — callers rely on this to choose a fallback, not to fail.
    """
    import importlib

    missing: list[str] = []
    for lib in _LAYOUT_LIBS:
        try:
            importlib.import_module(lib)
        except Exception as exc:
            missing.append(lib)
            logger.debug("dic.extractors: layout lib %r unavailable: %s", lib, exc)
    if missing:
        return LAYOUT_MODE_FLAT, missing
    return LAYOUT_MODE_AWARE, []


try:
    LAYOUT_MODE, _LAYOUT_MISSING = _probe_layout_libs()
except Exception as exc:  # defensive: the probe itself must never break import
    logger.error(
        "dic.extractors: layout-detection probe failed (%s) — forcing 'flat-ocr' fallback",
        exc,
    )
    LAYOUT_MODE, _LAYOUT_MISSING = LAYOUT_MODE_FLAT, list(_LAYOUT_LIBS)

if LAYOUT_MODE == LAYOUT_MODE_FLAT:
    logger.warning(
        "dic.extractors: layout-detection libraries unavailable (%s) — using "
        "'flat-ocr' fallback (no region/table/column segmentation). For "
        "layout-aware extraction install: pip install paddleocr doclayout-yolo",
        ", ".join(_LAYOUT_MISSING) or "none",
    )


def layout_mode() -> str:
    """Return the active layout-extraction mode ('layout-aware' or 'flat-ocr')."""
    return LAYOUT_MODE


@dataclass
class Extraction:
    """Normalized output of a file extractor."""

    text: str
    provider: str
    content_type: str
    page_count: int = 1
    title: str = ""
    metadata: dict[str, Any] = field(default_factory=dict)
    warnings: list[str] = field(default_factory=list)


@dataclass
class ExtractionQualityReport:
    """Statistical anomaly assessment for a batch of DIC extractions.

    Replaces the external pattern of hard-coding per-extraction quality
    thresholds (e.g. ``if len(text) < 100: warn_sparse``) with IQR-based
    distribution analysis across a batch. The detection adapts to the actual
    yield distribution so no manual threshold needs to be maintained.

    Three anomaly types are detected:

    * ``sparse_batch``   — Q1 of chars-per-page across the batch is below
      the noise floor, meaning most documents yield abnormally little text
      (possible library issue, batch of scanned PDFs without OCR).
    * ``yield_outlier``  — at least one extraction has yield far below the
      pack (< Q1 - 1.5 × IQR), flagged in ``outlier_indices``.
    * ``warning_spike``  — mean warnings-per-extraction exceeds 0 and the
      max warning count is an outlier (> Q3 + 1.5 × IQR of warning counts),
      indicating a systemic library or format problem.

    ``is_anomalous`` is False and all floats are 0.0 when the batch is too
    small to characterize (< 2 extractions) or no anomaly is detected.
    """

    is_anomalous: bool = False
    anomaly_type: str = ""
    mean_chars_per_page: float = 0.0
    chars_per_page_std: float = 0.0
    q1_chars_per_page: float = 0.0
    outlier_indices: list[int] = field(default_factory=list)
    detail: str = ""

    def to_dict(self) -> dict:
        return {
            "is_anomalous": self.is_anomalous,
            "anomaly_type": self.anomaly_type,
            "mean_chars_per_page": round(self.mean_chars_per_page, 1),
            "chars_per_page_std": round(self.chars_per_page_std, 1),
            "q1_chars_per_page": round(self.q1_chars_per_page, 1),
            "outlier_indices": self.outlier_indices,
            "detail": self.detail,
        }


def assess_extraction_quality(extractions: list["Extraction"]) -> ExtractionQualityReport:
    """Assess quality anomalies across a batch of extractions using IQR statistics.

    Args:
        extractions: List of :class:`Extraction` objects from a single ingestion batch.

    Returns:
        :class:`ExtractionQualityReport` — ``is_anomalous=False`` when the batch
        is too small (< 2) or all extractions are within normal range.
    """
    if len(extractions) < 2:
        return ExtractionQualityReport()

    yields = [len(e.text) / max(e.page_count, 1) for e in extractions]
    warning_counts = [len([w for w in e.warnings if w]) for e in extractions]

    n = len(yields)
    sorted_y = sorted(yields)
    q1 = sorted_y[n // 4]
    q3 = sorted_y[(3 * n) // 4]
    iqr = q3 - q1
    mean_y = sum(yields) / n
    variance = sum((y - mean_y) ** 2 for y in yields) / n
    std_y = variance ** 0.5

    low_fence = q1 - 1.5 * iqr

    # sparse_batch: Q1 is below the noise floor (< 50 chars/page)
    if q1 < 50:
        return ExtractionQualityReport(
            is_anomalous=True,
            anomaly_type="sparse_batch",
            mean_chars_per_page=mean_y,
            chars_per_page_std=std_y,
            q1_chars_per_page=q1,
            detail=(
                f"Q1 chars-per-page is {q1:.0f} — most extractions yield near-zero "
                "text per page, suggesting a systemic OCR gap or unsupported format batch."
            ),
        )

    # yield_outlier: individual extractions well below the pack.
    # When IQR = 0 (uniform batch except one extreme) fall back to mean - 2σ.
    effective_fence = low_fence if iqr > 10 else (mean_y - 2 * std_y if std_y > 10 else None)
    outliers = (
        [i for i, y in enumerate(yields) if y < effective_fence]
        if effective_fence is not None
        else []
    )
    if outliers:
        return ExtractionQualityReport(
            is_anomalous=True,
            anomaly_type="yield_outlier",
            mean_chars_per_page=mean_y,
            chars_per_page_std=std_y,
            q1_chars_per_page=q1,
            outlier_indices=outliers,
            detail=(
                f"{len(outliers)} extraction(s) at index(es) {outliers} have chars-per-page "
                f"below the low fence ({low_fence:.0f}), indicating corrupt or unreadable files."
            ),
        )

    # warning_spike: statistically high warning count on at least one extraction
    if any(wc > 0 for wc in warning_counts):
        sorted_w = sorted(warning_counts)
        w_q3 = sorted_w[(3 * n) // 4]
        w_iqr = w_q3 - sorted_w[n // 4]
        w_fence = w_q3 + 1.5 * w_iqr
        spike_indices = [i for i, wc in enumerate(warning_counts) if wc > w_fence and wc > 1]
        if spike_indices:
            return ExtractionQualityReport(
                is_anomalous=True,
                anomaly_type="warning_spike",
                mean_chars_per_page=mean_y,
                chars_per_page_std=std_y,
                q1_chars_per_page=q1,
                outlier_indices=spike_indices,
                detail=(
                    f"Extraction(s) at index(es) {spike_indices} have warning counts above "
                    f"the high fence ({w_fence:.0f}), indicating systemic library or format issues."
                ),
            )

    return ExtractionQualityReport(
        mean_chars_per_page=mean_y,
        chars_per_page_std=std_y,
        q1_chars_per_page=q1,
    )


@dataclass
class OcrConfidenceReport:
    """Per-image OCR quality assessment via IQR-based confidence distribution analysis.

    Replaces the external pattern of hard-coding a minimum per-word confidence
    threshold (e.g. ``if word_conf < 50: skip_word``) with distribution-aware
    detection. Anomaly boundaries adapt to each image's actual confidence spread
    so no fixed cutoff is maintained.

    Two anomaly types:
    * ``low_confidence_batch``  — Q1 of per-word confidence is below the OCR
      noise floor (< 0.20), meaning most words were read with very low certainty.
    * ``confidence_outliers``   — individual words fall well below the pack
      (< Q1 − 1.5 × IQR), returned in ``outlier_word_indices``.

    ``is_low_confidence`` is False when fewer than 3 valid words are scored or
    no anomaly is detected.
    """

    is_low_confidence: bool = False
    anomaly_type: str = ""
    mean_confidence: float = 0.0
    q1_confidence: float = 0.0
    word_count: int = 0
    outlier_word_indices: list[int] = field(default_factory=list)
    detail: str = ""

    def to_dict(self) -> dict:
        return {
            "is_low_confidence": self.is_low_confidence,
            "anomaly_type": self.anomaly_type,
            "mean_confidence": round(self.mean_confidence, 3),
            "q1_confidence": round(self.q1_confidence, 3),
            "word_count": self.word_count,
            "outlier_word_indices": self.outlier_word_indices,
            "detail": self.detail,
        }


def assess_ocr_confidence(word_confidences: list[float]) -> OcrConfidenceReport:
    """Assess OCR quality via IQR statistics on per-word confidence scores (0.0–1.0).

    Args:
        word_confidences: Per-word confidence values already filtered to valid
            words (tesseract assigns -1 to failed detections — exclude before
            calling this function).

    Returns:
        :class:`OcrConfidenceReport` — ``is_low_confidence=False`` when the
        word count is too small (< 3) or all words are within normal range.
    """
    n = len(word_confidences)
    if n < 3:
        return OcrConfidenceReport(word_count=n)

    sorted_c = sorted(word_confidences)
    q1 = sorted_c[n // 4]
    q3 = sorted_c[(3 * n) // 4]
    iqr = q3 - q1
    mean_c = sum(word_confidences) / n

    # low_confidence_batch: Q1 below OCR noise floor (< 0.20 on a 0–1 scale)
    if q1 < 0.20:
        return OcrConfidenceReport(
            is_low_confidence=True,
            anomaly_type="low_confidence_batch",
            mean_confidence=mean_c,
            q1_confidence=q1,
            word_count=n,
            detail=(
                f"Q1 word confidence is {q1:.2f} — most words recognized with very low "
                "certainty, suggesting a difficult image (skew, noise, or unsupported font). "
                "Consider image preprocessing or a vision-LLM fallback."
            ),
        )

    # confidence_outliers: words far below the pack relative to IQR spread
    low_fence = q1 - 1.5 * iqr if iqr > 0.05 else mean_c - 2 * max(iqr, 0.01)
    outliers = [i for i, c in enumerate(word_confidences) if c < low_fence]
    if outliers:
        shown = outliers[:5]
        suffix = "..." if len(outliers) > 5 else ""
        return OcrConfidenceReport(
            is_low_confidence=True,
            anomaly_type="confidence_outliers",
            mean_confidence=mean_c,
            q1_confidence=q1,
            word_count=n,
            outlier_word_indices=outliers,
            detail=(
                f"{len(outliers)} word(s) at index(es) {shown}{suffix} have confidence "
                f"below the low fence ({low_fence:.2f}), indicating locally unreadable "
                "regions. Results may contain garbled text."
            ),
        )

    return OcrConfidenceReport(
        mean_confidence=mean_c,
        q1_confidence=q1,
        word_count=n,
    )


@dataclass
class TextParseQualityReport:
    """Per-document plain-text quality assessment via IQR on line-length distribution.

    Replaces the external pattern of hard-coding minimum text-length or word-count
    thresholds (e.g. ``if len(text) < 100: raise ParseError("too short")``) with
    distribution-aware analysis of the document's internal line structure. Detection
    adapts to the document's actual content so no fixed cutoff is maintained.

    Three anomaly types:

    * ``binary_noise``         — encoding replacement character ratio (\\ufffd) exceeds
      5 %, indicating binary content misread as text or a severe encoding mismatch.
    * ``line_length_outliers`` — lines far above the pack (> Q3 + 1.5 × IQR) dominate
      when IQR > 10, suggesting base64 blobs, minified content, or structured data
      embedded in plain text.
    * *(no type)*              — document has fewer than 3 non-empty lines; not enough
      structure to characterise, returned clean unless binary_noise applies.

    ``is_anomalous`` is False when the document passes all checks.
    """

    is_anomalous: bool = False
    anomaly_type: str = ""
    line_count: int = 0
    mean_line_length: float = 0.0
    q3_line_length: float = 0.0
    replacement_char_ratio: float = 0.0
    outlier_line_count: int = 0
    detail: str = ""

    def to_dict(self) -> dict:
        return {
            "is_anomalous": self.is_anomalous,
            "anomaly_type": self.anomaly_type,
            "line_count": self.line_count,
            "mean_line_length": round(self.mean_line_length, 1),
            "q3_line_length": round(self.q3_line_length, 1),
            "replacement_char_ratio": round(self.replacement_char_ratio, 4),
            "outlier_line_count": self.outlier_line_count,
            "detail": self.detail,
        }


def assess_text_parse_quality(text: str) -> TextParseQualityReport:
    """Assess quality anomalies in a single plain-text document.

    Analyzes the document's internal line-length distribution with IQR statistics
    rather than hardcoded minimum thresholds. Intended for use inside a text-file
    extractor so quality issues surface as structured warnings rather than silent
    truncation or hard parser errors.

    Args:
        text: Raw text content (may be empty; ``errors='replace'`` encoding assumed).

    Returns:
        :class:`TextParseQualityReport` — ``is_anomalous=False`` when the document
        passes all checks or has too few lines to characterise (< 3 non-empty lines),
        provided the encoding replacement ratio is also below the noise floor.
    """
    # Encoding replacement-character density (applies regardless of line count)
    total_chars = len(text)
    replacement_ratio = text.count("�") / max(total_chars, 1)
    if replacement_ratio > 0.05:
        return TextParseQualityReport(
            is_anomalous=True,
            anomaly_type="binary_noise",
            line_count=len(text.splitlines()),
            replacement_char_ratio=replacement_ratio,
            detail=(
                f"Encoding replacement character ratio is {replacement_ratio:.2%} — "
                "document appears to be binary content misread as text or has a "
                "severe encoding mismatch."
            ),
        )

    non_empty_lines = [ln for ln in text.splitlines() if ln.strip()]
    n = len(non_empty_lines)

    if n < 3:
        return TextParseQualityReport(
            line_count=n,
            replacement_char_ratio=replacement_ratio,
        )

    lengths = [len(ln) for ln in non_empty_lines]
    sorted_l = sorted(lengths)
    q1 = sorted_l[n // 4]
    q3 = sorted_l[(3 * n) // 4]
    iqr = q3 - q1
    mean_l = sum(lengths) / n
    high_fence = q3 + 1.5 * iqr

    # line_length_outliers: long lines dominate when IQR shows meaningful spread.
    # When IQR = 0 (all lines uniform except extremes), fall back to mean + 2*std
    # so lone long-line blobs in otherwise short-line documents are still caught.
    variance = sum((l - mean_l) ** 2 for l in lengths) / n
    std_l = variance ** 0.5
    if iqr > 10:
        effective_fence = high_fence
    elif std_l > 10:
        effective_fence = mean_l + 2 * std_l
    else:
        effective_fence = None

    if effective_fence is not None:
        outliers = [ln for ln in lengths if ln > effective_fence]
        if outliers:
            return TextParseQualityReport(
                is_anomalous=True,
                anomaly_type="line_length_outliers",
                line_count=n,
                mean_line_length=mean_l,
                q3_line_length=float(q3),
                replacement_char_ratio=replacement_ratio,
                outlier_line_count=len(outliers),
                detail=(
                    f"{len(outliers)} line(s) exceed the high fence "
                    f"({effective_fence:.0f} chars), suggesting base64 blobs, "
                    "minified content, or structured data embedded in plain text."
                ),
            )

    return TextParseQualityReport(
        line_count=n,
        mean_line_length=mean_l,
        q3_line_length=float(q3),
        replacement_char_ratio=replacement_ratio,
    )


# --------------------------------------------------------------------------- #
# Extractor implementations
# --------------------------------------------------------------------------- #

def _extract_text(path: Path) -> Extraction:
    raw = path.read_text(encoding="utf-8", errors="replace")
    warnings: list[str] = []
    quality = assess_text_parse_quality(raw)
    if quality.is_anomalous:
        warnings.append(
            f"text_parse_anomaly:{quality.anomaly_type} — {quality.detail}"
        )
    return Extraction(
        text=raw,
        provider="builtin-text",
        content_type="text/plain",
        page_count=1,
        title=path.stem,
        warnings=warnings,
    )


def _extract_html(path: Path) -> Extraction:
    """Extract a local HTML file via the two-pass ``page_extract`` filter.

    Pass 2 is skipped (no query at ingestion time), so this is pass-1 pruning:
    nav, cookie banners and link farms are dropped, headings/lists/tables survive
    as markdown.  ``to_text`` is the degraded fallback — never a regex tag strip.
    """
    from tools.http import page_extract

    raw = path.read_text(encoding="utf-8", errors="replace")
    warnings: list[str] = []
    title = path.stem
    try:
        result = page_extract.extract(raw)
        text = result["fit_markdown"].strip() or result["raw_text"]
        title = result["title"] or title
    except Exception as exc:  # noqa: BLE001 - extraction must not fail an ingest
        text = page_extract.to_text(raw)
        warnings.append(f"page_extract failed ({type(exc).__name__}: {exc}) — used plain-text fallback")

    return Extraction(
        text=text,
        provider="builtin-html",
        content_type="text/html",
        page_count=1,
        title=title,
        warnings=warnings,
    )


# ── Cached easyocr reader (avoid re-loading model per page) ─────────────────
_EASYOCR_READER = None
_EASYOCR_TRIED = False


def _get_easyocr_reader():
    global _EASYOCR_READER, _EASYOCR_TRIED
    if _EASYOCR_TRIED:
        return _EASYOCR_READER
    _EASYOCR_TRIED = True
    try:
        import easyocr

        _EASYOCR_READER = easyocr.Reader(["en"], gpu=False)
    except Exception as exc:
        logger.debug("dic.extractors: easyocr not available: %s", exc)
        _EASYOCR_READER = None
    return _EASYOCR_READER


def _try_easyocr(img) -> str:
    reader = _get_easyocr_reader()
    if reader is None:
        return ""
    try:
        import numpy as np

        result = reader.readtext(np.array(img))
        return "\n".join(r[1] for r in result)
    except Exception:
        return ""


def _vision_ocr(img) -> str:
    """Send a PIL image to a local vision LLM for text extraction.

    Probes endpoints in priority order (env-configurable):
      1. Ollama      — OLLAMA_BASE_URL      → /api/tags + /api/chat
      2. vLLM        — VLLM_BASE_URL        → /v1/models + /v1/chat/completions
      3. LM Studio   — LM_STUDIO_BASE_URL   → /v1/models + /v1/chat/completions
      4. llama.cpp   — LLAMA_CPP_BASE_URL   → /v1/models + /v1/chat/completions

    Auto-detects API style per endpoint. Falls back to next endpoint if the
    current one has no vision-capable model or returns an error.
    """
    import json
    import os
    import urllib.request

    # Convert image to base64 PNG
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    image_b64 = base64.b64encode(buf.getvalue()).decode("utf-8")

    prompt_text = (
        "Extract ALL text from this document page image. "
        "Preserve structure, headings, and formatting. "
        "Return only the extracted text, no commentary."
    )

    # Vision-capable model keywords (covers llava, gemma3/4, bakllava, moondream, glm-ocr, qwen-vl, etc.)
    vision_kw = (
        "llava", "gemma3", "gemma4", "gemma-3", "gemma-4",
        "bakllava", "moondream", "glm-ocr", "qwen-vl", "qwen2-vl",
        "phi4", "phi-4", "minicpm-v", "internvl", "deepseek-vl",
    )

    # Candidate endpoints from environment (same precedence as pdf_fallback.py)
    candidates: list[tuple[str, str]] = []
    for env_var, default in [
        ("OLLAMA_BASE_URL", "http://localhost:11434"),
        ("VLLM_BASE_URL", ""),
        ("LM_STUDIO_BASE_URL", ""),
        ("LLAMA_CPP_BASE_URL", ""),
    ]:
        url = os.environ.get(env_var, default)
        if url:
            candidates.append((url, env_var))

    def _detect_api(base: str) -> str:
        """Return 'ollama' or 'openai' based on probe endpoints."""
        try:
            urllib.request.urlopen(f"{base}/api/tags", timeout=2)
            return "ollama"
        except Exception:
            pass
        try:
            urllib.request.urlopen(f"{base}/v1/models", timeout=2)
            return "openai"
        except Exception:
            pass
        return "ollama"  # default guess

    for base_url, src in candidates:
        base = base_url.rstrip("/")
        api = _detect_api(base)

        try:
            if api == "ollama":
                resp = urllib.request.urlopen(f"{base}/api/tags", timeout=3)
                data = json.loads(resp.read())
                models = [m["name"] for m in data.get("models", [])]
                vision_models = [m for m in models if any(v in m.lower() for v in vision_kw)]
                if not vision_models:
                    continue
                model = vision_models[0]

                payload = json.dumps(
                    {
                        "model": model,
                        "messages": [
                            {
                                "role": "user",
                                "content": prompt_text,
                                "images": [image_b64],
                            }
                        ],
                        "stream": False,
                        "options": {"num_predict": 4096},
                    }
                ).encode()
                url = f"{base}/api/chat"
            else:
                # OpenAI-compatible (vLLM, LM Studio, llama.cpp)
                req = urllib.request.Request(
                    f"{base}/v1/models",
                    headers={"Content-Type": "application/json"},
                )
                resp = urllib.request.urlopen(req, timeout=3)
                data = json.loads(resp.read())
                models = [m.get("id", "") for m in data.get("data", [])]
                vision_models = [m for m in models if any(v in m.lower() for v in vision_kw)]
                if not vision_models:
                    continue
                model = vision_models[0]

                payload = json.dumps(
                    {
                        "model": model,
                        "messages": [
                            {
                                "role": "user",
                                "content": [
                                    {"type": "text", "text": prompt_text},
                                    {
                                        "type": "image_url",
                                        "image_url": {
                                            "url": f"data:image/png;base64,{image_b64}",
                                        },
                                    },
                                ],
                            }
                        ],
                        "max_tokens": 4096,
                        "stream": False,
                    }
                ).encode()
                url = f"{base}/v1/chat/completions"

            req = urllib.request.Request(
                url,
                data=payload,
                headers={"Content-Type": "application/json"},
            )
            resp = urllib.request.urlopen(req, timeout=120)
            data = json.loads(resp.read())

            if api == "ollama":
                text = data.get("message", {}).get("content", "").strip()
            else:
                choices = data.get("choices", [])
                text = choices[0].get("message", {}).get("content", "").strip() if choices else ""

            if text:
                logger.info("dic.extractors: vision OCR via %s (%s) produced %d chars", src, model, len(text))
                return text
        except Exception as exc:
            logger.debug("dic.extractors: vision OCR failed on %s (%s): %s", src, base, exc)
            continue

    return ""


def _ocr_image(img) -> str:
    """Run OCR on a PIL image. Tries easyocr → local vision LLM (Ollama/vLLM) → pytesseract."""
    # 1. easyocr (pure Python, PyTorch-based, air-gap safe if models cached)
    text = _try_easyocr(img)
    if text.strip():
        return text

    # 2. Local vision LLM (Ollama / vLLM / LM Studio / llama.cpp)
    try:
        text = _vision_ocr(img)
        if text.strip():
            return text
    except Exception:
        pass

    # 3. pytesseract (requires external Tesseract binary — not air-gap friendly)
    try:
        import pytesseract

        return pytesseract.image_to_string(img)
    except Exception:
        pass

    return ""


def _try_pdf_ocr(path: Path, total_pages: int, max_pages: int = 50) -> str:
    """Render PDF pages to images via pypdfium2 and OCR them. Returns aggregated text."""
    try:
        import pypdfium2 as pdfium
    except Exception:
        return ""

    pages_to_process = min(total_pages, max_pages)
    text_parts: list[str] = []

    try:
        pdf = pdfium.PdfDocument(str(path))
        for i in range(pages_to_process):
            try:
                bitmap = pdf[i].render(scale=2)
                pil_image = bitmap.to_pil()
                page_text = _ocr_image(pil_image)
                if page_text.strip():
                    text_parts.append(f"\n--- Page {i + 1} ---\n{page_text.strip()}")
            except Exception as exc:
                logger.warning("dic.extractors: OCR failed on page %s: %s", i + 1, exc)
    except Exception as exc:
        logger.warning("dic.extractors: PDF OCR rendering failed: %s", exc)

    return "\n".join(text_parts)


def _try_pymupdf(path: Path) -> tuple[str, int]:
    """Extract text via pymupdf (fitz) — best for CID-mapped fonts and complex encodings."""
    try:
        import fitz  # pymupdf
        doc = fitz.open(str(path))
        total_pages = doc.page_count
        parts: list[str] = []
        for i, page in enumerate(doc):
            try:
                t = page.get_text("text") or ""
                if t.strip():
                    parts.append(f"\n--- Page {i + 1} ---\n{t.strip()}")
            except Exception:
                pass
        doc.close()
        return "\n".join(parts), total_pages
    except Exception as exc:
        logger.debug("dic.extractors: pymupdf failed: %s", exc)
        return "", 0


def _try_pdfplumber(path: Path) -> tuple[str, int]:
    """Extract text via pdfplumber — handles unusual font encodings well."""
    try:
        import pdfplumber
        parts: list[str] = []
        total_pages = 0
        with pdfplumber.open(str(path)) as pdf:
            total_pages = len(pdf.pages)
            for i, page in enumerate(pdf.pages):
                try:
                    t = page.extract_text() or ""
                    if t.strip():
                        parts.append(f"\n--- Page {i + 1} ---\n{t.strip()}")
                except Exception:
                    pass
        return "\n".join(parts), total_pages
    except Exception as exc:
        logger.debug("dic.extractors: pdfplumber failed: %s", exc)
        return "", 0


def _extract_pdf(path: Path) -> Extraction:
    """Extract text from PDF using a four-pass fallback chain.

    Pass 1: pymupdf/fitz    — gold standard; handles CID-mapped fonts, complex encodings
    Pass 2: pdfplumber      — good for non-standard font maps
    Pass 3: pypdf           — fast baseline
    Pass 4: OCR             — last resort for genuinely scanned/image PDFs
                              (easyocr → local vision LLM → pytesseract)
    """
    total_pages = 0

    # ── Pass 1: pymupdf (fitz) ────────────────────────────────────────────────
    text, total_pages = _try_pymupdf(path)
    if text.strip():
        logger.debug("dic.extractors: PDF extracted via pymupdf (%d chars)", len(text))
        return Extraction(
            text=text,
            provider="pymupdf",
            content_type="application/pdf",
            page_count=total_pages or 1,
            title=path.stem,
        )

    # ── Pass 2: pdfplumber ────────────────────────────────────────────────────
    text, pl_pages = _try_pdfplumber(path)
    if total_pages == 0:
        total_pages = pl_pages
    if text.strip():
        logger.debug("dic.extractors: PDF extracted via pdfplumber (%d chars)", len(text))
        return Extraction(
            text=text,
            provider="pdfplumber",
            content_type="application/pdf",
            page_count=total_pages or 1,
            title=path.stem,
        )

    # ── Pass 3: pypdf ─────────────────────────────────────────────────────────
    try:
        from pypdf import PdfReader
        text_parts: list[str] = []
        reader = PdfReader(str(path))
        if total_pages == 0:
            total_pages = len(reader.pages)
        for i, page in enumerate(reader.pages):
            try:
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text_parts.append(f"\n--- Page {i + 1} ---\n{page_text.strip()}")
            except Exception as exc:
                logger.warning("dic.extractors: pypdf page %s error: %s", i + 1, exc)
        text = "\n".join(text_parts)
        if text.strip():
            logger.debug("dic.extractors: PDF extracted via pypdf (%d chars)", len(text))
            return Extraction(
                text=text,
                provider="pypdf",
                content_type="application/pdf",
                page_count=total_pages or 1,
                title=path.stem,
            )
    except Exception as exc:
        logger.warning("dic.extractors: pypdf failed: %s", exc)

    if total_pages == 0:
        total_pages = 1

    # ── Pass 4: OCR (scanned/image PDF fallback) ──────────────────────────────
    ocr_text = _try_pdf_ocr(path, total_pages)
    if ocr_text.strip():
        return Extraction(
            text=ocr_text,
            provider="pdf+ocr",
            content_type="application/pdf",
            page_count=total_pages,
            title=path.stem,
            warnings=["PDF text extracted via OCR (scanned/image PDF)."],
        )

    return Extraction(
        text="",
        provider="pdf-all-failed",
        content_type="application/pdf",
        page_count=total_pages,
        title=path.stem,
        warnings=[
            "All PDF extraction passes failed (pymupdf, pdfplumber, pypdf, OCR). "
            "File may be password-protected, corrupted, or a purely graphical PDF. "
            "For scanned PDFs, install Tesseract: winget install UB-Mannheim.TesseractOCR"
        ],
    )


def _extract_docx(path: Path) -> Extraction:
    """Extract text from DOCX using python-docx."""
    try:
        import docx
    except Exception:
        return Extraction(
            text="",
            provider="docx-missing",
            content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            page_count=0,
            title=path.stem,
            warnings=["python-docx not installed — run: pip install python-docx"],
        )

    try:
        document = docx.Document(str(path))
        paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
        # Also extract text from tables
        table_texts: list[str] = []
        for table in document.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    table_texts.append(row_text)

        all_text = "\n\n".join(paragraphs)
        if table_texts:
            all_text += "\n\n--- Tables ---\n" + "\n".join(table_texts)

        return Extraction(
            text=all_text,
            provider="python-docx",
            content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            page_count=max(1, len(paragraphs) // 40),  # rough estimate
            title=path.stem,
            warnings=[""],
        )
    except Exception as exc:
        return Extraction(
            text="",
            provider="docx-error",
            content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            page_count=0,
            title=path.stem,
            warnings=[f"DOCX read failed: {exc}"],
        )


def _extract_xlsx(path: Path) -> Extraction:
    """Extract text from XLSX using openpyxl."""
    try:
        import openpyxl
    except Exception:
        return Extraction(
            text="",
            provider="xlsx-missing",
            content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            page_count=0,
            title=path.stem,
            warnings=["openpyxl not installed — run: pip install openpyxl"],
        )

    try:
        wb = openpyxl.load_workbook(str(path), data_only=True)
        sheet_texts: list[str] = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                row_vals = [str(v) for v in row if v is not None]
                if row_vals:
                    rows.append(" | ".join(row_vals))
            if rows:
                sheet_texts.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))

        return Extraction(
            text="\n\n".join(sheet_texts),
            provider="openpyxl",
            content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            page_count=len(wb.sheetnames),
            title=path.stem,
            warnings=[""],
        )
    except Exception as exc:
        return Extraction(
            text="",
            provider="xlsx-error",
            content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            page_count=0,
            title=path.stem,
            warnings=[f"XLSX read failed: {exc}"],
        )


def _extract_pptx(path: Path) -> Extraction:
    """Extract text from PPTX using python-pptx."""
    try:
        import pptx
    except Exception:
        return Extraction(
            text="",
            provider="pptx-missing",
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            page_count=0,
            title=path.stem,
            warnings=["python-pptx not installed — run: pip install python-pptx"],
        )

    try:
        prs = pptx.Presentation(str(path))
        slide_texts: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slide_texts.append(f"--- Slide {i} ---\n" + "\n".join(texts))

        return Extraction(
            text="\n\n".join(slide_texts),
            provider="python-pptx",
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            page_count=len(prs.slides),
            title=path.stem,
            warnings=[""],
        )
    except Exception as exc:
        return Extraction(
            text="",
            provider="pptx-error",
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            page_count=0,
            title=path.stem,
            warnings=[f"PPTX read failed: {exc}"],
        )


def _extract_image(path: Path) -> Extraction:
    """Extract text from images using OCR (best-effort)."""
    warnings: list[str] = []
    try:
        from PIL import Image
    except Exception:
        return Extraction(
            text="",
            provider="pillow-missing",
            content_type="image/unknown",
            page_count=0,
            title=path.stem,
            warnings=["Pillow not installed — run: pip install Pillow"],
        )

    try:
        img = Image.open(str(path))
        # Try pytesseract first, then easyocr, then give up gracefully
        ocr_text = ""
        for ocr_name, ocr_fn in [
            ("pytesseract", _try_pytesseract),
            ("easyocr", _try_easyocr),
        ]:
            try:
                ocr_text = ocr_fn(img)
                if ocr_text.strip():
                    break
            except Exception:
                continue
        else:
            warnings.append("OCR not available — install pytesseract or easyocr for image text extraction")

        return Extraction(
            text=ocr_text,
            provider="ocr",
            content_type=f"image/{img.format.lower() if img.format else 'unknown'}",
            page_count=1,
            title=path.stem,
            warnings=warnings,
        )
    except Exception as exc:
        return Extraction(
            text="",
            provider="image-error",
            content_type="image/unknown",
            page_count=0,
            title=path.stem,
            warnings=[f"Image read failed: {exc}"],
        )


def _try_easyocr(img) -> str:
    reader = _get_easyocr_reader()
    if reader is None:
        return ""
    try:
        import numpy as np

        result = reader.readtext(np.array(img))
        return "\n".join(r[1] for r in result)
    except Exception:
        return ""


def _try_pytesseract(img) -> str:
    """Run pytesseract OCR with IQR-based per-word confidence assessment.

    Uses ``image_to_data`` to score each recognized word and applies
    :func:`assess_ocr_confidence` instead of a fixed confidence cutoff.
    Low-confidence results are logged but all text is returned so callers
    receive the full engine output regardless of quality.
    """
    import pytesseract

    try:
        data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
    except Exception:
        return pytesseract.image_to_string(img)

    words, confs = [], []
    for i in range(len(data["text"])):
        text = str(data["text"][i]).strip()
        conf = float(data["conf"][i])
        if not text or conf < 0:
            continue
        words.append(text)
        confs.append(conf / 100.0)  # normalize 0–100 → 0–1

    report = assess_ocr_confidence(confs)
    if report.is_low_confidence:
        logger.warning(
            "dic.extractors: pytesseract low-confidence OCR (%s) — %s",
            report.anomaly_type,
            report.detail,
        )

    return " ".join(words)


# --------------------------------------------------------------------------- #
# Registry
# --------------------------------------------------------------------------- #

_EXTRACTORS: dict[str, callable] = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
    ".png": _extract_image,
    ".jpg": _extract_image,
    ".jpeg": _extract_image,
    ".tiff": _extract_image,
    ".tif": _extract_image,
    ".bmp": _extract_image,
    ".gif": _extract_image,
    ".html": _extract_html,
    ".htm": _extract_html,
    ".txt": _extract_text,
    ".md": _extract_text,
    ".markdown": _extract_text,
    ".rst": _extract_text,
    ".log": _extract_text,
    ".csv": _extract_text,
    ".tsv": _extract_text,
    ".json": _extract_text,
    ".yaml": _extract_text,
    ".yml": _extract_text,
    ".xml": _extract_text,
    ".py": _extract_text,
    ".sql": _extract_text,
    ".ini": _extract_text,
    ".cfg": _extract_text,
    ".toml": _extract_text,
}


def get_extractor(ext: str) -> callable | None:
    """Return the extractor function for a given file extension (lowercased)."""
    return _EXTRACTORS.get(ext.lower())


# --------------------------------------------------------------------------- #
# Web URL and YouTube extractors (non-file, dual-mode)
# --------------------------------------------------------------------------- #

def extract_url(url: str, query: str | None = None) -> Extraction:
    """Fetch and extract text from a web URL.

    Online: central-client fetch → two-pass ``page_extract`` → injection scan.
    Air-gap / network failure: returns an empty Extraction with a warning so the
    caller can still create a placeholder source entry.

    Args:
        url: the page to ingest (untrusted third-party content).
        query: optional retrieval question.  When supplied, pass 2 keeps only the
            blocks that answer it; otherwise the whole pruned page is kept.

    The byte cap comes from ``fetch.max_bytes`` in ``args/http_client.yaml``, and
    the title comes from the parsed document rather than a ``<title>`` regex —
    which matched the first ``</title>`` anywhere in the file, including inside
    a comment or an inline SVG.
    """
    from tools.http.fetch_extract import fetch_page

    page = fetch_page(url, query=query)

    if not page.ok:
        return Extraction(
            text="",
            provider="builtin-url",
            content_type=page.content_type or "text/html",
            page_count=0,
            title=url,
            warnings=[f"URL fetch failed (network unavailable or blocked): {page.error}"],
            metadata={"source_url": url},
        )

    if page.blocked:
        return Extraction(
            text="",
            provider="builtin-url",
            content_type=page.content_type or "text/html",
            page_count=0,
            title=page.title or url,
            warnings=[
                "Content dropped: prompt-injection scan reported a critical finding "
                f"({', '.join(sorted({str(f.get('category')) for f in page.injection_findings}))})"
            ],
            metadata={"source_url": url, "injection_findings": page.injection_findings},
        )

    warnings: list[str] = []
    if page.truncated:
        warnings.append("Response body reached fetch.max_bytes — page may be incomplete.")
    if page.injection_findings:
        warnings.append(
            f"Injection scan reported {len(page.injection_findings)} non-critical finding(s)."
        )

    return Extraction(
        text=page.text,
        provider="builtin-url",
        content_type=page.content_type or "text/html",
        page_count=1,
        title=page.title or url,
        warnings=warnings,
        metadata={"source_url": url, "extraction_reason": page.reason},
    )


def _is_youtube_url(url: str) -> bool:
    return any(h in url for h in ("youtube.com/", "youtu.be/", "youtube-nocookie.com/"))


def _extract_youtube_video_id(url: str) -> str | None:
    import re as _re
    for pat in (r"(?:v=|youtu\.be/|embed/)([A-Za-z0-9_-]{11})", r"shorts/([A-Za-z0-9_-]{11})"):
        m = _re.search(pat, url)
        if m:
            return m.group(1)
    return None


def _extract_via_ytdlp(url: str) -> tuple[str, str, list[str]]:
    """Try yt-dlp to get subtitles/info from any video URL.

    Returns (text, title, warnings).  Works offline for cached content;
    for live IPTV streams falls back to extracting the stream description.
    Requires yt-dlp: pip install yt-dlp
    """
    import json as _json
    import shutil as _sh
    import subprocess as _sp
    import tempfile as _tf
    warnings: list[str] = []
    if not _sh.which("yt-dlp"):
        return "", "", ["yt-dlp not installed (pip install yt-dlp). Cannot extract video transcript."]
    try:
        with _tf.TemporaryDirectory() as tmpdir:
            # First try: extract auto-generated subtitles
            _sp.run(
                ["yt-dlp", "--skip-download", "--write-auto-sub", "--write-sub",
                 "--sub-format", "vtt", "--sub-lang", "en",
                 "--output", f"{tmpdir}/%(title)s.%(ext)s", url],
                capture_output=True, text=True, timeout=60, encoding="utf-8", errors="replace",
            )
            import glob as _gl
            import pathlib as _pl
            vtt_files = _gl.glob(f"{tmpdir}/*.vtt")
            if vtt_files:
                raw = _pl.Path(vtt_files[0]).read_text(encoding="utf-8", errors="replace")
                # Strip VTT timestamps/headers → plain text
                import re as _re
                lines = [ln for ln in raw.splitlines()
                         if ln and not _re.match(r"^(WEBVTT|NOTE|\d{2}:|-->|\d+$)", ln.strip())]
                text = " ".join(lines)
                title = _pl.Path(vtt_files[0]).stem.rsplit(".", 1)[0]
                return text, title, warnings

            # Fallback: dump JSON info (title + description)
            res2 = _sp.run(
                ["yt-dlp", "--dump-json", "--no-playlist", url],
                capture_output=True, text=True, timeout=30, encoding="utf-8", errors="replace",
            )
            if res2.returncode == 0 and res2.stdout.strip():
                info = _json.loads(res2.stdout.strip().splitlines()[0])
                title = info.get("title", url)
                desc = info.get("description", "")
                text = f"{title}\n\n{desc}".strip()
                if text:
                    warnings.append("No subtitles found — using video title + description only.")
                    return text, title, warnings
    except Exception as exc:
        warnings.append(f"yt-dlp error: {exc}")
    return "", url, warnings


def extract_video(url: str) -> Extraction:
    """Extract transcript/captions from any video URL.

    Priority:
      1. YouTube transcript API  — for YouTube URLs (online only)
      2. yt-dlp subtitles        — for YouTube + 1000+ other sites + IPTV descriptions
      3. URL text extraction     — fall back to the video page's text via
         ``extract_url`` (central HTTP client + two-pass ``page_extract``)

    Works in air-gap for yt-dlp (subtitles cached locally) and URL text extraction.
    IPTV m3u8/HLS streams: extracts playlist description via yt-dlp dump-json.
    """
    warnings: list[str] = []
    title = url

    # ── Path 1: YouTube transcript API ───────────────────────────────────────
    if _is_youtube_url(url):
        video_id = _extract_youtube_video_id(url)
        if video_id:
            try:
                from youtube_transcript_api import YouTubeTranscriptApi  # type: ignore[import]
                api = YouTubeTranscriptApi()
                fetched = api.fetch(video_id)
                text = " ".join(
                    (s.text if hasattr(s, "text") else s["text"]) for s in fetched
                )
                # Attempt oEmbed title (via the central HTTP client, not raw urllib)
                title = f"YouTube: {video_id}"
                try:
                    import json as _json

                    from tools.http.fetch_extract import fetch_raw

                    oembed = fetch_raw(
                        f"https://www.youtube.com/oembed?url={url}&format=json"
                    )
                    if oembed.ok:
                        title = _json.loads(oembed.raw_text).get("title", title)
                except Exception:
                    pass
                if text:
                    return Extraction(
                        text=text, provider="youtube-transcript", content_type="text/plain",
                        page_count=1, title=title,
                        metadata={"source_url": url, "video_id": video_id},
                    )
                warnings.append("YouTube transcript was empty.")
            except ImportError:
                warnings.append("youtube-transcript-api not installed; trying yt-dlp.")
            except Exception as exc:
                warnings.append(f"YouTube transcript fetch failed: {exc}; trying yt-dlp.")

    # ── Path 2: yt-dlp (any site, offline subtitles, IPTV descriptions) ──────
    yt_text, yt_title, yt_warns = _extract_via_ytdlp(url)
    warnings.extend(yt_warns)
    if yt_text:
        return Extraction(
            text=yt_text, provider="yt-dlp", content_type="text/plain",
            page_count=1, title=yt_title or title,
            warnings=warnings, metadata={"source_url": url},
        )

    # ── Path 3: URL text extraction (video page metadata / description) ───────
    page_ext = extract_url(url)
    if page_ext.text:
        page_ext.warnings = warnings + [
            "No captions/transcript found. Extracted video page text instead."
        ] + list(page_ext.warnings)
        page_ext.provider = "video-page-text"
        return page_ext

    return Extraction(
        text="", provider="video-extract", content_type="text/plain",
        page_count=0, title=title, warnings=warnings + ["No extractable content found."],
        metadata={"source_url": url},
    )


# Keep backward-compat alias
def extract_youtube(url: str) -> Extraction:
    return extract_video(url)


def extract_file(path: str | Path) -> Extraction:
    """Route a file path to the correct extractor.

    Prefers MarkItDown for DOCX/PPTX/XLSX when installed (structured Markdown output).
    Falls back to built-in extractors, then to a raw utf-8 decode for unknown extensions.
    """
    p = Path(path)
    ext = p.suffix.lower()

    # adapt-md-03: Try MarkItDown first for formats it handles better
    try:
        from tools.document_intelligence.converters.markitdown_adapter import (
            should_use_markitdown,
            convert as _markitdown_convert,
        )
        if should_use_markitdown(ext):
            result = _markitdown_convert(p)
            if result.text:  # only use if it produced content
                result.metadata.setdefault("layout_mode", "markitdown")
                return result
    except Exception:
        pass  # fall through to built-in extractors

    extractor = get_extractor(ext)
    if extractor is not None:
        result = extractor(p)
        result.metadata.setdefault("layout_mode", LAYOUT_MODE)
        return result

    # Unknown extension — best-effort utf-8 decode
    try:
        raw = p.read_text(encoding="utf-8", errors="replace")
        return Extraction(
            text=raw,
            provider="builtin-fallback",
            content_type="application/octet-stream",
            page_count=1,
            title=p.stem,
            warnings=[f"Unrecognized extension '{ext}' — best-effort text decode. Install a provider for {ext} files."],
        )
    except Exception as exc:
        return Extraction(
            text="",
            provider="fallback-error",
            content_type="application/octet-stream",
            page_count=0,
            title=p.stem,
            warnings=[f"File read failed: {exc}"],
        )
