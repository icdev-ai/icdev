# CUI // SP-CTI
"""Cell-level extraction that keeps the formulas and the coordinates.

Why this exists alongside tools/document_intelligence/extractors.py rather than
inside it:

  * ``_extract_xlsx`` loads with ``data_only=True`` — the formulas are gone — and
    joins cells with ``" | "`` — the coordinates are gone.
  * ``_extract_pptx`` keeps only shapes that have ``.text``. A table is a
    ``GraphicFrame`` and has none, so PPTX tables are dropped entirely.

Both are the right choices for feeding a RAG index, and both are fatal here.
Without formulas you cannot tell a licence that has been counted twice from one
you are buying two of — the difference lives entirely in which SUM() consumes
which cell. Without coordinates you cannot answer "where did this number come
from?" with anything better than a filename, and a number you cannot source is a
number a CFO is right to ignore.

Bending ``extract_file()`` to carry cell provenance would change the contract
every existing DIC and RAG consumer depends on. So this reads the grid, and
``extract_file()`` keeps reading the prose. Two extractors, two jobs.

Public API::

    extract_grid(path) -> GridExtraction
"""
from __future__ import annotations

import csv
import re
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from tools.bom.constants import REPRESENTATION_FIDELITY

# Excel serialises money as a bare float and puts the currency in the number
# format; text sources embed it in the string. Both have to land on the same
# numeric.
_MONEY_CHARS = re.compile(r"[,$€£\s]")
_PAREN_NEGATIVE = re.compile(r"^\((.*)\)$")


@dataclass
class Cell:
    """One cell, with everything that lets us cite it and audit it later."""

    sheet: str
    locator: str          # "A9", or "slide9!tbl1!r4c3", or "p3!t1!r5c2"
    row: int
    col: int
    value_text: str = ""
    value_num: float | None = None
    formula: str = ""     # "=C9*D9" — the whole reason this module exists
    number_format: str = ""

    @property
    def is_formula(self) -> bool:
        return bool(self.formula)

    @property
    def is_blank(self) -> bool:
        return not self.value_text and self.value_num is None


@dataclass
class GridExtraction:
    path: str
    filename: str
    media_type: str
    representation: str          # keys of REPRESENTATION_FIDELITY
    sheets: list[str] = field(default_factory=list)
    cells: list[Cell] = field(default_factory=list)
    # drawio only: a diagram claims components, it does not price them.
    nodes: list[dict] = field(default_factory=list)
    edges: list[dict] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)
    # Real-world messiness is reported, never swallowed. A source we could only
    # partly read must say so, or a missing sheet reads as an empty one.
    warnings: list[str] = field(default_factory=list)

    @property
    def fidelity(self) -> int:
        return REPRESENTATION_FIDELITY.get(self.representation, 0)

    @property
    def has_formulas(self) -> bool:
        return any(c.is_formula for c in self.cells)

    def sheet_cells(self, sheet: str) -> list[Cell]:
        return [c for c in self.cells if c.sheet == sheet]


def _to_num(value: Any) -> float | None:
    """Parse a number out of whatever a document happened to put there.

    Handles the three things real spreadsheets and PDFs actually do: a native
    float, a formatted string like "$1,234.00", and "(1,200)" — which is how
    accountants write negative twelve hundred.
    """
    if value is None or isinstance(value, bool):
        return None
    if isinstance(value, (int, float)):
        return float(value)

    text = str(value).strip()
    if not text:
        return None

    negative = False
    m = _PAREN_NEGATIVE.match(text)
    if m:
        negative, text = True, m.group(1)

    text = _MONEY_CHARS.sub("", text)
    if text.endswith("%"):
        # A percentage is a number, but silently dividing it by 100 here would
        # turn a 40% discount into $0.40 somewhere downstream. Leave it to the
        # caller who knows what the column means.
        text = text[:-1]

    if not text or not re.fullmatch(r"-?\d*\.?\d+", text):
        return None
    try:
        n = float(text)
    except ValueError:
        return None
    return -n if negative else n


def _extract_xlsx(path: Path, out: GridExtraction) -> None:
    """Load the workbook TWICE — once for values, once for formulas.

    openpyxl gives you one or the other, never both: ``data_only=True`` returns
    the last value Excel cached, ``data_only=False`` returns the formula string.
    The double load is the only way to hold "this cell shows 0" and "this cell
    says =C9*D9" at the same time — and that pair is what exposes a quantity
    multiplied by a price nobody ever filled in.
    """
    try:
        import openpyxl
    except ImportError:  # pragma: no cover
        out.warnings.append("openpyxl not installed; cannot read .xlsx")
        return

    wb_val = openpyxl.load_workbook(str(path), data_only=True)
    wb_fml = openpyxl.load_workbook(str(path), data_only=False)

    out.sheets = list(wb_val.sheetnames)
    out.metadata["sheet_states"] = {
        ws.title: ws.sheet_state for ws in wb_val.worksheets
    }

    saw_formula = False
    for name in wb_val.sheetnames:
        ws_v = wb_val[name]
        ws_f = wb_fml[name] if name in wb_fml.sheetnames else None

        for row in ws_v.iter_rows():
            for cv in row:
                cf = ws_f.cell(cv.row, cv.column) if ws_f is not None else None
                raw_f = cf.value if cf is not None else None
                formula = raw_f if isinstance(raw_f, str) and raw_f.startswith("=") else ""
                if formula:
                    saw_formula = True

                # A cell that is blank in BOTH loads carries nothing. A cell that
                # is blank in the value load but holds a formula does — that is
                # exactly the zeroed line we are hunting.
                if cv.value is None and not formula:
                    continue

                out.cells.append(
                    Cell(
                        sheet=name,
                        locator=cv.coordinate,
                        row=cv.row,
                        col=cv.column,
                        value_text="" if cv.value is None else str(cv.value).strip(),
                        value_num=_to_num(cv.value),
                        formula=formula,
                        number_format=cv.number_format or "",
                    )
                )

    out.representation = "xlsx_formulas" if saw_formula else "xlsx"

    props = wb_val.properties
    out.metadata.update(
        {
            "creator": props.creator or "",
            "last_modified_by": props.lastModifiedBy or "",
            "revision": props.revision or "",
            "created": str(props.created) if props.created else "",
            "modified": str(props.modified) if props.modified else "",
        }
    )


def _cluster(values: list[int], tolerance: int) -> list[int]:
    """Collapse near-equal coordinates into shared track positions.

    Shapes that a human reads as "the same row" are never at pixel-identical
    tops. Sort, then start a new track whenever the gap exceeds the tolerance.
    """
    if not values:
        return []
    tracks = [values[0]]
    for v in sorted(values)[1:]:
        if v - tracks[-1] > tolerance:
            tracks.append(v)
    return tracks


def _snap(value: int, tracks: list[int], tolerance: int) -> int | None:
    for i, t in enumerate(tracks):
        if abs(value - t) <= tolerance:
            return i
    return None


def _reconstruct_grid(shapes: list[tuple[int, int, int, int, str]]) -> list[tuple[int, int, str]]:
    """Rebuild a table from text boxes that only LOOK like one.

    A cost-summary slide is very often not a table at all. It is a few dozen
    separate text boxes positioned in a grid, and ``shape.has_table`` returns
    False for every one of them.

    That is not an exotic edge case — it is what you get from any deck a script
    generated, and from plenty that humans made by duplicating a text box. An
    extractor that only understands real tables reads those slides as loose
    words, and the single most important number on them goes missing without a
    warning.

    So: cluster the shapes' top edges into rows and their left edges into
    columns, and read the grid back out. Geometry is all that is left to go on
    once the semantics have been thrown away.
    """
    if len(shapes) < 4:
        return []

    heights = sorted(h for _, _, _, h, _ in shapes)
    widths = sorted(w for _, _, w, _, _ in shapes)
    row_tol = max(heights[len(heights) // 2] // 2, 1)
    col_tol = max(widths[len(widths) // 2] // 4, 1)

    row_tracks = _cluster([t for _, t, _, _, _ in shapes], row_tol)
    col_tracks = _cluster([left for left, _, _, _, _ in shapes], col_tol)

    # One row, or one column, is a list or a caption — not a table. Demanding at
    # least a 2x2 keeps us from turning a bulleted slide into fictional cells.
    if len(row_tracks) < 2 or len(col_tracks) < 2:
        return []

    out: list[tuple[int, int, str]] = []
    for left, top, _, _, text in shapes:
        r = _snap(top, row_tracks, row_tol)
        c = _snap(left, col_tracks, col_tol)
        if r is None or c is None:
            continue
        out.append((r + 1, c + 1, text))
    return out


def _extract_pptx(path: Path, out: GridExtraction) -> None:
    """Read real tables AND the ones drawn out of loose text boxes.

    A PowerPoint table is a ``GraphicFrame`` with no ``.text``, so any extractor
    filtering on ``has_text_frame`` loses it. And a machine-generated deck often
    has no real tables at all — just text boxes arranged to look like one. Both
    failure modes end the same way: the roll-up on the money slide silently is
    not there.
    """
    try:
        from pptx import Presentation
    except ImportError:  # pragma: no cover
        out.warnings.append("python-pptx not installed; cannot read .pptx")
        return

    prs = Presentation(str(path))
    saw_table = False

    for s_idx, slide in enumerate(prs.slides, start=1):
        sheet = f"slide{s_idx}"
        out.sheets.append(sheet)
        t_idx = 0
        loose: list[tuple[int, int, int, int, str]] = []

        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                t_idx += 1
                saw_table = True
                for r, row in enumerate(shape.table.rows, start=1):
                    for c, cell in enumerate(row.cells, start=1):
                        text = (cell.text or "").strip()
                        if not text:
                            continue
                        out.cells.append(
                            Cell(sheet=sheet, locator=f"tbl{t_idx}!r{r}c{c}",
                                 row=r, col=c, value_text=text, value_num=_to_num(text))
                        )
                continue

            if not shape.has_text_frame:
                continue
            text = (shape.text_frame.text or "").strip()
            if not text or shape.top is None or shape.left is None:
                continue
            loose.append((
                int(shape.left), int(shape.top),
                int(shape.width or 0), int(shape.height or 0),
                text,
            ))

        for r, c, text in _reconstruct_grid(loose):
            saw_table = True
            out.cells.append(
                Cell(sheet=sheet, locator=f"grid!r{r}c{c}", row=r, col=c,
                     value_text=text, value_num=_to_num(text))
            )

    out.representation = "pptx_tables" if saw_table else "pptx"

    props = prs.core_properties
    out.metadata.update({
        "creator": props.author or "",
        "last_modified_by": props.last_modified_by or "",
        "revision": props.revision or "",
        "comments": props.comments or "",
    })


def _extract_pdf(path: Path, out: GridExtraction) -> None:
    """Tables only. Prose is extract_file()'s job, not ours.

    A PDF has no formulas — it is a photograph of a spreadsheet, taken after the
    arithmetic already happened. That is precisely why a PDF ranks low in
    REPRESENTATION_FIDELITY: everything that would let us audit the number has
    already been flattened out of it.
    """
    try:
        import pdfplumber
    except ImportError:  # pragma: no cover
        out.warnings.append("pdfplumber not installed; cannot read tables from .pdf")
        out.representation = "pdf"
        return

    out.representation = "pdf"
    with pdfplumber.open(str(path)) as pdf:
        out.metadata.update(
            {k: str(v) for k, v in (pdf.metadata or {}).items()}
        )
        for p_idx, page in enumerate(pdf.pages, start=1):
            sheet = f"p{p_idx}"
            try:
                tables = page.extract_tables() or []
            except Exception as exc:  # pragma: no cover - pdfplumber edge cases
                out.warnings.append(f"page {p_idx}: table extraction failed ({exc})")
                continue
            if tables:
                out.sheets.append(sheet)
            for t_idx, table in enumerate(tables, start=1):
                for r, row in enumerate(table, start=1):
                    for c, val in enumerate(row, start=1):
                        text = (val or "").strip()
                        if not text:
                            continue
                        out.cells.append(
                            Cell(
                                sheet=sheet,
                                locator=f"t{t_idx}!r{r}c{c}",
                                row=r,
                                col=c,
                                value_text=text,
                                value_num=_to_num(text),
                            )
                        )


def _extract_csv(path: Path, out: GridExtraction) -> None:
    out.representation = "csv"
    sheet = path.stem
    out.sheets.append(sheet)
    with path.open("r", encoding="utf-8", errors="replace", newline="") as fh:
        for r, row in enumerate(csv.reader(fh), start=1):
            for c, val in enumerate(row, start=1):
                text = (val or "").strip()
                if not text:
                    continue
                out.cells.append(
                    Cell(sheet=sheet, locator=f"r{r}c{c}", row=r, col=c,
                         value_text=text, value_num=_to_num(text))
                )


def _drawio_xml(path: Path) -> tuple[str, list[str]]:
    """Return the diagram XML, inflating draw.io's compressed payload if needed.

    draw.io routinely writes ``<diagram>`` as deflate+base64. ``parse_drawio``
    takes a plain XML string, so handing it a compressed file yields zero nodes
    and no error at all — a silently empty architecture, which is worse than a
    crash because it looks like a diagram with nothing in it.
    """
    warnings: list[str] = []
    raw = path.read_text(encoding="utf-8", errors="replace")
    if "<mxCell" in raw or "<mxGraphModel" in raw:
        return raw, warnings

    import base64
    import zlib
    from urllib.parse import unquote
    from xml.etree import ElementTree as ET

    try:
        root = ET.fromstring(raw)
    except ET.ParseError as exc:
        warnings.append(f"drawio: not parseable as XML ({exc})")
        return raw, warnings

    parts: list[str] = []
    for diagram in root.iter("diagram"):
        payload = (diagram.text or "").strip()
        if not payload:
            continue
        try:
            inflated = zlib.decompress(base64.b64decode(payload), -15)
            parts.append(unquote(inflated.decode("utf-8")))
        except Exception as exc:
            warnings.append(f"drawio: could not inflate <diagram> payload ({exc})")
    if not parts:
        warnings.append("drawio: no readable diagram content")
        return raw, warnings
    return "<root>" + "".join(parts) + "</root>", warnings


def _extract_drawio(path: Path, out: GridExtraction) -> None:
    """A diagram CLAIMS components. It never prices them.

    Nodes land in ``out.nodes``, not ``out.cells``, and that separation is
    deliberate: a rack elevation drawing twelve servers is an assertion that
    twelve servers *should be there*. It is not evidence that twelve exist, and
    it is not a line item. Treating it as either is how a drawing quietly becomes
    a budget.

    Each tab is parsed SEPARATELY and tagged with its name. A .drawio file
    routinely holds a floor plan, a logical topology and a rack elevation, and
    those are three different kinds of claim: the rack elevation is where the
    unit counts live, and merging it into the floor plan would throw away the
    only place the diagram says *how many*.
    """
    from xml.etree import ElementTree as ET

    from tools.simulation.parsers.drawio_parser import parse_drawio

    out.representation = "drawio"
    xml, warnings = _drawio_xml(path)
    out.warnings.extend(warnings)

    diagrams: list[tuple[str, str]] = []  # (tab name, mxGraphModel xml)
    try:
        root = ET.fromstring(xml)
        for i, diagram in enumerate(root.iter("diagram"), start=1):
            model = diagram.find("mxGraphModel")
            if model is None:
                continue
            name = diagram.get("name") or f"diagram{i}"
            diagrams.append((name, ET.tostring(model, encoding="unicode")))
    except ET.ParseError as exc:
        out.warnings.append(f"drawio: could not split diagrams ({exc})")

    # A bare <mxGraphModel> (or an inflated payload we already unwrapped) has no
    # <diagram> wrapper — hand the whole thing over as a single unnamed tab.
    if not diagrams:
        diagrams = [("", xml)]

    for name, model_xml in diagrams:
        parsed = parse_drawio(model_xml) or {}
        for node in parsed.get("nodes", []) or []:
            node["diagram"] = name
            out.nodes.append(node)
        for edge in parsed.get("edges", []) or []:
            edge["diagram"] = name
            out.edges.append(edge)
        if name:
            out.sheets.append(name)

    out.metadata["diagrams"] = [n for n, _ in diagrams if n]

    if not out.nodes:
        out.warnings.append(
            "drawio: parsed zero nodes — the diagram is either empty or the "
            "payload could not be inflated. An empty architecture reads as "
            "'nothing was agreed', which is almost never true, so this is "
            "reported rather than passed on as a diagram with nothing in it."
        )


_DISPATCH = {
    ".xlsx": _extract_xlsx,
    ".xlsm": _extract_xlsx,
    ".pptx": _extract_pptx,
    ".pdf": _extract_pdf,
    ".csv": _extract_csv,
    ".tsv": _extract_csv,
    ".drawio": _extract_drawio,
    ".xml": _extract_drawio,
}

_MEDIA = {
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".xlsm": "application/vnd.ms-excel.sheet.macroEnabled.12",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".pdf": "application/pdf",
    ".csv": "text/csv",
    ".drawio": "application/vnd.jgraph.mxfile",
}


def _msip_label(path: Path) -> str:
    """Pull a Microsoft sensitivity label out of an OOXML package, if present.

    Not security theatre: a file stamped "Proprietary" that is about to be pasted
    into a cloud LLM is a decision someone should get to make on purpose.
    """
    if path.suffix.lower() not in (".xlsx", ".xlsm", ".pptx", ".docx"):
        return ""
    try:
        with zipfile.ZipFile(path) as zf:
            for name in zf.namelist():
                if name.endswith("custom.xml"):
                    blob = zf.read(name).decode("utf-8", errors="replace")
                    m = re.search(r"MSIP_Label_[0-9a-f-]+_Name[^>]*>\s*<[^>]*>([^<]+)", blob)
                    if m:
                        return m.group(1).strip()
    except (zipfile.BadZipFile, KeyError, OSError):
        return ""
    return ""


def extract_grid(path: str | Path) -> GridExtraction:
    """Read a document as a grid of cells, keeping formulas and coordinates.

    Unsupported formats return an empty extraction with a warning rather than
    raising: one unreadable file in a corpus of twenty should degrade that file,
    not the run.
    """
    p = Path(path)
    suffix = p.suffix.lower()

    out = GridExtraction(
        path=str(p),
        filename=p.name,
        media_type=_MEDIA.get(suffix, "application/octet-stream"),
        representation=suffix.lstrip(".") or "unknown",
    )

    if not p.exists():
        out.warnings.append(f"file not found: {p}")
        return out

    handler = _DISPATCH.get(suffix)
    if handler is None:
        out.warnings.append(
            f"no grid extractor for '{suffix}' — this file carries no line items. "
            "Prose is read separately by extract_file()."
        )
        return out

    try:
        handler(p, out)
    except Exception as exc:  # noqa: BLE001 - one bad file must not kill the run
        out.warnings.append(f"grid extraction failed: {type(exc).__name__}: {exc}")
        return out

    label = _msip_label(p)
    if label:
        out.metadata["sensitivity_label"] = label

    return out


def main() -> int:  # pragma: no cover - CLI convenience
    import argparse
    import json

    ap = argparse.ArgumentParser(description=__doc__.split("\n")[0])
    ap.add_argument("path")
    ap.add_argument("--json", action="store_true")
    ap.add_argument("--formulas", action="store_true", help="show only formula cells")
    args = ap.parse_args()

    g = extract_grid(args.path)
    cells = [c for c in g.cells if c.is_formula] if args.formulas else g.cells

    if args.json:
        print(json.dumps({
            "filename": g.filename,
            "representation": g.representation,
            "fidelity": g.fidelity,
            "has_formulas": g.has_formulas,
            "sheets": g.sheets,
            "cell_count": len(g.cells),
            "node_count": len(g.nodes),
            "metadata": g.metadata,
            "warnings": g.warnings,
            "cells": [vars(c) for c in cells[:200]],
        }, indent=2, default=str))
        return 0

    print(f"{g.filename}  [{g.representation}, fidelity {g.fidelity}]")
    print(f"  sheets: {', '.join(g.sheets) or '-'}")
    print(f"  cells: {len(g.cells)}  formulas: {sum(c.is_formula for c in g.cells)}"
          f"  nodes: {len(g.nodes)}")
    for w in g.warnings:
        print(f"  WARN {w}")
    for c in cells[:40]:
        f = f"  [{c.formula}]" if c.formula else ""
        print(f"  {c.sheet}!{c.locator} = {c.value_text[:44]!r}{f}")
    return 0


if __name__ == "__main__":  # pragma: no cover
    raise SystemExit(main())
