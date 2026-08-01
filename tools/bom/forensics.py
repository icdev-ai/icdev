# CUI // SP-CTI
"""What the documents are hiding.

Every artifact here is *in* the file and *not on the screen*. That is not a
metaphor, and none of it is exotic — it is what ordinary business documents look
like once you stop reading only the parts that render:

  * **A worksheet that reports its dimensions as ``A1:A1``** and still carries an
    anchored image. Excel says the sheet is empty. It opens empty. Nobody
    scrolling the workbook ever finds it. It is where somebody pasted a
    screenshot of a constraint that never made it into a cell — routinely a
    compliance limit or an unpriced cost, and now invisible to every reader and
    every tool that reads only cells.

  * **A gap in the sheetId sequence.** Excel never reuses a sheetId, so a hole in
    the run is a sheet that was deleted. In a costing workbook that is a category
    somebody removed, and the surviving subtotals will not tell you whether its
    money went with it. The give-away is usually a section header whose number no
    longer matches its own subtotal row.

  * **Speaker notes.** The talk track holds the constraint the slide was too
    polite to show — the long-lead item that gates the schedule, the caveat on
    the estimate, the thing still marked TBD. It renders nowhere.

  * **Documents written by a script.** ``openpyxl`` and ``python-pptx`` sign
    their work. A generated estimate is a different kind of evidence from one a
    person authored, and that difference belongs in the credibility assessment
    rather than nowhere.

Also: cell comments, links out to the real source, tracked changes, embedded
photographs, sensitivity labels, and values left unfinished in a document that is
being costed.
"""
from __future__ import annotations

import re
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

# Namespaces, spelled out once. OOXML is a swamp of these.
_NS = {
    "main": "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
    "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
    "xdr": "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}

# A source that stamps itself as machine-written is telling you something about
# how much its numbers are worth. openpyxl and python-pptx both sign their work.
_GENERATOR_MARKERS = ("openpyxl", "python-pptx", "reportlab", "Steve Canny")

# Words that mean "this number is not finished". Configurable vocabulary lives in
# args/bom_credibility.yaml; these are the structural ones that are language, not
# policy.
_PLACEHOLDER_RE = re.compile(
    # "???" is spelled outside the \b group on purpose: '?' is not a word
    # character, so \b???\b can never match anything. A guard that silently never
    # fires is worse than no guard, because you stop looking.
    r"(?:\b(?:TBD|TBC|FIXME|XXX|placeholder|pending"
    r"|to\s+be\s+(?:determined|confirmed)|to\s+work)\b|\?{3,})",
    re.IGNORECASE,
)


@dataclass
class Artifact:
    """One thing found in a file that a reader would not have seen."""

    kind: str            # speaker_notes | hidden_image | hidden_sheet | deleted_sheet | ...
    sheet: str = ""
    locator: str = ""
    text: str = ""
    detail: str = ""
    # OCR is a guess. Anything read out of pixels says so, so a reader can weigh
    # it — quoting an OCR error at a CFO is its own kind of disaster.
    confidence: float = 1.0
    data: dict[str, Any] = field(default_factory=dict)


@dataclass
class Forensics:
    path: str
    filename: str
    artifacts: list[Artifact] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)
    warnings: list[str] = field(default_factory=list)

    def of_kind(self, kind: str) -> list[Artifact]:
        return [a for a in self.artifacts if a.kind == kind]

    @property
    def machine_generated(self) -> bool:
        blob = " ".join(str(v) for v in self.metadata.values())
        return any(m.lower() in blob.lower() for m in _GENERATOR_MARKERS)


# ── OCR ──────────────────────────────────────────────────────────────────────

def _ocr(image_bytes: bytes) -> tuple[str, float]:
    """Read text out of an image. Returns ("", 0.0) if we cannot.

    Reuses ICDEV's cached easyocr reader rather than loading a second copy of the
    model. Falls back to pytesseract. If neither is available we still report the
    image as hidden content — knowing that a blank-looking sheet contains a
    picture is most of the value, even when we cannot read it.
    """
    try:
        import io

        from PIL import Image
    except ImportError:  # pragma: no cover
        return "", 0.0

    try:
        img = Image.open(io.BytesIO(image_bytes))
    except Exception:  # noqa: BLE001
        return "", 0.0

    try:
        from tools.document_intelligence.extractors import _get_easyocr_reader

        reader = _get_easyocr_reader()
        if reader is not None:
            import numpy as np

            results = reader.readtext(np.array(img.convert("RGB")))
            if results:
                text = " ".join(r[1] for r in results)
                conf = sum(float(r[2]) for r in results) / len(results)
                return text.strip(), conf
    except Exception:  # noqa: BLE001 - OCR is best-effort by nature
        pass

    try:
        import pytesseract

        text = pytesseract.image_to_string(img)
        return text.strip(), 0.6 if text.strip() else 0.0
    except Exception:  # noqa: BLE001
        return "", 0.0


# ── XLSX ─────────────────────────────────────────────────────────────────────

def _xlsx_sheet_order(zf: zipfile.ZipFile) -> list[dict]:
    """Read workbook.xml for sheet names, ids and visibility."""
    try:
        root = ET.fromstring(zf.read("xl/workbook.xml"))
    except (KeyError, ET.ParseError):
        return []
    sheets = []
    for s in root.iter(f"{{{_NS['main']}}}sheet"):
        sheets.append({
            "name": s.get("name", ""),
            "sheetId": int(s.get("sheetId", "0") or 0),
            "state": s.get("state", "visible"),
            "rId": s.get(f"{{{_NS['r']}}}id", ""),
        })
    return sheets


def _xlsx_deleted_sheets(sheets: list[dict]) -> list[int]:
    """Gaps in the sheetId sequence are sheets that used to exist.

    Excel never reuses a sheetId. A hole in the run is a deletion, and a deleted
    sheet in a costing workbook is a category somebody removed — possibly along
    with its money, possibly not, and the surviving subtotals will not tell you
    which.
    """
    ids = sorted(s["sheetId"] for s in sheets if s["sheetId"] > 0)
    if not ids:
        return []
    return [i for i in range(1, max(ids) + 1) if i not in set(ids)]


def _resolve_target(target: str, base: str) -> str:
    """Resolve an OOXML relationship Target to a path inside the package.

    There are two spellings in the wild and you meet both:

      Excel writes RELATIVE   — "../drawings/drawing1.xml", resolved against the
                                directory holding the .rels file.
      openpyxl writes ABSOLUTE — "/xl/drawings/drawing1.xml", rooted at the
                                package.

    Handling only the relative form works perfectly on every workbook a human
    saved and silently finds nothing in every workbook a script produced — and
    script-produced workbooks are exactly the ones whose numbers most need
    checking. The bug is invisible on a corpus of hand-saved files, which is
    precisely why the fixtures build both.
    """
    target = target.strip()
    if target.startswith("/"):
        return target.lstrip("/")

    parts = base.split("/")[:-1]           # directory of the .rels' owner
    for segment in target.split("/"):
        if segment in ("", "."):
            continue
        if segment == "..":
            if parts:
                parts.pop()
        else:
            parts.append(segment)
    return "/".join(parts)


def _xlsx_images(zf: zipfile.ZipFile, sheets: list[dict]) -> list[tuple[str, str, bytes]]:
    """Find images anchored onto worksheets. Returns (sheet, anchor, bytes).

    This is the one that matters. A sheet can report zero cells and still carry a
    drawing, and openpyxl's normal read path will never mention it.
    """
    found: list[tuple[str, str, bytes]] = []
    names = set(zf.namelist())

    for idx, sheet in enumerate(sheets, start=1):
        owner = f"xl/worksheets/sheet{idx}.xml"
        rels_path = f"xl/worksheets/_rels/sheet{idx}.xml.rels"
        if rels_path not in names or owner not in names:
            continue

        try:
            rels = ET.fromstring(zf.read(rels_path))
        except ET.ParseError:
            continue

        for r in rels.iter(f"{{{_NS['rel']}}}Relationship"):
            if not r.get("Type", "").endswith("/drawing"):
                continue

            drawing_path = _resolve_target(r.get("Target", ""), owner)
            if drawing_path not in names:
                continue

            head, _, tail = drawing_path.rpartition("/")
            drels_path = f"{head}/_rels/{tail}.rels"
            if drels_path not in names:
                continue
            try:
                drels = ET.fromstring(zf.read(drels_path))
            except ET.ParseError:
                continue

            for ir in drels.iter(f"{{{_NS['rel']}}}Relationship"):
                if not ir.get("Type", "").endswith("/image"):
                    continue
                media = _resolve_target(ir.get("Target", ""), drawing_path)
                if media not in names:
                    continue
                found.append((sheet["name"], media.rsplit("/", 1)[-1], zf.read(media)))
    return found


def _xlsx_comments_and_links(zf: zipfile.ZipFile, sheets: list[dict]) -> list[Artifact]:
    out: list[Artifact] = []
    names = set(zf.namelist())

    for idx, sheet in enumerate(sheets, start=1):
        cpath = f"xl/comments{idx}.xml"
        if cpath in names:
            try:
                root = ET.fromstring(zf.read(cpath))
                for c in root.iter(f"{{{_NS['main']}}}comment"):
                    text = " ".join(t.text or "" for t in c.iter(f"{{{_NS['main']}}}t")).strip()
                    if text:
                        out.append(Artifact(
                            kind="cell_comment", sheet=sheet["name"],
                            locator=c.get("ref", ""), text=text,
                            detail="A comment on a cell. Never rendered in a printout.",
                        ))
            except ET.ParseError:
                pass

        rels_path = f"xl/worksheets/_rels/sheet{idx}.xml.rels"
        if rels_path in names:
            try:
                rels = ET.fromstring(zf.read(rels_path))
                for r in rels.iter(f"{{{_NS['rel']}}}Relationship"):
                    if r.get("Type", "").endswith("/hyperlink"):
                        out.append(Artifact(
                            kind="hyperlink", sheet=sheet["name"],
                            text=r.get("Target", ""),
                            detail="A link out of the workbook — often to the real source.",
                        ))
            except ET.ParseError:
                pass
    return out


def _forensics_xlsx(path: Path, out: Forensics) -> None:
    with zipfile.ZipFile(path) as zf:
        sheets = _xlsx_sheet_order(zf)
        out.metadata["sheets"] = [(s["name"], s["state"], s["sheetId"]) for s in sheets]

        for s in sheets:
            if s["state"] != "visible":
                out.artifacts.append(Artifact(
                    kind="hidden_sheet", sheet=s["name"],
                    detail=f"Sheet is {s['state']}. It does not render, and it still has numbers in it.",
                ))

        for missing in _xlsx_deleted_sheets(sheets):
            out.artifacts.append(Artifact(
                kind="deleted_sheet",
                detail=(
                    f"sheetId {missing} is missing from the sequence. Excel never "
                    "reuses a sheetId, so a sheet was deleted. In a costing "
                    "workbook that is a category somebody removed — and the "
                    "surviving subtotals will not tell you whether its money went "
                    "with it."
                ),
                data={"sheet_id": missing},
            ))

        for sheet_name, anchor, blob in _xlsx_images(zf, sheets):
            text, conf = _ocr(blob)
            out.artifacts.append(Artifact(
                kind="hidden_image", sheet=sheet_name, locator=anchor,
                text=text, confidence=conf,
                detail=(
                    "An image anchored onto the worksheet. It carries content that "
                    "no cell holds, and a sheet can look completely empty while "
                    "containing one."
                ),
                data={"bytes": len(blob), "ocr": bool(text)},
            ))

        out.artifacts.extend(_xlsx_comments_and_links(zf, sheets))

    import openpyxl

    wb = openpyxl.load_workbook(str(path), data_only=True)
    for ws in wb.worksheets:
        if ws.max_row == 1 and ws.max_column == 1 and ws["A1"].value is None:
            has_image = any(
                a.kind == "hidden_image" and a.sheet == ws.title for a in out.artifacts
            )
            if has_image:
                out.artifacts.append(Artifact(
                    kind="empty_sheet_with_content", sheet=ws.title,
                    detail=(
                        "This sheet reports zero cells. Excel says it is empty and "
                        "it opens empty — and it is carrying an image. Nobody "
                        "scrolling this workbook would ever find it."
                    ),
                ))
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and _PLACEHOLDER_RE.search(cell.value):
                    out.artifacts.append(Artifact(
                        kind="placeholder", sheet=ws.title, locator=cell.coordinate,
                        text=cell.value.strip(),
                        detail="An unfinished value, left in a document that is being costed.",
                    ))

    props = wb.properties
    out.metadata.update({
        "creator": props.creator or "",
        "last_modified_by": props.lastModifiedBy or "",
        "revision": props.revision or "",
    })


# ── PPTX ─────────────────────────────────────────────────────────────────────

def _forensics_pptx(path: Path, out: Forensics) -> None:
    from pptx import Presentation

    prs = Presentation(str(path))
    for i, slide in enumerate(prs.slides, start=1):
        if not slide.has_notes_slide:
            continue
        notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        if not notes:
            continue
        out.artifacts.append(Artifact(
            kind="speaker_notes", sheet=f"slide{i}", text=notes,
            detail=(
                "The talk track. Routinely holds the constraint, the caveat or the "
                "critical path that the slide itself is too polite to show."
            ),
        ))
        for m in _PLACEHOLDER_RE.finditer(notes):
            out.artifacts.append(Artifact(
                kind="placeholder", sheet=f"slide{i}", text=m.group(0),
                detail="An unfinished item, hidden in the speaker notes.",
            ))

    props = prs.core_properties
    out.metadata.update({
        "creator": props.author or "",
        "last_modified_by": props.last_modified_by or "",
        "revision": props.revision or "",
        "comments": props.comments or "",
        "created": str(props.created) if props.created else "",
    })

    slides_with_notes = len(out.of_kind("speaker_notes"))
    out.metadata["slides"] = len(prs.slides._sldIdLst)  # noqa: SLF001
    out.metadata["slides_with_notes"] = slides_with_notes


# ── DOCX ─────────────────────────────────────────────────────────────────────

def _forensics_docx(path: Path, out: Forensics) -> None:
    with zipfile.ZipFile(path) as zf:
        names = set(zf.namelist())

        if "word/document.xml" in names:
            try:
                root = ET.fromstring(zf.read("word/document.xml"))
                ins = list(root.iter(f"{{{_NS['w']}}}ins"))
                dele = list(root.iter(f"{{{_NS['w']}}}del"))
                if ins or dele:
                    out.artifacts.append(Artifact(
                        kind="tracked_changes",
                        detail=(
                            f"{len(ins)} insertions and {len(dele)} deletions are still "
                            "tracked in this document. The edits — and who made them — "
                            "are part of the file."
                        ),
                        data={"insertions": len(ins), "deletions": len(dele)},
                    ))
            except ET.ParseError:
                pass

        if "word/comments.xml" in names:
            try:
                root = ET.fromstring(zf.read("word/comments.xml"))
                for c in root.iter(f"{{{_NS['w']}}}comment"):
                    text = " ".join(t.text or "" for t in c.iter(f"{{{_NS['w']}}}t")).strip()
                    if text:
                        out.artifacts.append(Artifact(
                            kind="document_comment", text=text,
                            data={"author": c.get(f"{{{_NS['w']}}}author", "")},
                            detail="A reviewer comment. Not printed.",
                        ))
            except ET.ParseError:
                pass

        media = [n for n in names if n.startswith("word/media/")]
        if media:
            total = sum(zf.getinfo(n).file_size for n in media)
            out.artifacts.append(Artifact(
                kind="embedded_media",
                detail=(
                    f"{len(media)} embedded images ({total // 1024}KB). Photographs "
                    "in an engineering document are usually evidence of something."
                ),
                data={"count": len(media), "bytes": total},
            ))


# ── Package-level ────────────────────────────────────────────────────────────

def _sensitivity_label(zf: zipfile.ZipFile) -> str:
    for name in zf.namelist():
        if name.endswith("custom.xml"):
            blob = zf.read(name).decode("utf-8", errors="replace")
            m = re.search(r"MSIP_Label_[0-9a-f-]+_Name[^>]*>\s*<[^>]*>([^<]+)", blob)
            if m:
                return m.group(1).strip()
    return ""


_DISPATCH = {
    ".xlsx": _forensics_xlsx,
    ".xlsm": _forensics_xlsx,
    ".pptx": _forensics_pptx,
    ".docx": _forensics_docx,
}


def analyze(path: str | Path) -> Forensics:
    """Find what a document is carrying that nobody would see by opening it."""
    p = Path(path)
    out = Forensics(path=str(p), filename=p.name)

    if not p.exists():
        out.warnings.append(f"file not found: {p}")
        return out

    handler = _DISPATCH.get(p.suffix.lower())
    if handler is None:
        return out

    try:
        handler(p, out)
    except Exception as exc:  # noqa: BLE001 - one bad file degrades that file only
        out.warnings.append(f"forensics failed: {type(exc).__name__}: {exc}")
        return out

    try:
        with zipfile.ZipFile(p) as zf:
            label = _sensitivity_label(zf)
            if label:
                out.metadata["sensitivity_label"] = label
                out.artifacts.append(Artifact(
                    kind="sensitivity_label", text=label,
                    detail=(
                        f"Stamped '{label}'. Worth knowing before any of it is pasted "
                        "into a cloud model."
                    ),
                ))
    except (zipfile.BadZipFile, OSError):
        pass

    if out.machine_generated:
        out.artifacts.append(Artifact(
            kind="machine_generated",
            detail=(
                "This document was written by a script, not a person "
                f"({out.metadata.get('creator') or out.metadata.get('last_modified_by')}). "
                "That is a fact about how much its numbers are worth, and it "
                "belongs in the credibility assessment rather than in a footnote."
            ),
        ))

    return out


def main() -> int:  # pragma: no cover
    import argparse
    import json

    ap = argparse.ArgumentParser(description="Find what a document is hiding.")
    ap.add_argument("path")
    ap.add_argument("--json", action="store_true")
    args = ap.parse_args()

    f = analyze(args.path)
    if args.json:
        print(json.dumps({
            "filename": f.filename,
            "machine_generated": f.machine_generated,
            "metadata": f.metadata,
            "warnings": f.warnings,
            "artifacts": [vars(a) for a in f.artifacts],
        }, indent=2, default=str))
        return 0

    print(f"{f.filename}")
    for w in f.warnings:
        print(f"  WARN {w}")
    for a in f.artifacts:
        where = f"{a.sheet}!{a.locator}" if a.sheet or a.locator else "-"
        print(f"  [{a.kind}] {where}")
        if a.text:
            preview = a.text[:200].replace("\n", " ")
            conf = f" (ocr conf {a.confidence:.2f})" if a.confidence < 1.0 else ""
            print(f"      {preview!r}{conf}")
    return 0


if __name__ == "__main__":  # pragma: no cover
    raise SystemExit(main())
