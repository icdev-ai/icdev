# CUI // SP-CTI
"""Native python-pptx renderer for the Viz Kernel.

Adds *editable* PowerPoint charts and tables (not images) to a slide, so the
recipient can restyle/retype data in PowerPoint. Fully air-gap (no network).
Diagrams are not native-drawable in pptx, so callers embed a PNG from
``render_png.diagram_to_png`` instead.

Public API (all mutate the given ``slide`` in place):
  add_chart(slide, spec, left, top, width, height, theme)
  add_table(slide, spec, left, top, width, height, theme)
"""
from __future__ import annotations

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Emu, Pt

from tools.viz.palette import get_palette
from tools.viz.spec import ChartSpec, TableSpec

_CHART_TYPE_MAP = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "area": XL_CHART_TYPE.AREA,
    "pie": XL_CHART_TYPE.PIE,
    "donut": XL_CHART_TYPE.DOUGHNUT,
    "gauge": XL_CHART_TYPE.DOUGHNUT,  # gauge approximated as doughnut in native pptx
}


def add_chart(slide, spec: ChartSpec, left, top, width, height,
              theme: str = "midnight_executive"):
    """Add a native, editable chart to *slide*. Returns the chart object."""
    pal = get_palette(theme)
    chart_data = CategoryChartData()

    cats = spec.categories or [
        str(i + 1) for i in range(max((len(s.values) for s in spec.series), default=0))
    ]
    chart_data.categories = cats

    if spec.chart_type in ("pie", "donut", "gauge"):
        s0 = spec.series[0] if spec.series else None
        vals = (s0.values[:len(cats)] if s0 else [0] * len(cats))
        chart_data.add_series(s0.name if s0 else "Series 1", vals)
    else:
        for s in spec.series:
            chart_data.add_series(s.name, s.values[:len(cats)])

    xl_type = _CHART_TYPE_MAP.get(spec.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    gframe = slide.shapes.add_chart(xl_type, left, top, width, height, chart_data)
    chart = gframe.chart

    # Title
    if spec.title:
        chart.has_title = True
        chart.chart_title.text_frame.text = spec.title
        try:
            run = chart.chart_title.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = pal.rgbcolor("accent")
        except (IndexError, AttributeError):
            pass
    else:
        chart.has_title = False

    # Legend (multi-series only)
    multi = len(spec.series) > 1 or spec.chart_type in ("pie", "donut")
    chart.has_legend = multi
    if multi:
        try:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.color.rgb = pal.rgbcolor("text")
            chart.legend.font.size = Pt(9)
        except (AttributeError, ValueError):
            pass

    _color_series(chart, spec, pal)
    return chart


def _color_series(chart, spec: ChartSpec, pal) -> None:
    """Apply the theme series-colour cycle to the plotted data."""
    try:
        plot = chart.plots[0]
    except (IndexError, AttributeError):
        return

    if spec.chart_type in ("pie", "donut", "gauge"):
        # colour individual points of the single series
        try:
            series = plot.series[0]
            for i, point in enumerate(series.points):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = pal.series_rgbcolor(i)
        except (IndexError, AttributeError):
            pass
        return

    for i, series in enumerate(plot.series):
        try:
            if spec.chart_type in ("line",):
                series.format.line.color.rgb = pal.series_rgbcolor(i)
                series.format.line.width = Pt(2.25)
            else:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = pal.series_rgbcolor(i)
        except (AttributeError, ValueError):
            pass


def add_table(slide, spec: TableSpec, left, top, width, height,
              theme: str = "midnight_executive"):
    """Add a native, editable table to *slide*. Returns the table graphic frame."""
    pal = get_palette(theme)
    n_cols = len(spec.headers) or (len(spec.rows[0]) if spec.rows else 1)
    n_rows = len(spec.rows) + (1 if spec.headers else 0)
    n_rows = max(n_rows, 1)

    gframe = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gframe.table

    r0 = 0
    # Header row
    if spec.headers:
        for c, htext in enumerate(spec.headers[:n_cols]):
            cell = table.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = pal.rgbcolor("accent")
            _set_cell_text(cell, htext, pal.rgbcolor("bg"), bold=True, size=12)
        r0 = 1

    # Body rows (zebra striping)
    for ri, row in enumerate(spec.rows):
        tr = r0 + ri
        if tr >= n_rows:
            break
        stripe = pal.rgbcolor("dark") if ri % 2 == 0 else pal.rgbcolor("bg")
        for c in range(n_cols):
            cell = table.cell(tr, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = stripe
            text = str(row[c]) if c < len(row) else ""
            _set_cell_text(cell, text, pal.rgbcolor("text"), bold=False, size=11)

    return gframe


def _set_cell_text(cell, text: str, color, bold: bool, size: int) -> None:
    cell.margin_left = Emu(45720)   # 0.05"
    cell.margin_right = Emu(45720)
    cell.margin_top = Emu(18288)
    cell.margin_bottom = Emu(18288)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
