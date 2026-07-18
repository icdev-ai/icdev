# CUI // SP-CTI
"""PPTX export for Boundary Design Canvas designs (task bdr-feat-5).

Sibling of the existing boundary-canvas exporters (vsdx / json / markdown /
csv / drawio / svg). Produces a native, editable PowerPoint deck:

  1. Title slide       — design name, metadata (node/edge counts, timestamps)
                          and the classification marking when present.
  2. Boundary diagram  — the SAME geometry the SVG/drawio exporters consume
                          (``tools.canvas.export_utils.export_svg`` over the
                          design ``graph_json``), rendered as native DrawingML
                          shapes via ``tools.viz.svg_to_pptx``.
  3. Compliance /       — a native table (``tools.viz.render_pptx.add_table``)
     readiness slide      populated from the design's latest ``bd_assessments``
                          row and the cATO readiness composite when available.

Design notes
------------
* Pure-Python. ``python-pptx`` is already a platform dependency; no npm/node.
* The exporter takes plain dicts (design + optional assessment + optional
  readiness) so it is fully unit-testable without a DB. The Flask route does
  the DB reads and hands the dicts in.
* Everything degrades gracefully: an empty graph yields a placeholder diagram
  slide; missing assessment AND readiness yields a placeholder compliance
  slide — never an exception that costs the caller a deck.
"""

from __future__ import annotations

import io
from datetime import datetime, timezone
from typing import Any, Dict, Optional

from tools.logging.icdev_logger import get_logger

logger = get_logger("icdev.boundary_canvas.export_pptx")

_CUI_BANNER = "CUI // SP-CTI"
_DEFAULT_THEME = "midnight_executive"

# 16:9 canvas in inches (matches the bi_dashboard exporter).
_SLIDE_W_IN = 13.333
_SLIDE_H_IN = 7.5


def _now() -> str:
    return datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")


def _parse_graph(graph_json: Any) -> Dict[str, Any]:
    """Normalise *graph_json* (dict | JSON string | None) to a dict."""
    if graph_json is None:
        return {}
    if isinstance(graph_json, str):
        import json

        try:
            graph_json = json.loads(graph_json)
        except (ValueError, TypeError):
            return {}
    return graph_json if isinstance(graph_json, dict) else {}


def _readiness_rows(readiness: Optional[Dict[str, Any]]) -> list[list[str]]:
    """Metric/value rows for the cATO readiness composite (may be empty)."""
    if not isinstance(readiness, dict):
        return []
    score = readiness.get("score")
    if score is None:
        return []
    rows: list[list[str]] = [
        ["cATO Readiness Score", f"{score}"],
        ["Readiness Band", str(readiness.get("band", "unknown")).upper()],
    ]
    components = readiness.get("components")
    if isinstance(components, dict):
        for name, detail in components.items():
            comp_score = detail.get("score") if isinstance(detail, dict) else None
            label = name.replace("_", " ").title()
            rows.append([f"  {label}", "n/a" if comp_score is None else f"{comp_score}"])
    return rows


def _assessment_rows(assessment: Optional[Dict[str, Any]]) -> list[list[str]]:
    """Metric/value rows for the latest bd_assessments row (may be empty)."""
    if not isinstance(assessment, dict) or not assessment:
        return []
    rows: list[list[str]] = [
        ["Assessment Type", str(assessment.get("assessment_type", "full"))],
        ["Grade", str(assessment.get("grade", "N/A"))],
        ["Score", f"{assessment.get('score', 0)}"],
        ["CAT I Findings", f"{assessment.get('cat1_findings', 0)}"],
        ["CAT II Findings", f"{assessment.get('cat2_findings', 0)}"],
        ["CAT III Findings", f"{assessment.get('cat3_findings', 0)}"],
    ]
    assessed_at = assessment.get("created_at")
    if assessed_at:
        rows.append(["Assessed At", str(assessed_at)])
    return rows


def design_to_pptx(
    design: Dict[str, Any],
    *,
    assessment: Optional[Dict[str, Any]] = None,
    readiness: Optional[Dict[str, Any]] = None,
    theme: str = _DEFAULT_THEME,
) -> bytes:
    """Render a boundary *design* dict to PowerPoint (.pptx) bytes.

    Args:
        design:     Dict with at least ``name`` and ``graph_json`` (dict or JSON
                    string). ``classification``/``description``/timestamps are
                    used when present.
        assessment: Optional latest ``bd_assessments`` row as a dict. When
                    absent the compliance slide falls back to readiness / a
                    clean placeholder.
        readiness:  Optional cATO readiness composite ``{score, band,
                    components}`` (see ``cato_readiness.compute_readiness``).
                    Missing/None ``score`` degrades gracefully.
        theme:      Viz-kernel palette name for the native table.

    Returns:
        The ``.pptx`` file as bytes.
    """
    from pptx import Presentation
    from pptx.util import Inches

    from tools.canvas.export_utils import export_svg
    from tools.viz.render_pptx import add_table
    from tools.viz.spec import TableSpec

    name = str(design.get("name") or "Boundary Design")
    classification = str(design.get("classification") or "CUI").upper()
    graph = _parse_graph(design.get("graph_json"))
    nodes = graph.get("nodes") or []
    edges = graph.get("edges") or []

    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_W_IN)
    prs.slide_height = Inches(_SLIDE_H_IN)

    # ── Slide 1: title ────────────────────────────────────────────────────
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = name
    subtitle_lines = [
        f"{classification} · {_CUI_BANNER}",
        f"Boundary Design Canvas · {len(nodes)} nodes · {len(edges)} edges",
        f"Generated {_now()}",
    ]
    description = str(design.get("description") or "").strip()
    if description:
        subtitle_lines.insert(1, description[:180])
    if len(cover.placeholders) > 1:
        cover.placeholders[1].text = "\n".join(subtitle_lines)

    # ── Slide 2: boundary diagram (same geometry as SVG/drawio export) ─────
    diagram = prs.slides.add_slide(prs.slide_layouts[5])
    diagram.shapes.title.text = "Boundary Diagram"
    box = (Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.2))
    if nodes:
        try:
            svg_source = export_svg(name, graph, "BDC")
            from tools.viz.svg_to_pptx import render_svg_into_slide

            shapes = render_svg_into_slide(diagram, svg_source, *box)
            if not shapes:
                diagram.shapes.add_textbox(*box).text_frame.text = (
                    "Diagram could not be rendered from the design geometry."
                )
        except Exception as exc:  # noqa: BLE001 — never lose the deck over one slide
            logger.warning("boundary diagram render failed for %r: %s", name, exc)
            diagram.shapes.add_textbox(*box).text_frame.text = f"Diagram render failed: {exc}"
    else:
        diagram.shapes.add_textbox(*box).text_frame.text = (
            "This design has no nodes yet — nothing to diagram."
        )

    # ── Slide 3: compliance / readiness ───────────────────────────────────
    compliance = prs.slides.add_slide(prs.slide_layouts[5])
    compliance.shapes.title.text = "Compliance & cATO Readiness"
    rows = _assessment_rows(assessment) + _readiness_rows(readiness)
    if rows:
        spec = TableSpec(title="Compliance & cATO Readiness", headers=["Metric", "Value"], rows=rows)
        try:
            add_table(compliance, spec, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.2), theme=theme)
        except Exception as exc:  # noqa: BLE001 — placeholder rather than 500
            logger.warning("compliance table render failed for %r: %s", name, exc)
            compliance.shapes.add_textbox(
                Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.2)
            ).text_frame.text = f"Compliance table render failed: {exc}"
    else:
        compliance.shapes.add_textbox(
            Inches(0.7), Inches(1.6), Inches(11.9), Inches(2.0)
        ).text_frame.text = (
            "No compliance assessment or cATO readiness data is available for "
            "this design yet. Run a Diagram Grade assessment and populate the "
            "cATO readiness inputs to see them here."
        )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
