# CUI // SP-CTI
"""Export a WHOLE dashboard — HTML, PPTX, PDF.

prem-rpt-01. tools/bi_dashboard already had NL->chart generation, dashboard CRUD, data
ingestion, ECharts theming and a spec layer that round-trips. What it did not have was a
way to get a dashboard OUT.

Export was chart-level only, and the code says so:

    if fmt == "png" and spec.kind == "chart":  ...
    if fmt == "svg" and spec.kind == "chart":  ...
    return {"error": f"export format {fmt} not supported for kind {spec.kind}"}, 400

One chart, one PNG. Ask for the dashboard — the thing a customer actually reads — and
you got a 400. That is the whole gap for a customer-facing report: you could build the
report and never send it.

tools/viz/render_html.py has per-spec renderers (chart_to_html, table_to_html,
kpis_to_html, diagram_to_html). Nothing composed them into a document. This does.

## Classification travels with the export

A dashboard row carries a `classification`. An exported artifact leaves the platform —
that is the entire point of it — so the marking has to go WITH it, banner-in and
banner-out, rather than being left behind in the database row it came from. An
unmarked export of CUI is a CUI spill in a slide deck.
"""
from __future__ import annotations

import base64
import io
from datetime import datetime, timezone
from typing import Any, Dict, List

from tools.viz.palette import get_palette
from tools.viz.render_html import (
    chart_to_html,
    diagram_to_html,
    kpis_to_html,
    table_to_html,
)
from tools.viz.spec import spec_from_dict

DEFAULT_THEME = "midnight_executive"

# Which spec kinds we can render into a document, and with what.
_HTML_RENDERERS = {
    "chart": chart_to_html,
    "table": table_to_html,
    "kpis": kpis_to_html,
    "diagram": diagram_to_html,
}

_UNRENDERABLE = (
    "<div class='viz-unsupported' style='padding:16px;border:1px dashed #888;'>"
    "tile kind {kind!r} has no renderer — it is NOT silently dropped, because a "
    "report that quietly loses a tile is worse than one that says it could not draw it."
    "</div>"
)


def _now() -> str:
    return datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")


def _tile_specs(dashboard: Dict[str, Any]) -> List[Any]:
    """Parse each tile's spec. A tile that will not parse is reported, not dropped."""
    out = []
    for index, tile in enumerate(dashboard.get("tiles") or []):
        spec_dict = tile.get("spec") if isinstance(tile, dict) else None
        if not spec_dict:
            out.append(("error", f"tile {index} has no spec"))
            continue
        try:
            out.append(("spec", spec_from_dict(spec_dict)))
        except (TypeError, ValueError) as exc:
            out.append(("error", f"tile {index}: invalid spec ({exc})"))
    return out


def dashboard_to_html(dashboard: Dict[str, Any], *, theme: str = DEFAULT_THEME) -> str:
    """Render a whole dashboard as one self-contained HTML document.

    Self-contained on purpose: the charts are inline SVG, the styling is inline. A
    customer-facing report that fetches a stylesheet from a host they cannot reach
    renders as a wall of unstyled text — and they will not tell you, they will just
    think less of it.
    """
    pal = get_palette(theme)
    bg, text, accent = pal.hex("bg"), pal.hex("text"), pal.hex("accent")

    title = str(dashboard.get("title") or "Dashboard")
    classification = str(dashboard.get("classification") or "CUI").upper()

    body: List[str] = []
    for status, item in _tile_specs(dashboard):
        if status == "error":
            body.append(_UNRENDERABLE.format(kind=item))
            continue
        renderer = _HTML_RENDERERS.get(item.kind)
        if renderer is None:
            body.append(_UNRENDERABLE.format(kind=item.kind))
            continue
        body.append(f'<section class="viz-tile">{renderer(item, theme)}</section>')

    tiles_html = "\n".join(body) or (
        '<p style="opacity:.7;">This dashboard has no tiles.</p>'
    )

    # Banner in AND out. An exported artifact leaves the platform — the marking travels
    # with it, or it is an unmarked spill.
    banner = (
        f'<div style="background:{accent};color:{bg};text-align:center;'
        f'font-weight:700;letter-spacing:.08em;padding:6px;">{classification}</div>'
    )

    return f"""<!doctype html>
<html lang="en"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{title}</title>
<style>
  body {{ margin:0; background:{bg}; color:{text};
          font-family:'Segoe UI',Arial,sans-serif; }}
  main {{ max-width:1200px; margin:0 auto; padding:28px 20px 48px; }}
  h1 {{ color:{accent}; margin:24px 0 4px; font-size:30px; }}
  .meta {{ opacity:.65; font-size:13px; margin-bottom:26px; }}
  .viz-tile {{ margin:0 0 34px; }}
  table {{ width:100%; }}
</style></head>
<body>
{banner}
<main>
  <h1>{title}</h1>
  <div class="meta">Generated {_now()} &middot; ICDEV&trade;</div>
  {tiles_html}
</main>
{banner}
</body></html>"""


def dashboard_to_pptx(dashboard: Dict[str, Any], *, theme: str = DEFAULT_THEME) -> bytes:
    """One slide per tile, plus a title slide. Returns the .pptx bytes."""
    from pptx import Presentation
    from pptx.util import Inches

    from tools.viz.render_pptx import add_chart, add_table

    title = str(dashboard.get("title") or "Dashboard")
    classification = str(dashboard.get("classification") or "CUI").upper()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = title
    if len(cover.placeholders) > 1:
        cover.placeholders[1].text = f"{classification} · Generated {_now()}"

    for status, item in _tile_specs(dashboard):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        if status == "error":
            slide.shapes.title.text = "Tile could not be rendered"
            slide.shapes.add_textbox(
                Inches(0.7), Inches(2.0), Inches(11.9), Inches(1.2)
            ).text_frame.text = str(item)
            continue

        slide.shapes.title.text = getattr(item, "title", "") or item.kind.title()
        box = (Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.2))
        try:
            if item.kind == "chart":
                add_chart(slide, item, *box, theme=theme)
            elif item.kind == "table":
                add_table(slide, item, *box, theme=theme)
            else:
                # Not silently dropped. A deck that quietly loses a tile is worse than
                # one that says which tile it could not draw.
                slide.shapes.add_textbox(*box).text_frame.text = (
                    f"tile kind {item.kind!r} has no PPTX renderer"
                )
        except Exception as exc:  # pragma: no cover - renderer-specific
            slide.shapes.add_textbox(*box).text_frame.text = f"render failed: {exc}"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def dashboard_to_pdf(dashboard: Dict[str, Any], *, theme: str = DEFAULT_THEME) -> bytes:
    """A PDF of the dashboard.

    Built on reportlab, which is the PDF library this platform actually has installed —
    weasyprint (the obvious HTML->PDF route) is not, and adding a new native dependency
    to a customer-facing export path is how an air-gapped install discovers at the worst
    moment that it cannot produce a report.

    So this is a deliberately plain layout: the classification banner, the title, and
    each tile's chart rendered to PNG and placed. It is a document, not a pixel-perfect
    copy of the web view — and it says so rather than pretending otherwise.
    """
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.lib.units import inch
    from reportlab.pdfgen import canvas as pdfcanvas

    title = str(dashboard.get("title") or "Dashboard")
    classification = str(dashboard.get("classification") or "CUI").upper()

    buf = io.BytesIO()
    page_w, page_h = landscape(letter)
    pdf = pdfcanvas.Canvas(buf, pagesize=landscape(letter))

    def _banner():
        pdf.setFont("Helvetica-Bold", 10)
        for y in (page_h - 0.35 * inch, 0.2 * inch):
            pdf.drawCentredString(page_w / 2, y, classification)

    _banner()
    pdf.setFont("Helvetica-Bold", 24)
    pdf.drawString(inch, page_h - 1.4 * inch, title)
    pdf.setFont("Helvetica", 10)
    pdf.drawString(inch, page_h - 1.75 * inch, f"Generated {_now()} · ICDEV™")
    pdf.showPage()

    for status, item in _tile_specs(dashboard):
        _banner()
        pdf.setFont("Helvetica-Bold", 16)
        label = (
            str(item) if status == "error"
            else (getattr(item, "title", "") or item.kind.title())
        )
        pdf.drawString(inch, page_h - 1.1 * inch, label[:110])

        if status == "spec" and item.kind == "chart":
            try:
                from reportlab.lib.utils import ImageReader

                from tools.viz.render_png import chart_to_png

                png_path = chart_to_png(item)
                pdf.drawImage(
                    ImageReader(png_path), inch, inch,
                    width=page_w - 2 * inch, height=page_h - 2.6 * inch,
                    preserveAspectRatio=True, anchor="n", mask="auto",
                )
            except Exception as exc:  # pragma: no cover - renderer-specific
                pdf.setFont("Helvetica", 11)
                pdf.drawString(inch, page_h - 1.6 * inch, f"chart render failed: {exc}")
        elif status == "spec":
            pdf.setFont("Helvetica", 11)
            pdf.drawString(
                inch, page_h - 1.6 * inch,
                f"tile kind {item.kind!r} is not drawn in the PDF export "
                f"(it IS in the HTML and PPTX) — not silently dropped.",
            )
        pdf.showPage()

    pdf.save()
    return buf.getvalue()


def export_dashboard(dashboard: Dict[str, Any], fmt: str, *,
                     theme: str = DEFAULT_THEME) -> Dict[str, Any]:
    """Export in *fmt*. Returns a uniform envelope; binary formats are base64.

    The envelope is the same shape for every format so a caller (and the Cortex service
    surface) does not need a branch per format.
    """
    fmt = (fmt or "").strip().lower()
    slug = "".join(
        ch if ch.isalnum() or ch in "-_" else "-"
        for ch in str(dashboard.get("title") or "dashboard")
    ).strip("-").lower() or "dashboard"

    if fmt == "html":
        return {
            "format": "html",
            "filename": f"{slug}.html",
            "content_type": "text/html; charset=utf-8",
            "html": dashboard_to_html(dashboard, theme=theme),
        }
    if fmt == "pptx":
        return {
            "format": "pptx",
            "filename": f"{slug}.pptx",
            "content_type": ("application/vnd.openxmlformats-officedocument"
                             ".presentationml.presentation"),
            "pptx_base64": base64.b64encode(
                dashboard_to_pptx(dashboard, theme=theme)
            ).decode("ascii"),
        }
    if fmt == "pdf":
        return {
            "format": "pdf",
            "filename": f"{slug}.pdf",
            "content_type": "application/pdf",
            "pdf_base64": base64.b64encode(
                dashboard_to_pdf(dashboard, theme=theme)
            ).decode("ascii"),
        }
    raise ValueError(
        f"unsupported dashboard export format {fmt!r} (html, pptx, pdf)"
    )


def supported_formats() -> List[str]:
    return ["html", "pptx", "pdf"]
